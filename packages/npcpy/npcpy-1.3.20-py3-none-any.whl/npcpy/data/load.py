import fitz
import pandas as pd
import json
import io
from PIL import Image
import numpy as np
from typing import Optional, List
import os
import tempfile
import subprocess

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None

def load_csv(file_path):
    df = pd.read_csv(file_path)
    return df

def load_json(file_path):
    with open(file_path, "r", encoding='utf-8') as f:
        data = json.load(f)
    return data

def load_txt(file_path):
    with open(file_path, "r", encoding='utf-8') as f:
        text = f.read()
    return text

def load_excel(file_path):
    df = pd.read_excel(file_path)
    return df

def load_image(file_path):
    img = Image.open(file_path)
    img_array = np.array(img)
    df = pd.DataFrame(
        {
            "image_array": [img_array.tobytes()],
            "shape": [img_array.shape],
            "dtype": [img_array.dtype.str],
        }
    )
    return df

def load_pdf(file_path):
    pdf_document = fitz.open(file_path)
    full_text = ""
    for page in pdf_document:
        full_text += page.get_text() + "\n"
    return full_text

def load_docx(file_path):
    if Document is None:
        raise ImportError("Please install python-docx to load .docx files.")
    doc = Document(file_path)
    full_text = "\n".join([para.text for para in doc.paragraphs])
    return full_text

def load_pptx(file_path):
    if Presentation is None:
        raise ImportError("Please install python-pptx to load .pptx files.")
    prs = Presentation(file_path)
    full_text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text += shape.text + "\n"
    return full_text

def load_html(file_path):
    if BeautifulSoup is None:
        raise ImportError("Please install beautifulsoup4 to load .html files.")
    with open(file_path, 'r', encoding='utf-8') as f:
        soup = BeautifulSoup(f, 'html.parser')
    return soup.get_text(separator='\n', strip=True)

extension_map = {
    "PNG": "images",
    "JPG": "images",
    "JPEG": "images",
    "GIF": "images",
    "SVG": "images",
    "WEBP": "images",
    "BMP": "images",
    "TIFF": "images",
    "MP4": "videos",
    "AVI": "videos",
    "MOV": "videos",
    "WMV": "videos",
    "MPG": "videos",
    "MPEG": "videos",
    "WEBM": "videos",
    "MKV": "videos",
    "DOCX": "documents",
    "PPTX": "documents",
    "PDF": "documents",
    "XLSX": "documents",
    "TXT": "documents",
    "CSV": "documents",
    "MD": "documents",
    "HTML": "documents",
    "HTM": "documents",
    "MP3": "audio",
    "WAV": "audio",
    "M4A": "audio",
    "AAC": "audio",
    "FLAC": "audio",
    "OGG": "audio",
    "ZIP": "archives",
    "RAR": "archives",
    "7Z": "archives",
    "TAR": "archives",
    "GZ": "archives",
}

def _chunk_text(full_content: str, chunk_size: int) -> List[str]:
    """Split long content into reasonably sized chunks for model input."""
    chunks = []
    for i in range(0, len(full_content), chunk_size):
        chunk = full_content[i:i+chunk_size].strip()
        if chunk:
            chunks.append(chunk)
    return chunks

def _transcribe_audio(file_path: str, language: Optional[str] = None) -> str:
    """
    Best-effort audio transcription using optional dependencies.
    Tries faster-whisper, then openai/whisper. Falls back to metadata only.
    """
    # Prefer the existing audio module helper if present
    try:
        from npcpy.data.audio import transcribe_audio_file  # type: ignore
        text = transcribe_audio_file(file_path, language=language)
        if text:
            return text
    except Exception:
        pass

    # Try faster-whisper first
    try:
        from faster_whisper import WhisperModel
        try:
            import torch
            device = "cuda" if torch.cuda.is_available() else "cpu"
        except Exception:
            device = "cpu"
        model = WhisperModel("small", device=device)
        segments, _ = model.transcribe(file_path, language=language, beam_size=5)
        return " ".join(seg.text.strip() for seg in segments if seg.text).strip()
    except Exception:
        pass

    # Fallback: openai/whisper
    try:
        import whisper
        model = whisper.load_model("small")
        result = model.transcribe(file_path, language=language)
        return result.get("text", "").strip()
    except Exception:
        pass

    # Last resort metadata message
    return f"[Audio file at {file_path}; install faster-whisper or whisper for transcription]"

def load_audio(file_path: str, language: Optional[str] = None) -> str:
    """Load and transcribe an audio file into text."""
    transcript = _transcribe_audio(file_path, language=language)
    if transcript:
        return transcript
    return f"[Audio file at {file_path}; no transcript available]"

def _extract_audio_from_video(file_path: str, max_duration: int = 600) -> Optional[str]:
    """
    Use ffmpeg to dump the audio track from a video into a temp wav for transcription.
    Returns the temp path or None.
    """
    try:
        temp_audio = tempfile.NamedTemporaryFile(delete=False, suffix=".wav")
        temp_audio.close()
        cmd = [
            "ffmpeg",
            "-y",
            "-i",
            file_path,
            "-vn",
            "-ac",
            "1",
            "-ar",
            "16000",
            "-t",
            str(max_duration),
            temp_audio.name,
        ]
        subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        return temp_audio.name
    except Exception:
        return None

def load_video(file_path: str, language: Optional[str] = None, max_audio_seconds: int = 600) -> str:
    """
    Summarize a video by reporting metadata and (optionally) transcribing its audio track.
    """
    # Prefer the video module helper if present
    try:
        from npcpy.data.video import summarize_video_file  # type: ignore
        return summarize_video_file(file_path, language=language, max_audio_seconds=max_audio_seconds)
    except Exception:
        pass

    # Fallback to minimal summary/transcription
    meta_bits = []
    try:
        import cv2
        video = cv2.VideoCapture(file_path)
        fps = video.get(cv2.CAP_PROP_FPS)
        frame_count = int(video.get(cv2.CAP_PROP_FRAME_COUNT))
        width = int(video.get(cv2.CAP_PROP_FRAME_WIDTH))
        height = int(video.get(cv2.CAP_PROP_FRAME_HEIGHT))
        duration = frame_count / fps if fps else 0
        meta_bits.append(
            f"Video file: {os.path.basename(file_path)} | {width}x{height} | {fps:.2f} fps | {frame_count} frames | ~{duration:.1f}s"
        )
        video.release()
    except Exception:
        meta_bits.append(f"Video file: {os.path.basename(file_path)}")

    audio_path = _extract_audio_from_video(file_path, max_duration=max_audio_seconds)
    transcript = ""
    if audio_path:
        try:
            transcript = _transcribe_audio(audio_path, language=language)
        finally:
            try:
                os.remove(audio_path)
            except Exception:
                pass

    if transcript:
        meta_bits.append("Audio transcript:")
        meta_bits.append(transcript)
    else:
        meta_bits.append("[No transcript extracted; ensure ffmpeg and faster-whisper/whisper are installed]")

    return "\n".join(meta_bits)

def load_file_contents(file_path, chunk_size=None):
    file_ext = os.path.splitext(file_path)[1].upper().lstrip('.')
    full_content = ""
    if not isinstance(chunk_size, int):
        chunk_size=8000
    try:
        if file_ext == 'PDF':
            full_content = load_pdf(file_path)
        elif file_ext == 'DOCX':
            full_content = load_docx(file_path)
        elif file_ext == 'PPTX':
            full_content = load_pptx(file_path)
        elif file_ext in ['HTML', 'HTM']:
            full_content = load_html(file_path)
        elif file_ext == 'CSV':
            df = load_csv(file_path)
            full_content = df.to_string()
        elif file_ext in ['XLS', 'XLSX']:
            df = load_excel(file_path)
            full_content = df.to_string()
        elif file_ext in ['TXT', 'MD', 'PY', 'JSX', 'TSX', 'TS', 'JS', 'JSON', 'SQL', 'NPC', 'JINX', 'LINE', 'YAML', 'DART', 'JAVA']:
            full_content = load_txt(file_path)
        elif file_ext == 'JSON':
            data = load_json(file_path)
            full_content = json.dumps(data, indent=2)
        elif file_ext in ['MP3', 'WAV', 'M4A', 'AAC', 'FLAC', 'OGG']:
            full_content = load_audio(file_path)
        elif file_ext in ['MP4', 'AVI', 'MOV', 'WMV', 'MPG', 'MPEG', 'WEBM', 'MKV']:
            full_content = load_video(file_path)
        else:
            return [f"Unsupported file format for content loading: {file_ext}"]
        
        if not full_content:
            return []

        return _chunk_text(full_content, chunk_size)
            
    except Exception as e:
        return [f"Error loading file {file_path}: {str(e)}"]
