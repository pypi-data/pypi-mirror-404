"""Template class for AI-friendly template analysis.

Provides comprehensive analysis of PowerPoint templates with
AI-readable descriptions of layouts, placeholders, and theme colors.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation as PptxPresentation

from .layout import LayoutDescription, LayoutRecommendation, analyze_layout, recommend_layout


class Template:
    """AI-friendly wrapper for PowerPoint templates.

    Provides template analysis and creates presentations with
    semantic, high-level APIs designed for AI agents.

    Example:
        >>> template = Template("corporate.pptx")
        >>> layouts = template.describe()
        >>> print(layouts[0]["name"], layouts[0]["description"])
        >>> pres = template.create_presentation()
    """

    def __init__(self, template_path: str | Path) -> None:
        """Load and analyze a template.

        Args:
            template_path: Path to the .pptx template file
        """
        self._path = Path(template_path)
        if not self._path.exists():
            raise FileNotFoundError(f"Template not found: {template_path}")

        # Load the template with python-pptx
        self._pptx = PptxPresentation(self._path)

        # Analyze layouts
        self._layouts: list[LayoutDescription] = []
        self._analyze_layouts()

        # Extract theme
        self._colors: dict[str, str] = {}
        self._fonts: dict[str, str] = {}
        self._extract_theme()

    def _analyze_layouts(self) -> None:
        """Analyze all layouts and build descriptions."""
        for idx, layout in enumerate(self._pptx.slide_layouts):
            # Convert placeholder info to dict format
            placeholders = []
            for ph in layout.placeholders:
                try:
                    placeholders.append({
                        "type": self._get_placeholder_type(ph),
                        "idx": ph.placeholder_format.idx,
                        "name": ph.name,
                        "x": ph.left,
                        "y": ph.top,
                        "cx": ph.width,
                        "cy": ph.height,
                    })
                except Exception:
                    continue

            layout_desc = analyze_layout(
                name=layout.name,
                index=idx,
                placeholders=placeholders,
            )
            self._layouts.append(layout_desc)

    def _get_placeholder_type(self, placeholder) -> str:
        """Get placeholder type as string."""
        from pptx.enum.shapes import PP_PLACEHOLDER

        type_map = {
            PP_PLACEHOLDER.TITLE: "title",
            PP_PLACEHOLDER.CENTER_TITLE: "ctrTitle",
            PP_PLACEHOLDER.SUBTITLE: "subTitle",
            PP_PLACEHOLDER.BODY: "body",
            PP_PLACEHOLDER.OBJECT: "obj",
            PP_PLACEHOLDER.CHART: "chart",
            PP_PLACEHOLDER.TABLE: "tbl",
            PP_PLACEHOLDER.PICTURE: "pic",
            PP_PLACEHOLDER.FOOTER: "ftr",
            PP_PLACEHOLDER.DATE: "dt",
            PP_PLACEHOLDER.SLIDE_NUMBER: "sldNum",
        }
        ph_type = placeholder.placeholder_format.type
        return type_map.get(ph_type, "body")

    def _extract_theme(self) -> None:
        """Extract theme colors and fonts from the slide master's theme."""
        from lxml import etree

        ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

        try:
            # Access theme through slide master's relationships
            sm = self._pptx.slide_master
            for rel in sm.part.rels.values():
                if "theme" in rel.reltype:
                    theme_elem = etree.fromstring(rel.target_part.blob)

                    # Extract colors
                    clr_scheme = theme_elem.find(f".//{ns}clrScheme")
                    if clr_scheme is not None:
                        for child in clr_scheme:
                            name = child.tag.split("}")[-1]
                            srgb = child.find(f"{ns}srgbClr")
                            if srgb is not None:
                                self._colors[name] = f"#{srgb.get('val')}"
                            else:
                                sys_clr = child.find(f"{ns}sysClr")
                                if sys_clr is not None:
                                    last_clr = sys_clr.get("lastClr")
                                    if last_clr:
                                        self._colors[name] = f"#{last_clr}"

                    # Extract fonts
                    font_scheme = theme_elem.find(f".//{ns}fontScheme")
                    if font_scheme is not None:
                        major = font_scheme.find(f"{ns}majorFont")
                        minor = font_scheme.find(f"{ns}minorFont")
                        if major is not None:
                            latin = major.find(f"{ns}latin")
                            if latin is not None:
                                self._fonts["heading"] = latin.get("typeface", "Calibri Light")
                        if minor is not None:
                            latin = minor.find(f"{ns}latin")
                            if latin is not None:
                                self._fonts["body"] = latin.get("typeface", "Calibri")
                    break
        except Exception:
            pass

        # Defaults if not found
        if "heading" not in self._fonts:
            self._fonts["heading"] = "Calibri Light"
        if "body" not in self._fonts:
            self._fonts["body"] = "Calibri"

    @property
    def path(self) -> Path:
        """Get the template file path."""
        return self._path

    @property
    def colors(self) -> dict[str, str]:
        """Get theme colors.

        Returns:
            Dict mapping color names to hex values.
            Common keys: dk1, lt1, dk2, lt2, accent1-6, hlink, folHlink

        Example:
            >>> colors = template.colors
            >>> primary = colors.get("accent1", "#000000")
        """
        return self._colors.copy()

    @property
    def fonts(self) -> dict[str, str]:
        """Get theme fonts.

        Returns:
            Dict with 'heading' and 'body' font names.

        Example:
            >>> fonts = template.fonts
            >>> print(f"Use {fonts['heading']} for titles")
        """
        return self._fonts.copy()

    def describe(self) -> list[dict[str, Any]]:
        """Get AI-friendly descriptions of all layouts.

        Returns a list of layout descriptions that can be easily
        consumed by AI agents to understand what layouts are available
        and how to use them.

        Returns:
            List of layout description dicts with:
            - name: Layout name (e.g., "Title Slide")
            - index: Layout index for selection
            - type: Classified type (title_slide, content, etc.)
            - description: Human-readable description
            - placeholders: Dict of semantic placeholder names
            - best_for: List of content types this layout suits

        Example:
            >>> layouts = template.describe()
            >>> for layout in layouts:
            ...     print(f"{layout['name']}: {layout['description']}")
            ...     print(f"  Best for: {', '.join(layout['best_for'])}")
        """
        return [layout.to_dict() for layout in self._layouts]

    def describe_as_text(self) -> str:
        """Get a text description of the template for AI prompts.

        Returns a formatted string that can be included in AI system
        prompts to give context about the template.

        Returns:
            Multi-line string describing the template.

        Example:
            >>> description = template.describe_as_text()
            >>> # Include in AI prompt:
            >>> prompt = f"Template info:\\n{description}\\n\\nCreate a presentation..."
        """
        lines = [
            f"Template: {self._path.name}",
            f"Theme Colors: {', '.join(f'{k}={v}' for k, v in list(self._colors.items())[:6])}",
            f"Fonts: heading={self._fonts.get('heading')}, body={self._fonts.get('body')}",
            "",
            "Available Layouts:",
        ]

        for layout in self._layouts:
            ph_names = list(layout.placeholders.keys())
            lines.append(f"  {layout.index}: {layout.name} ({layout.layout_type.value})")
            lines.append(f"      Placeholders: {', '.join(ph_names)}")
            lines.append(f"      Best for: {', '.join(layout.best_for)}")

        return "\n".join(lines)

    def get_layout(self, name_or_index: str | int) -> LayoutDescription | None:
        """Get a specific layout by name or index.

        Args:
            name_or_index: Layout name (fuzzy matched) or index

        Returns:
            LayoutDescription or None if not found

        Example:
            >>> layout = template.get_layout("title")
            >>> if layout:
            ...     print(layout.placeholders)
        """
        if isinstance(name_or_index, int):
            for layout in self._layouts:
                if layout.index == name_or_index:
                    return layout
            return None

        # Fuzzy name match
        name_lower = name_or_index.lower()
        for layout in self._layouts:
            if name_lower in layout.name.lower() or layout.name.lower() in name_lower:
                return layout
        return None

    def recommend_layout(
        self,
        content_type: str,
        has_image: bool = False,
        bullet_count: int = 0,
    ) -> list[dict[str, Any]]:
        """Get layout recommendations for content type.

        Args:
            content_type: Type of content ("bullets", "comparison", etc.)
            has_image: Whether content includes an image
            bullet_count: Number of bullet points

        Returns:
            List of recommendations sorted by confidence

        Example:
            >>> recs = template.recommend_layout("comparison")
            >>> best = recs[0]
            >>> print(f"Use '{best['name']}' ({best['confidence']:.0%} confidence)")
        """
        recommendations = recommend_layout(
            layouts=self._layouts,
            content_type=content_type,
            has_image=has_image,
            bullet_count=bullet_count,
        )

        return [
            {
                "name": r.layout_name,
                "index": r.layout_index,
                "confidence": r.confidence,
                "reason": r.reason,
            }
            for r in recommendations
        ]

    def create_presentation(self) -> "Presentation":
        """Create a new presentation from this template.

        Returns a Presentation object with high-level methods for
        adding slides with semantic APIs.

        Returns:
            Presentation object

        Example:
            >>> pres = template.create_presentation()
            >>> pres.add_title_slide("My Title", "Subtitle")
            >>> pres.save("output.pptx")
        """
        from .presentation import Presentation

        return Presentation(self)

    def get_layout_names(self) -> list[str]:
        """Get list of all layout names.

        Returns:
            List of layout names in order
        """
        return [layout.name for layout in self._layouts]

    def __repr__(self) -> str:
        return f"Template({self._path.name!r}, {len(self._layouts)} layouts)"
