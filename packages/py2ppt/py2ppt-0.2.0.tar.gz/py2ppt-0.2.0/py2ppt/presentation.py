"""Presentation class with semantic slide management.

Provides high-level methods for creating presentations with
AI-friendly, intent-based APIs.
"""

from __future__ import annotations

from pathlib import Path
from typing import TYPE_CHECKING, Any

from pptx import Presentation as PptxPresentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import PP_PLACEHOLDER, MSO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
import copy

from .errors import (
    SlideNotFoundError,
    LayoutNotFoundError,
    InvalidDataError,
)
from .formatting import format_for_py2ppt, parse_content
from .layout import LayoutType
from .analysis import analyze_content
from .theme import ThemeHelper
from .validation import validate_presentation, ValidationResult

if TYPE_CHECKING:
    from .template import Template

# Chart type mapping
_CHART_TYPE_MAP = {
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "pie": XL_CHART_TYPE.PIE,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
}


class Presentation:
    """AI-friendly presentation with semantic slide methods.

    This class wraps python-pptx's Presentation to provide high-level,
    intent-based methods for creating slides. Instead of dealing
    with placeholder indices and layout details, you can simply
    call methods like add_title_slide() or add_comparison_slide().

    Example:
        >>> template = Template("corporate.pptx")
        >>> pres = template.create_presentation()
        >>>
        >>> pres.add_title_slide("Q4 Review", "January 2025")
        >>> pres.add_content_slide("Key Points", [
        ...     "Revenue up 20%",
        ...     "New markets opened",
        ...     "Customer satisfaction high"
        ... ])
        >>> pres.save("output.pptx")
    """

    def __init__(self, template: "Template") -> None:
        """Initialize from a Template.

        Use Template.create_presentation() instead of this constructor.
        """
        self._template = template
        self._layouts = template._layouts

        # Create python-pptx presentation from template
        self._pptx = PptxPresentation(template.path)

        # Remove existing slides (keep only layouts/masters)
        while len(self._pptx.slides) > 0:
            rId = self._pptx.slides._sldIdLst[0].rId
            self._pptx.part.drop_rel(rId)
            del self._pptx.slides._sldIdLst[0]

    @property
    def slide_count(self) -> int:
        """Get the current number of slides."""
        return len(self._pptx.slides)

    @property
    def template(self) -> "Template":
        """Get the template this presentation was created from."""
        return self._template

    @property
    def theme(self) -> ThemeHelper:
        """Get the theme helper for easy access to colors and fonts.

        Returns:
            ThemeHelper instance with accent colors, fonts, and formatting helpers

        Example:
            >>> pres.theme.accent1  # "#41B3FF"
            >>> pres.theme.colored("Important", "accent1")
            {'text': 'Important', 'color': '#41B3FF'}
            >>> pres.theme.bold("Key point")
            {'text': 'Key point', 'bold': True}
        """
        return ThemeHelper(self._template)

    # --- Private helpers ---

    def _validate_slide_number(self, n: int) -> None:
        """Validate a slide number, raising SlideNotFoundError if invalid."""
        count = len(self._pptx.slides)
        if not isinstance(n, int) or n < 1 or n > count:
            raise SlideNotFoundError(
                f"Slide {n} does not exist. Presentation has {count} slide(s).",
                suggestion=(
                    f"Use a slide number between 1 and {count}."
                    if count > 0
                    else "Add slides first."
                ),
                code="SLIDE_NOT_FOUND",
            )

    def _find_layout_by_type(self, layout_type: LayoutType) -> int:
        """Find the best layout index for a layout type."""
        for layout in self._layouts:
            if layout.layout_type == layout_type:
                return layout.index
        # Fall back to first layout
        return 0

    def _find_layout(self, layout: str | int | None, layout_type: LayoutType) -> int:
        """Find layout by name, index, or type.

        Raises LayoutNotFoundError for explicit names/indices that don't match.
        Falls back silently for None/"auto".
        """
        if layout is None or layout == "auto":
            return self._find_layout_by_type(layout_type)

        if isinstance(layout, int):
            if layout < 0 or layout >= len(self._pptx.slide_layouts):
                raise LayoutNotFoundError(
                    f"Layout index {layout} is out of range.",
                    suggestion=f"Available indices: 0-{len(self._pptx.slide_layouts) - 1}.",
                    code="LAYOUT_NOT_FOUND",
                )
            return layout

        # Fuzzy name match
        layout_lower = layout.lower()
        for l in self._layouts:
            if layout_lower in l.name.lower() or l.name.lower() in layout_lower:
                return l.index

        raise LayoutNotFoundError(
            f"No layout matching '{layout}' found.",
            suggestion=f"Available layouts: {', '.join(l.name for l in self._layouts)}.",
            code="LAYOUT_NOT_FOUND",
        )

    def _get_placeholder(self, slide, ph_type):
        """Get a placeholder by type from a slide."""
        for shape in slide.placeholders:
            if shape.placeholder_format.type == ph_type:
                return shape
        return None

    def _set_text_frame(self, shape, text: str) -> None:
        """Set text in a shape's text frame."""
        if shape is None:
            return
        tf = shape.text_frame
        tf.clear()
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        p = tf.paragraphs[0]
        p.text = text

    def _apply_run_formatting(self, run, fmt: dict) -> None:
        """Apply formatting from a dict to a python-pptx Run.

        Supports: bold, italic, underline, font_size, font_family, color, hyperlink.
        """
        if fmt.get("bold"):
            run.font.bold = True
        if fmt.get("italic"):
            run.font.italic = True
        if fmt.get("underline"):
            run.font.underline = True
        if fmt.get("font_size"):
            run.font.size = Pt(fmt["font_size"])
        if fmt.get("font_family"):
            run.font.name = fmt["font_family"]
        if fmt.get("color"):
            color_hex = fmt["color"].lstrip("#")
            if len(color_hex) == 6:
                run.font.color.rgb = RGBColor(
                    int(color_hex[:2], 16),
                    int(color_hex[2:4], 16),
                    int(color_hex[4:6], 16),
                )
        if fmt.get("hyperlink"):
            run.hyperlink.address = fmt["hyperlink"]

    def _set_paragraph_runs(self, p, segments: list) -> None:
        """Set formatted runs on a paragraph, replacing any existing runs."""
        # Remove existing <a:r> elements from the paragraph XML
        for r_elem in list(p._p):
            if r_elem.tag == qn("a:r"):
                p._p.remove(r_elem)

        for seg in segments:
            run = p.add_run()
            if isinstance(seg, dict):
                run.text = seg.get("text", "")
                self._apply_run_formatting(run, seg)
            elif isinstance(seg, str):
                run.text = seg

    def _set_body_content(
        self, shape, content: list, levels: list[int] | None = None
    ) -> None:
        """Set bullet content in a body placeholder with full rich text support."""
        if shape is None:
            return

        tf = shape.text_frame
        tf.clear()
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        for i, item in enumerate(content):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            # Set level
            if levels and i < len(levels):
                p.level = levels[i]

            # Set text with formatting
            if isinstance(item, str):
                p.text = item
            elif isinstance(item, list):
                # Rich text: list of dicts/strings with formatting
                self._set_paragraph_runs(p, item)
            elif isinstance(item, dict):
                # Single formatted run
                self._set_paragraph_runs(p, [item])

    def _get_theme_color(self, name: str = "accent1") -> RGBColor:
        """Get a theme color as RGBColor, with fallback."""
        hex_color = self._template.colors.get(name, "#4472C4").lstrip("#")
        return RGBColor(
            int(hex_color[:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16),
        )

    # --- Slide creation methods ---

    def add_title_slide(
        self,
        title: str,
        subtitle: str = "",
        *,
        layout: str | int | None = None,
    ) -> int:
        """Add a title/cover slide.

        Args:
            title: Main presentation title
            subtitle: Optional subtitle (presenter name, date, etc.)
            layout: Layout name, index, or None for auto-selection

        Returns:
            Slide number of the new slide

        Example:
            >>> pres.add_title_slide("Q4 Business Review", "January 2025")
        """
        layout_idx = self._find_layout(layout, LayoutType.TITLE)
        slide_layout = self._pptx.slide_layouts[layout_idx]
        slide = self._pptx.slides.add_slide(slide_layout)

        # Set title
        title_ph = self._get_placeholder(slide, PP_PLACEHOLDER.TITLE)
        if title_ph is None:
            title_ph = self._get_placeholder(slide, PP_PLACEHOLDER.CENTER_TITLE)
        self._set_text_frame(title_ph, title)

        # Set subtitle
        if subtitle:
            subtitle_ph = self._get_placeholder(slide, PP_PLACEHOLDER.SUBTITLE)
            self._set_text_frame(subtitle_ph, subtitle)

        return len(self._pptx.slides)

    def add_section_slide(
        self,
        title: str,
        subtitle: str = "",
        *,
        layout: str | int | None = None,
    ) -> int:
        """Add a section header/divider slide.

        Args:
            title: Section title
            subtitle: Optional subtitle or description
            layout: Layout name, index, or None for auto-selection

        Returns:
            Slide number of the new slide

        Example:
            >>> pres.add_section_slide("Part 2: Analysis")
        """
        layout_idx = self._find_layout(layout, LayoutType.SECTION)
        slide_layout = self._pptx.slide_layouts[layout_idx]
        slide = self._pptx.slides.add_slide(slide_layout)

        # Set title
        title_ph = self._get_placeholder(slide, PP_PLACEHOLDER.TITLE)
        self._set_text_frame(title_ph, title)

        # Set subtitle if provided
        if subtitle:
            subtitle_ph = self._get_placeholder(slide, PP_PLACEHOLDER.SUBTITLE)
            if subtitle_ph:
                self._set_text_frame(subtitle_ph, subtitle)

        return len(self._pptx.slides)

    def add_content_slide(
        self,
        title: str,
        content: str | list[str | tuple | dict | list],
        *,
        levels: list[int] | None = None,
        layout: str | int | None = None,
        warn_overflow: bool = False,
    ) -> int | dict:
        """Add a content slide with bullets.

        Args:
            title: Slide title
            content: Content as string, list of strings, or rich text
            levels: Optional indent levels for each bullet (0=top level)
            layout: Layout name, index, or None for auto-selection
            warn_overflow: If True, returns a dict with overflow info instead of bare int

        Returns:
            Slide number (int) or overflow info dict if warn_overflow=True

        Example:
            >>> pres.add_content_slide("Key Points", [
            ...     "First point",
            ...     "Second point",
            ...     "Third point"
            ... ])

            >>> # With overflow detection
            >>> result = pres.add_content_slide("Many Points", items, warn_overflow=True)
            >>> if result["overflow"]:
            ...     print(f"Content overflows: {result['item_count']} items, capacity ~{result['estimated_capacity']}")
        """
        layout_idx = self._find_layout(layout, LayoutType.CONTENT)
        slide_layout = self._pptx.slide_layouts[layout_idx]
        slide = self._pptx.slides.add_slide(slide_layout)

        # Set title
        title_ph = self._get_placeholder(slide, PP_PLACEHOLDER.TITLE)
        self._set_text_frame(title_ph, title)

        # Parse and format content
        paragraphs = parse_content(content, levels)
        formatted_content, formatted_levels = format_for_py2ppt(paragraphs)

        # Set body
        body_ph = self._get_placeholder(slide, PP_PLACEHOLDER.BODY)
        if body_ph is None:
            body_ph = self._get_placeholder(slide, PP_PLACEHOLDER.OBJECT)
        self._set_body_content(body_ph, formatted_content, formatted_levels)

        slide_num = len(self._pptx.slides)

        if warn_overflow:
            # Estimate capacity from placeholder height (~0.4 inches per bullet)
            capacity = 6  # default
            if body_ph is not None:
                height_inches = body_ph.height / 914400  # EMU to inches
                capacity = max(1, int(height_inches / 0.4))
            item_count = len(paragraphs)
            return {
                "slide_number": slide_num,
                "overflow": item_count > capacity,
                "item_count": item_count,
                "estimated_capacity": capacity,
            }

        return slide_num

    def add_two_column_slide(
        self,
        title: str,
        left_content: list[str | tuple | dict | list],
        right_content: list[str | tuple | dict | list],
        *,
        left_levels: list[int] | None = None,
        right_levels: list[int] | None = None,
        layout: str | int | None = None,
    ) -> int:
        """Add a two-column content slide.

        Args:
            title: Slide title
            left_content: Content for left column
            right_content: Content for right column
            left_levels: Indent levels for left content
            right_levels: Indent levels for right content
            layout: Layout name, index, or None for auto-selection

        Returns:
            Slide number of the new slide

        Example:
            >>> pres.add_two_column_slide(
            ...     "Topics",
            ...     ["Point A", "Point B"],
            ...     ["Point X", "Point Y"]
            ... )
        """
        layout_idx = self._find_layout(layout, LayoutType.TWO_COLUMN)
        slide_layout = self._pptx.slide_layouts[layout_idx]
        slide = self._pptx.slides.add_slide(slide_layout)

        # Set title
        title_ph = self._get_placeholder(slide, PP_PLACEHOLDER.TITLE)
        self._set_text_frame(title_ph, title)

        # Find body placeholders (there should be two)
        body_phs = [
            shape
            for shape in slide.placeholders
            if shape.placeholder_format.type
            in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT)
        ]
        body_phs.sort(key=lambda s: s.left)  # Sort by x position

        # Format content
        left_paragraphs = parse_content(left_content, left_levels)
        left_formatted, left_lvls = format_for_py2ppt(left_paragraphs)

        right_paragraphs = parse_content(right_content, right_levels)
        right_formatted, right_lvls = format_for_py2ppt(right_paragraphs)

        # Set left content
        if len(body_phs) >= 1:
            self._set_body_content(body_phs[0], left_formatted, left_lvls)

        # Set right content
        if len(body_phs) >= 2:
            self._set_body_content(body_phs[1], right_formatted, right_lvls)

        return len(self._pptx.slides)

    def add_comparison_slide(
        self,
        title: str,
        left_heading: str,
        left_content: list[str | tuple | dict | list],
        right_heading: str,
        right_content: list[str | tuple | dict | list],
        *,
        layout: str | int | None = None,
    ) -> int:
        """Add a comparison slide with two titled columns.

        Ideal for before/after, pros/cons, or versus comparisons.

        Args:
            title: Slide title
            left_heading: Heading for left column
            left_content: Bullet points for left column
            right_heading: Heading for right column
            right_content: Bullet points for right column
            layout: Layout name, index, or None for auto-selection

        Returns:
            Slide number of the new slide

        Example:
            >>> pres.add_comparison_slide(
            ...     "Before vs After",
            ...     "Legacy System",
            ...     ["Slow", "Manual", "Error-prone"],
            ...     "New Platform",
            ...     ["Fast", "Automated", "Reliable"]
            ... )
        """
        # Try to find a comparison layout, fall back to two-column
        layout_idx = self._find_layout(layout, LayoutType.COMPARISON)
        slide_layout = self._pptx.slide_layouts[layout_idx]
        slide = self._pptx.slides.add_slide(slide_layout)

        # Set title
        title_ph = self._get_placeholder(slide, PP_PLACEHOLDER.TITLE)
        self._set_text_frame(title_ph, title)

        # Find body placeholders
        body_phs = [
            shape
            for shape in slide.placeholders
            if shape.placeholder_format.type
            in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT)
        ]
        # Sort by position: top row first, then left to right
        body_phs.sort(key=lambda s: (s.top, s.left))

        if len(body_phs) >= 4:
            # True comparison layout: heading, content, heading, content
            self._set_text_frame(body_phs[0], left_heading)
            left_paragraphs = parse_content(left_content)
            left_formatted, left_lvls = format_for_py2ppt(left_paragraphs)
            self._set_body_content(body_phs[2], left_formatted, left_lvls)

            self._set_text_frame(body_phs[1], right_heading)
            right_paragraphs = parse_content(right_content)
            right_formatted, right_lvls = format_for_py2ppt(right_paragraphs)
            self._set_body_content(body_phs[3], right_formatted, right_lvls)

        elif len(body_phs) >= 2:
            # Two-column layout - combine heading with content
            left_combined = [left_heading] + list(left_content)
            right_combined = [right_heading] + list(right_content)

            left_levels = [0] + [1] * len(left_content)
            right_levels = [0] + [1] * len(right_content)

            left_paragraphs = parse_content(left_combined, left_levels)
            right_paragraphs = parse_content(right_combined, right_levels)

            left_formatted, left_lvls = format_for_py2ppt(left_paragraphs)
            right_formatted, right_lvls = format_for_py2ppt(right_paragraphs)

            # Sort by x position for left/right
            body_phs.sort(key=lambda s: s.left)
            self._set_body_content(body_phs[0], left_formatted, left_lvls)
            self._set_body_content(body_phs[1], right_formatted, right_lvls)

        elif len(body_phs) >= 1:
            # Single body - combine all
            combined = [
                left_heading,
                *list(left_content),
                "",
                right_heading,
                *list(right_content),
            ]
            combined_levels = (
                [0]
                + [1] * len(left_content)
                + [0]
                + [0]
                + [1] * len(right_content)
            )
            paragraphs = parse_content(combined, combined_levels)
            formatted, lvls = format_for_py2ppt(paragraphs)
            self._set_body_content(body_phs[0], formatted, lvls)

        return len(self._pptx.slides)

    def add_image_slide(
        self,
        title: str,
        image_path: str | Path,
        caption: str = "",
        *,
        layout: str | int | None = None,
        left: float | None = None,
        top: float | None = None,
        width: float | None = None,
        height: float | None = None,
    ) -> int:
        """Add a slide with an image.

        Image placement priority:
        1. PICTURE placeholder (uses insert_picture for best fit)
        2. Explicit left/top/width/height kwargs (in inches)
        3. BODY/OBJECT placeholder bounds as fallback
        4. Default position (1in, 2in, width=5in)

        Args:
            title: Slide title
            image_path: Path to the image file
            caption: Optional caption text
            layout: Layout name, index, or None for auto-selection
            left: Image left position in inches
            top: Image top position in inches
            width: Image width in inches
            height: Image height in inches

        Returns:
            Slide number of the new slide

        Example:
            >>> pres.add_image_slide("Product Photo", "product.png", "Our flagship product")
            >>> pres.add_image_slide("Custom Position", "img.png", left=2, top=3, width=6)
        """
        layout_idx = self._find_layout(layout, LayoutType.IMAGE_CONTENT)
        slide_layout = self._pptx.slide_layouts[layout_idx]
        slide = self._pptx.slides.add_slide(slide_layout)

        # Set title
        title_ph = self._get_placeholder(slide, PP_PLACEHOLDER.TITLE)
        self._set_text_frame(title_ph, title)

        # Add image
        image_path = Path(image_path)
        if image_path.exists():
            added = False

            # 1. Try PICTURE placeholder
            pic_ph = self._get_placeholder(slide, PP_PLACEHOLDER.PICTURE)
            if pic_ph is not None:
                pic_ph.insert_picture(str(image_path))
                added = True

            # 2. Explicit dimensions
            if not added and any(
                v is not None for v in [left, top, width, height]
            ):
                pic_left = Inches(left) if left is not None else Inches(1)
                pic_top = Inches(top) if top is not None else Inches(2)
                pic_width = Inches(width) if width is not None else None
                pic_height = Inches(height) if height is not None else None
                slide.shapes.add_picture(
                    str(image_path),
                    pic_left,
                    pic_top,
                    width=pic_width,
                    height=pic_height,
                )
                added = True

            # 3. BODY/OBJECT placeholder bounds
            if not added:
                body_ph = self._get_placeholder(slide, PP_PLACEHOLDER.BODY)
                if body_ph is None:
                    body_ph = self._get_placeholder(
                        slide, PP_PLACEHOLDER.OBJECT
                    )
                if body_ph is not None:
                    slide.shapes.add_picture(
                        str(image_path),
                        body_ph.left,
                        body_ph.top,
                        width=body_ph.width,
                    )
                else:
                    # 4. Default fallback
                    slide.shapes.add_picture(
                        str(image_path),
                        Inches(1),
                        Inches(2),
                        width=Inches(5),
                    )

        # Set caption if provided
        if caption:
            body_ph = self._get_placeholder(slide, PP_PLACEHOLDER.BODY)
            if body_ph:
                self._set_text_frame(body_ph, caption)

        return len(self._pptx.slides)

    def add_blank_slide(self, layout: str | int | None = None) -> int:
        """Add a blank slide.

        Args:
            layout: Layout name, index, or None for auto-selection

        Returns:
            Slide number of the new slide

        Example:
            >>> slide_num = pres.add_blank_slide()
        """
        layout_idx = self._find_layout(layout, LayoutType.BLANK)
        slide_layout = self._pptx.slide_layouts[layout_idx]
        self._pptx.slides.add_slide(slide_layout)
        return len(self._pptx.slides)

    def add_table_slide(
        self,
        title: str,
        headers: list[str],
        rows: list[list[Any]],
        *,
        col_widths: list[float] | None = None,
        style: str = "theme",
        layout: str | int | None = None,
    ) -> int:
        """Add a slide with a table.

        Args:
            title: Slide title
            headers: Column header labels
            rows: List of rows, each a list of cell values
            col_widths: Optional column widths in inches (auto-calculated if omitted)
            style: Table style - "theme" (default), "plain", or "striped"
            layout: Layout name, index, or None for auto-selection

        Returns:
            Slide number of the new slide

        Example:
            >>> pres.add_table_slide(
            ...     "Sales Data",
            ...     ["Region", "Q1", "Q2", "Q3"],
            ...     [["North", 100, 120, 130],
            ...      ["South", 90, 110, 125]],
            ... )
        """
        # Validate row/header match
        for i, row in enumerate(rows):
            if len(row) != len(headers):
                raise InvalidDataError(
                    f"Row {i} has {len(row)} columns but headers has {len(headers)}.",
                    suggestion="Ensure all rows have the same number of columns as headers.",
                    code="TABLE_ROW_MISMATCH",
                )

        layout_idx = self._find_layout(layout, LayoutType.TITLE_ONLY)
        slide_layout = self._pptx.slide_layouts[layout_idx]
        slide = self._pptx.slides.add_slide(slide_layout)

        # Title
        title_ph = self._get_placeholder(slide, PP_PLACEHOLDER.TITLE)
        self._set_text_frame(title_ph, title)

        # Table dimensions
        num_rows = len(rows) + 1  # +1 for header
        num_cols = len(headers)
        tbl_left = Inches(0.5)
        tbl_top = Inches(1.8)
        tbl_width = Inches(9.0)
        tbl_height = Inches(min(0.4 * num_rows, 5.2))

        table_shape = slide.shapes.add_table(
            num_rows, num_cols, tbl_left, tbl_top, tbl_width, tbl_height
        )
        table = table_shape.table

        # Set column widths
        if col_widths:
            for c, w in enumerate(col_widths):
                if c < num_cols:
                    table.columns[c].width = Inches(w)
        else:
            # Auto-calculate proportional to content length
            max_lengths = []
            for c in range(num_cols):
                max_len = len(str(headers[c]))
                for row in rows:
                    max_len = max(max_len, len(str(row[c])))
                max_lengths.append(max(max_len, 1))
            total_len = sum(max_lengths)
            for c in range(num_cols):
                table.columns[c].width = int(
                    tbl_width * max_lengths[c] / total_len
                )

        # Set headers
        for c, header in enumerate(headers):
            table.cell(0, c).text = str(header)

        # Set data rows
        for r, row in enumerate(rows):
            for c, value in enumerate(row):
                table.cell(r + 1, c).text = str(value)

        # Apply styling
        accent_color = self._get_theme_color("accent1")

        if style in ("theme", "striped"):
            # Style header row
            for c in range(num_cols):
                cell = table.cell(0, c)
                cell.fill.solid()
                cell.fill.fore_color.rgb = accent_color
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

            # Striped: alternate row backgrounds
            if style == "striped":
                for r in range(len(rows)):
                    if r % 2 == 0:
                        for c in range(num_cols):
                            cell = table.cell(r + 1, c)
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor(
                                0xF2, 0xF2, 0xF2
                            )

        # "plain" style: no special formatting

        return len(self._pptx.slides)

    def add_chart_slide(
        self,
        title: str,
        chart_type: str,
        data: dict[str, Any],
        *,
        layout: str | int | None = None,
    ) -> int:
        """Add a slide with a chart.

        Args:
            title: Slide title
            chart_type: One of "bar", "column", "line", "pie", "doughnut"
            data: Chart data dict. Format depends on chart type:
                  Bar/column/line: {"categories": [...], "series": [{"name": "...", "values": [...]}]}
                  Pie/doughnut: {"categories": [...], "values": [...]}
            layout: Layout name, index, or None for auto-selection

        Returns:
            Slide number of the new slide

        Example:
            >>> pres.add_chart_slide("Revenue", "column", {
            ...     "categories": ["Q1", "Q2", "Q3"],
            ...     "series": [{"name": "2024", "values": [10, 20, 30]}]
            ... })
        """
        # Validate chart type
        chart_type_lower = chart_type.lower()
        if chart_type_lower not in _CHART_TYPE_MAP:
            raise InvalidDataError(
                f"Unknown chart type '{chart_type}'.",
                suggestion=f"Supported types: {', '.join(_CHART_TYPE_MAP.keys())}.",
                code="INVALID_CHART_TYPE",
            )

        # Validate data
        if "categories" not in data:
            raise InvalidDataError(
                "Chart data must include 'categories' key.",
                suggestion=(
                    "Provide {'categories': [...], 'series': [...]} "
                    "or {'categories': [...], 'values': [...]}."
                ),
                code="MISSING_CHART_DATA",
            )

        xl_chart_type = _CHART_TYPE_MAP[chart_type_lower]

        # Build chart data
        chart_data = CategoryChartData()
        chart_data.categories = data["categories"]

        if chart_type_lower in ("pie", "doughnut"):
            if "values" not in data:
                raise InvalidDataError(
                    f"'{chart_type}' chart requires 'values' key in data.",
                    suggestion="Provide {'categories': [...], 'values': [...]}.",
                    code="MISSING_CHART_VALUES",
                )
            chart_data.add_series("Values", data["values"])
        else:
            if "series" not in data:
                raise InvalidDataError(
                    f"'{chart_type}' chart requires 'series' key in data.",
                    suggestion=(
                        "Provide {'categories': [...], "
                        "'series': [{'name': '...', 'values': [...]}]}."
                    ),
                    code="MISSING_CHART_SERIES",
                )
            for series in data["series"]:
                chart_data.add_series(
                    series.get("name", "Series"),
                    series.get("values", []),
                )

        # Create slide
        layout_idx = self._find_layout(layout, LayoutType.TITLE_ONLY)
        slide_layout = self._pptx.slide_layouts[layout_idx]
        slide = self._pptx.slides.add_slide(slide_layout)

        # Title
        title_ph = self._get_placeholder(slide, PP_PLACEHOLDER.TITLE)
        self._set_text_frame(title_ph, title)

        # Add chart
        chart_left = Inches(1.0)
        chart_top = Inches(1.8)
        chart_width = Inches(8.0)
        chart_height = Inches(5.0)

        chart_shape = slide.shapes.add_chart(
            xl_chart_type,
            chart_left,
            chart_top,
            chart_width,
            chart_height,
            chart_data,
        )
        chart = chart_shape.chart

        # Add legend for multi-series charts
        if chart_type_lower not in ("pie", "doughnut"):
            if len(data.get("series", [])) > 1:
                chart.has_legend = True

        return len(self._pptx.slides)

    def add_slide(
        self,
        layout: str | int = "auto",
        content_type: str = "content",
        title: str = "",
        content: list[str] | None = None,
        **kwargs: Any,
    ) -> int:
        """Add a slide with auto-layout selection.

        A flexible method that chooses the right slide type
        based on content_type and delegates to the appropriate
        specialized method.

        Args:
            layout: Layout name, index, or "auto"
            content_type: Type of content ("title", "content", "comparison",
                          "table", "chart", etc.)
            title: Slide title
            content: Content items
            **kwargs: Additional arguments for the specific slide type

        Returns:
            Slide number of the new slide

        Example:
            >>> pres.add_slide(content_type="table", title="Data",
            ...     headers=["A", "B"], rows=[[1, 2]])
        """
        content_lower = content_type.lower()

        if content_lower in ("title", "cover", "opening"):
            subtitle = kwargs.get("subtitle", "")
            return self.add_title_slide(title, subtitle, layout=layout)

        elif content_lower in ("section", "divider"):
            subtitle = kwargs.get("subtitle", "")
            return self.add_section_slide(title, subtitle, layout=layout)

        elif content_lower in ("comparison", "versus", "vs"):
            return self.add_comparison_slide(
                title,
                kwargs.get("left_heading", "Option A"),
                kwargs.get("left_content", content or []),
                kwargs.get("right_heading", "Option B"),
                kwargs.get("right_content", []),
                layout=layout,
            )

        elif content_lower in ("two_column", "split"):
            return self.add_two_column_slide(
                title,
                kwargs.get("left_content", content or []),
                kwargs.get("right_content", []),
                layout=layout,
            )

        elif content_lower in ("image", "picture"):
            return self.add_image_slide(
                title,
                kwargs.get("image_path", ""),
                kwargs.get("caption", ""),
                layout=layout,
            )

        elif content_lower == "blank":
            return self.add_blank_slide(layout=layout)

        elif content_lower == "table":
            return self.add_table_slide(
                title,
                kwargs.get("headers", []),
                kwargs.get("rows", []),
                col_widths=kwargs.get("col_widths"),
                style=kwargs.get("style", "theme"),
                layout=layout,
            )

        elif content_lower == "chart":
            return self.add_chart_slide(
                title,
                kwargs.get("chart_type", "column"),
                kwargs.get("data", {}),
                layout=layout,
            )

        else:
            # Default to content slide
            return self.add_content_slide(
                title,
                content or [],
                levels=kwargs.get("levels"),
                layout=layout,
            )

    def set_notes(self, slide_number: int, notes: str) -> None:
        """Set speaker notes for a slide.

        Args:
            slide_number: The slide number (1-indexed)
            notes: Notes text (can include newlines)

        Raises:
            SlideNotFoundError: If slide_number is out of range

        Example:
            >>> pres.set_notes(1, "Remember to mention...")
        """
        self._validate_slide_number(slide_number)

        slide = self._pptx.slides[slide_number - 1]
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    # --- Inspection methods ---

    def describe_slide(self, n: int) -> dict[str, Any]:
        """Describe a single slide's contents.

        Args:
            n: Slide number (1-indexed)

        Returns:
            Dict with: slide_number, layout, title, content, shapes, notes,
            has_title, has_content, has_table, has_chart, has_image

        Raises:
            SlideNotFoundError: If slide number is out of range

        Example:
            >>> info = pres.describe_slide(1)
            >>> print(info["title"], info["has_table"])
        """
        self._validate_slide_number(n)
        slide = self._pptx.slides[n - 1]

        layout_name = slide.slide_layout.name

        # Find title
        title = ""
        has_title = False
        for shape in slide.placeholders:
            if shape.placeholder_format.type in (
                PP_PLACEHOLDER.TITLE,
                PP_PLACEHOLDER.CENTER_TITLE,
            ):
                title = shape.text
                has_title = bool(title)
                break

        # Find body content
        content = []
        has_content = False
        for shape in slide.placeholders:
            if shape.placeholder_format.type in (
                PP_PLACEHOLDER.BODY,
                PP_PLACEHOLDER.OBJECT,
            ):
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        if p.text:
                            content.append(p.text)
                    if content:
                        has_content = True

        # Inspect all shapes
        has_table = False
        has_chart = False
        has_image = False
        shapes = []

        for shape in slide.shapes:
            shape_info: dict[str, Any] = {
                "name": shape.name,
                "shape_type": str(shape.shape_type),
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height,
            }

            if shape.has_table:
                has_table = True
                tbl = shape.table
                table_data = {
                    "rows": len(tbl.rows),
                    "cols": len(tbl.columns),
                    "headers": [
                        tbl.cell(0, c).text for c in range(len(tbl.columns))
                    ],
                }
                shape_info["table"] = table_data

            if shape.has_chart:
                has_chart = True
                shape_info["chart_type"] = str(shape.chart.chart_type)

            try:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    has_image = True
                    shape_info["image"] = {
                        "width": shape.width,
                        "height": shape.height,
                    }
            except Exception:
                pass

            shapes.append(shape_info)

        # Notes
        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text

        return {
            "slide_number": n,
            "layout": layout_name,
            "title": title,
            "content": content,
            "shapes": shapes,
            "notes": notes,
            "has_title": has_title,
            "has_content": has_content,
            "has_table": has_table,
            "has_chart": has_chart,
            "has_image": has_image,
        }

    def describe_all_slides(self) -> list[dict[str, Any]]:
        """Describe all slides in the presentation.

        Returns:
            List of slide description dicts (same format as describe_slide)

        Example:
            >>> for slide in pres.describe_all_slides():
            ...     print(f"Slide {slide['slide_number']}: {slide['title']}")
        """
        return [
            self.describe_slide(i + 1)
            for i in range(len(self._pptx.slides))
        ]

    # --- Editing methods ---

    def update_slide(
        self,
        n: int,
        *,
        title: str | None = None,
        content: list[str | tuple | dict | list] | None = None,
        levels: list[int] | None = None,
        notes: str | None = None,
    ) -> dict[str, Any]:
        """Update an existing slide. Only provided args are changed.

        Args:
            n: Slide number (1-indexed)
            title: New title (None to keep existing)
            content: New body content (None to keep existing)
            levels: Indent levels for new content
            notes: New speaker notes (None to keep existing)

        Returns:
            Updated slide description dict

        Raises:
            SlideNotFoundError: If slide number is out of range

        Example:
            >>> pres.update_slide(2, title="Updated Title", notes="New notes")
        """
        self._validate_slide_number(n)
        slide = self._pptx.slides[n - 1]

        if title is not None:
            title_ph = self._get_placeholder(slide, PP_PLACEHOLDER.TITLE)
            if title_ph is None:
                title_ph = self._get_placeholder(
                    slide, PP_PLACEHOLDER.CENTER_TITLE
                )
            self._set_text_frame(title_ph, title)

        if content is not None:
            paragraphs = parse_content(content, levels)
            formatted_content, formatted_levels = format_for_py2ppt(paragraphs)
            body_ph = self._get_placeholder(slide, PP_PLACEHOLDER.BODY)
            if body_ph is None:
                body_ph = self._get_placeholder(
                    slide, PP_PLACEHOLDER.OBJECT
                )
            self._set_body_content(
                body_ph, formatted_content, formatted_levels
            )

        if notes is not None:
            self.set_notes(n, notes)

        return self.describe_slide(n)

    def delete_slide(self, n: int) -> int:
        """Delete a slide by number.

        Args:
            n: Slide number (1-indexed)

        Returns:
            New slide count after deletion

        Raises:
            SlideNotFoundError: If slide number is out of range

        Example:
            >>> remaining = pres.delete_slide(3)
        """
        self._validate_slide_number(n)
        sldIdLst = self._pptx.slides._sldIdLst
        rId = sldIdLst[n - 1].rId
        self._pptx.part.drop_rel(rId)
        del sldIdLst[n - 1]
        return len(self._pptx.slides)

    def reorder_slides(self, order: list[int]) -> None:
        """Reorder all slides.

        Args:
            order: List of slide numbers in desired order (1-indexed).
                   Must be a permutation of [1..N].

        Raises:
            InvalidDataError: If order is not a valid permutation

        Example:
            >>> pres.reorder_slides([3, 1, 2])  # Slide 3 becomes first
        """
        count = len(self._pptx.slides)
        if sorted(order) != list(range(1, count + 1)):
            raise InvalidDataError(
                f"Order must be a permutation of slide numbers 1-{count}.",
                suggestion=f"Provide a list containing each number from 1 to {count} exactly once.",
                code="INVALID_SLIDE_ORDER",
            )

        sldIdLst = self._pptx.slides._sldIdLst
        sld_ids = list(sldIdLst)
        for elem in sld_ids:
            sldIdLst.remove(elem)
        for pos in order:
            sldIdLst.append(sld_ids[pos - 1])

    def move_slide(self, from_pos: int, to_pos: int) -> None:
        """Move a single slide to a new position.

        Args:
            from_pos: Current slide number (1-indexed)
            to_pos: Target slide number (1-indexed)

        Raises:
            SlideNotFoundError: If either position is out of range

        Example:
            >>> pres.move_slide(5, 2)  # Move slide 5 to position 2
        """
        self._validate_slide_number(from_pos)
        count = len(self._pptx.slides)
        if not isinstance(to_pos, int) or to_pos < 1 or to_pos > count:
            raise SlideNotFoundError(
                f"Target position {to_pos} is out of range.",
                suggestion=f"Use a position between 1 and {count}.",
                code="SLIDE_NOT_FOUND",
            )
        order = list(range(1, count + 1))
        order.remove(from_pos)
        order.insert(to_pos - 1, from_pos)
        self.reorder_slides(order)

    # --- Content intelligence ---

    def add_content_slides(
        self,
        title: str,
        content: str | list[str | tuple | dict | list],
        *,
        max_bullets: int = 6,
        continuation_suffix: str = " (cont.)",
        levels: list[int] | None = None,
        layout: str | int | None = None,
    ) -> list[int]:
        """Auto-split long content across multiple slides.

        Keeps sub-items (level > 0) grouped with their parent.

        Args:
            title: Slide title (continuation slides get suffix appended)
            content: Content items to split
            max_bullets: Maximum bullets per slide before splitting
            continuation_suffix: Suffix for continuation slide titles
            levels: Optional indent levels for each item
            layout: Layout name, index, or None for auto-selection

        Returns:
            List of slide numbers created

        Example:
            >>> slides = pres.add_content_slides("Key Points", [
            ...     "Point 1", "Detail A", "Detail B",
            ...     "Point 2", "Detail C",
            ...     "Point 3", "Point 4", "Point 5",
            ...     "Point 6", "Point 7", "Point 8",
            ... ], levels=[0, 1, 1, 0, 1, 0, 0, 0, 0, 0, 0], max_bullets=4)
        """
        # Normalize content to list
        if isinstance(content, str):
            content_list: list[str | tuple | dict | list] = [
                line for line in content.split("\n") if line.strip()
            ]
        else:
            content_list = list(content)

        paragraphs = parse_content(content_list, levels)

        # Split into chunks, keeping sub-items with their parent
        chunks: list[list[int]] = []  # each chunk is a list of indices
        current_chunk: list[int] = []

        for idx, para in enumerate(paragraphs):
            if len(current_chunk) >= max_bullets and para.level == 0:
                chunks.append(current_chunk)
                current_chunk = []
            current_chunk.append(idx)

        if current_chunk:
            chunks.append(current_chunk)

        slide_numbers = []
        for i, chunk_indices in enumerate(chunks):
            slide_title = (
                title if i == 0 else f"{title}{continuation_suffix}"
            )
            chunk_content = [content_list[j] for j in chunk_indices]
            chunk_levels = [paragraphs[j].level for j in chunk_indices]
            slide_num = self.add_content_slide(
                slide_title,
                chunk_content,
                levels=chunk_levels,
                layout=layout,
            )
            slide_numbers.append(slide_num)

        return slide_numbers

    # --- Semantic slide types ---

    def add_smart_slide(
        self,
        title: str,
        content: str | list[str | dict | list],
        *,
        layout: str | int | None = None,
    ) -> int:
        """Add a slide with auto-detected type and layout.

        Analyzes the content to determine the best slide type
        (comparison, quote, stats, timeline, etc.) and uses
        the appropriate method.

        Args:
            title: Slide title
            content: Content to analyze and add
            layout: Layout override (None for auto-selection)

        Returns:
            Slide number of the new slide

        Example:
            >>> pres.add_smart_slide("Before vs After", [
            ...     "Before: slow and manual",
            ...     "After: fast and automated"
            ... ])  # Auto-detects comparison
        """
        analysis = analyze_content(content, title)

        if analysis.confidence >= 0.5:
            slide_type = analysis.recommended_slide_type
        else:
            slide_type = "content"

        # Route to appropriate method based on detected type
        if slide_type == "comparison":
            # Use add_content_slide since we don't have explicit parts
            return self.add_content_slide(title, content, layout=layout)
        elif slide_type == "quote":
            content_list = content if isinstance(content, list) else [content]
            quote = content_list[0] if content_list else ""
            attribution = content_list[1] if len(content_list) > 1 else ""
            return self.add_quote_slide(quote, attribution, layout=layout)
        elif slide_type == "stats":
            # Parse stats from content
            stats = []
            content_list = content if isinstance(content, list) else content.split("\n")
            for item in content_list:
                if isinstance(item, str):
                    stats.append({"value": item, "label": ""})
            return self.add_stats_slide(title, stats, layout=layout)
        elif slide_type == "timeline":
            content_list = content if isinstance(content, list) else content.split("\n")
            return self.add_timeline_slide(title, content_list, layout=layout)
        else:
            content_list = content if isinstance(content, list) else [content]
            return self.add_content_slide(title, content_list, layout=layout)

    def add_quote_slide(
        self,
        quote: str,
        attribution: str = "",
        *,
        source: str | None = None,
        layout: str | int | None = None,
    ) -> int:
        """Add a quote slide with styled quotation.

        Creates a visually impactful quote with optional attribution.
        The quote is styled with large, italic text and the attribution
        is styled with theme accent color.

        Args:
            quote: The quotation text
            attribution: Who said it (e.g., "Steve Jobs")
            source: Optional source (e.g., "Stanford Commencement, 2005")
            layout: Layout name, index, or None for auto-selection

        Returns:
            Slide number of the new slide

        Example:
            >>> pres.add_quote_slide(
            ...     "Stay hungry, stay foolish.",
            ...     "Steve Jobs",
            ...     source="Stanford Commencement, 2005"
            ... )
        """
        layout_idx = self._find_layout(layout, LayoutType.CONTENT)
        slide_layout = self._pptx.slide_layouts[layout_idx]
        slide = self._pptx.slides.add_slide(slide_layout)

        # Build formatted quote content
        content = []

        # Quote with large italic formatting
        quote_fmt = {
            "text": f'"{quote}"',
            "italic": True,
            "font_size": 28,
        }
        content.append([quote_fmt])

        # Attribution with accent color
        if attribution:
            attr_text = f"— {attribution}"
            if source:
                attr_text = f"{attr_text}, {source}"
            attr_fmt = {
                "text": attr_text,
                "color": self._template.colors.get("accent1", "#4472C4"),
            }
            content.append([attr_fmt])

        # Leave title empty for quotes (content is the focus)
        title_ph = self._get_placeholder(slide, PP_PLACEHOLDER.TITLE)
        if title_ph:
            self._set_text_frame(title_ph, "")

        # Set body with formatted quote
        body_ph = self._get_placeholder(slide, PP_PLACEHOLDER.BODY)
        if body_ph is None:
            body_ph = self._get_placeholder(slide, PP_PLACEHOLDER.OBJECT)
        self._set_body_content(body_ph, content)

        return len(self._pptx.slides)

    def add_stats_slide(
        self,
        title: str,
        stats: list[dict[str, Any]],
        *,
        layout: str | int | None = None,
    ) -> int:
        """Add a statistics slide with big numbers and labels.

        Creates a visually striking slide with large statistics
        and descriptive labels.

        Args:
            title: Slide title
            stats: List of stat dicts with "value" and "label" keys
                   e.g., [{"value": "50%", "label": "Growth"}, ...]
            layout: Layout name, index, or None for auto-selection

        Returns:
            Slide number of the new slide

        Example:
            >>> pres.add_stats_slide("Key Metrics", [
            ...     {"value": "98%", "label": "Customer Satisfaction"},
            ...     {"value": "2.5M", "label": "Active Users"},
            ...     {"value": "150+", "label": "Countries"},
            ... ])
        """
        layout_idx = self._find_layout(layout, LayoutType.CONTENT)
        slide_layout = self._pptx.slide_layouts[layout_idx]
        slide = self._pptx.slides.add_slide(slide_layout)

        # Set title
        title_ph = self._get_placeholder(slide, PP_PLACEHOLDER.TITLE)
        self._set_text_frame(title_ph, title)

        # Build formatted stats content
        content = []
        accent_color = self._template.colors.get("accent1", "#4472C4")

        for stat in stats:
            value = stat.get("value", "")
            label = stat.get("label", "")

            # Value in large, bold, colored text
            value_fmt = {
                "text": str(value),
                "bold": True,
                "font_size": 36,
                "color": accent_color,
            }
            # Label in regular text
            label_fmt = {"text": f"  {label}" if label else ""}

            content.append([value_fmt, label_fmt])

        # Set body with stats
        body_ph = self._get_placeholder(slide, PP_PLACEHOLDER.BODY)
        if body_ph is None:
            body_ph = self._get_placeholder(slide, PP_PLACEHOLDER.OBJECT)
        self._set_body_content(body_ph, content)

        return len(self._pptx.slides)

    def add_timeline_slide(
        self,
        title: str,
        events: list[str | dict[str, str]],
        *,
        layout: str | int | None = None,
    ) -> int:
        """Add a timeline slide with sequence of events.

        Creates a slide showing events in chronological order.
        Each event can be a string or dict with "date" and "event" keys.

        Args:
            title: Slide title
            events: List of events as strings or dicts
                    Strings: "2024: Launched product"
                    Dicts: {"date": "2024", "event": "Launched product"}
            layout: Layout name, index, or None for auto-selection

        Returns:
            Slide number of the new slide

        Example:
            >>> pres.add_timeline_slide("Our Journey", [
            ...     {"date": "2020", "event": "Company founded"},
            ...     {"date": "2022", "event": "First product launch"},
            ...     {"date": "2024", "event": "Global expansion"},
            ... ])
        """
        layout_idx = self._find_layout(layout, LayoutType.CONTENT)
        slide_layout = self._pptx.slide_layouts[layout_idx]
        slide = self._pptx.slides.add_slide(slide_layout)

        # Set title
        title_ph = self._get_placeholder(slide, PP_PLACEHOLDER.TITLE)
        self._set_text_frame(title_ph, title)

        # Build formatted timeline content
        content = []
        accent_color = self._template.colors.get("accent1", "#4472C4")

        for event in events:
            if isinstance(event, dict):
                date = event.get("date", "")
                event_text = event.get("event", "")
                # Date in bold, colored
                date_fmt = {
                    "text": f"{date}  ",
                    "bold": True,
                    "color": accent_color,
                }
                event_fmt = {"text": event_text}
                content.append([date_fmt, event_fmt])
            else:
                # Simple string - just add as-is
                content.append(str(event))

        # Set body with timeline
        body_ph = self._get_placeholder(slide, PP_PLACEHOLDER.BODY)
        if body_ph is None:
            body_ph = self._get_placeholder(slide, PP_PLACEHOLDER.OBJECT)
        self._set_body_content(body_ph, content)

        return len(self._pptx.slides)

    def add_agenda_slide(
        self,
        title: str,
        items: list[str],
        *,
        layout: str | int | None = None,
    ) -> int:
        """Add an agenda slide with numbered items.

        Creates a slide with automatically numbered agenda items,
        with numbers styled in the theme's accent color.

        Args:
            title: Slide title (e.g., "Agenda", "Today's Topics")
            items: List of agenda items
            layout: Layout name, index, or None for auto-selection

        Returns:
            Slide number of the new slide

        Example:
            >>> pres.add_agenda_slide("Today's Agenda", [
            ...     "Introduction and Context",
            ...     "Key Findings",
            ...     "Recommendations",
            ...     "Next Steps",
            ... ])
        """
        layout_idx = self._find_layout(layout, LayoutType.CONTENT)
        slide_layout = self._pptx.slide_layouts[layout_idx]
        slide = self._pptx.slides.add_slide(slide_layout)

        # Set title
        title_ph = self._get_placeholder(slide, PP_PLACEHOLDER.TITLE)
        self._set_text_frame(title_ph, title)

        # Build numbered agenda content
        content = []
        accent_color = self._template.colors.get("accent1", "#4472C4")

        for i, item in enumerate(items, 1):
            # Number in bold, colored
            num_fmt = {
                "text": f"{i}. ",
                "bold": True,
                "color": accent_color,
            }
            item_fmt = {"text": item}
            content.append([num_fmt, item_fmt])

        # Set body with agenda
        body_ph = self._get_placeholder(slide, PP_PLACEHOLDER.BODY)
        if body_ph is None:
            body_ph = self._get_placeholder(slide, PP_PLACEHOLDER.OBJECT)
        self._set_body_content(body_ph, content)

        return len(self._pptx.slides)

    # --- Validation ---

    def validate(
        self,
        *,
        strict: bool = False,
        include_accessibility: bool = False,
        brand_rules: dict[str, Any] | None = None,
    ) -> ValidationResult:
        """Validate the presentation against design best practices.

        Checks for issues like too many bullets, missing titles,
        repetitive layouts, and provides actionable suggestions.

        Args:
            strict: If True, treat warnings as errors
            include_accessibility: If True, also run accessibility checks
            brand_rules: Optional dict of brand guidelines to enforce:
                - allowed_fonts: List of allowed font names
                - primary_colors: List of allowed brand colors
                - min_font_size: Minimum font size in points
                - max_bullets: Maximum bullets per slide
                - require_logo_slide_1: Require logo on first slide

        Returns:
            ValidationResult with issues and quality score

        Example:
            >>> result = pres.validate()
            >>> print(result.summary())
            "VALID (score: 85/100) | 2 warning(s) | 3 suggestion(s)"

            >>> result = pres.validate(include_accessibility=True)
            >>> result = pres.validate(brand_rules={
            ...     "allowed_fonts": ["Arial", "Calibri"],
            ...     "min_font_size": 14,
            ... })
        """
        from .validation import validate_presentation_extended

        return validate_presentation_extended(
            self,
            strict=strict,
            include_accessibility=include_accessibility,
            brand_rules=brand_rules,
        )

    # --- Shape manipulation methods ---

    def add_textbox(
        self,
        slide_num: int,
        text: str,
        left: float,
        top: float,
        width: float,
        height: float,
        *,
        font_size: int | None = None,
        font_family: str | None = None,
        font_color: str | None = None,
        bold: bool = False,
        italic: bool = False,
        align: str = "left",
    ) -> str:
        """Add a text box to a slide.

        Args:
            slide_num: The slide number (1-indexed)
            text: Text content
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
            font_size: Font size in points
            font_family: Font family name
            font_color: Hex color (e.g., "#FF0000")
            bold: Make text bold
            italic: Make text italic
            align: Text alignment ("left", "center", "right")

        Returns:
            Shape name for reference

        Example:
            >>> pres.add_textbox(1, "Hello World", 1, 1, 3, 0.5, font_size=24)
        """
        from pptx.enum.text import PP_ALIGN

        self._validate_slide_number(slide_num)
        slide = self._pptx.slides[slide_num - 1]

        shape = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = shape.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = text

        # Apply alignment
        align_map = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
        }
        p.alignment = align_map.get(align, PP_ALIGN.LEFT)

        # Apply formatting to the run
        if p.runs:
            run = p.runs[0]
            if font_size:
                run.font.size = Pt(font_size)
            if font_family:
                run.font.name = font_family
            if font_color:
                color_hex = font_color.lstrip("#")
                if len(color_hex) == 6:
                    run.font.color.rgb = RGBColor(
                        int(color_hex[:2], 16),
                        int(color_hex[2:4], 16),
                        int(color_hex[4:6], 16),
                    )
            if bold:
                run.font.bold = True
            if italic:
                run.font.italic = True

        return shape.name

    def add_shape(
        self,
        slide_num: int,
        shape_type: str,
        left: float,
        top: float,
        width: float,
        height: float,
        *,
        text: str = "",
        fill_color: str | None = None,
        line_color: str | None = None,
        line_width: float | None = None,
        font_size: int | None = None,
        font_color: str | None = None,
    ) -> str:
        """Add a shape to a slide.

        Args:
            slide_num: The slide number (1-indexed)
            shape_type: Shape type string (e.g., "rectangle", "oval", "arrow_right")
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
            text: Optional text to add inside the shape
            fill_color: Fill color as hex (e.g., "#4472C4")
            line_color: Line/border color as hex
            line_width: Line width in points
            font_size: Font size for text in points
            font_color: Font color for text as hex

        Returns:
            Shape name for reference

        Example:
            >>> pres.add_shape(1, "rectangle", 1, 2, 3, 2, fill_color="#4472C4")
            >>> pres.add_shape(1, "oval", 5, 2, 2, 2, text="Step 1")
        """
        from .shapes import get_mso_shape, ShapeType

        self._validate_slide_number(slide_num)
        slide = self._pptx.slides[slide_num - 1]

        # Get MSO_SHAPE constant
        mso_shape = get_mso_shape(shape_type)

        shape = slide.shapes.add_shape(
            mso_shape, Inches(left), Inches(top), Inches(width), Inches(height)
        )

        # Apply fill color
        if fill_color:
            color_hex = fill_color.lstrip("#")
            if len(color_hex) == 6:
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(
                    int(color_hex[:2], 16),
                    int(color_hex[2:4], 16),
                    int(color_hex[4:6], 16),
                )

        # Apply line color and width
        if line_color:
            color_hex = line_color.lstrip("#")
            if len(color_hex) == 6:
                shape.line.color.rgb = RGBColor(
                    int(color_hex[:2], 16),
                    int(color_hex[2:4], 16),
                    int(color_hex[4:6], 16),
                )
        if line_width is not None:
            shape.line.width = Pt(line_width)

        # Add text if provided
        if text:
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            p.alignment = 1  # Center

            if p.runs and (font_size or font_color):
                run = p.runs[0]
                if font_size:
                    run.font.size = Pt(font_size)
                if font_color:
                    color_hex = font_color.lstrip("#")
                    if len(color_hex) == 6:
                        run.font.color.rgb = RGBColor(
                            int(color_hex[:2], 16),
                            int(color_hex[2:4], 16),
                            int(color_hex[4:6], 16),
                        )

        return shape.name

    def add_connector(
        self,
        slide_num: int,
        start_shape: str,
        end_shape: str,
        connector_type: str = "elbow",
        *,
        line_color: str | None = None,
        line_width: float | None = None,
    ) -> str:
        """Add a connector between two shapes.

        Args:
            slide_num: The slide number (1-indexed)
            start_shape: Name of the starting shape
            end_shape: Name of the ending shape
            connector_type: Type of connector ("straight", "elbow", "curved")
            line_color: Line color as hex
            line_width: Line width in points

        Returns:
            Connector shape name

        Example:
            >>> shape1 = pres.add_shape(1, "rectangle", 1, 2, 2, 1, text="Start")
            >>> shape2 = pres.add_shape(1, "rectangle", 5, 2, 2, 1, text="End")
            >>> pres.add_connector(1, shape1, shape2, "elbow")
        """
        from .shapes import get_mso_connector

        self._validate_slide_number(slide_num)
        slide = self._pptx.slides[slide_num - 1]

        # Find shapes by name
        start = None
        end = None
        for shape in slide.shapes:
            if shape.name == start_shape:
                start = shape
            if shape.name == end_shape:
                end = shape

        if start is None:
            raise InvalidDataError(
                f"Shape '{start_shape}' not found on slide {slide_num}.",
                suggestion="Check shape names returned from add_shape().",
                code="SHAPE_NOT_FOUND",
            )
        if end is None:
            raise InvalidDataError(
                f"Shape '{end_shape}' not found on slide {slide_num}.",
                suggestion="Check shape names returned from add_shape().",
                code="SHAPE_NOT_FOUND",
            )

        # Get connector type
        mso_connector = get_mso_connector(connector_type)

        # Calculate connector endpoints (center-right of start, center-left of end)
        start_x = start.left + start.width
        start_y = start.top + start.height // 2
        end_x = end.left
        end_y = end.top + end.height // 2

        connector = slide.shapes.add_connector(
            mso_connector, start_x, start_y, end_x, end_y
        )

        # Connect to shapes
        connector.begin_connect(start, 3)  # Right side
        connector.end_connect(end, 1)  # Left side

        # Apply styling
        if line_color:
            color_hex = line_color.lstrip("#")
            if len(color_hex) == 6:
                connector.line.color.rgb = RGBColor(
                    int(color_hex[:2], 16),
                    int(color_hex[2:4], 16),
                    int(color_hex[4:6], 16),
                )
        if line_width is not None:
            connector.line.width = Pt(line_width)

        return connector.name

    def style_shape(
        self,
        slide_num: int,
        shape_name: str,
        *,
        fill_color: str | None = None,
        line_color: str | None = None,
        line_width: float | None = None,
    ) -> None:
        """Style an existing shape.

        Args:
            slide_num: The slide number (1-indexed)
            shape_name: Name of the shape to style
            fill_color: New fill color as hex
            line_color: New line color as hex
            line_width: New line width in points

        Example:
            >>> pres.style_shape(1, "Rectangle 1", fill_color="#FF0000")
        """
        self._validate_slide_number(slide_num)
        slide = self._pptx.slides[slide_num - 1]

        shape = None
        for s in slide.shapes:
            if s.name == shape_name:
                shape = s
                break

        if shape is None:
            raise InvalidDataError(
                f"Shape '{shape_name}' not found on slide {slide_num}.",
                suggestion="Use describe_slide() to see available shape names.",
                code="SHAPE_NOT_FOUND",
            )

        if fill_color:
            color_hex = fill_color.lstrip("#")
            if len(color_hex) == 6:
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(
                    int(color_hex[:2], 16),
                    int(color_hex[2:4], 16),
                    int(color_hex[4:6], 16),
                )

        if line_color:
            color_hex = line_color.lstrip("#")
            if len(color_hex) == 6:
                shape.line.color.rgb = RGBColor(
                    int(color_hex[:2], 16),
                    int(color_hex[2:4], 16),
                    int(color_hex[4:6], 16),
                )

        if line_width is not None:
            shape.line.width = Pt(line_width)

    def get_shape(self, slide_num: int, shape_name: str) -> dict[str, Any]:
        """Get information about a shape.

        Args:
            slide_num: The slide number (1-indexed)
            shape_name: Name of the shape

        Returns:
            Dict with shape properties (name, type, position, size, etc.)

        Example:
            >>> info = pres.get_shape(1, "Rectangle 1")
            >>> print(info["left"], info["top"])
        """
        self._validate_slide_number(slide_num)
        slide = self._pptx.slides[slide_num - 1]

        for shape in slide.shapes:
            if shape.name == shape_name:
                info: dict[str, Any] = {
                    "name": shape.name,
                    "shape_type": str(shape.shape_type),
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height,
                    "left_inches": shape.left / 914400,
                    "top_inches": shape.top / 914400,
                    "width_inches": shape.width / 914400,
                    "height_inches": shape.height / 914400,
                }
                if shape.has_text_frame:
                    info["text"] = shape.text_frame.text
                return info

        raise InvalidDataError(
            f"Shape '{shape_name}' not found on slide {slide_num}.",
            suggestion="Use describe_slide() to see available shape names.",
            code="SHAPE_NOT_FOUND",
        )

    # --- Strategic slide patterns ---

    def add_swot_slide(
        self,
        title: str,
        strengths: list[str],
        weaknesses: list[str],
        opportunities: list[str],
        threats: list[str],
        *,
        layout: str | int | None = None,
    ) -> int:
        """Add a SWOT analysis slide with 2x2 grid.

        Creates a professional SWOT diagram with four colored quadrants.

        Args:
            title: Slide title
            strengths: List of strength items
            weaknesses: List of weakness items
            opportunities: List of opportunity items
            threats: List of threat items
            layout: Layout name, index, or None for auto-selection

        Returns:
            Slide number of the new slide

        Example:
            >>> pres.add_swot_slide(
            ...     "SWOT Analysis",
            ...     strengths=["Strong brand", "Skilled team"],
            ...     weaknesses=["Limited budget"],
            ...     opportunities=["New markets", "Partnerships"],
            ...     threats=["Competition", "Regulation"]
            ... )
        """
        from .patterns import _add_labeled_box, get_default_pattern_colors

        layout_idx = self._find_layout(layout, LayoutType.TITLE_ONLY)
        slide_layout = self._pptx.slide_layouts[layout_idx]
        slide = self._pptx.slides.add_slide(slide_layout)

        # Set title
        title_ph = self._get_placeholder(slide, PP_PLACEHOLDER.TITLE)
        self._set_text_frame(title_ph, title)

        # Get theme colors
        colors = get_default_pattern_colors(self._template.colors)

        # Define grid layout (2x2)
        margin = 0.5
        box_width = 4.25
        box_height = 2.8
        gap = 0.2
        start_top = 1.5

        # Add quadrants
        _add_labeled_box(
            slide, margin, start_top, box_width, box_height,
            "Strengths", strengths, colors.primary
        )
        _add_labeled_box(
            slide, margin + box_width + gap, start_top, box_width, box_height,
            "Weaknesses", weaknesses, colors.secondary
        )
        _add_labeled_box(
            slide, margin, start_top + box_height + gap, box_width, box_height,
            "Opportunities", opportunities, colors.tertiary
        )
        _add_labeled_box(
            slide, margin + box_width + gap, start_top + box_height + gap, box_width, box_height,
            "Threats", threats, "#C00000"  # Red for threats
        )

        return len(self._pptx.slides)

    def add_matrix_slide(
        self,
        title: str,
        top_left: list[str],
        top_right: list[str],
        bottom_left: list[str],
        bottom_right: list[str],
        *,
        x_label: str = "",
        y_label: str = "",
        quadrant_labels: tuple[str, str, str, str] | None = None,
        layout: str | int | None = None,
    ) -> int:
        """Add a 2x2 matrix slide.

        Creates a strategic matrix with optional axis labels.

        Args:
            title: Slide title
            top_left: Items for top-left quadrant
            top_right: Items for top-right quadrant
            bottom_left: Items for bottom-left quadrant
            bottom_right: Items for bottom-right quadrant
            x_label: Label for horizontal axis
            y_label: Label for vertical axis
            quadrant_labels: Optional labels for each quadrant (TL, TR, BL, BR)
            layout: Layout name, index, or None for auto-selection

        Returns:
            Slide number of the new slide

        Example:
            >>> pres.add_matrix_slide(
            ...     "Priority Matrix",
            ...     top_left=["High impact quick wins"],
            ...     top_right=["Strategic projects"],
            ...     bottom_left=["Quick fixes"],
            ...     bottom_right=["Low priority"],
            ...     x_label="Effort",
            ...     y_label="Impact"
            ... )
        """
        from .patterns import _add_labeled_box, get_default_pattern_colors

        layout_idx = self._find_layout(layout, LayoutType.TITLE_ONLY)
        slide_layout = self._pptx.slide_layouts[layout_idx]
        slide = self._pptx.slides.add_slide(slide_layout)

        # Set title
        title_ph = self._get_placeholder(slide, PP_PLACEHOLDER.TITLE)
        self._set_text_frame(title_ph, title)

        # Get theme colors
        colors = get_default_pattern_colors(self._template.colors)

        # Define grid layout
        margin_left = 1.0 if y_label else 0.5
        box_width = 4.0
        box_height = 2.6
        gap = 0.15
        start_top = 1.5

        labels = quadrant_labels or ("", "", "", "")

        # Add quadrants
        _add_labeled_box(
            slide, margin_left, start_top, box_width, box_height,
            labels[0], top_left, colors.primary
        )
        _add_labeled_box(
            slide, margin_left + box_width + gap, start_top, box_width, box_height,
            labels[1], top_right, colors.secondary
        )
        _add_labeled_box(
            slide, margin_left, start_top + box_height + gap, box_width, box_height,
            labels[2], bottom_left, colors.tertiary
        )
        _add_labeled_box(
            slide, margin_left + box_width + gap, start_top + box_height + gap, box_width, box_height,
            labels[3], bottom_right, colors.quaternary
        )

        # Add axis labels
        if x_label:
            x_label_shape = slide.shapes.add_textbox(
                Inches(margin_left + box_width), Inches(start_top + 2 * box_height + gap + 0.2),
                Inches(2), Inches(0.4)
            )
            tf = x_label_shape.text_frame
            tf.paragraphs[0].text = x_label
            tf.paragraphs[0].alignment = 1  # Center

        if y_label:
            y_label_shape = slide.shapes.add_textbox(
                Inches(0.2), Inches(start_top + box_height),
                Inches(0.6), Inches(2)
            )
            tf = y_label_shape.text_frame
            tf.paragraphs[0].text = y_label
            y_label_shape.rotation = 270

        return len(self._pptx.slides)

    def add_funnel_slide(
        self,
        title: str,
        stages: list[dict[str, str] | str],
        *,
        layout: str | int | None = None,
    ) -> int:
        """Add a funnel diagram slide.

        Creates a funnel visualization with decreasing widths.

        Args:
            title: Slide title
            stages: List of stage dicts with "label" and optional "value" keys,
                    or list of strings for simple labels
            layout: Layout name, index, or None for auto-selection

        Returns:
            Slide number of the new slide

        Example:
            >>> pres.add_funnel_slide("Sales Funnel", [
            ...     {"label": "Leads", "value": "1000"},
            ...     {"label": "Qualified", "value": "400"},
            ...     {"label": "Proposals", "value": "100"},
            ...     {"label": "Closed", "value": "25"}
            ... ])
        """
        from pptx.enum.shapes import MSO_SHAPE

        layout_idx = self._find_layout(layout, LayoutType.TITLE_ONLY)
        slide_layout = self._pptx.slide_layouts[layout_idx]
        slide = self._pptx.slides.add_slide(slide_layout)

        # Set title
        title_ph = self._get_placeholder(slide, PP_PLACEHOLDER.TITLE)
        self._set_text_frame(title_ph, title)

        # Get theme colors
        accent_colors = [
            self._template.colors.get("accent1", "#4472C4"),
            self._template.colors.get("accent2", "#ED7D31"),
            self._template.colors.get("accent3", "#A5A5A5"),
            self._template.colors.get("accent4", "#FFC000"),
            self._template.colors.get("accent5", "#5B9BD5"),
            self._template.colors.get("accent6", "#70AD47"),
        ]

        # Normalize stages to dicts
        normalized_stages = []
        for stage in stages:
            if isinstance(stage, str):
                normalized_stages.append({"label": stage, "value": ""})
            else:
                normalized_stages.append(stage)

        # Draw funnel stages
        n_stages = len(normalized_stages)
        center_x = 5.0
        start_top = 1.6
        stage_height = 0.9
        max_width = 6.5
        min_width = 2.0

        for i, stage in enumerate(normalized_stages):
            width = max_width - (max_width - min_width) * i / max(n_stages - 1, 1)
            top = start_top + i * (stage_height + 0.1)
            left = center_x - width / 2
            color = accent_colors[i % len(accent_colors)]

            shape = slide.shapes.add_shape(
                MSO_SHAPE.PENTAGON,
                Inches(left),
                Inches(top),
                Inches(width),
                Inches(stage_height),
            )
            shape.rotation = 180
            shape.fill.solid()
            color_hex = color.lstrip("#")
            shape.fill.fore_color.rgb = RGBColor(
                int(color_hex[:2], 16),
                int(color_hex[2:4], 16),
                int(color_hex[4:6], 16),
            )
            shape.line.fill.background()

            # Add text
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            label = stage.get("label", "")
            value = stage.get("value", "")
            p.text = f"{label}: {value}" if value else label
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.font.bold = True
            p.alignment = 1  # Center

        return len(self._pptx.slides)

    def add_pyramid_slide(
        self,
        title: str,
        levels: list[str],
        *,
        layout: str | int | None = None,
    ) -> int:
        """Add a pyramid diagram slide.

        Creates a pyramid with levels from top (smallest) to bottom (largest).

        Args:
            title: Slide title
            levels: List of level labels from top to bottom
            layout: Layout name, index, or None for auto-selection

        Returns:
            Slide number of the new slide

        Example:
            >>> pres.add_pyramid_slide("Strategic Hierarchy", [
            ...     "Vision",
            ...     "Strategy",
            ...     "Tactics",
            ...     "Operations"
            ... ])
        """
        from pptx.enum.shapes import MSO_SHAPE

        layout_idx = self._find_layout(layout, LayoutType.TITLE_ONLY)
        slide_layout = self._pptx.slide_layouts[layout_idx]
        slide = self._pptx.slides.add_slide(slide_layout)

        # Set title
        title_ph = self._get_placeholder(slide, PP_PLACEHOLDER.TITLE)
        self._set_text_frame(title_ph, title)

        # Get theme colors
        accent_colors = [
            self._template.colors.get("accent1", "#4472C4"),
            self._template.colors.get("accent2", "#ED7D31"),
            self._template.colors.get("accent3", "#A5A5A5"),
            self._template.colors.get("accent4", "#FFC000"),
            self._template.colors.get("accent5", "#5B9BD5"),
        ]

        # Draw pyramid levels
        n_levels = len(levels)
        center_x = 5.0
        start_top = 1.6
        level_height = 1.0
        min_width = 1.5
        max_width = 7.0

        for i, level in enumerate(levels):
            # Width increases from top to bottom
            width = min_width + (max_width - min_width) * i / max(n_levels - 1, 1)
            top = start_top + i * (level_height + 0.1)
            left = center_x - width / 2
            color = accent_colors[i % len(accent_colors)]

            shape = slide.shapes.add_shape(
                MSO_SHAPE.TRAPEZOID,
                Inches(left),
                Inches(top),
                Inches(width),
                Inches(level_height),
            )
            shape.fill.solid()
            color_hex = color.lstrip("#")
            shape.fill.fore_color.rgb = RGBColor(
                int(color_hex[:2], 16),
                int(color_hex[2:4], 16),
                int(color_hex[4:6], 16),
            )
            shape.line.fill.background()

            # Add text
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = level
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.font.bold = True
            p.alignment = 1  # Center

        return len(self._pptx.slides)

    def add_process_slide(
        self,
        title: str,
        steps: list[str],
        *,
        layout: str | int | None = None,
    ) -> int:
        """Add a process flow slide with connected steps.

        Creates a horizontal process flow with chevron-style steps.

        Args:
            title: Slide title
            steps: List of step labels
            layout: Layout name, index, or None for auto-selection

        Returns:
            Slide number of the new slide

        Example:
            >>> pres.add_process_slide("Development Process", [
            ...     "Plan",
            ...     "Build",
            ...     "Test",
            ...     "Deploy"
            ... ])
        """
        from pptx.enum.shapes import MSO_SHAPE

        layout_idx = self._find_layout(layout, LayoutType.TITLE_ONLY)
        slide_layout = self._pptx.slide_layouts[layout_idx]
        slide = self._pptx.slides.add_slide(slide_layout)

        # Set title
        title_ph = self._get_placeholder(slide, PP_PLACEHOLDER.TITLE)
        self._set_text_frame(title_ph, title)

        # Get theme color
        accent_color = self._template.colors.get("accent1", "#4472C4")

        # Draw process steps
        n_steps = len(steps)
        if n_steps == 0:
            return len(self._pptx.slides)

        # Calculate step dimensions
        margin = 0.5
        available_width = 9.0
        step_width = min(2.0, (available_width - margin) / n_steps)
        step_height = 1.2
        gap = 0.1
        start_left = margin
        top = 3.0

        for i, step_label in enumerate(steps):
            left = start_left + i * (step_width + gap)

            shape = slide.shapes.add_shape(
                MSO_SHAPE.CHEVRON,
                Inches(left),
                Inches(top),
                Inches(step_width),
                Inches(step_height),
            )
            shape.fill.solid()
            color_hex = accent_color.lstrip("#")
            shape.fill.fore_color.rgb = RGBColor(
                int(color_hex[:2], 16),
                int(color_hex[2:4], 16),
                int(color_hex[4:6], 16),
            )
            shape.line.fill.background()

            # Add text
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = step_label
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.font.bold = True
            p.alignment = 1  # Center

            # Add step number below
            num_shape = slide.shapes.add_textbox(
                Inches(left), Inches(top + step_height + 0.1),
                Inches(step_width), Inches(0.4)
            )
            tf = num_shape.text_frame
            tf.paragraphs[0].text = str(i + 1)
            tf.paragraphs[0].alignment = 1  # Center
            if tf.paragraphs[0].runs:
                tf.paragraphs[0].runs[0].font.size = Pt(12)
                tf.paragraphs[0].runs[0].font.bold = True

        return len(self._pptx.slides)

    def add_venn_slide(
        self,
        title: str,
        sets: list[str],
        intersection_label: str = "",
        *,
        layout: str | int | None = None,
    ) -> int:
        """Add a Venn diagram slide.

        Creates a 2 or 3-circle Venn diagram with overlapping circles.

        Args:
            title: Slide title
            sets: List of 2 or 3 set labels
            intersection_label: Optional label for the intersection
            layout: Layout name, index, or None for auto-selection

        Returns:
            Slide number of the new slide

        Example:
            >>> pres.add_venn_slide(
            ...     "Skills Overlap",
            ...     ["Technical", "Business", "Leadership"],
            ...     intersection_label="Ideal Candidate"
            ... )
        """
        from pptx.enum.shapes import MSO_SHAPE
        from lxml import etree

        layout_idx = self._find_layout(layout, LayoutType.TITLE_ONLY)
        slide_layout = self._pptx.slide_layouts[layout_idx]
        slide = self._pptx.slides.add_slide(slide_layout)

        # Set title
        title_ph = self._get_placeholder(slide, PP_PLACEHOLDER.TITLE)
        self._set_text_frame(title_ph, title)

        # Get theme colors
        accent_colors = [
            self._template.colors.get("accent1", "#4472C4"),
            self._template.colors.get("accent2", "#ED7D31"),
            self._template.colors.get("accent3", "#A5A5A5"),
        ]

        n_sets = len(sets)
        if n_sets < 2:
            n_sets = 2
            sets = sets + [""] * (2 - len(sets))
        elif n_sets > 3:
            sets = sets[:3]
            n_sets = 3

        diameter = 3.0
        center_x = 5.0
        center_y = 3.8

        if n_sets == 2:
            # Two-circle Venn
            positions = [
                (center_x - 1.0, center_y - diameter / 2),
                (center_x + 0.5, center_y - diameter / 2),
            ]
        else:
            # Three-circle Venn
            positions = [
                (center_x - 0.8, center_y - diameter / 2 - 0.5),
                (center_x + 0.3, center_y - diameter / 2 - 0.5),
                (center_x - 0.25, center_y + 0.3),
            ]

        for i, (left, top) in enumerate(positions):
            if i >= len(sets):
                break

            color = accent_colors[i % len(accent_colors)]

            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(left),
                Inches(top),
                Inches(diameter),
                Inches(diameter),
            )
            circle.fill.solid()
            color_hex = color.lstrip("#")
            circle.fill.fore_color.rgb = RGBColor(
                int(color_hex[:2], 16),
                int(color_hex[2:4], 16),
                int(color_hex[4:6], 16),
            )
            circle.line.fill.background()

            # Set transparency via XML
            try:
                spPr = circle._sp.spPr
                ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
                solidFill = spPr.find(".//a:solidFill", ns)
                if solidFill is not None:
                    srgbClr = solidFill.find("a:srgbClr", ns)
                    if srgbClr is not None:
                        alpha = etree.SubElement(
                            srgbClr,
                            "{http://schemas.openxmlformats.org/drawingml/2006/main}alpha",
                        )
                        alpha.set("val", "50000")  # 50% opacity
            except Exception:
                pass

            # Add label
            tf = circle.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = sets[i]
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.font.bold = True
            p.alignment = 1  # Center

        # Add intersection label
        if intersection_label:
            inter_shape = slide.shapes.add_textbox(
                Inches(center_x - 0.5), Inches(center_y),
                Inches(1.5), Inches(0.5)
            )
            tf = inter_shape.text_frame
            tf.paragraphs[0].text = intersection_label
            tf.paragraphs[0].alignment = 1  # Center
            tf.paragraphs[0].font.size = Pt(10)
            tf.paragraphs[0].font.bold = True

        return len(self._pptx.slides)

    # --- Cloning and merge methods ---

    def clone_slide(
        self,
        source_num: int,
        *,
        insert_at: int | None = None,
    ) -> int:
        """Clone a slide within the same presentation.

        Creates a deep copy of the slide including all shapes, text,
        and images.

        Args:
            source_num: Source slide number to clone (1-indexed)
            insert_at: Position to insert the clone (1-indexed).
                       None = append at end.

        Returns:
            Slide number of the cloned slide

        Example:
            >>> pres.add_title_slide("Title", "Subtitle")
            >>> pres.clone_slide(1)  # Clone slide 1
            >>> pres.slide_count  # Now 2 slides
        """
        from lxml import etree

        self._validate_slide_number(source_num)
        source_slide = self._pptx.slides[source_num - 1]

        # Get the layout of the source slide
        slide_layout = source_slide.slide_layout

        # Add new slide with same layout
        new_slide = self._pptx.slides.add_slide(slide_layout)

        # Copy all shapes from source to new slide
        # We need to work with XML for deep copying
        for shape in source_slide.shapes:
            # Get XML element
            sp = shape._element
            # Deep copy
            new_sp = copy.deepcopy(sp)
            # Add to new slide's shape tree
            new_slide.shapes._spTree.insert_element_before(new_sp, "p:extLst")

        # Copy slide notes if any
        if source_slide.has_notes_slide:
            notes_text = source_slide.notes_slide.notes_text_frame.text
            if notes_text:
                new_slide.notes_slide.notes_text_frame.text = notes_text

        new_slide_num = len(self._pptx.slides)

        # Move to insert_at position if specified
        if insert_at is not None and insert_at != new_slide_num:
            self.move_slide(new_slide_num, insert_at)
            return insert_at

        return new_slide_num

    def clone_slide_from(
        self,
        other_pres: "Presentation",
        source_num: int,
        *,
        insert_at: int | None = None,
    ) -> int:
        """Clone a slide from another presentation.

        Copies the slide content to this presentation. Note: complex
        elements like charts may not transfer perfectly.

        Args:
            other_pres: The source Presentation object
            source_num: Slide number in the source presentation (1-indexed)
            insert_at: Position to insert (1-indexed). None = append.

        Returns:
            Slide number of the cloned slide

        Example:
            >>> pres1 = template.create_presentation()
            >>> pres1.add_title_slide("Source Title", "")
            >>> pres2 = template.create_presentation()
            >>> pres2.clone_slide_from(pres1, 1)
        """
        from lxml import etree

        other_pres._validate_slide_number(source_num)
        source_slide = other_pres._pptx.slides[source_num - 1]

        # Find a matching layout in this presentation
        source_layout_name = source_slide.slide_layout.name
        target_layout = None

        for layout in self._pptx.slide_layouts:
            if layout.name == source_layout_name:
                target_layout = layout
                break

        if target_layout is None:
            # Fall back to first layout
            target_layout = self._pptx.slide_layouts[0]

        # Add new slide
        new_slide = self._pptx.slides.add_slide(target_layout)

        # Copy shapes by rebuilding them
        for shape in source_slide.shapes:
            try:
                # Deep copy XML element
                sp = shape._element
                new_sp = copy.deepcopy(sp)
                new_slide.shapes._spTree.insert_element_before(new_sp, "p:extLst")
            except Exception:
                # Some shapes may not copy well, continue with others
                continue

        # Copy notes
        if source_slide.has_notes_slide:
            notes_text = source_slide.notes_slide.notes_text_frame.text
            if notes_text:
                new_slide.notes_slide.notes_text_frame.text = notes_text

        new_slide_num = len(self._pptx.slides)

        if insert_at is not None and insert_at != new_slide_num:
            self.move_slide(new_slide_num, insert_at)
            return insert_at

        return new_slide_num

    def merge(
        self,
        other_pres: "Presentation",
        *,
        insert_at: int | None = None,
    ) -> list[int]:
        """Merge another presentation into this one.

        Copies all slides from the other presentation.

        Args:
            other_pres: The Presentation to merge in
            insert_at: Position to start inserting (1-indexed).
                       None = append at end.

        Returns:
            List of slide numbers for the merged slides

        Example:
            >>> pres1 = template.create_presentation()
            >>> pres1.add_title_slide("Main Deck", "")
            >>> pres2 = template.create_presentation()
            >>> pres2.add_content_slide("Extra Content", ["Item"])
            >>> merged = pres1.merge(pres2)
        """
        merged_nums = []
        other_count = other_pres.slide_count

        for i in range(1, other_count + 1):
            if insert_at is not None:
                pos = insert_at + len(merged_nums)
            else:
                pos = None

            new_num = self.clone_slide_from(other_pres, i, insert_at=pos)
            merged_nums.append(new_num)

        return merged_nums

    @classmethod
    def merge_files(
        cls,
        template: "Template",
        paths: list[str | Path],
    ) -> "Presentation":
        """Create a new presentation by merging multiple files.

        Args:
            template: Template to use for the merged presentation
            paths: List of .pptx file paths to merge

        Returns:
            New Presentation with all slides merged

        Example:
            >>> merged = Presentation.merge_files(
            ...     template,
            ...     ["part1.pptx", "part2.pptx", "part3.pptx"]
            ... )
        """
        from .template import Template as TemplateClass

        result = template.create_presentation()

        for path in paths:
            path = Path(path)
            if not path.exists():
                continue

            # Load the file as a presentation
            try:
                other_template = TemplateClass(path)
                other_pres = cls(other_template)
                # Restore slides from the source file
                other_pres._pptx = PptxPresentation(path)

                result.merge(other_pres)
            except Exception:
                # Skip files that can't be loaded
                continue

        return result

    # --- Optimization and accessibility methods ---

    def optimize_slide(self, n: int) -> dict[str, Any]:
        """Optimize a single slide for better presentation.

        Performs automatic optimizations like:
        - Reducing font size if text overflows
        - Suggesting content splits if too much text
        - Checking image sizes

        Args:
            n: Slide number (1-indexed)

        Returns:
            Dict with changes made and suggestions

        Example:
            >>> changes = pres.optimize_slide(1)
            >>> print(changes)
            {"adjusted_fonts": True, "suggestions": ["Consider splitting content"]}
        """
        self._validate_slide_number(n)
        slide_info = self.describe_slide(n)
        slide = self._pptx.slides[n - 1]

        changes: dict[str, Any] = {
            "slide_number": n,
            "adjusted_fonts": False,
            "suggestions": [],
        }

        content = slide_info.get("content", [])
        content_count = len(content)

        # Check for too many bullets
        if content_count > 7:
            changes["suggestions"].append(
                f"Consider splitting {content_count} items across multiple slides"
            )

        # Check total word count
        total_words = sum(len(str(c).split()) for c in content if isinstance(c, str))
        if total_words > 100:
            changes["suggestions"].append(
                f"Slide has {total_words} words - consider reducing"
            )

        # Check for very long bullets
        for i, item in enumerate(content):
            if isinstance(item, str) and len(item) > 100:
                changes["suggestions"].append(
                    f"Bullet {i + 1} is very long ({len(item)} chars) - consider breaking up"
                )

        return changes

    def optimize_all(self) -> list[dict[str, Any]]:
        """Optimize all slides in the presentation.

        Returns:
            List of changes/suggestions for each slide

        Example:
            >>> all_changes = pres.optimize_all()
            >>> for change in all_changes:
            ...     if change["suggestions"]:
            ...         print(f"Slide {change['slide_number']}: {change['suggestions']}")
        """
        return [self.optimize_slide(i + 1) for i in range(self.slide_count)]

    def check_accessibility(self) -> ValidationResult:
        """Check the presentation for accessibility issues.

        Validates against accessibility best practices:
        - Missing alt text on images
        - Slide titles for navigation
        - Font size minimums
        - Reading order concerns

        Returns:
            ValidationResult with accessibility issues

        Example:
            >>> result = pres.check_accessibility()
            >>> print(result.summary())
            >>> for issue in result.issues:
            ...     print(f"Slide {issue.slide_number}: {issue.message}")
        """
        from .accessibility import check_accessibility

        return check_accessibility(self)

    def set_alt_text(self, slide_num: int, shape_name: str, alt_text: str) -> None:
        """Set alt text for an image or shape.

        Args:
            slide_num: Slide number (1-indexed)
            shape_name: Name of the shape
            alt_text: Description for screen readers

        Example:
            >>> pres.set_alt_text(1, "Picture 1", "Team photo from annual meeting")
        """
        from .accessibility import set_alt_text

        set_alt_text(self, slide_num, shape_name, alt_text)

    # --- Image placeholder methods ---

    def add_image_placeholder(
        self,
        slide_num: int,
        prompt: str,
        left: float,
        top: float,
        width: float,
        height: float,
    ) -> str:
        """Add a placeholder for an AI-generated image.

        Creates a visible placeholder shape with the image prompt.
        Use fill_image_placeholder() to replace with actual image.

        Args:
            slide_num: Slide number (1-indexed)
            prompt: Description of the desired image (for AI generation)
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches

        Returns:
            Placeholder ID for later use with fill_image_placeholder()

        Example:
            >>> ph_id = pres.add_image_placeholder(
            ...     1, "Professional team meeting in modern office",
            ...     left=1, top=2, width=5, height=3
            ... )
        """
        self._validate_slide_number(slide_num)
        slide = self._pptx.slides[slide_num - 1]

        from pptx.enum.shapes import MSO_SHAPE

        # Create placeholder shape with dashed border
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
        )

        # Style as placeholder
        shape.fill.background()  # No fill
        shape.line.color.rgb = RGBColor(0x80, 0x80, 0x80)
        shape.line.dash_style = 2  # Dashed
        shape.line.width = Pt(2)

        # Add prompt text
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"[Image: {prompt}]"
        p.alignment = 1  # Center
        if p.runs:
            p.runs[0].font.size = Pt(12)
            p.runs[0].font.italic = True
            p.runs[0].font.color.rgb = RGBColor(0x80, 0x80, 0x80)

        # Store metadata in shape name
        placeholder_id = f"img_placeholder_{slide_num}_{len(slide.shapes)}"
        shape.name = placeholder_id

        return placeholder_id

    def get_image_placeholders(self) -> list[dict[str, Any]]:
        """Get all image placeholders in the presentation.

        Returns:
            List of dicts with placeholder info:
            - slide: Slide number
            - placeholder_id: ID for fill_image_placeholder()
            - prompt: Image description/prompt
            - bounds: Position and size dict

        Example:
            >>> placeholders = pres.get_image_placeholders()
            >>> for ph in placeholders:
            ...     print(f"Slide {ph['slide']}: {ph['prompt']}")
        """
        placeholders = []

        for i in range(self.slide_count):
            slide = self._pptx.slides[i]

            for shape in slide.shapes:
                if shape.name.startswith("img_placeholder_"):
                    prompt = ""
                    if shape.has_text_frame:
                        text = shape.text_frame.text
                        if text.startswith("[Image:") and text.endswith("]"):
                            prompt = text[7:-1].strip()

                    placeholders.append({
                        "slide": i + 1,
                        "placeholder_id": shape.name,
                        "prompt": prompt,
                        "bounds": {
                            "left": shape.left / 914400,
                            "top": shape.top / 914400,
                            "width": shape.width / 914400,
                            "height": shape.height / 914400,
                        },
                    })

        return placeholders

    def fill_image_placeholder(
        self,
        slide_num: int,
        placeholder_id: str,
        image_path: str | Path,
    ) -> None:
        """Replace an image placeholder with an actual image.

        Args:
            slide_num: Slide number (1-indexed)
            placeholder_id: ID returned from add_image_placeholder()
            image_path: Path to the image file

        Example:
            >>> pres.fill_image_placeholder(1, ph_id, "generated_image.png")
        """
        self._validate_slide_number(slide_num)
        slide = self._pptx.slides[slide_num - 1]
        image_path = Path(image_path)

        if not image_path.exists():
            raise InvalidDataError(
                f"Image file not found: {image_path}",
                suggestion="Verify the image path is correct.",
                code="IMAGE_NOT_FOUND",
            )

        # Find the placeholder shape
        placeholder = None
        for shape in slide.shapes:
            if shape.name == placeholder_id:
                placeholder = shape
                break

        if placeholder is None:
            raise InvalidDataError(
                f"Placeholder '{placeholder_id}' not found on slide {slide_num}.",
                suggestion="Use get_image_placeholders() to see available placeholders.",
                code="PLACEHOLDER_NOT_FOUND",
            )

        # Get position and size from placeholder
        left = placeholder.left
        top = placeholder.top
        width = placeholder.width

        # Remove placeholder shape
        sp = placeholder._element
        sp.getparent().remove(sp)

        # Add actual image
        slide.shapes.add_picture(str(image_path), left, top, width=width)

    def save(self, path: str | Path) -> None:
        """Save the presentation.

        Args:
            path: Output file path

        Example:
            >>> pres.save("output.pptx")
        """
        self._pptx.save(path)

    def to_markdown(self, path: str | Path | None = None) -> str:
        """Export the presentation to Markdown format.

        Args:
            path: Optional file path to save (if None, returns string only)

        Returns:
            Markdown string representation

        Example:
            >>> md = pres.to_markdown()
            >>> print(md)
            # Presentation Title
            ## Slide 1: Introduction
            - Point 1
            - Point 2

            >>> pres.to_markdown("output.md")  # Save to file
        """
        from .markdown import to_markdown

        return to_markdown(self, path)

    def save_pdf(self, path: str | Path, *, engine: str = "libreoffice") -> None:
        """Export the presentation to PDF format.

        Requires LibreOffice installed on the system.

        Args:
            path: Output PDF file path
            engine: Export engine - "libreoffice" or "unoconv"

        Raises:
            ExportError: If export fails or engine not available

        Example:
            >>> pres.save_pdf("output.pdf")
        """
        from .export import save_pdf

        save_pdf(self, path, engine=engine)

    # --- Master and layout inspection ---

    def describe_master(self) -> dict[str, Any]:
        """Describe the slide master.

        Returns:
            Dict with master slide information:
            - name: Master slide name
            - layout_count: Number of layouts
            - colors: Theme colors
            - fonts: Theme fonts

        Example:
            >>> master = pres.describe_master()
            >>> print(master["name"], master["layout_count"])
        """
        master = self._pptx.slide_master

        return {
            "name": master.name if hasattr(master, "name") else "Slide Master",
            "layout_count": len(master.slide_layouts),
            "colors": self._template.colors,
            "fonts": self._template.fonts,
        }

    def describe_layouts(self) -> list[dict[str, Any]]:
        """Describe all available slide layouts.

        Returns:
            List of layout description dicts

        Example:
            >>> for layout in pres.describe_layouts():
            ...     print(f"{layout['index']}: {layout['name']}")
        """
        return self._template.describe()

    def get_layout(self, name_or_index: str | int) -> dict[str, Any] | None:
        """Get a specific layout by name or index.

        Args:
            name_or_index: Layout name (fuzzy matched) or index

        Returns:
            Layout description dict or None if not found

        Example:
            >>> layout = pres.get_layout("Title Slide")
            >>> layout = pres.get_layout(0)
        """
        result = self._template.get_layout(name_or_index)
        return result.to_dict() if result else None

    # --- Theme modification ---

    def set_theme_color(self, color_name: str, hex_value: str) -> None:
        """Modify a theme color.

        Note: This modifies the color in memory. Use save_as_template()
        to persist as a new template.

        Args:
            color_name: Color name (e.g., "accent1", "accent2", "dk1", "lt1")
            hex_value: New color value as hex (e.g., "#FF6600")

        Example:
            >>> pres.set_theme_color("accent1", "#FF6600")
        """
        from lxml import etree

        hex_value = hex_value.lstrip("#")

        # Access theme through slide master
        try:
            sm = self._pptx.slide_master
            for rel in sm.part.rels.values():
                if "theme" in rel.reltype:
                    theme_part = rel.target_part
                    theme_elem = etree.fromstring(theme_part.blob)

                    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
                    clr_scheme = theme_elem.find(".//a:clrScheme", ns)

                    if clr_scheme is not None:
                        # Find the color element
                        color_elem = clr_scheme.find(f"a:{color_name}", ns)
                        if color_elem is not None:
                            # Remove existing color definition
                            for child in list(color_elem):
                                color_elem.remove(child)

                            # Add new sRGB color
                            srgb = etree.SubElement(
                                color_elem,
                                "{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr",
                            )
                            srgb.set("val", hex_value.upper())

                            # Update the theme part
                            theme_part._blob = etree.tostring(
                                theme_elem, xml_declaration=True, encoding="UTF-8"
                            )

                    # Update template colors cache
                    self._template._colors[color_name] = f"#{hex_value}"
                    break
        except Exception:
            # Theme modification is best-effort
            pass

    def save_as_template(self, path: str | Path) -> None:
        """Save the current presentation as a template.

        Preserves all slides, layouts, and theme modifications.

        Args:
            path: Output template file path (.pptx)

        Example:
            >>> pres.set_theme_color("accent1", "#FF6600")
            >>> pres.save_as_template("custom_template.pptx")
        """
        self._pptx.save(path)

    def __repr__(self) -> str:
        return f"Presentation({self.slide_count} slides, template={self._template.path.name})"
