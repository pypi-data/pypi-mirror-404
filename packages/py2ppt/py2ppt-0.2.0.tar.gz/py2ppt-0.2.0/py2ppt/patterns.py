"""Strategic slide pattern helpers for py2ppt.

Provides pre-built slide patterns like SWOT, matrix, funnel,
pyramid, process flow, and Venn diagrams using shapes.
"""

from __future__ import annotations

from dataclasses import dataclass
from typing import TYPE_CHECKING, Any

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

if TYPE_CHECKING:
    from pptx.slide import Slide


@dataclass
class PatternColors:
    """Color scheme for pattern slides."""

    primary: str
    secondary: str
    tertiary: str
    quaternary: str
    text_dark: str = "#000000"
    text_light: str = "#FFFFFF"


def get_default_pattern_colors(template_colors: dict[str, str]) -> PatternColors:
    """Get pattern colors from template theme colors.

    Args:
        template_colors: Dict of template colors (accent1, accent2, etc.)

    Returns:
        PatternColors with theme-appropriate values
    """
    return PatternColors(
        primary=template_colors.get("accent1", "#4472C4"),
        secondary=template_colors.get("accent2", "#ED7D31"),
        tertiary=template_colors.get("accent3", "#A5A5A5"),
        quaternary=template_colors.get("accent4", "#FFC000"),
        text_dark=template_colors.get("dk1", "#000000"),
        text_light=template_colors.get("lt1", "#FFFFFF"),
    )


def _parse_color(color: str) -> RGBColor:
    """Parse hex color to RGBColor."""
    color_hex = color.lstrip("#")
    return RGBColor(
        int(color_hex[:2], 16),
        int(color_hex[2:4], 16),
        int(color_hex[4:6], 16),
    )


def _add_labeled_box(
    slide: "Slide",
    left: float,
    top: float,
    width: float,
    height: float,
    label: str,
    items: list[str],
    fill_color: str,
    text_color: str = "#FFFFFF",
    label_size: int = 16,
    item_size: int = 12,
) -> None:
    """Add a labeled box with bullet items to a slide.

    Used for SWOT quadrants and matrix cells.
    """
    from pptx.enum.shapes import MSO_SHAPE

    # Create the box
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = _parse_color(fill_color)
    box.line.fill.background()

    # Add label
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.1)

    p = tf.paragraphs[0]
    p.text = label
    p.font.bold = True
    p.font.size = Pt(label_size)
    p.font.color.rgb = _parse_color(text_color)
    p.alignment = PP_ALIGN.CENTER

    # Add items
    for item in items:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(item_size)
        p.font.color.rgb = _parse_color(text_color)
        p.alignment = PP_ALIGN.LEFT


def _add_funnel_stage(
    slide: "Slide",
    center_x: float,
    top: float,
    width: float,
    height: float,
    label: str,
    value: str | None,
    fill_color: str,
    text_color: str = "#FFFFFF",
) -> None:
    """Add a funnel stage (trapezoid shape)."""
    from pptx.enum.shapes import MSO_SHAPE

    left = center_x - width / 2

    stage = slide.shapes.add_shape(
        MSO_SHAPE.TRAPEZOID,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    stage.fill.solid()
    stage.fill.fore_color.rgb = _parse_color(fill_color)
    stage.line.fill.background()

    # Rotate to make it a funnel shape (wider at top)
    stage.rotation = 180

    # Add text
    tf = stage.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]

    if value:
        p.text = f"{label}: {value}"
    else:
        p.text = label

    p.font.size = Pt(14)
    p.font.color.rgb = _parse_color(text_color)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER


def _add_pyramid_level(
    slide: "Slide",
    center_x: float,
    top: float,
    width: float,
    height: float,
    label: str,
    fill_color: str,
    text_color: str = "#FFFFFF",
) -> None:
    """Add a pyramid level (trapezoid)."""
    from pptx.enum.shapes import MSO_SHAPE

    left = center_x - width / 2

    level = slide.shapes.add_shape(
        MSO_SHAPE.TRAPEZOID,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    level.fill.solid()
    level.fill.fore_color.rgb = _parse_color(fill_color)
    level.line.fill.background()

    tf = level.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(14)
    p.font.color.rgb = _parse_color(text_color)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER


def _add_process_step(
    slide: "Slide",
    left: float,
    top: float,
    width: float,
    height: float,
    label: str,
    step_num: int,
    fill_color: str,
    text_color: str = "#FFFFFF",
    is_last: bool = False,
) -> str:
    """Add a process step with chevron shape.

    Returns:
        Shape name for connector reference
    """
    from pptx.enum.shapes import MSO_SHAPE

    step = slide.shapes.add_shape(
        MSO_SHAPE.CHEVRON if not is_last else MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    step.fill.solid()
    step.fill.fore_color.rgb = _parse_color(fill_color)
    step.line.fill.background()

    tf = step.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(12)
    p.font.color.rgb = _parse_color(text_color)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    return step.name


def _add_venn_circle(
    slide: "Slide",
    left: float,
    top: float,
    diameter: float,
    label: str,
    fill_color: str,
    text_color: str = "#FFFFFF",
    transparency: float = 0.5,
) -> None:
    """Add a Venn diagram circle with transparency."""
    from pptx.enum.shapes import MSO_SHAPE
    from lxml import etree

    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left),
        Inches(top),
        Inches(diameter),
        Inches(diameter),
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = _parse_color(fill_color)
    circle.line.fill.background()

    # Set transparency via XML (python-pptx doesn't expose this directly)
    spPr = circle._sp.spPr
    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    solidFill = spPr.find(".//a:solidFill", ns)
    if solidFill is not None:
        srgbClr = solidFill.find("a:srgbClr", ns)
        if srgbClr is not None:
            alpha = etree.SubElement(
                srgbClr,
                "{http://schemas.openxmlformats.org/drawingml/2006/main}alpha",
            )
            alpha.set("val", str(int((1 - transparency) * 100000)))

    tf = circle.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(14)
    p.font.color.rgb = _parse_color(text_color)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
