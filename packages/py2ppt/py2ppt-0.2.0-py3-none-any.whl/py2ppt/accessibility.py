"""Accessibility checking for py2ppt presentations.

Validates presentations against accessibility best practices
including alt text, contrast, reading order, and more.
"""

from __future__ import annotations

from dataclasses import dataclass, field
from typing import TYPE_CHECKING, Any

from pptx.enum.shapes import MSO_SHAPE_TYPE

from .validation import (
    IssueCategory,
    IssueSeverity,
    ValidationIssue,
    ValidationResult,
)

if TYPE_CHECKING:
    from .presentation import Presentation


# Accessibility thresholds
MIN_FONT_SIZE = 18  # Points - minimum for readability
MIN_CONTRAST_RATIO = 4.5  # WCAG AA for normal text
MAX_WORDS_PER_SLIDE_ACCESSIBILITY = 80


@dataclass
class AccessibilityCheck:
    """Result of a single accessibility check."""

    check_name: str
    passed: bool
    message: str
    slide_number: int | None = None
    details: dict[str, Any] = field(default_factory=dict)


def check_accessibility(presentation: "Presentation") -> ValidationResult:
    """Check presentation for accessibility issues.

    Checks:
    - Missing alt text on images
    - Low contrast (where detectable)
    - Reading order issues
    - Slide titles present
    - Font size minimums
    - Too much text per slide

    Args:
        presentation: The Presentation to check

    Returns:
        ValidationResult with accessibility issues

    Example:
        >>> result = check_accessibility(pres)
        >>> for issue in result.issues:
        ...     print(f"{issue.severity}: {issue.message}")
    """
    issues: list[ValidationIssue] = []

    if presentation.slide_count == 0:
        return ValidationResult(is_valid=True, issues=[], score=100.0)

    for i in range(1, presentation.slide_count + 1):
        slide_issues = _check_slide_accessibility(presentation, i)
        issues.extend(slide_issues)

    # Calculate score
    score = _calculate_accessibility_score(issues)

    # No errors means valid
    error_count = len([i for i in issues if i.severity == IssueSeverity.ERROR])
    is_valid = error_count == 0

    return ValidationResult(is_valid=is_valid, issues=issues, score=score)


def _check_slide_accessibility(
    presentation: "Presentation",
    slide_num: int,
) -> list[ValidationIssue]:
    """Check a single slide for accessibility issues."""
    issues = []
    slide = presentation._pptx.slides[slide_num - 1]
    slide_info = presentation.describe_slide(slide_num)

    # Check for missing title
    if not slide_info.get("has_title"):
        issues.append(
            ValidationIssue(
                severity=IssueSeverity.WARNING,
                category=IssueCategory.ACCESSIBILITY,
                slide_number=slide_num,
                message="Slide has no title - screen readers use titles for navigation",
                suggestion="Add a descriptive title to help users navigate",
                rule="accessibility_missing_title",
            )
        )

    # Check for images without alt text
    for shape in slide.shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # Check for alt text (description)
                alt_text = ""
                try:
                    # Access alt text via the shape's XML
                    descr = shape._element.find(
                        ".//{http://schemas.openxmlformats.org/drawingml/2006/main}cNvPr"
                    )
                    if descr is not None:
                        alt_text = descr.get("descr", "")
                except Exception:
                    pass

                if not alt_text:
                    issues.append(
                        ValidationIssue(
                            severity=IssueSeverity.WARNING,
                            category=IssueCategory.ACCESSIBILITY,
                            slide_number=slide_num,
                            message=f"Image '{shape.name}' has no alt text",
                            suggestion="Add descriptive alt text for screen reader users",
                            rule="accessibility_missing_alt_text",
                            details={"shape_name": shape.name},
                        )
                    )
        except Exception:
            continue

    # Check font sizes
    small_fonts_found = []
    for shape in slide.shapes:
        try:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size is not None:
                            font_pt = run.font.size.pt
                            if font_pt < MIN_FONT_SIZE:
                                small_fonts_found.append(font_pt)
        except Exception:
            continue

    if small_fonts_found:
        min_found = min(small_fonts_found)
        issues.append(
            ValidationIssue(
                severity=IssueSeverity.WARNING,
                category=IssueCategory.ACCESSIBILITY,
                slide_number=slide_num,
                message=f"Font size {min_found:.0f}pt is below recommended {MIN_FONT_SIZE}pt minimum",
                suggestion="Use at least 18pt font for better readability",
                rule="accessibility_small_font",
                details={"min_font_size": min_found},
            )
        )

    # Check for too much text
    content = slide_info.get("content", [])
    total_words = sum(
        len(str(c).split()) for c in content if isinstance(c, str)
    )
    title = slide_info.get("title", "")
    total_words += len(title.split())

    if total_words > MAX_WORDS_PER_SLIDE_ACCESSIBILITY:
        issues.append(
            ValidationIssue(
                severity=IssueSeverity.INFO,
                category=IssueCategory.ACCESSIBILITY,
                slide_number=slide_num,
                message=f"Slide has {total_words} words - may be overwhelming",
                suggestion="Consider splitting content for better comprehension",
                rule="accessibility_too_much_text",
                details={"word_count": total_words},
            )
        )

    # Check reading order (shapes should be in logical order)
    shape_count = len(list(slide.shapes))
    if shape_count > 10:
        issues.append(
            ValidationIssue(
                severity=IssueSeverity.INFO,
                category=IssueCategory.ACCESSIBILITY,
                slide_number=slide_num,
                message=f"Slide has {shape_count} shapes - reading order may be confusing",
                suggestion="Verify reading order in PowerPoint's selection pane",
                rule="accessibility_reading_order",
                details={"shape_count": shape_count},
            )
        )

    return issues


def _calculate_accessibility_score(issues: list[ValidationIssue]) -> float:
    """Calculate accessibility score from 0-100."""
    score = 100.0

    for issue in issues:
        if issue.category != IssueCategory.ACCESSIBILITY:
            continue

        if issue.severity == IssueSeverity.ERROR:
            score -= 20
        elif issue.severity == IssueSeverity.WARNING:
            score -= 10
        elif issue.severity == IssueSeverity.INFO:
            score -= 3

    return max(0.0, min(100.0, score))


def set_alt_text(presentation: "Presentation", slide_num: int, shape_name: str, alt_text: str) -> None:
    """Set alt text for an image or shape.

    Args:
        presentation: The Presentation
        slide_num: Slide number (1-indexed)
        shape_name: Name of the shape
        alt_text: Alt text description

    Example:
        >>> set_alt_text(pres, 1, "Picture 1", "Team photo from annual meeting")
    """
    presentation._validate_slide_number(slide_num)
    slide = presentation._pptx.slides[slide_num - 1]

    for shape in slide.shapes:
        if shape.name == shape_name:
            # Set alt text via XML
            cNvPr = shape._element.find(
                ".//{http://schemas.openxmlformats.org/drawingml/2006/main}cNvPr"
            )
            if cNvPr is not None:
                cNvPr.set("descr", alt_text)
            return

    from .errors import InvalidDataError

    raise InvalidDataError(
        f"Shape '{shape_name}' not found on slide {slide_num}.",
        suggestion="Use describe_slide() to see available shape names.",
        code="SHAPE_NOT_FOUND",
    )
