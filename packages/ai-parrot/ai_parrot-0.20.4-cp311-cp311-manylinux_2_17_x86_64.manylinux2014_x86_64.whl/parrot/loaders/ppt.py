from typing import List, Union, Optional, Literal
from pathlib import PurePath
from collections.abc import Callable
import re
from ..stores.models import Document
from .abstract import AbstractLoader

# Optional dependencies
try:
    from markitdown import MarkItDown
    MARKITDOWN_AVAILABLE = True
except ImportError:
    MARKITDOWN_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class PowerPointLoader(AbstractLoader):
    """
    Enhanced PowerPoint loader with multiple backends.

    Supports:
    1. MarkItDown backend for rich markdown extraction (primary)
    2. python-pptx backend for detailed control and fallback

    Features:
    - Slide-by-slide processing with proper markdown formatting
    - Automatic slide title detection
    - Bullet point preservation
    - Slide notes extraction
    - Image-only slide detection and filtering
    - Configurable output formats
    """

    extensions: List[str] = ['.pptx', '.ppt']

    def __init__(
        self,
        source: Optional[Union[str, PurePath, List[PurePath]]] = None,
        *,
        tokenizer: Union[str, Callable] = None,
        text_splitter: Union[str, Callable] = None,
        source_type: str = 'file',

        # Backend selection
        backend: str = "auto",  # "markitdown", "pptx", "auto"

        # Output format
        output_format: Literal["markdown", "plain"] = "markdown",

        # Processing options
        skip_image_only_slides: bool = True,
        skip_empty_slides: bool = True,
        extract_slide_notes: bool = True,
        preserve_slide_structure: bool = True,

        # Slide filtering
        min_slide_content_length: int = 10,

        # Content processing
        clean_whitespace: bool = True,
        merge_consecutive_headers: bool = True,

        **kwargs
    ):
        super().__init__(
            source,
            tokenizer=tokenizer,
            text_splitter=text_splitter,
            source_type=source_type,
            **kwargs
        )

        # Backend configuration
        self.backend = self._select_backend(backend)
        self.output_format = output_format

        # Processing options
        self.skip_image_only_slides = skip_image_only_slides
        self.skip_empty_slides = skip_empty_slides
        self.extract_slide_notes = extract_slide_notes
        self.preserve_slide_structure = preserve_slide_structure
        self.min_slide_content_length = min_slide_content_length

        # Content processing
        self.clean_whitespace = clean_whitespace
        self.merge_consecutive_headers = merge_consecutive_headers

        # Initialize backend
        self._setup_backend()

    def _select_backend(self, preferred: str) -> str:
        """Select the best available backend."""
        if preferred == "auto":
            if MARKITDOWN_AVAILABLE:
                return "markitdown"
            elif PPTX_AVAILABLE:
                return "pptx"
            else:
                raise ImportError(
                    "No PowerPoint processing backend available. Install 'markitdown' or 'python-pptx'"
                )
        elif preferred == "markitdown" and MARKITDOWN_AVAILABLE:
            return "markitdown"
        elif preferred == "pptx" and PPTX_AVAILABLE:
            return "pptx"
        else:
            self.logger.warning(
                f"Backend '{preferred}' not available, falling back"
            )
            return self._select_backend("auto")

    def _setup_backend(self):
        """Initialize the selected backend."""
        if self.backend == "markitdown":
            self.md_converter = MarkItDown()
            self.logger.info("Using MarkItDown backend for PowerPoint processing")
        else:
            self.logger.info("Using python-pptx backend for PowerPoint processing")

    def _clean_content(self, content: str) -> str:
        """Clean and normalize content."""
        if not content:
            return ""

        if self.clean_whitespace:
            # Normalize whitespace while preserving markdown structure
            lines = content.split('\n')
            cleaned_lines = []
            for line in lines:
                cleaned_line = ' '.join(line.split())
                cleaned_lines.append(cleaned_line)
            content = '\n'.join(cleaned_lines)

        return content.strip()

    def _extract_slides_from_markdown(self, markdown_content: str) -> List[dict]:
        """Extract individual slides from MarkItDown markdown output."""
        slides = []

        # Split by slide separators (MarkItDown typically uses headers or page breaks)
        # Try multiple patterns for slide separation
        slide_patterns = [
            r'\n(?=#{1,2}\s)',  # Level 1-2 headers (typical slide titles)
            r'\n---+\n',       # Horizontal rules
            r'\n\*{3,}\n',     # Multiple asterisks
            r'(?:\n\s*){3,}',  # Multiple blank lines
        ]

        slide_sections = [markdown_content]  # Start with full content

        for pattern in slide_patterns:
            new_sections = []
            for section in slide_sections:
                parts = re.split(pattern, section)
                new_sections.extend([part.strip() for part in parts if part.strip()])
            if len(new_sections) > len(slide_sections):
                slide_sections = new_sections
                break

        # Process each section as a potential slide
        for i, section in enumerate(slide_sections):
            if len(section) < self.min_slide_content_length:
                continue

            # Extract title (first header if present)
            title_match = re.match(r'^(#{1,3})\s*(.+)$', section, re.MULTILINE)
            title = title_match.group(2) if title_match else f"Slide {i+1}"

            # Extract content (everything after title or full content if no title)
            if title_match:
                content_start = section.find('\n', title_match.end())
                content = section[content_start:].strip() if content_start != -1 else ""
            else:
                content = section.strip()

            slides.append({
                "slide_number": i + 1,
                "title": title,
                "content": content,
                "full_content": section,
                "has_title": bool(title_match)
            })

        return slides

    def _process_markitdown_content(self, path: Union[str, PurePath]) -> List[dict]:
        """Process PowerPoint using MarkItDown backend."""
        try:
            result = self.md_converter.convert(str(path))
            if not result or not result.text_content:
                self.logger.warning("MarkItDown returned empty content")
                return []

            markdown_content = result.text_content
            slides = self._extract_slides_from_markdown(markdown_content)

            self.logger.info(f"MarkItDown extracted {len(slides)} slides")
            return slides

        except Exception as e:
            self.logger.error(f"MarkItDown processing failed: {e}")
            return []

    # Original python-pptx methods (preserved as fallback)
    def extract_slide_text(self, slide):
        """Extract all text from a slide as a single string."""
        text_chunks = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_chunks.append(shape.text.strip())
        return "\n\n".join(text_chunks).strip()

    def slide_has_text(self, slide) -> bool:
        """Determine if a slide contains any text."""
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                return True
        return False

    def slide_has_images_only(self, slide) -> bool:
        """Return True if slide has images and no text."""
        has_image = False
        for shape in slide.shapes:
            if shape.shape_type == 13:  # PICTURE shape type in python-pptx
                has_image = True
            if hasattr(shape, "text") and shape.text.strip():
                return False
        return has_image

    def _extract_slide_title(self, slide) -> str:
        """Extract slide title from python-pptx slide object."""
        # Try to get title from title placeholder
        try:
            if slide.shapes.title and slide.shapes.title.text.strip():
                return slide.shapes.title.text.strip()
        except:
            pass

        # Look for first text shape that looks like a title
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                # Simple heuristic: short text, single line, likely a title
                if len(text) < 100 and '\n' not in text:
                    return text
                break

        return ""

    def _format_slide_as_markdown(self, slide_data: dict, slide_text: str, slide_notes: str = "") -> str:
        """Format slide content as markdown."""
        markdown_parts = []

        # Add title
        if slide_data.get("title"):
            markdown_parts.append(f"# {slide_data['title']}")
        elif not slide_data.get("has_title", False):
            markdown_parts.append(f"# Slide {slide_data['slide_number']}")

        # Add main content
        if slide_text:
            # Convert plain text to markdown if needed
            if self.output_format == "markdown" and not slide_data.get("full_content"):
                # Basic markdown conversion for bullet points
                content_lines = []
                for line in slide_text.split('\n'):
                    line = line.strip()
                    if line:
                        # Convert indented text to bullet points
                        if line.startswith('•') or line.startswith('-'):
                            content_lines.append(f"- {line[1:].strip()}")
                        elif line.startswith('  ') or line.startswith('\t'):
                            content_lines.append(f"- {line.strip()}")
                        else:
                            content_lines.append(line)
                markdown_parts.append('\n'.join(content_lines))
            else:
                markdown_parts.append(slide_text)

        # Add notes if present
        if slide_notes and self.extract_slide_notes:
            markdown_parts.append("## Notes")
            markdown_parts.append(slide_notes)

        return '\n\n'.join(markdown_parts)

    def _process_pptx_content(self, path: Union[str, PurePath]) -> List[dict]:
        """Process PowerPoint using python-pptx backend (original implementation enhanced)."""
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx not available for fallback processing")

        try:
            prs = Presentation(str(path))
            slides = []
            slide_count = len(prs.slides)

            for i, slide in enumerate(prs.slides):
                # Skip image-only slides if configured
                if self.skip_image_only_slides and self.slide_has_images_only(slide):
                    self.logger.debug(f"Slide {i+1}/{slide_count}: only images, skipping.")
                    continue

                # Extract slide text
                slide_text = self.extract_slide_text(slide)

                # Skip empty slides if configured
                if self.skip_empty_slides and (not slide_text or len(slide_text) < self.min_slide_content_length):
                    self.logger.debug(f"Slide {i+1}/{slide_count}: no sufficient text content, skipping.")
                    continue

                # Extract slide title
                slide_title = self._extract_slide_title(slide)

                # Extract slide notes
                slide_notes = ""
                if self.extract_slide_notes and slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    slide_notes = slide.notes_slide.notes_text_frame.text.strip()

                slides.append({
                    "slide_number": i + 1,
                    "slide_id": slide.slide_id,
                    "title": slide_title,
                    "content": slide_text,
                    "notes": slide_notes,
                    "has_title": bool(slide_title)
                })

            self.logger.info(f"python-pptx extracted {len(slides)} slides from {slide_count} total slides")
            return slides

        except Exception as e:
            self.logger.error(f"python-pptx processing failed: {e}")
            return []

    async def _load(self, path: Union[str, PurePath, List[PurePath]], **kwargs) -> List[Document]:
        """
        Load PowerPoint presentation with enhanced markdown support.

        Args:
            path: Path to the PowerPoint file

        Returns:
            List of Document objects, one per slide
        """
        self.logger.info(f"Loading PowerPoint file: {path}")
        docs = []

        # Try primary backend
        if self.backend == "markitdown":
            slides_data = self._process_markitdown_content(path)

            # Fallback to python-pptx if MarkItDown fails or returns no slides
            if not slides_data and PPTX_AVAILABLE:
                self.logger.info("MarkItDown failed or returned no slides, falling back to python-pptx")
                slides_data = self._process_pptx_content(path)
        else:
            slides_data = self._process_pptx_content(path)

        if not slides_data:
            self.logger.warning(f"No slides extracted from {path}")
            return docs

        # Create documents for each slide
        for slide_data in slides_data:
            # Format content based on output format and backend
            if self.backend == "markitdown" and self.output_format == "markdown":
                if slide_data.get("full_content"):
                    content = slide_data["full_content"]
                else:
                    content = self._format_slide_as_markdown(
                        slide_data,
                        slide_data.get("content", ""),
                        slide_data.get("notes", "")
                    )
            elif self.output_format == "markdown":
                content = self._format_slide_as_markdown(
                    slide_data,
                    slide_data.get("content", ""),
                    slide_data.get("notes", "")
                )
            else:
                # Plain text format
                parts = []
                if slide_data.get("title"):
                    parts.append(f"Title: {slide_data['title']}")
                if slide_data.get("content"):
                    parts.append(slide_data["content"])
                if slide_data.get("notes") and self.extract_slide_notes:
                    parts.append(f"Notes: {slide_data['notes']}")
                content = "\n\n".join(parts)

            content = self._clean_content(content)

            if not content or len(content) < self.min_slide_content_length:
                continue

            # Create metadata
            slide_meta = {
                "slide_number": slide_data["slide_number"],
                "slide_title": slide_data.get("title", ""),
                "has_notes": bool(slide_data.get("notes", "")),
                "content_length": len(content),
            }

            # Add backend-specific metadata
            if "slide_id" in slide_data:
                slide_meta["slide_id"] = slide_data["slide_id"]

            metadata = self.create_metadata(
                path=path,
                doctype="pptx",
                source_type="powerpoint",
                doc_metadata={
                    **slide_meta,
                    "extraction_backend": self.backend,
                    "output_format": self.output_format,
                },
            )

            # Create context header if preserve_slide_structure is True
            if self.preserve_slide_structure:
                context_parts = [
                    f"File Name: {path.name if hasattr(path, 'name') else str(path).split('/')[-1]}",
                    f"Slide Number: {slide_data['slide_number']}",
                    f"Document Type: pptx",
                    f"Source Type: powerpoint",
                ]

                if slide_data.get("slide_id"):
                    context_parts.append(f"Slide ID: {slide_data['slide_id']}")

                context_str = "\n".join(context_parts) + "\n======\n\n"
                full_content = context_str + content
            else:
                full_content = content

            doc = self.create_document(
                content=full_content,
                path=path,
                metadata=metadata
            )
            docs.append(doc)

        self.logger.info(f"Created {len(docs)} documents from PowerPoint slides")
        return docs

    def get_supported_backends(self) -> List[str]:
        """Get list of available backends."""
        backends = []

        if MARKITDOWN_AVAILABLE:
            backends.append("markitdown")
        if PPTX_AVAILABLE:
            backends.append("pptx")

        return backends

    def get_backend_info(self) -> dict:
        """Get information about current backend configuration."""
        return {
            "current_backend": self.backend,
            "available_backends": self.get_supported_backends(),
            "output_format": self.output_format,
            "settings": {
                "skip_image_only_slides": self.skip_image_only_slides,
                "skip_empty_slides": self.skip_empty_slides,
                "extract_slide_notes": self.extract_slide_notes,
                "preserve_slide_structure": self.preserve_slide_structure,
                "min_slide_content_length": self.min_slide_content_length,
            }
        }
