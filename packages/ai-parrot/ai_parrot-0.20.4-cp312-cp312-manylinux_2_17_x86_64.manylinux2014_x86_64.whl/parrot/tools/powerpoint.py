"""
PowerPoint Tool migrated to use AbstractDocumentTool framework.
"""
import re
from pathlib import Path
import tempfile
from typing import Dict, List, Optional, Any, Union
import io
import traceback
import asyncio
import aiohttp
from pptx import Presentation
from pptx.util import Pt, Inches, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
from jinja2 import Environment, FileSystemLoader
from pydantic import Field, field_validator
import markdown
from bs4 import BeautifulSoup, NavigableString
from navconfig import BASE_DIR
from .document import (
    AbstractDocumentTool,
    DocumentGenerationArgs
)

class PowerPointArgs(DocumentGenerationArgs):
    """Arguments schema for PowerPoint presentation generation."""

    template_name: Optional[str] = Field(
        None,
        description="Name of the HTML template (e.g., 'presentation.html') to render before conversion"
    )
    template_vars: Optional[Dict[str, Any]] = Field(
        None,
        description="Variables to pass to the HTML template (e.g., title, author, date)"
    )
    pptx_template: Optional[str] = Field(
        None,
        description="Filename of PowerPoint template file (.pptx or .potx) to use as base"
    )
    pptx_template_path: Optional[Path] = Field(
        None,
        description="Path where the PowerPoint template file is located"
    )
    slide_layout: int = Field(
        1,
        description="Default slide layout index (0=Title Slide, 1=Title and Content, etc.)",
        ge=0,
        le=15
    )
    title_styles: Optional[Dict[str, Any]] = Field(
        None,
        description="Styles to apply to slide titles (font_name, font_size, bold, italic, font_color, alignment)"
    )
    content_styles: Optional[Dict[str, Any]] = Field(
        None,
        description="Styles to apply to slide content (font_name, font_size, bold, italic, font_color, alignment)"
    )
    max_slides: int = Field(
        50,
        description="Maximum number of slides to generate",
        ge=1,
        le=100
    )
    split_by_headings: bool = Field(
        True,
        description="Whether to split content into slides based on headings (H1, H2, etc.)"
    )
    enable_images: bool = Field(
        True,
        description="Whether to download and include images from URLs in markdown"
    )
    image_width: Optional[float] = Field(
        None,
        description="Default image width in inches (None for auto-sizing)"
    )
    image_height: Optional[float] = Field(
        None,
        description="Default image height in inches (None for auto-sizing)"
    )
    max_image_size: float = Field(
        6.0,
        description="Maximum image size in inches (width or height)"
    )
    image_quality: str = Field(
        "high",
        description="Image quality: 'high', 'medium', 'low'"
    )
    image_timeout: int = Field(
        30,
        description="Timeout for image downloads in seconds"
    )

    @field_validator('template_name')
    @classmethod
    def validate_template_name(cls, v):
        if v and not v.endswith('.html'):
            v = f"{v}.html"
        return v

class PowerPointTool(AbstractDocumentTool):
    """
    PowerPoint Presentation Generator Tool.

    This tool converts text content (including Markdown and HTML) into professionally
    formatted PowerPoint presentations. It automatically splits content into slides
    based on headings and supports custom templates, styling, and layout options.

    Features:
    - Automatic slide splitting based on headings (H1, H2, H3, etc.)
    - Markdown to PowerPoint conversion with proper formatting
    - HTML to PowerPoint conversion support
    - Custom PowerPoint template support
    - Jinja2 HTML template processing
    - Configurable slide layouts and styling
    - Table, list, and content formatting
    - Professional presentation generation

    Slide Splitting Logic:
    - H1 (# Title) → Title slide (layout 0)
    - H2 (## Section) → Content slide (layout 1)
    - H3 (### Subsection) → Content slide (layout 1)
    - Content between headings → Added to the slide
    """

    name = "powerpoint_generator"
    description = (
        "Generate PowerPoint presentations from text, Markdown, or HTML content. "
        "Automatically splits content into slides based on headings. "
        "Supports custom templates, styling, and professional presentation formatting."
    )
    args_schema = PowerPointArgs

    # Document type configuration
    document_type = "presentation"
    default_extension = "pptx"
    supported_extensions = [".pptx", ".potx"]

    def __init__(
        self,
        templates_dir: Optional[Path] = None,
        output_dir: Optional[Union[str, Path]] = None,
        pptx_template_path: Optional[Path] = None,
        default_html_template: Optional[str] = None,
        **kwargs
    ):
        """
        Initialize the PowerPoint Tool.

        Args:
            templates_dir: Directory containing HTML and PowerPoint templates
            output_dir: Directory where generated presentations will be saved
            default_html_template: Default HTML template for content processing
            **kwargs: Additional arguments for AbstractDocumentTool
        """
        # Set up output directory before calling super().__init__
        if output_dir:
            kwargs['output_dir'] = Path(output_dir)

        super().__init__(templates_dir=templates_dir, **kwargs)

        self.default_html_template = default_html_template
        self.pptx_template_path = pptx_template_path or BASE_DIR.joinpath('presentations')

        # Initialize Jinja2 environment for HTML templates
        if self.templates_dir:
            self.html_env = Environment(
                loader=FileSystemLoader(str(self.templates_dir)),
                autoescape=True
            )
        else:
            self.html_env = None

    def _render_html_template(self, content: str, template_name: Optional[str], template_vars: Optional[Dict[str, Any]]) -> str:
        """Render content through Jinja2 HTML template if provided."""
        if not template_name or not self.html_env:
            return content

        try:
            template = self.html_env.get_template(template_name)
            vars_dict = template_vars or {}

            # Add default variables
            vars_dict.setdefault('content', content)
            vars_dict.setdefault('date', self._get_current_date())
            vars_dict.setdefault('timestamp', self._get_current_timestamp())

            rendered = template.render(**vars_dict)
            self.logger.info(f"Rendered content through HTML template: {template_name}")
            return rendered

        except Exception as e:
            self.logger.error(f"HTML template rendering failed: {e}")
            return content

    def _preprocess_markdown(self, text: str) -> str:
        """Preprocess markdown to handle common issues."""
        # Replace placeholder variables with empty strings
        text = re.sub(r'\{[a-zA-Z0-9_]+\}', '', text)

        # Handle f-strings that weren't evaluated
        text = re.sub(r'f"""(.*?)"""', r'\1', text, flags=re.DOTALL)
        text = re.sub(r"f'''(.*?)'''", r'\1', text, flags=re.DOTALL)

        # Remove triple backticks and language indicators
        text = re.sub(r'```[a-zA-Z]*\n', '', text)
        text = re.sub(r'```', '', text)

        # Fix heading issues (ensure space after #) - this should work correctly
        text = re.sub(r'(#+)([^ \n])', r'\1 \2', text)

        # Fix escaped newlines if any
        text = text.replace('\\n', '\n')

        # Clean up extra whitespace but preserve line structure
        lines = text.split('\n')
        cleaned_lines = []
        for line in lines:
            # Strip leading/trailing whitespace but preserve the line
            cleaned_line = line.strip()
            cleaned_lines.append(cleaned_line)

        return '\n'.join(cleaned_lines)

    def _markdown_to_html(self, markdown_text: str) -> str:
        """Convert markdown to HTML."""
        try:
            self.logger.debug(f"Converting markdown to HTML. Input preview: {markdown_text[:100]}...")

            html = markdown.markdown(
                markdown_text,
                extensions=['extra', 'codehilite', 'tables']  # Removed 'toc' to avoid issues
            )

            self.logger.debug(f"HTML conversion result preview: {html[:200]}...")
            return html

        except Exception as e:
            self.logger.error(f"Markdown conversion failed: {e}")
            # Fallback: wrap in paragraphs
            paragraphs = markdown_text.split('\n\n')
            html_paragraphs = [f'<p>{p.replace(chr(10), "<br>")}</p>' for p in paragraphs if p.strip()]
            fallback_html = '\n'.join(html_paragraphs)
            self.logger.debug(f"Using fallback HTML: {fallback_html[:200]}...")
            return fallback_html

    def _extract_slides_from_html(self, html_content: str, max_slides: int) -> List[Dict[str, Any]]:
        """Extract slides from HTML content based on headings."""
        soup = BeautifulSoup(html_content, 'html.parser')
        slides = []

        # Find all heading elements
        headings = soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6'])

        if not headings:
            # If no headings, create a single slide with all content
            slides.append({
                'title': 'Presentation',
                'content': self._extract_content_elements(soup),
                'level': 1,
                'layout': None
            })
            return slides

        # Get all elements in the document to process sequentially
        all_elements = soup.find_all(
            ['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'ul', 'ol', 'table', 'blockquote', 'div', 'img']
        )

        current_slide = None

        for element in all_elements:
            if len(slides) >= max_slides:
                self.logger.warning(
                    f"Reached maximum slides limit ({max_slides})"
                )
                break

            # If this is a heading, start a new slide
            if element.name in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
                if current_slide is not None and current_slide['content']:
                    slides.append(current_slide)

                heading_level = int(element.name[1])
                heading_text = self._get_text_content(element)
                heading_text = heading_text.strip()

                if not heading_text:
                    continue

                current_slide = {
                    'title': heading_text,
                    'level': heading_level,
                    'content': [],
                    # 'layout': 0 if (len(slides) == 0 and heading_level == 1) else None
                    'layout': None
                }
            # If this is content and we have a current slide, add it
            elif element.name in ['p', 'ul', 'ol', 'table', 'blockquote', 'div', 'img'] and current_slide is not None:
                if element.name == 'img':
                    # Always add images
                    current_slide['content'].append(element)
                    img_alt = element.get('alt', 'Image')
                    self.logger.debug(f"Added image to slide '{current_slide['title']}': {img_alt}")
                else:
                    # Add other content if not empty
                    content_text = self._get_text_content(element).strip()
                    if content_text:
                        current_slide['content'].append(element)

        # Don't forget the last slide
        if current_slide is not None and current_slide['content']:
            slides.append(current_slide)

        self.logger.info(
            f"Extracted {len(slides)} slides from HTML content"
        )
        return slides

    def _extract_content_elements(self, soup) -> List:
        """Extract content elements from soup."""
        content_elements = []
        # Get all content elements, including images
        for element in soup.find_all(['p', 'ul', 'ol', 'table', 'blockquote', 'div', 'img']):
            # Skip if this div only contains headings
            if element.name == 'div':
                child_tags = [child.name for child in element.find_all() if hasattr(child, 'name')]
                if all(tag in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6'] for tag in child_tags):
                    continue

            # Special handling for paragraphs that contain images
            if element.name == 'p':
                # Check if paragraph contains images
                images = element.find_all('img')
                if images:
                    # Add text content if any
                    text_content = self._get_text_content(element).strip()
                    if text_content:
                        content_elements.append(element)
                    # Add images separately
                    for img in images:
                        content_elements.append(img)
                else:
                    content_elements.append(element)
            else:
                content_elements.append(element)

        return content_elements

    def _create_presentation(self, template_path: Optional[str] = None) -> Presentation:
        """Create or load PowerPoint presentation."""
        if template_path:
            pptx_template = self._get_template_path(template_path)
            if pptx_template and pptx_template.exists():
                self.logger.info(f"Loading PowerPoint template: {pptx_template}")
                return Presentation(str(pptx_template))

        # Create new presentation
        return Presentation()

    async def _create_slides(
        self,
        prs: Presentation,
        slides_data: List[Dict[str, Any]],
        slide_layout: int,
        title_styles: Optional[Dict[str, Any]],
        content_styles: Optional[Dict[str, Any]],
        **kwargs: Any
    ) -> None:
        """Create slides from extracted data with dynamic layout detection (async for image downloads)."""

        # Debug available layouts
        self.logger.info(
            f"Available slide layouts: {len(prs.slide_layouts)}"
        )
        for i, layout in enumerate(prs.slide_layouts):
            self.logger.debug(
                f"Layout {i}: {layout.name}"
            )
        if len(prs.slides) > 0:
            # Remove the empty slide if it exists
            # slide_to_remove = prs.slides[0]
            # slide_id = slide_to_remove.slide_id
            prs.part.drop_rel(prs.slides._sldIdLst[0].rId)
            del prs.slides._sldIdLst[0]

        for i, slide_data in enumerate(slides_data):
            try:
                if len(prs.slide_layouts) == 1:
                    # Template has only one layout - use it for everything
                    layout_idx = 0
                elif len(prs.slide_layouts) >= 2:
                    # Template has multiple layouts - use layout 1 for content, 0 for title-only
                    if slide_data['content']:
                        layout_idx = 1  # Content slide
                    else:
                        layout_idx = 0  # Title-only slide
                    self.logger.debug(
                        f"Using layout {layout_idx} for slide with content: {bool(slide_data['content'])}"
                    )
                else:
                    # Fallback - shouldn't happen but just in case
                    layout_idx = 0
                # Additional safety check
                if layout_idx >= len(prs.slide_layouts):
                    layout_idx = 0
                    self.logger.warning(
                        f"Layout index {layout_idx} out of range, using layout 0"
                    )

                # Add slide
                slide_layout_obj = prs.slide_layouts[layout_idx]
                slide = prs.slides.add_slide(slide_layout_obj)

                # FIXED: Better title cleaning
                if slide.shapes.title and slide_data['title']:
                    # Remove markdown symbols more thoroughly
                    _title = slide_data['title']
                    _title = re.sub(r'^#+\s*', '', _title)  # Remove leading # symbols
                    _title = _title.lstrip('#').strip()    # Remove any remaining # symbols
                    _title = _title.strip()                # Final cleanup

                    slide.shapes.title.text = _title
                    self.logger.debug(f"Set slide title: '{_title}'")

                    if title_styles:
                        self._apply_text_styles(slide.shapes.title, title_styles)

                # Add content only if there is content
                if slide_data['content']:
                    await self._add_slide_content(
                        slide, slide_data['content'], content_styles, **kwargs
                    )

                self.logger.debug(
                    f"Created slide: '{_title}' using layout {layout_idx}"
                )

            except Exception as e:
                self.logger.error(
                    f"Error creating slide '{slide_data.get('title', 'Unknown')}': {e}"
                )
                # Continue to next slide instead of failing completely
                continue

    def _find_content_placeholder(self, slide):
        """Find the appropriate content placeholder on a slide."""
        try:
            # Debug: Log all available placeholders
            self.logger.debug(
                f"Available placeholders: {len(slide.shapes.placeholders)}"
            )
            for i, placeholder in enumerate(slide.shapes.placeholders):
                self.logger.debug(
                    f"Placeholder {i}: {placeholder.placeholder_format.type}"
                )

            # Try to find content placeholder
            # Common content placeholder types: BODY (2), OBJECT (7), CONTENT (13)
            content_placeholder_types = [2, 7, 13, 14, 15, 16, 17, 18]

            for placeholder in slide.shapes.placeholders:
                try:
                    if hasattr(placeholder.placeholder_format, 'type'):
                        if placeholder.placeholder_format.type in content_placeholder_types:
                            if hasattr(placeholder, 'text_frame'):
                                self.logger.debug(
                                    f"Found content placeholder: {placeholder.placeholder_format.type}"
                                )
                                return placeholder
                except Exception as e:
                    self.logger.error(
                        f"Error finding content placeholder: {e}"
                    )
                    continue

            # Fallback: Try to find any placeholder that's not the title (idx 0)
            if len(slide.shapes.placeholders) > 1:
                # Skip title placeholder (usually index 0) and try others
                for i in range(1, len(slide.shapes.placeholders)):
                    placeholder = slide.shapes.placeholders[i]
                    if hasattr(placeholder, 'text_frame'):
                        self.logger.debug(f"Using placeholder at index {i} as content placeholder")
                        return placeholder

            self.logger.warning(
                "No suitable content placeholder found"
            )
            return None

        except Exception as e:
            self.logger.error(f"Error finding content placeholder: {e}")
            return None

    def _add_text_as_textbox(self, slide, content_elements, content_styles):
        """Add text content as a text box when no content placeholder is available."""
        try:
            # FIXED: Better positioning that allows movement
            left = Inches(0.5)
            top = Inches(1.8)    # Start below title area
            width = Inches(9)    # Full width minus margins
            height = Inches(5)   # Reasonable height

            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame

            # FIXED: Better text frame settings
            text_frame.clear()
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            text_frame.margin_top = Inches(0.1)
            text_frame.margin_bottom = Inches(0.1)
            text_frame.word_wrap = True

            # Add content to text box with better formatting
            first_paragraph = True
            for element in content_elements:
                if element.name == 'p':
                    paragraph_text = self._get_text_content(element)

                    # Split by line breaks for better formatting
                    lines = [line.strip() for line in paragraph_text.split('\n') if line.strip()]

                    for line in lines:
                        if first_paragraph and not text_frame.paragraphs[0].text:
                            p = text_frame.paragraphs[0]
                            first_paragraph = False
                        else:
                            p = text_frame.add_paragraph()

                        p.text = line
                        p.space_after = Pt(6)  # Add some spacing between lines

                        if content_styles:
                            self._apply_paragraph_styles(p, content_styles)

                elif element.name in ['ul', 'ol']:
                    for li in element.find_all('li', recursive=False):
                        if first_paragraph and not text_frame.paragraphs[0].text:
                            p = text_frame.paragraphs[0]
                            first_paragraph = False
                        else:
                            p = text_frame.add_paragraph()

                        p.text = f"• {self._get_text_content(li)}"
                        p.space_after = Pt(3)

                        if content_styles:
                            self._apply_paragraph_styles(p, content_styles)

                elif element.name == 'table':
                    table_text = self._extract_table_text(element)
                    if first_paragraph and not text_frame.paragraphs[0].text:
                        p = text_frame.paragraphs[0]
                        first_paragraph = False
                    else:
                        p = text_frame.add_paragraph()

                    p.text = table_text
                    p.space_after = Pt(6)

                    if content_styles:
                        self._apply_paragraph_styles(p, content_styles)

                elif element.name == 'blockquote':
                    if first_paragraph and not text_frame.paragraphs[0].text:
                        p = text_frame.paragraphs[0]
                        first_paragraph = False
                    else:
                        p = text_frame.add_paragraph()

                    p.text = f'"{self._get_text_content(element)}"'
                    p.space_after = Pt(6)

                    if content_styles:
                        self._apply_paragraph_styles(p, content_styles)

            self.logger.debug("Added content as movable textbox")

        except Exception as e:
            self.logger.error(f"Error adding text as textbox: {e}")

    async def _add_slide_content(
        self,
        slide: Any,
        content_elements: List,
        content_styles: Optional[Dict[str, Any]],
        **kwargs
    ) -> None:
        """Add content to a slide placeholder (async for image downloads)."""

        enable_images = kwargs.get('enable_images', True)
        image_width = kwargs.get('image_width')
        image_height = kwargs.get('image_height')
        max_image_size = kwargs.get('max_image_size', 6.0)

        # Separate images from other content
        images = [elem for elem in content_elements if elem.name == 'img']
        other_content = [elem for elem in content_elements if elem.name != 'img']

        # Handle text content first
        if other_content:
            content_placeholder = self._find_content_placeholder(slide)
            if content_placeholder:
                try:
                    text_frame = content_placeholder.text_frame
                    text_frame.clear()

                    for element in other_content:
                        if element.name == 'p':
                            if len(text_frame.paragraphs) == 1 and not text_frame.paragraphs[0].text:
                                p = text_frame.paragraphs[0]
                            else:
                                p = text_frame.add_paragraph()
                            p.text = self._get_text_content(element)
                            if content_styles:
                                self._apply_paragraph_styles(p, content_styles)

                        elif element.name in ['ul', 'ol']:
                            for li in element.find_all('li', recursive=False):
                                p = text_frame.add_paragraph()
                                p.text = self._get_text_content(li)
                                p.level = 1
                                if content_styles:
                                    self._apply_paragraph_styles(p, content_styles)

                        elif element.name == 'table':
                            table_text = self._extract_table_text(element)
                            p = text_frame.add_paragraph()
                            p.text = table_text
                            if content_styles:
                                self._apply_paragraph_styles(p, content_styles)

                        elif element.name == 'blockquote':
                            p = text_frame.add_paragraph()
                            p.text = f'"{self._get_text_content(element)}"'
                            if content_styles:
                                self._apply_paragraph_styles(p, content_styles)

                except Exception as e:
                    self.logger.error(f"Error adding text content to placeholder: {e}")
                    # Fallback: Add text as a text box
                    self._add_text_as_textbox(slide, other_content, content_styles)
            else:
                # No content placeholder found, add text as a text box
                self.logger.info("No content placeholder found, adding text as textbox")
                self._add_text_as_textbox(slide, other_content, content_styles)

        # Handle images
        if images and enable_images:
            await self._add_images_to_slide(
                slide,
                images,
                image_width,
                image_height,
                max_image_size
            )

    async def _download_image(self, img_url: str, timeout: int = 30) -> Optional[bytes]:
        """
        Download image from URL using async aiohttp.
        
        This is non-blocking and works correctly when downloading from other
        agents running on the same event loop.
        """
        try:
            self.logger.debug(f"Downloading image: {img_url}")

            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
            }

            client_timeout = aiohttp.ClientTimeout(total=timeout)
            async with aiohttp.ClientSession(timeout=client_timeout) as session:
                async with session.get(img_url, headers=headers) as response:
                    response.raise_for_status()

                    # Check content type
                    content_type = response.headers.get('content-type', '').lower()
                    if not any(img_type in content_type for img_type in ['image/', 'jpeg', 'jpg', 'png', 'gif', 'bmp']):
                        self.logger.warning(f"URL does not appear to be an image: {content_type}")

                    # Read image data
                    image_data = await response.read()

                    if len(image_data) == 0:
                        self.logger.error(f"Downloaded image is empty: {img_url}")
                        return None

                    self.logger.debug(f"Successfully downloaded image: {len(image_data)} bytes")
                    return image_data

        except aiohttp.ClientError as e:
            self.logger.error(f"Failed to download image {img_url}: {e}")
            return None
        except asyncio.TimeoutError:
            self.logger.error(f"Timeout downloading image {img_url}")
            return None
        except Exception as e:
            self.logger.error(f"Unexpected error downloading image {img_url}: {e}")
            return None

    def _calculate_image_position(
        self,
        image_index: int,
        total_images: int,
        image_width: Optional[float],
        image_height: Optional[float],
        max_image_size: float
    ) -> tuple:
        """Calculate position and size for image on slide."""

        # Default dimensions
        slide_width = Inches(10)  # Standard slide width
        slide_height = Inches(7.5)  # Standard slide height

        # Calculate default size
        if image_width and image_height:
            width = Inches(image_width)
            height = Inches(image_height)
        else:
            # Auto-size with max constraints
            max_width = Inches(max_image_size)
            max_height = Inches(max_image_size * 0.75)  # Maintain aspect ratio

            # For multiple images, make them smaller
            if total_images > 1:
                max_width = Inches(max_image_size * 0.7)
                max_height = Inches(max_image_size * 0.5)

            width = max_width
            height = max_height

        # Calculate position
        if total_images == 1:
            # Center single image
            left = (slide_width - width) / 2
            top = Inches(2)  # Below title
        else:
            # Arrange multiple images
            images_per_row = 2 if total_images <= 4 else 3
            row = image_index // images_per_row
            col = image_index % images_per_row

            margin = Inches(0.5)
            available_width = slide_width - (2 * margin)
            available_height = slide_height - Inches(3)  # Account for title

            img_width = (available_width - (images_per_row - 1) * Inches(0.2)) / images_per_row
            img_height = min(height, available_height / ((total_images - 1) // images_per_row + 1))

            left = margin + col * (img_width + Inches(0.2))
            top = Inches(2) + row * (img_height + Inches(0.2))

            width = img_width
            height = img_height

        return left, top, width, height

    async def _add_images_to_slide(
        self,
        slide,
        images: List,
        image_width: Optional[float],
        image_height: Optional[float],
        max_image_size: float
    ) -> None:
        """Add images to a PowerPoint slide (async for non-blocking downloads)."""

        for i, img_element in enumerate(images):
            try:
                img_src = img_element.get('src')
                img_alt = img_element.get('alt', f'Image {i+1}')

                if not img_src:
                    self.logger.warning(f"Image has no src attribute: {img_alt}")
                    continue

                # Download image asynchronously
                image_data = await self._download_image(img_src)
                if not image_data:
                    self.logger.error(f"Failed to download image: {img_src}")
                    continue

                # Calculate position and size
                left, top, width, height = self._calculate_image_position(
                    i, len(images), image_width, image_height, max_image_size
                )

                # Add image to slide
                with tempfile.NamedTemporaryFile(delete=False, suffix='.jpg') as tmp_file:
                    tmp_file.write(image_data)
                    tmp_file_path = tmp_file.name

                try:
                    picture = slide.shapes.add_picture(tmp_file_path, left, top, width, height)
                    self.logger.debug(f"Added image to slide: {img_alt}")
                finally:
                    # Clean up temporary file
                    Path(tmp_file_path).unlink(missing_ok=True)

            except Exception as e:
                self.logger.error(f"Error adding image '{img_alt}': {e}")

    def _extract_table_text(self, table_element) -> str:
        """Extract text from table element."""
        rows = table_element.find_all('tr')
        table_lines = []

        for row in rows:
            cells = row.find_all(['td', 'th'])
            row_text = ' | '.join([self._get_text_content(cell) for cell in cells])
            table_lines.append(row_text)

        return '\n'.join(table_lines)

    def _get_text_content(self, element) -> str:
        """Extract clean text content from HTML element with preserved line breaks."""
        if isinstance(element, NavigableString):
            return str(element).strip()

        # For HTML elements, get the text content with line break preservation
        if hasattr(element, 'get_text'):
            # Use separator to preserve line breaks between elements
            text = element.get_text(separator='\n', strip=True)

            # Clean up excessive newlines but preserve intentional line breaks
            # Replace multiple consecutive newlines with single newlines
            text = re.sub(r'\n{3,}', '\n\n', text)

            # Handle specific markdown formatting that should have line breaks
            # Convert patterns like "**Key:** value" to have proper line breaks
            text = re.sub(r'(\*\*[^*]+\*\*[^*\n]*?)(\*\*[^*]+\*\*)', r'\1\n\2', text)

            return text

        # Fallback method for manual text extraction
        text_parts = []
        for content in element.contents:
            if isinstance(content, NavigableString):
                text_parts.append(str(content).strip())
            else:
                text_parts.append(self._get_text_content(content))

        result = '\n'.join([part for part in text_parts if part.strip()])

        # Additional cleanup: remove any remaining markdown symbols
        result = re.sub(r'^#+\s*', '', result)  # Remove leading hashtags
        result = re.sub(r'\*\*([^*]+)\*\*', r'\1', result)  # Remove bold markers
        result = re.sub(r'\*([^*]+)\*', r'\1', result)  # Remove italic markers

        return result

    def _apply_text_styles(self, shape, styles: Dict[str, Any]) -> None:
        """Apply styles to a text shape."""
        if not shape.has_text_frame:
            return

        try:
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                self._apply_paragraph_styles(paragraph, styles)
        except Exception as e:
            self.logger.error(f"Error applying text styles: {e}")

    def _apply_paragraph_styles(self, paragraph, styles: Dict[str, Any]) -> None:
        """Apply styles to a paragraph."""
        try:
            # Font styling
            if 'font_name' in styles:
                paragraph.font.name = styles['font_name']
            if 'font_size' in styles:
                paragraph.font.size = Pt(styles['font_size'])
            if 'bold' in styles:
                paragraph.font.bold = styles['bold']
            if 'italic' in styles:
                paragraph.font.italic = styles['italic']
            if 'font_color' in styles:
                # Convert hex color to RGB
                color_hex = styles['font_color'].lstrip('#')
                r, g, b = tuple(int(color_hex[i:i+2], 16) for i in (0, 2, 4))
                paragraph.font.color.rgb = RGBColor(r, g, b)

            # Alignment
            if 'alignment' in styles:
                alignment_map = {
                    'left': PP_ALIGN.LEFT,
                    'center': PP_ALIGN.CENTER,
                    'right': PP_ALIGN.RIGHT,
                    'justify': PP_ALIGN.JUSTIFY
                }
                paragraph.alignment = alignment_map.get(styles['alignment'], PP_ALIGN.LEFT)

        except Exception as e:
            self.logger.error(f"Error applying paragraph styles: {e}")

    def debug_content_parsing(self, content: str) -> Dict[str, Any]:
        """
        Debug method to see how content is being parsed.

        Args:
            content: Input content to debug

        Returns:
            Dictionary with debug information
        """
        try:
            # Process the content the same way as in generation
            processed_content = self._preprocess_markdown(content)
            html_content = self._markdown_to_html(processed_content)

            # Parse with BeautifulSoup
            soup = BeautifulSoup(html_content, 'html.parser')
            headings = soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6'])

            # Extract slide information
            slides_data = self._extract_slides_from_html(html_content, 50)

            debug_info = {
                "original_content_length": len(content),
                "original_content_preview": content[:300] + "..." if len(content) > 300 else content,
                "processed_content_preview": processed_content[:300] + "..." if len(processed_content) > 300 else processed_content,
                "html_content": html_content,  # Show full HTML to debug
                "headings_found": [
                    {
                        "tag": h.name,
                        "level": int(h.name[1]),
                        "raw_html": str(h),
                        "extracted_text": self._get_text_content(h),
                        "inner_text": h.get_text() if hasattr(h, 'get_text') else "N/A"
                    } for h in headings
                ],
                "slides_extracted": [
                    {
                        "title": slide['title'],
                        "level": slide['level'],
                        "layout": slide['layout'],
                        "content_count": len(slide['content']),
                        "content_preview": [
                            {
                                "tag": elem.name,
                                "text": self._get_text_content(elem)[:100]
                            } for elem in slide['content'][:3]
                        ]
                    } for slide in slides_data
                ],
                "total_slides": len(slides_data)
            }

            return debug_info

        except Exception as e:
            return {
                "error": str(e),
                "traceback": traceback.format_exc(),
                "message": "Debug parsing failed"
            }

    async def _generate_document_content(self, content: str, **kwargs) -> bytes:
        """
        Generate PowerPoint presentation content from input.

        Args:
            content: Input content (text, markdown, or HTML)
            **kwargs: Additional arguments from PowerPointArgs

        Returns:
            PowerPoint presentation as bytes
        """
        try:
            # Extract arguments
            template_name = kwargs.get('template_name')
            template_vars = kwargs.get('template_vars')
            pptx_template = kwargs.get('pptx_template')
            slide_layout = kwargs.pop('slide_layout', 1)
            title_styles = kwargs.pop('title_styles', None)
            content_styles = kwargs.pop('content_styles', None)
            max_slides = kwargs.pop('max_slides', 50)
            split_by_headings = kwargs.pop('split_by_headings', True)

            # Process content through HTML template if provided
            processed_content = self._render_html_template(content, template_name, template_vars)

            if pptx_template:
                pptx_template = self.pptx_template_path.joinpath(pptx_template)

            # Preprocess markdown
            cleaned_content = self._preprocess_markdown(processed_content)

            # Convert to HTML
            html_content = self._markdown_to_html(cleaned_content)

            # Extract slides from HTML
            if split_by_headings:
                slides_data = self._extract_slides_from_html(
                    html_content,
                    max_slides
                )
            else:
                # Create single slide with all content
                soup = BeautifulSoup(html_content, 'html.parser')
                slides_data = [{
                    'title': 'Presentation',
                    'content': self._extract_content_elements(soup),
                    'level': 1,
                    'layout': 0
                }]

            self.logger.info(
                f"Generated {len(slides_data)} slides from content"
            )

            # Create PowerPoint presentation
            prs = self._create_presentation(
                pptx_template
            )

            # Create slides
            create_slides_kwargs = {
                k: v for k, v in kwargs.items()
                if k not in [
                    'template_name', 'template_vars', 'pptx_template',
                    'slide_layout', 'title_styles', 'content_styles',
                    'max_slides', 'split_by_headings'
                ]
            }
            await self._create_slides(
                prs=prs,
                slides_data=slides_data,
                slide_layout=slide_layout,
                title_styles=title_styles,
                content_styles=content_styles,
                **create_slides_kwargs
            )

            # Save presentation to bytes
            ppt_bytes = io.BytesIO()
            prs.save(ppt_bytes)
            ppt_bytes.seek(0)

            return ppt_bytes.getvalue()

        except Exception as e:
            self.logger.error(f"Error generating PowerPoint presentation: {e}")
            raise

    async def _execute(
        self,
        content: str,
        output_filename: Optional[str] = None,
        file_prefix: str = "presentation",
        output_dir: Optional[str] = None,
        overwrite_existing: bool = False,
        template_name: Optional[str] = None,
        template_vars: Optional[Dict[str, Any]] = None,
        pptx_template: Optional[str] = None,
        slide_layout: int = 1,
        title_styles: Optional[Dict[str, Any]] = None,
        content_styles: Optional[Dict[str, Any]] = None,
        max_slides: int = 50,
        split_by_headings: bool = True,
        **kwargs
    ) -> Dict[str, Any]:
        """
        Execute PowerPoint presentation generation (AbstractTool interface).

        Args:
            content: Content to convert to PowerPoint presentation
            output_filename: Custom filename (without extension)
            file_prefix: Prefix for auto-generated filenames
            output_dir: Custom output directory
            overwrite_existing: Whether to overwrite existing files
            template_name: HTML template name for content processing
            template_vars: Variables for HTML template
            pptx_template: PowerPoint template file path
            slide_layout: Default slide layout index
            title_styles: Styles for slide titles
            content_styles: Styles for slide content
            max_slides: Maximum number of slides to generate
            split_by_headings: Whether to split by headings
            **kwargs: Additional arguments

        Returns:
            Dictionary with presentation generation results
        """
        try:
            self.logger.info(f"Starting PowerPoint generation with {len(content)} characters of content")

            # Use the safe document creation workflow
            result = await self._create_document_safely(
                content=content,
                output_filename=output_filename,
                file_prefix=file_prefix,
                output_dir=output_dir,
                overwrite_existing=overwrite_existing,
                extension="pptx",
                template_name=template_name,
                template_vars=template_vars,
                pptx_template=pptx_template,
                slide_layout=slide_layout,
                title_styles=title_styles,
                content_styles=content_styles,
                max_slides=max_slides,
                split_by_headings=split_by_headings
            )

            if result['status'] == 'success':
                # Add presentation-specific metadata
                result['presentation_info'] = {
                    'max_slides_limit': max_slides,
                    'split_by_headings': split_by_headings,
                    'slide_layout_used': slide_layout
                }

                self.logger.debug(
                    f"PowerPoint presentation created successfully: {result['metadata']['filename']}"
                )

            return result

        except Exception as e:
            self.logger.error(f"Error in PowerPoint generation: {e}")
            raise
