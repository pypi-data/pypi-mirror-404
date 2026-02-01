"""PowerPoint file reader using python-pptx (optional dependency)."""

from __future__ import annotations

from pathlib import Path

from middleware.filesystem.read.types import FileType, ReadLimits, ReadResult

try:
    from pptx import Presentation

    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


def read_pptx(
    path: Path,
    limits: ReadLimits,
    start_slide: int | None = None,
    limit_slides: int | None = None,
) -> ReadResult:
    """
    Read PowerPoint file with slide-based pagination.

    Args:
        path: Absolute path to PPTX file
        limits: ReadLimits configuration (max_chars applies to total output)
        start_slide: Start slide (1-indexed, None = start from slide 1)
        limit_slides: Number of slides to read (None = default 20 slides)

    Returns:
        ReadResult with extracted text and metadata
    """
    if not HAS_PPTX:
        return _no_pptx_result(path)

    stat = path.stat()
    result = ReadResult(
        file_path=str(path),
        file_type=FileType.DOCUMENT,
        total_size=stat.st_size,
    )

    try:
        prs = Presentation(path)
    except Exception as e:
        result.error = f"Error opening PPTX: {e}"
        return result

    slides = list(prs.slides)
    total_slides = len(slides)
    result.total_pages = total_slides

    start_idx = (start_slide - 1) if start_slide and start_slide > 0 else 0
    if start_idx >= total_slides:
        result.error = f"Start slide {start_slide} exceeds total slides {total_slides}"
        return result

    effective_limit = limit_slides if limit_slides else 20
    end_idx = min(start_idx + effective_limit, total_slides)

    output_parts: list[str] = []
    total_chars = 0
    truncated = False
    truncation_reason: str | None = None
    slides_read = 0

    for slide_num in range(start_idx, end_idx):
        slide = slides[slide_num]
        slide_text = _extract_slide_text(slide, slide_num + 1, total_slides)

        if total_chars + len(slide_text) > limits.max_chars:
            truncated = True
            truncation_reason = f"max_chars={limits.max_chars}"
            remaining = limits.max_chars - total_chars
            if remaining > 100:
                output_parts.append(slide_text[:remaining] + "\n... [truncated]")
                slides_read += 1
            break

        output_parts.append(slide_text)
        total_chars += len(slide_text)
        slides_read += 1

    if end_idx < total_slides and not truncated:
        truncated = True
        truncation_reason = f"limit_slides={effective_limit}"

    result.content = "\n".join(output_parts)
    result.start_page = start_idx + 1
    result.end_page = start_idx + slides_read
    result.truncated = truncated
    result.truncation_reason = truncation_reason

    return result


def _extract_slide_text(slide, slide_num: int, total_slides: int) -> str:
    """Extract text content from a slide."""
    parts = [f"\n{'='*60}", f"Slide {slide_num}/{total_slides}", "=" * 60]

    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            parts.append(shape.text.strip())

    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
        notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            parts.append(f"\n[Speaker Notes]: {notes}")

    return "\n".join(parts)


def _no_pptx_result(path: Path) -> ReadResult:
    """Return result when python-pptx is not installed."""
    stat = path.stat()
    content = (
        f"PowerPoint file: {path.name}\n"
        f"  Size: {stat.st_size:,} bytes\n"
        f"\n"
        f"python-pptx is not installed. To read PPTX files:\n"
        f"  uv pip install python-pptx"
    )
    return ReadResult(
        file_path=str(path),
        file_type=FileType.DOCUMENT,
        content=content,
        total_size=stat.st_size,
    )
