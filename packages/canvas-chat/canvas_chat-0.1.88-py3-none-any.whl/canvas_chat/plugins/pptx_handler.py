"""PowerPoint (PPTX) File Upload Handler Plugin

Handles PowerPoint file uploads and extraction on the backend.

Pipeline:
- Validate file size / extension
- Use LibreOffice headless to render slides to PNG
- Transcode PNG -> WebP (and generate thumbnails) for frontend display/storage
- Use python-pptx to extract per-slide titles and text content
"""

from __future__ import annotations

import base64
import logging
import os
import re
import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import Any

import pymupdf
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation

from canvas_chat.file_upload_handler_plugin import FileUploadHandlerPlugin
from canvas_chat.file_upload_registry import PRIORITY, FileUploadRegistry

logger = logging.getLogger(__name__)

# Maximum PPT/PPTX file size (50 MB)
MAX_PPTX_SIZE = 50 * 1024 * 1024

# LibreOffice conversion timeout (seconds)
LIBREOFFICE_TIMEOUT_S = 60

# Render width (pixels) for stored slide images / thumbnails
SLIDE_RENDER_WIDTH = 1024
THUMB_RENDER_WIDTH = 240


def _find_libreoffice_cmd() -> str | None:
    """Find an available LibreOffice headless command."""
    return shutil.which("soffice") or shutil.which("libreoffice")


def _safe_stem(filename: str) -> str:
    """Create a safe filesystem stem from a filename."""
    stem = Path(filename).stem or "presentation"
    # Keep alphanumerics, dash, underscore; replace others with underscore
    stem = re.sub(r"[^A-Za-z0-9_-]+", "_", stem).strip("_")
    return stem or "presentation"


def _read_b64(path: Path) -> str:
    return base64.b64encode(path.read_bytes()).decode("utf-8")


def _image_to_webp_b64(img: Image.Image, *, quality: int) -> str:
    """Encode PIL image to WebP base64."""
    import io

    buf = io.BytesIO()
    # Preserve alpha if present; otherwise convert to RGB for better compression.
    if img.mode not in ("RGB", "RGBA"):
        img = img.convert("RGBA")
    img.save(buf, format="WEBP", quality=quality, method=6)
    return base64.b64encode(buf.getvalue()).decode("utf-8")


def _resize_to_width(img: Image.Image, width: int) -> Image.Image:
    if img.width <= width:
        return img
    ratio = width / float(img.width)
    height = max(1, int(img.height * ratio))
    return img.resize((width, height), Image.LANCZOS)


def _extract_slide_text(slide) -> str:
    parts: list[str] = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = (shape.text or "").strip()
        if text:
            parts.append(text)
    return "\n".join(parts).strip()


def _extract_slide_title(slide) -> str | None:
    # 1) Prefer PowerPoint's title placeholder if present
    try:
        title_shape = slide.shapes.title
        if title_shape is not None:
            title_text = (title_shape.text or "").strip()
            if title_text:
                return title_text
    except Exception:
        # python-pptx can throw for certain templates; fall back below
        pass

    # 2) Heuristic: first non-empty line from any text shape
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = (shape.text or "").strip()
        if not text:
            continue
        first_line = text.splitlines()[0].strip()
        if first_line:
            return first_line[:120]
    return None


def _png_sort_key(p: Path, base_stem: str) -> tuple[int, str]:
    """Sort LibreOffice PNG outputs in slide order, robustly.

    LibreOffice output naming varies by platform/version. We've observed:
    - <stem>.png
    - <stem>1.png
    - <stem>_1.png
    - <stem>-1.png
    - <stem>-Slide1.png (rare)
    """
    stem = p.stem

    if stem == base_stem:
        return (-1, p.name)  # usually slide 1

    if stem.startswith(base_stem):
        rest = stem[len(base_stem) :]
        rest = rest.lstrip("_-")

        m = re.match(r"(\d+)$", rest)
        if m:
            return (int(m.group(1)), p.name)

        # Last resort: first digit sequence anywhere in the remainder
        m = re.search(r"(\d+)", rest)
        if m:
            return (int(m.group(1)), p.name)

    return (10_000, p.name)


def _convert_with_libreoffice(
    input_path: Path,
    *,
    out_dir: Path,
    timeout_s: int,
    convert_to: str,
) -> subprocess.CompletedProcess[str]:
    cmd = _find_libreoffice_cmd()
    if not cmd:
        raise ValueError("LibreOffice (soffice) not found.")

    out_dir.mkdir(parents=True, exist_ok=True)
    # LibreOffice sometimes requires HOME to be writable for profile/temp data.
    env = os.environ.copy()
    env.setdefault("HOME", str(out_dir))

    args = [
        cmd,
        "--headless",
        "--nologo",
        "--nolockcheck",
        "--nodefault",
        "--nofirststartwizard",
        "--convert-to",
        convert_to,
        "--outdir",
        str(out_dir),
        str(input_path),
    ]

    try:
        proc = subprocess.run(
            args,
            check=True,
            capture_output=True,
            text=True,
            timeout=timeout_s,
            env=env,
        )
        return proc
    except subprocess.TimeoutExpired as e:
        raise ValueError(
            f"PowerPoint conversion timed out after {timeout_s}s. "
            "Try a smaller deck or fewer embedded media assets."
        ) from e
    except subprocess.CalledProcessError as e:
        stderr = (e.stderr or "").strip()
        stdout = (e.stdout or "").strip()
        details = stderr or stdout or "Unknown LibreOffice error"
        raise ValueError(
            f"Failed to convert PowerPoint via LibreOffice: {details}"
        ) from e


def _render_pdf_to_pngs(
    pdf_path: Path, *, out_dir: Path, target_width: int
) -> list[Path]:
    out_dir.mkdir(parents=True, exist_ok=True)

    doc = pymupdf.open(pdf_path)
    png_paths: list[Path] = []
    try:
        for i in range(len(doc)):
            page = doc[i]
            rect = page.rect
            if rect.width <= 0:
                zoom = 1.0
            else:
                zoom = max(0.1, target_width / float(rect.width))

            mat = pymupdf.Matrix(zoom, zoom)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            out_path = out_dir / f"page_{i + 1:04d}.png"
            out_path.write_bytes(pix.tobytes("png"))
            png_paths.append(out_path)
    finally:
        doc.close()

    return png_paths


def _render_pptx_to_pngs(
    pptx_path: Path,
    *,
    out_dir: Path,
    timeout_s: int,
    expected_slide_count: int | None = None,
) -> tuple[list[Path], str]:
    cmd = _find_libreoffice_cmd()
    if not cmd:
        # Development/test-friendly fallback: generate placeholder PNGs so the UI
        # can still function even if LibreOffice isn't installed.
        logger.warning(
            "LibreOffice not found; generating placeholder slide images (filename=%s)",
            pptx_path.name,
        )
        pres = Presentation(str(pptx_path))
        base_stem = pptx_path.stem
        out_dir.mkdir(parents=True, exist_ok=True)

        def wrap_text(text: str, width: int = 60) -> str:
            words = (text or "").split()
            lines = []
            line = []
            count = 0
            for w in words:
                if count + len(w) + (1 if line else 0) > width:
                    lines.append(" ".join(line))
                    line = [w]
                    count = len(w)
                else:
                    line.append(w)
                    count += len(w) + (1 if line else 0)
            if line:
                lines.append(" ".join(line))
            return "\n".join(lines)

        try:
            font = ImageFont.load_default()
        except Exception:
            font = None

        png_paths: list[Path] = []
        for idx, slide in enumerate(pres.slides):
            title = _extract_slide_title(slide) or f"Slide {idx + 1}"
            body = _extract_slide_text(slide) or "(No text content)"

            img = Image.new("RGB", (1280, 720), color=(255, 255, 255))
            draw = ImageDraw.Draw(img)
            draw.text((40, 40), title, fill=(0, 0, 0), font=font)
            draw.text((40, 120), wrap_text(body), fill=(30, 30, 30), font=font)
            draw.text(
                (40, 680),
                "Placeholder rendering (install LibreOffice for accurate slides)",
                fill=(120, 120, 120),
                font=font,
            )

            out_path = out_dir / f"{base_stem}_{idx}.png"
            img.save(out_path, format="PNG")
            png_paths.append(out_path)

        if not png_paths:
            raise ValueError("No slides found in PowerPoint.")
        return png_paths, "placeholder"

    # Prefer explicit Impress export filter; some LO builds only render first slide
    # with plain "png" conversion.
    proc = _convert_with_libreoffice(
        pptx_path,
        out_dir=out_dir,
        timeout_s=timeout_s,
        convert_to="png:impress_png_Export",
    )

    logger.debug(
        "LibreOffice convert stdout=%s stderr=%s",
        (proc.stdout or "").strip(),
        (proc.stderr or "").strip(),
    )

    base_stem = pptx_path.stem
    pngs = sorted(out_dir.glob("*.png"), key=lambda p: _png_sort_key(p, base_stem))

    if not pngs:
        raise ValueError("LibreOffice conversion produced no PNG outputs.")

    # Workaround: if we got fewer images than slides (common on some platforms),
    # export to PDF and rasterize each page with PyMuPDF.
    if (
        expected_slide_count
        and expected_slide_count > 1
        and len(pngs) < expected_slide_count
    ):
        logger.warning(
            (
                "LibreOffice PNG export produced %s images for %s slides; "
                "falling back to PDF rasterization."
            ),
            len(pngs),
            expected_slide_count,
        )

        # Convert to PDF (single multipage PDF)
        _convert_with_libreoffice(
            pptx_path,
            out_dir=out_dir,
            timeout_s=timeout_s,
            convert_to="pdf:impress_pdf_Export",
        )
        pdfs = sorted(out_dir.glob("*.pdf"))
        if not pdfs:
            return pngs, "libreoffice_png_partial"

        pdf_path = pdfs[0]
        pdf_png_dir = out_dir / "pdf_pages"
        pdf_pngs = _render_pdf_to_pngs(
            pdf_path, out_dir=pdf_png_dir, target_width=SLIDE_RENDER_WIDTH
        )
        if pdf_pngs and len(pdf_pngs) >= expected_slide_count:
            return pdf_pngs, "libreoffice_pdf"

        return pdf_pngs or pngs, "libreoffice_pdf_partial"

    return pngs, "libreoffice_png"


class PptxFileUploadHandler(FileUploadHandlerPlugin):
    """Handler for PowerPoint (PPTX) file uploads."""

    async def process_file(self, file_bytes: bytes, filename: str) -> dict[str, Any]:
        # Validate file size
        self.validate_file_size(file_bytes, MAX_PPTX_SIZE, "PowerPoint")

        filename_lower = (filename or "").lower()
        if filename_lower.endswith(".ppt"):
            raise ValueError(
                "Only .pptx files are supported (legacy .ppt is not supported)."
            )

        title = filename.rsplit(".", 1)[0] if "." in filename else filename

        # Parse PPTX with python-pptx (for slide count, titles, text)
        try:
            with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
                tmp.write(file_bytes)
                tmp_path = Path(tmp.name)
        except Exception as e:
            raise ValueError(
                "Failed to write uploaded PowerPoint to a temp file."
            ) from e

        try:
            pres = Presentation(str(tmp_path))
            slide_count = len(pres.slides)
        except Exception as e:
            raise ValueError(
                "Failed to read PowerPoint file (is it a valid .pptx?)."
            ) from e

        # Convert to images via LibreOffice in a dedicated temp dir
        has_libreoffice = _find_libreoffice_cmd() is not None
        safe_name = _safe_stem(filename)
        with tempfile.TemporaryDirectory(prefix="canvas-chat-pptx-") as td:
            workdir = Path(td)
            pptx_path = workdir / f"{safe_name}.pptx"
            pptx_path.write_bytes(file_bytes)

            render_dir = workdir / "rendered"
            png_paths, rendering_mode = _render_pptx_to_pngs(
                pptx_path,
                out_dir=render_dir,
                timeout_s=LIBREOFFICE_TIMEOUT_S,
                expected_slide_count=slide_count,
            )

            # Best-effort alignment to slide count
            if slide_count and len(png_paths) != slide_count:
                logger.warning(
                    "PPTX conversion produced %s images for %s slides (filename=%s)",
                    len(png_paths),
                    slide_count,
                    filename,
                )

            slides: list[dict[str, Any]] = []
            for idx, slide in enumerate(pres.slides):
                slide_image: dict[str, Any] = {"mimeType": "image/webp"}
                if idx < len(png_paths):
                    png_path = png_paths[idx]
                    try:
                        img = Image.open(png_path)
                        img = _resize_to_width(img, SLIDE_RENDER_WIDTH)
                        thumb = _resize_to_width(img.copy(), THUMB_RENDER_WIDTH)

                        image_webp = _image_to_webp_b64(img, quality=90)
                        thumb_webp = _image_to_webp_b64(thumb, quality=70)
                        slide_image = {
                            "mimeType": "image/webp",
                            "image_webp": image_webp,
                            "thumb_webp": thumb_webp,
                        }
                    except Exception as e:
                        # Fallback to PNG base64 if WebP encode fails in this
                        # environment
                        logger.warning("WebP encode failed, falling back to PNG: %s", e)
                        slide_image = {
                            "mimeType": "image/png",
                            "image_png": _read_b64(png_path),
                            "thumb_png": _read_b64(png_path),
                        }
                else:
                    # Keep slide metadata even if an image is missing, so the UI
                    # can still
                    # navigate and show a clear "missing image" state.
                    logger.warning(
                        (
                            "Missing rendered image for slide %s "
                            "(have %s images, %s slides): %s"
                        ),
                        idx + 1,
                        len(png_paths),
                        slide_count,
                        filename,
                    )

                slide_title = _extract_slide_title(slide)
                slide_text = _extract_slide_text(slide)

                slides.append(
                    {
                        "index": idx,
                        "title": slide_title,
                        "text_content": slide_text,
                        **slide_image,
                    }
                )

            # Build a text export for context in chat
            parts: list[str] = []
            for s in slides:
                header = f"## Slide {s['index'] + 1}"
                if s.get("title"):
                    header += f": {s['title']}"
                parts.append(header)
                if s.get("text_content"):
                    parts.append(s["text_content"])
                else:
                    parts.append("(No text content found on this slide)")
                parts.append("")

            content = "\n".join(parts).strip() if parts else "(No slide content found)"

            logger.info(
                "Successfully processed PowerPoint: %s (%s slides)",
                filename,
                slide_count,
            )

            return {
                "content": content,
                "title": title,
                "slide_count": slide_count,
                "slides": slides,
                "metadata": {
                    "content_type": "powerpoint",
                    "slide_count": slide_count,
                    "source": "upload",
                    "rendering_mode": rendering_mode
                    if has_libreoffice
                    else "placeholder",
                },
            }
        # Ensure temp file is deleted
        try:
            tmp_path.unlink(missing_ok=True)
        except Exception:
            pass


# Register PowerPoint handler
FileUploadRegistry.register(
    id="pptx",
    mime_types=[
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint",  # legacy mime type
    ],
    extensions=[".pptx", ".ppt"],
    handler=PptxFileUploadHandler,
    priority=PRIORITY["BUILTIN"],
)

logger.info("PowerPoint file upload handler plugin loaded")
