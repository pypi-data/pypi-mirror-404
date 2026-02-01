"""Tests for PowerPoint file upload handler.

These are unit tests that avoid requiring LibreOffice at test time by
monkeypatching the LibreOffice conversion step.
"""

import asyncio
import io
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation

from canvas_chat.plugins import pptx_handler


def _make_sample_pptx_bytes() -> bytes:
    pres = Presentation()

    # Slide 1 (title slide)
    slide1 = pres.slides.add_slide(pres.slide_layouts[0])
    slide1.shapes.title.text = "Intro"
    if len(slide1.placeholders) > 1:
        slide1.placeholders[1].text = "Hello world"

    # Slide 2
    slide2 = pres.slides.add_slide(pres.slide_layouts[0])
    slide2.shapes.title.text = "Results"
    if len(slide2.placeholders) > 1:
        slide2.placeholders[1].text = "Key numbers"

    buf = io.BytesIO()
    pres.save(buf)
    return buf.getvalue()


class TestPptxHandler:
    def test_rejects_legacy_ppt(self):
        handler = pptx_handler.PptxFileUploadHandler()
        data = _make_sample_pptx_bytes()

        with pytest.raises(ValueError, match=r"Only \.pptx"):
            asyncio.run(handler.process_file(data, "legacy.ppt"))

    def test_process_file_extracts_slides_titles_text_and_images(
        self, monkeypatch, tmp_path: Path
    ):
        data = _make_sample_pptx_bytes()

        # Patch LibreOffice render step to produce deterministic PNGs
        def fake_render(
            pptx_path: Path,
            *,
            out_dir: Path,
            timeout_s: int,
            expected_slide_count=None,  # noqa: ARG001
        ):
            out_dir.mkdir(parents=True, exist_ok=True)
            out_paths = []
            # Use python-pptx to know slide count
            pres = Presentation(str(pptx_path))
            for i in range(len(pres.slides)):
                p = out_dir / f"{pptx_path.stem}_{i}.png"
                img = Image.new("RGB", (800, 600), color=(255, 255, 255))
                img.save(p, format="PNG")
                out_paths.append(p)
            return out_paths, "placeholder"

        monkeypatch.setattr(pptx_handler, "_render_pptx_to_pngs", fake_render)

        # Patch WebP encoder to avoid environment-specific WEBP support issues
        monkeypatch.setattr(
            pptx_handler, "_image_to_webp_b64", lambda _img, *, quality: "d2VicA=="
        )

        handler = pptx_handler.PptxFileUploadHandler()
        result = asyncio.run(handler.process_file(data, "deck.pptx"))

        assert result["title"] == "deck"
        assert result["slide_count"] == 2
        assert result["metadata"]["content_type"] == "powerpoint"

        slides = result["slides"]
        assert len(slides) == 2

        assert slides[0]["title"] == "Intro"
        assert "Hello world" in slides[0]["text_content"]
        assert slides[0]["mimeType"] in ("image/webp", "image/png")
        # Prefer webp keys when encoder is patched
        assert "image_webp" in slides[0]
        assert "thumb_webp" in slides[0]

        assert slides[1]["title"] == "Results"
        assert "Key numbers" in slides[1]["text_content"]
