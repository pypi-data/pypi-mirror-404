"""PowerPoint presentation export for Oscura reports.

This module provides PPTX generation for stakeholder presentations
with automated slide layouts, embedded plots, and speaker notes.


Note:
    This is a stub implementation. For full functionality, install python-pptx:
    pip install python-pptx

Example:
    >>> from oscura.reporting import export_pptx
    >>> data = {"title": "Analysis Report", "findings": [...]}
    >>> export_pptx(data, "presentation.pptx")
"""

from __future__ import annotations

import logging
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)

# Check if python-pptx is available
try:
    from pptx import Presentation
    from pptx.util import Inches

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not available. Install with: pip install python-pptx")


@dataclass
class PPTXSlide:
    """PowerPoint slide configuration.

    Attributes:
        title: Slide title
        content: Slide content (text, bullet points, or image path)
        layout: Slide layout type
        notes: Speaker notes
        chart_data: Optional chart data to embed
    """

    title: str
    content: str | list[str] | Path = ""
    layout: str = "title_content"  # title, title_content, content, blank
    notes: str = ""
    chart_data: dict[str, Any] | None = None


@dataclass
class PPTXPresentation:
    """PowerPoint presentation configuration.

    Attributes:
        title: Presentation title
        subtitle: Presentation subtitle
        author: Author name
        slides: List of slides
        template: Optional template file path

    References:
        REPORT-023: PowerPoint/PPTX Export
    """

    title: str
    subtitle: str = ""
    author: str = ""
    slides: list[PPTXSlide] = field(default_factory=list)
    template: Path | None = None

    def add_slide(
        self,
        title: str,
        content: str | list[str] | Path = "",
        *,
        layout: str = "title_content",
        notes: str = "",
    ) -> PPTXSlide:
        """Add slide to presentation.

        Args:
            title: Slide title
            content: Slide content
            layout: Layout type
            notes: Speaker notes

        Returns:
            Created slide object
        """
        slide = PPTXSlide(title=title, content=content, layout=layout, notes=notes)
        self.slides.append(slide)
        return slide


def export_pptx(
    report_data: dict[str, Any],
    output_path: str | Path,
    *,
    title: str = "Analysis Report",
    subtitle: str = "",
    author: str = "",
    template: Path | None = None,
) -> Path:
    """Export report data to PowerPoint presentation.

    Args:
        report_data: Report data dictionary containing:
            - 'summary': Executive summary text
            - 'findings': List of key findings
            - 'measurements': Measurement results
            - 'plots': List of plot image paths
        output_path: Output PPTX file path
        title: Presentation title
        subtitle: Presentation subtitle
        author: Author name
        template: Optional template PPTX file

    Returns:
        Path to generated PPTX file

    Example:
        >>> data = {
        ...     'summary': 'All tests passed',
        ...     'findings': ['Rise time: 2.3ns', 'Fall time: 2.1ns'],
        ...     'plots': [Path('plot1.png'), Path('plot2.png')]
        ... }
        >>> export_pptx(data, 'report.pptx', title='Signal Analysis')

    References:
        REPORT-023: PowerPoint/PPTX Export
    """
    output_path = Path(output_path)

    if not PPTX_AVAILABLE:
        # Create stub file
        logger.warning("Creating stub PPTX file (python-pptx not installed)")
        _create_stub_pptx(output_path, title, report_data)
        return output_path

    # Create presentation
    prs = Presentation(str(template)) if template else Presentation()

    # Title slide
    _add_title_slide(prs, title, subtitle, author)

    # Summary slide
    summary = report_data.get("summary", "")
    if summary:
        _add_summary_slide(prs, summary)

    # Key findings slides
    findings = report_data.get("findings", [])
    if findings:
        _add_findings_slide(prs, findings)

    # Measurement results slides
    measurements = report_data.get("measurements", [])
    if measurements:
        _add_measurement_slides(prs, measurements)

    # Plot slides
    plots = report_data.get("plots", [])
    for i, plot_path in enumerate(plots, 1):
        _add_plot_slide(prs, plot_path, f"Figure {i}")

    # Save presentation
    prs.save(str(output_path))
    logger.info("Exported PPTX presentation to %s", output_path)

    return output_path


def _add_title_slide(
    prs: Any,
    title: str,
    subtitle: str,
    author: str,
) -> None:
    """Add title slide to presentation."""
    if not PPTX_AVAILABLE:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout

    # Set title
    if slide.shapes.title:
        slide.shapes.title.text = title

    # Set subtitle (placeholder 1)
    if len(slide.placeholders) > 1:
        subtitle_text = subtitle
        if author:
            subtitle_text += f"\n{author}"
        slide.placeholders[1].text = subtitle_text


def _add_summary_slide(prs: Any, summary: str) -> None:
    """Add executive summary slide."""
    if not PPTX_AVAILABLE:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content layout

    if slide.shapes.title:
        slide.shapes.title.text = "Executive Summary"

    # Add summary text
    if len(slide.placeholders) > 1:
        text_frame = slide.placeholders[1].text_frame
        text_frame.text = summary


def _add_findings_slide(prs: Any, findings: list[str]) -> None:
    """Add key findings slide with bullet points."""
    if not PPTX_AVAILABLE:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content layout

    if slide.shapes.title:
        slide.shapes.title.text = "Key Findings"

    # Add bullet points
    if len(slide.placeholders) > 1:
        text_frame = slide.placeholders[1].text_frame
        text_frame.clear()

        for finding in findings:
            p = text_frame.add_paragraph()
            p.text = finding
            p.level = 0


def _add_measurement_slides(prs: Any, measurements: list[dict[str, Any]]) -> None:
    """Add measurement results slides."""
    if not PPTX_AVAILABLE:
        return

    # Group measurements into slides (5 per slide)
    chunk_size = 5
    for i in range(0, len(measurements), chunk_size):
        chunk = measurements[i : i + chunk_size]

        slide = prs.slides.add_slide(prs.slide_layouts[1])

        if slide.shapes.title:
            slide.shapes.title.text = f"Measurement Results ({i + 1}-{i + len(chunk)})"

        # Add measurements as bullet points
        if len(slide.placeholders) > 1:
            text_frame = slide.placeholders[1].text_frame
            text_frame.clear()

            for meas in chunk:
                p = text_frame.add_paragraph()
                name = meas.get("name", "Unknown")
                value = meas.get("value", "")
                unit = meas.get("unit", "")
                status = meas.get("status", "")

                p.text = f"{name}: {value} {unit} {status}"
                p.level = 0


def _add_plot_slide(prs: Any, plot_path: Path, caption: str) -> None:
    """Add slide with embedded plot image."""
    if not PPTX_AVAILABLE:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Add title
    if slide.shapes.title:
        slide.shapes.title.text = caption

    # Add image (centered)
    if plot_path.exists():
        left = Inches(1.5)
        top = Inches(2)
        width = Inches(7)

        slide.shapes.add_picture(str(plot_path), left, top, width=width)


def _create_stub_pptx(
    output_path: Path,
    title: str,
    report_data: dict[str, Any],
) -> None:
    """Create stub PPTX file when python-pptx not available."""
    # Write text file with .pptx extension as placeholder
    with output_path.open("w") as f:
        f.write("PowerPoint Export Stub\n")
        f.write("======================\n\n")
        f.write(f"Title: {title}\n\n")
        f.write("Install python-pptx for full PPTX export:\n")
        f.write("  pip install python-pptx\n\n")
        f.write("Report Data Summary:\n")
        f.write(f"  - Summary: {report_data.get('summary', 'N/A')}\n")
        f.write(f"  - Findings: {len(report_data.get('findings', []))} items\n")
        f.write(f"  - Measurements: {len(report_data.get('measurements', []))} items\n")
        f.write(f"  - Plots: {len(report_data.get('plots', []))} items\n")


def generate_presentation_from_report(
    report: dict[str, Any],
    output_path: str | Path,
    *,
    presentation_config: PPTXPresentation | None = None,
) -> Path:
    """Generate PowerPoint presentation from report structure.

    Args:
        report: Report dictionary with standard structure
        output_path: Output file path
        presentation_config: Optional presentation configuration

    Returns:
        Path to generated presentation

    References:
        REPORT-023: PowerPoint/PPTX Export
    """
    if presentation_config is None:
        presentation_config = PPTXPresentation(
            title=report.get("title", "Analysis Report"),
            subtitle=report.get("subtitle", ""),
            author=report.get("author", ""),
        )

    # Extract data from report
    report_data = {
        "summary": report.get("executive_summary", ""),
        "findings": report.get("key_findings", []),
        "measurements": report.get("measurements", []),
        "plots": report.get("plot_paths", []),
    }

    return export_pptx(
        report_data,
        output_path,
        title=presentation_config.title,
        subtitle=presentation_config.subtitle,
        author=presentation_config.author,
        template=presentation_config.template,
    )


__all__ = [
    "PPTXPresentation",
    "PPTXSlide",
    "export_pptx",
    "generate_presentation_from_report",
]
