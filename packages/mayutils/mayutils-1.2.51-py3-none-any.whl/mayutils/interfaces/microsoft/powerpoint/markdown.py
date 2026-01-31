"""Markdown to PowerPoint text_frame conversion utilities."""

from __future__ import annotations

from typing import Any, Literal, Optional, cast

from lxml import etree  # type: ignore[import-untyped]
from pptx.oxml import register_element_cls
from pptx.oxml.ns import qn
from pptx.oxml.simpletypes import ST_Coordinate32
from pptx.oxml.text import CT_TextCharacterProperties
from pptx.oxml.xmlchemy import OptionalAttribute
from pptx.text.text import TextFrame, _Paragraph, _Run
from pptx.util import Pt

from mayutils.interfaces.markdown import EMOJI_MAP, create_markdown_parser
from mayutils.objects.colours import Colour


class CT_TextCharacterPropertiesExtended(CT_TextCharacterProperties):
    """Extended CT_TextCharacterProperties with baseline support for superscript/subscript."""

    baseline = OptionalAttribute("baseline", ST_Coordinate32)


register_element_cls("a:rPr", CT_TextCharacterPropertiesExtended)


# Header sizes relative to base font size
HEADER_SIZES = {
    1: 2.0,
    2: 1.5,
    3: 1.25,
    4: 1.1,
    5: 1.0,
    6: 0.9,
}

# Code font family
CODE_FONT = "Consolas"


def _get_or_add_first_paragraph(
    text_frame: TextFrame,
) -> _Paragraph:
    """Get the first paragraph from a text_frame, creating one if none exist."""
    if len(text_frame.paragraphs) == 0:
        return text_frame.add_paragraph()
    return text_frame.paragraphs[0]


def _set_run_formatting(
    run: _Run,
    bold: Optional[bool] = None,
    italic: Optional[bool] = None,
    underline: Optional[bool] = None,
    strikethrough: Optional[bool] = None,
    font_size: Optional[int] = None,
    font_family: Optional[str] = None,
    font_colour: Optional[Colour] = None,
    highlight_colour: Optional[Colour] = None,
    hyperlink: Optional[str] = None,
    hyperlink_colour: Optional[Colour] = None,
    superscript: bool = False,
    subscript: bool = False,
) -> None:
    """Apply formatting to a run."""
    if bold is not None:
        run.font.bold = bold
    if italic is not None:
        run.font.italic = italic
    if underline is not None:
        run.font.underline = underline
    if strikethrough:
        run.font._rPr.attrib["strike"] = "sngStrike"
    if font_size is not None:
        run.font.size = Pt(font_size)
    if font_family is not None:
        run.font.name = font_family
    if font_colour is not None:
        run.font.color.rgb = font_colour.pptx_colour  # type: ignore[attr-defined]
    if highlight_colour is not None:
        rPr = run.font._rPr
        highlight = etree.SubElement(rPr, qn("a:highlight"))
        srgbClr = etree.SubElement(highlight, qn("a:srgbClr"))
        srgbClr.set("val", highlight_colour.to_str(method="hex").lstrip("#"))
    if hyperlink is not None:
        run.hyperlink.address = hyperlink
        if hyperlink_colour is not None:
            run.font.color.rgb = hyperlink_colour.pptx_colour  # type: ignore[attr-defined]
    if superscript:
        run.font._rPr.baseline = 30000  # 30% above baseline
    if subscript:
        run.font._rPr.baseline = -25000  # 25% below baseline


def _process_inline_tokens(
    paragraph: _Paragraph,
    tokens: list[dict[str, Any]],
    formatting_stack: dict[str, Any],
    font_size: Optional[int] = None,
    font_family: Optional[str] = None,
    font_colour: Optional[Colour] = None,
    highlight_colour: Optional[Colour] = None,
    hyperlink_colour: Optional[Colour] = None,
) -> None:
    """Process inline tokens and add runs to paragraph."""
    for token in tokens:
        token_type = token.get("type")

        if token_type == "text":
            run = paragraph.add_run()
            run.text = token.get("raw", "")
            _set_run_formatting(
                run=run,
                bold=formatting_stack.get("bold"),
                italic=formatting_stack.get("italic"),
                underline=formatting_stack.get("underline"),
                strikethrough=formatting_stack.get("strikethrough"),
                font_size=font_size,
                font_family=formatting_stack.get("code_font") or font_family,
                font_colour=font_colour,
                highlight_colour=formatting_stack.get("highlight") or highlight_colour,
                hyperlink=formatting_stack.get("hyperlink"),
                hyperlink_colour=hyperlink_colour,
                superscript=formatting_stack.get("superscript", False),
                subscript=formatting_stack.get("subscript", False),
            )

        elif token_type == "emphasis":
            new_stack = formatting_stack.copy()
            new_stack["italic"] = True
            _process_inline_tokens(
                paragraph=paragraph,
                tokens=token.get("children", []),
                formatting_stack=new_stack,
                font_size=font_size,
                font_family=font_family,
                font_colour=font_colour,
                highlight_colour=highlight_colour,
                hyperlink_colour=hyperlink_colour,
            )

        elif token_type == "strong":
            new_stack = formatting_stack.copy()
            new_stack["bold"] = True
            _process_inline_tokens(
                paragraph=paragraph,
                tokens=token.get("children", []),
                formatting_stack=new_stack,
                font_size=font_size,
                font_family=font_family,
                font_colour=font_colour,
                highlight_colour=highlight_colour,
                hyperlink_colour=hyperlink_colour,
            )

        elif token_type == "underline":
            new_stack = formatting_stack.copy()
            new_stack["underline"] = True
            _process_inline_tokens(
                paragraph=paragraph,
                tokens=token.get("children", []),
                formatting_stack=new_stack,
                font_size=font_size,
                font_family=font_family,
                font_colour=font_colour,
                highlight_colour=highlight_colour,
                hyperlink_colour=hyperlink_colour,
            )

        elif token_type == "strikethrough":
            new_stack = formatting_stack.copy()
            new_stack["strikethrough"] = True
            _process_inline_tokens(
                paragraph=paragraph,
                tokens=token.get("children", []),
                formatting_stack=new_stack,
                font_size=font_size,
                font_family=font_family,
                font_colour=font_colour,
                highlight_colour=highlight_colour,
                hyperlink_colour=hyperlink_colour,
            )

        elif token_type == "mark":
            new_stack = formatting_stack.copy()
            new_stack["highlight"] = highlight_colour or Colour.parse("#FFFF00")
            _process_inline_tokens(
                paragraph=paragraph,
                tokens=token.get("children", []),
                formatting_stack=new_stack,
                font_size=font_size,
                font_family=font_family,
                font_colour=font_colour,
                highlight_colour=highlight_colour,
                hyperlink_colour=hyperlink_colour,
            )

        elif token_type == "superscript":
            new_stack = formatting_stack.copy()
            new_stack["superscript"] = True
            _process_inline_tokens(
                paragraph=paragraph,
                tokens=token.get("children", []),
                formatting_stack=new_stack,
                font_size=font_size,
                font_family=font_family,
                font_colour=font_colour,
                highlight_colour=highlight_colour,
                hyperlink_colour=hyperlink_colour,
            )

        elif token_type == "subscript":
            new_stack = formatting_stack.copy()
            new_stack["subscript"] = True
            _process_inline_tokens(
                paragraph=paragraph,
                tokens=token.get("children", []),
                formatting_stack=new_stack,
                font_size=font_size,
                font_family=font_family,
                font_colour=font_colour,
                highlight_colour=highlight_colour,
                hyperlink_colour=hyperlink_colour,
            )

        elif token_type == "codespan":
            run = paragraph.add_run()
            run.text = token.get("raw", "")
            _set_run_formatting(
                run=run,
                bold=formatting_stack.get("bold"),
                italic=formatting_stack.get("italic"),
                font_size=font_size,
                font_family=CODE_FONT,
                font_colour=font_colour,
                highlight_colour=Colour.parse("#E8E8E8"),
            )

        elif token_type == "link":
            new_stack = formatting_stack.copy()
            new_stack["hyperlink"] = token.get("attrs", {}).get("url", "")
            _process_inline_tokens(
                paragraph=paragraph,
                tokens=token.get("children", []),
                formatting_stack=new_stack,
                font_size=font_size,
                font_family=font_family,
                font_colour=font_colour,
                highlight_colour=highlight_colour,
                hyperlink_colour=hyperlink_colour,
            )

        elif token_type == "softbreak":
            run = paragraph.add_run()
            run.text = " "
            _set_run_formatting(
                run=run,
                font_size=font_size,
                font_family=font_family,
                font_colour=font_colour,
            )

        elif token_type == "linebreak":
            run = paragraph.add_run()
            run.text = "\n"

        else:
            # Handle unknown inline types by processing children if present
            if "children" in token:
                _process_inline_tokens(
                    paragraph=paragraph,
                    tokens=token.get("children", []),
                    formatting_stack=formatting_stack,
                    font_size=font_size,
                    font_family=font_family,
                    font_colour=font_colour,
                    highlight_colour=highlight_colour,
                    hyperlink_colour=hyperlink_colour,
                )
            elif "raw" in token:
                run = paragraph.add_run()
                run.text = token.get("raw", "")
                _set_run_formatting(
                    run=run,
                    font_size=font_size,
                    font_family=font_family,
                    font_colour=font_colour,
                )


def _add_bullet_subelement(
    parent,
    tagname: str,
    **kwargs,
) -> etree._Element:
    """Helper to add OxmlElement for bullet formatting."""
    element = etree.SubElement(parent, qn(tagname))
    for key, value in kwargs.items():
        element.set(key, value)
    return element


def _set_paragraph_bullet(
    paragraph: _Paragraph,
    level: int = 0,
    numbered: bool = False,
    bullet_char: Optional[str] = None,
    space_before: Optional[int] = 10,
    space_after: Optional[int] = 0,
    bullet_size_pct: int = 100000,  # 100% = 100000
    from_placeholder: bool = False,
) -> None:
    """
    Set bullet point formatting on a paragraph.

    Based on: https://stackoverflow.com/questions/49384134/add-bullet-list-using-python-pptx

    Args:
        paragraph: The paragraph to format.
        level: Indentation level (0-based).
        numbered: Whether to use numbered list.
        bullet_char: Custom bullet character (default: filled square).
        space_before: Space before paragraph in points.
        space_after: Space after paragraph in points.
        bullet_size_pct: Bullet size as percentage (100000 = 100%).
        from_placeholder: If True, only set level and let placeholder handle bullet style.
    """
    # Set the indentation level using python-pptx's built-in property
    paragraph.level = level

    # Set spacing between bullet points
    if space_before is not None:
        paragraph.space_before = Pt(space_before)
    if space_after is not None:
        paragraph.space_after = Pt(space_after)

    # If from_placeholder, just set the level and let the placeholder handle bullets
    if from_placeholder:
        return

    # Access the underlying XML to set bullet style
    pPr = paragraph._p.get_or_add_pPr()

    # Remove any existing bullet-related elements
    for tag in ("a:buNone", "a:buChar", "a:buFont", "a:buAutoNum", "a:buSzPct"):
        existing = pPr.find(qn(tag))
        if existing is not None:
            pPr.remove(existing)

    # Set indentation for proper bullet formatting (in EMUs)
    # marL = left margin (positive), indent = hanging indent (negative for bullets)
    base_margin = 342900  # ~0.375 inch
    indent_per_level = 342900  # ~0.375 inch per level
    bullet_indent = 171450  # ~0.1875 inch hanging indent

    pPr.set("marL", str(base_margin + (indent_per_level * level)))
    pPr.set("indent", str(-bullet_indent))

    if numbered:
        # Add auto-numbering
        _add_bullet_subelement(pPr, "a:buAutoNum", type="arabicPeriod")
    else:
        # Add bullet size percentage
        _add_bullet_subelement(pPr, "a:buSzPct", val=str(bullet_size_pct))

        # Add bullet font
        _add_bullet_subelement(
            pPr,
            "a:buFont",
            typeface="Arial",
            panose="020B0604020202020204",
            pitchFamily="34",
            charset="0",
        )

        # Add bullet character (default: filled square)
        char = bullet_char if bullet_char is not None else "•"
        _add_bullet_subelement(pPr, "a:buChar", char=char)


def _set_paragraph_task(
    paragraph: _Paragraph,
    checked: bool = False,
    level: int = 0,
    space_before: Optional[int] = 6,
    space_after: Optional[int] = 0,
) -> None:
    """Set task checkbox formatting on a paragraph."""
    # Set the indentation level
    paragraph.level = level

    # Set spacing between items
    if space_before is not None:
        paragraph.space_before = Pt(space_before)
    if space_after is not None:
        paragraph.space_after = Pt(space_after)

    pPr = paragraph._p.get_or_add_pPr()

    buNone = pPr.find(qn("a:buNone"))
    if buNone is not None:
        pPr.remove(buNone)

    # Remove any existing auto-numbering
    buAutoNum = pPr.find(qn("a:buAutoNum"))
    if buAutoNum is not None:
        pPr.remove(buAutoNum)

    buChar = pPr.find(qn("a:buChar"))
    if buChar is None:
        buChar = etree.SubElement(pPr, qn("a:buChar"))
    buChar.set("char", EMOJI_MAP["checkmark"] if checked else "☐")


def _process_list_item(
    text_frame: TextFrame,
    token: dict[str, Any],
    level: int,
    numbered: bool,
    first_paragraph: bool,
    font_size: Optional[int],
    font_family: Optional[str],
    font_colour: Optional[Colour],
    highlight_colour: Optional[Colour],
    hyperlink_colour: Optional[Colour],
    bullet_char: Optional[str],
    from_placeholder: bool,
) -> bool:
    """Process a list item token, returning whether we used first_paragraph."""
    children = token.get("children", [])

    for child in children:
        child_type = child.get("type")

        if child_type in ("paragraph", "block_text"):
            if first_paragraph:
                paragraph = _get_or_add_first_paragraph(text_frame)
                first_paragraph = False
            else:
                paragraph = text_frame.add_paragraph()

            # Check if this is a task list item
            is_task = False
            checked = False
            inline_children = child.get("children", [])
            if inline_children and inline_children[0].get("type") == "task_list_marker":
                is_task = True
                checked = inline_children[0].get("attrs", {}).get("checked", False)
                inline_children = inline_children[1:]

            # First add the text content
            _process_inline_tokens(
                paragraph=paragraph,
                tokens=inline_children,
                formatting_stack={},
                font_size=font_size,
                font_family=font_family,
                font_colour=font_colour,
                highlight_colour=highlight_colour,
                hyperlink_colour=hyperlink_colour,
            )

            # Then set bullet formatting (after text is added)
            if is_task:
                _set_paragraph_task(
                    paragraph=paragraph,
                    checked=checked,
                    level=level,
                )
            else:
                _set_paragraph_bullet(
                    paragraph=paragraph,
                    level=level,
                    numbered=numbered,
                    bullet_char=bullet_char,
                    from_placeholder=from_placeholder,
                )

        elif child_type in ("list", "bullet_list", "ordered_list"):
            first_paragraph = _process_list(
                text_frame=text_frame,
                token=child,
                level=level + 1,
                first_paragraph=first_paragraph,
                font_size=font_size,
                font_family=font_family,
                font_colour=font_colour,
                highlight_colour=highlight_colour,
                hyperlink_colour=hyperlink_colour,
                bullet_char=bullet_char,
                from_placeholder=from_placeholder,
            )

    return first_paragraph


def _process_list(
    text_frame: TextFrame,
    token: dict[str, Any],
    level: int,
    first_paragraph: bool,
    font_size: Optional[int],
    font_family: Optional[str],
    font_colour: Optional[Colour],
    highlight_colour: Optional[Colour],
    hyperlink_colour: Optional[Colour],
    bullet_char: Optional[str],
    from_placeholder: bool,
) -> bool:
    """Process a list token (bullet or ordered)."""
    token_type = token.get("type")
    numbered = token_type == "list" and token.get("attrs", {}).get("ordered", False)

    for item in token.get("children", []):
        if item.get("type") == "list_item":
            first_paragraph = _process_list_item(
                text_frame=text_frame,
                token=item,
                level=level,
                numbered=numbered,
                first_paragraph=first_paragraph,
                font_size=font_size,
                font_family=font_family,
                font_colour=font_colour,
                highlight_colour=highlight_colour,
                hyperlink_colour=hyperlink_colour,
                bullet_char=bullet_char,
                from_placeholder=from_placeholder,
            )

    return first_paragraph


def add_markdown_to_text_frame(
    text_frame: TextFrame,
    markdown: str,
    font_family: Optional[str] = None,
    font_size: Optional[int] = None,
    line_spacing: Optional[float] = None,
    text_alignment: Optional[Literal["left", "center", "right", "justify"]] = None,
    font_colour: Optional[Colour | str] = None,
    highlight_colour: Optional[Colour | str] = None,
    hyperlink_colour: Optional[Colour | str] = None,
    bullet_char: Optional[str] = None,
    from_placeholder: bool = False,
) -> TextFrame:
    """
    Add markdown-formatted text to a PowerPoint text_frame.

    Supports the following markdown syntax:
    - `**text**` for bold
    - `*text*` for italic
    - `__text__` for underline
    - `~~text~~` for strikethrough
    - `` `code` `` for inline code
    - `- item` for bullet points (with indentation for levels)
    - `1. item` for numbered lists (with indentation for levels)
    - `> quote` for blockquotes
    - ``` for code blocks
    - `[text](url)` for hyperlinks
    - `# Header` for headers (h1-h6)
    - `:emoji:` for emojis
    - `==text==` for highlighting
    - `X^2^` for superscript
    - `X~2~` for subscript
    - `- [ ]` and `- [x]` for task lists
    - Footnotes

    Args:
        text_frame: The PowerPoint text_frame to add content to.
        markdown: The markdown string to parse and add.
        font_family: Default font family (e.g., "Arial", "Calibri").
        font_size: Default font size in points.
        line_spacing: Line spacing multiplier (e.g., 1.5 for 150%).
        text_alignment: Text alignment ("left", "center", "right", "justify").
        font_colour: Default font colour (Colour object or parseable string).
        highlight_colour: Default highlight colour for marked text.
        hyperlink_colour: Colour for hyperlinks.
        bullet_char: Custom bullet character for unordered lists (default: filled square).
        from_placeholder: If True, assumes the text_frame already has bullet formatting
            from a placeholder. Only sets paragraph levels without adding bullet chars.
            Validates that markdown only contains list content.

    Returns:
        The modified text_frame.
    """
    # Parse colours if strings
    if font_colour is not None and not isinstance(font_colour, Colour):
        font_colour = Colour.parse(colour=font_colour)
    if highlight_colour is not None and not isinstance(highlight_colour, Colour):
        highlight_colour = Colour.parse(colour=highlight_colour)
    if hyperlink_colour is not None and not isinstance(hyperlink_colour, Colour):
        hyperlink_colour = Colour.parse(colour=hyperlink_colour)

    # Parse markdown to AST
    parser = create_markdown_parser(renderer=None)
    tokens = cast(list[dict[str, Any]], parser(markdown))

    # Validate that markdown only contains lists when from_placeholder=True
    if from_placeholder:
        allowed_types = {"list", "blank_line"}
        for token in tokens:
            token_type = token.get("type")
            if token_type not in allowed_types:
                raise ValueError(
                    f"When from_placeholder=True, markdown must only contain lists. "
                    f"Found: {token_type}"
                )

    # Track whether we've used the first paragraph
    first_paragraph = True
    footnotes_content: dict[str, list] = {}

    # Process block-level tokens
    for token in tokens:
        token_type = token.get("type")

        if token_type == "paragraph":
            if first_paragraph:
                paragraph = _get_or_add_first_paragraph(text_frame)
                first_paragraph = False
            else:
                paragraph = text_frame.add_paragraph()

            if line_spacing is not None:
                paragraph.line_spacing = line_spacing
            if text_alignment is not None:
                from pptx.enum.text import PP_ALIGN

                alignment_map = {
                    "left": PP_ALIGN.LEFT,
                    "center": PP_ALIGN.CENTER,
                    "right": PP_ALIGN.RIGHT,
                    "justify": PP_ALIGN.JUSTIFY,
                }
                paragraph.alignment = alignment_map.get(text_alignment)

            _process_inline_tokens(
                paragraph=paragraph,
                tokens=token.get("children", []),
                formatting_stack={},
                font_size=font_size,
                font_family=font_family,
                font_colour=font_colour,
                highlight_colour=highlight_colour,
                hyperlink_colour=hyperlink_colour,
            )

        elif token_type == "heading":
            if first_paragraph:
                paragraph = _get_or_add_first_paragraph(text_frame)
                first_paragraph = False
            else:
                paragraph = text_frame.add_paragraph()

            if line_spacing is not None:
                paragraph.line_spacing = line_spacing
            if text_alignment is not None:
                from pptx.enum.text import PP_ALIGN

                alignment_map = {
                    "left": PP_ALIGN.LEFT,
                    "center": PP_ALIGN.CENTER,
                    "right": PP_ALIGN.RIGHT,
                    "justify": PP_ALIGN.JUSTIFY,
                }
                paragraph.alignment = alignment_map.get(text_alignment)

            level = token.get("attrs", {}).get("level", 1)
            header_size = (
                int(font_size * HEADER_SIZES.get(level, 1.0))
                if font_size
                else int(24 * HEADER_SIZES.get(level, 1.0))
            )

            _process_inline_tokens(
                paragraph=paragraph,
                tokens=token.get("children", []),
                formatting_stack={"bold": True},
                font_size=header_size,
                font_family=font_family,
                font_colour=font_colour,
                highlight_colour=highlight_colour,
                hyperlink_colour=hyperlink_colour,
            )

        elif token_type == "list":
            first_paragraph = _process_list(
                text_frame=text_frame,
                token=token,
                level=0,
                first_paragraph=first_paragraph,
                font_size=font_size,
                font_family=font_family,
                font_colour=font_colour,
                highlight_colour=highlight_colour,
                hyperlink_colour=hyperlink_colour,
                bullet_char=bullet_char,
                from_placeholder=from_placeholder,
            )

        elif token_type == "block_quote":
            # Process blockquote content with indentation
            for child in token.get("children", []):
                if child.get("type") == "paragraph":
                    if first_paragraph:
                        paragraph = _get_or_add_first_paragraph(text_frame)
                        first_paragraph = False
                    else:
                        paragraph = text_frame.add_paragraph()

                    # Add quote indicator
                    run = paragraph.add_run()
                    run.text = "│ "
                    _set_run_formatting(
                        run=run,
                        font_size=font_size,
                        font_family=font_family,
                        font_colour=Colour.parse("#666666"),
                    )

                    _process_inline_tokens(
                        paragraph=paragraph,
                        tokens=child.get("children", []),
                        formatting_stack={"italic": True},
                        font_size=font_size,
                        font_family=font_family,
                        font_colour=Colour.parse("#666666")
                        if font_colour is None
                        else font_colour,
                        highlight_colour=highlight_colour,
                        hyperlink_colour=hyperlink_colour,
                    )

        elif token_type == "block_code":
            if first_paragraph:
                paragraph = _get_or_add_first_paragraph(text_frame)
                first_paragraph = False
            else:
                paragraph = text_frame.add_paragraph()

            code_text = token.get("raw", "").rstrip("\n")
            run = paragraph.add_run()
            run.text = code_text
            _set_run_formatting(
                run=run,
                font_size=font_size if font_size else 10,
                font_family=CODE_FONT,
                font_colour=font_colour,
                highlight_colour=Colour.parse("#F5F5F5"),
            )

        elif token_type == "footnote_list":
            # Store footnotes for later processing
            for footnote in token.get("children", []):
                if footnote.get("type") == "footnote_item":
                    key = footnote.get("attrs", {}).get("key", "")
                    footnotes_content[key] = footnote.get("children", [])

        elif token_type == "thematic_break":
            if first_paragraph:
                paragraph = _get_or_add_first_paragraph(text_frame)
                first_paragraph = False
            else:
                paragraph = text_frame.add_paragraph()

            run = paragraph.add_run()
            run.text = "─" * 40
            _set_run_formatting(
                run=run,
                font_size=font_size,
                font_family=font_family,
                font_colour=Colour.parse("#CCCCCC"),
            )

        elif token_type == "blank_line":
            # Skip blank lines (they're just spacing in markdown)
            pass

    # Add footnotes at the end if any exist
    if footnotes_content:
        # Add separator
        paragraph = text_frame.add_paragraph()
        run = paragraph.add_run()
        run.text = "─" * 20
        _set_run_formatting(
            run=run,
            font_size=font_size if font_size else 8,
            font_family=font_family,
            font_colour=Colour.parse("#999999"),
        )

        # Add each footnote
        for key, content in footnotes_content.items():
            paragraph = text_frame.add_paragraph()
            run = paragraph.add_run()
            run.text = f"[{key}] "
            _set_run_formatting(
                run=run,
                font_size=(font_size - 2) if font_size else 8,
                font_family=font_family,
                font_colour=Colour.parse("#666666"),
                superscript=True,
            )

            for child in content:
                if child.get("type") == "paragraph":
                    _process_inline_tokens(
                        paragraph=paragraph,
                        tokens=child.get("children", []),
                        formatting_stack={},
                        font_size=(font_size - 2) if font_size else 8,
                        font_family=font_family,
                        font_colour=Colour.parse("#666666"),
                        highlight_colour=highlight_colour,
                        hyperlink_colour=hyperlink_colour,
                    )

    return text_frame
