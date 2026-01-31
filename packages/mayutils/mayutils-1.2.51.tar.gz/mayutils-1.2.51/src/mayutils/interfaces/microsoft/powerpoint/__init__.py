from __future__ import annotations
from pathlib import Path
from typing import Optional, Self
from pptx import Presentation as Init
from pptx.shapes.autoshape import Shape
from pptx.slide import SlideLayouts, SlideLayout, Slides, Slide
from pptx.util import Pt, Length as BaseLength
from pptx.dml.color import RGBColor
from mayutils.objects.colours import Colour
import six
from copy import deepcopy
from mayutils.export.images import IMAGES_FOLDER
import subprocess
import shutil


class Length(BaseLength):
    @classmethod
    def from_float(
        cls,
        value: float,
    ) -> Self:
        return cls(value)  # type: ignore


class SlideContext:
    def __init__(
        self,
        presentation: Presentation,
        layout: Optional[SlideLayout] = None,
    ) -> None:
        self.presentation = presentation
        self.layout = layout if layout is not None else presentation.blank_layout
        self.slide = self.presentation.new_slide(layout=self.layout).slides[-1]

    def __enter__(
        self,
    ) -> Slide:
        return self.slide

    def __exit__(
        self,
        exc_type,
        exc_value,
        traceback,
    ) -> None:
        if exc_type is not None:
            # self.presentation.delete_slide(slide_number=len(self.presentation.slides))
            raise exc_value


class Presentation:
    def __init__(
        self,
        template: Path | str,
    ) -> None:
        self.template = Path(template)
        if not self.template.exists():
            raise FileNotFoundError(f"Template file {self.template} does not exist.")
        elif self.template.is_dir():
            raise ValueError(
                f"Template file {self.template} is a directory, not a file."
            )
        elif not self.template.is_file():
            raise ValueError(f"Template file {self.template} is not a valid file.")
        elif self.template.suffix.lower() not in [".pptx", ".ppt"]:
            raise ValueError(
                f"Template file {self.template} is not a valid PowerPoint file."
            )
        elif self.template.suffix.lower() == ".ppt":
            raise ValueError(
                f"Template file {self.template} is a legacy PowerPoint file (.ppt). "
                "Please convert it to .pptx format."
            )

        self.internal = Init(pptx=str(self.template))

        self.blank_layout = self.internal.slide_layouts[
            len(self.internal.slide_layouts) - 1
        ]

    # TODO: Implement display/__repr__/_repr_mimebundle_
    # def __repr__(
    #     self,
    #     text: str,
    # ) -> Self:

    @property
    def layouts(
        self,
    ) -> SlideLayouts:
        return self.internal.slide_layouts

    @property
    def height(
        self,
    ) -> Length:
        return (
            Length(emu=self.internal.slide_height)
            if self.internal.slide_height is not None
            else Length.from_float(value=7.5 * BaseLength._EMUS_PER_INCH)
        )

    @height.setter
    def height(
        self,
        value: Length,
    ) -> None:
        self.internal.slide_height = value

    @property
    def width(
        self,
    ) -> Length:
        return (
            Length(emu=self.internal.slide_width)
            if self.internal.slide_width is not None
            else Length.from_float(value=13.33 * BaseLength._EMUS_PER_INCH)
        )

    @width.setter
    def width(
        self,
        value: Length,
    ) -> None:
        self.internal.slide_width = value

    # TODO: Implement
    # @property
    # def title(
    #     self,
    # ) -> str:

    @property
    def slides(
        self,
    ) -> Slides:
        return self.internal.slides

    def slide(
        self,
        slide_number: int,
    ) -> Slide:
        if slide_number < 1 or slide_number > len(self.slides):
            raise IndexError(
                f"Slide number {slide_number} is out of range. Presentation has {len(self.slides)} slides."
            )

        return self.slides[slide_number - 1]

    def new_slide(
        self,
        layout: Optional[SlideLayout] = None,
    ) -> Self:
        self.slides.add_slide(
            slide_layout=layout if layout is not None else self.blank_layout
        )

        return self

    def enter_new_slide(
        self,
        layout: Optional[SlideLayout] = None,
    ) -> SlideContext:
        return SlideContext(
            presentation=self,
            layout=layout,
        )

    def empty(
        self,
    ) -> Self:
        for i in range(len(self.slides) - 1, -1, -1):
            rId = self.slides._sldIdLst[i].rId
            self.internal.part.drop_rel(rId=rId)
            del self.slides._sldIdLst[i]

        return self

    def delete_slide(
        self,
        slide_number: int,
    ) -> Self:
        raise NotImplementedError("Deleting slides is not implemented yet.")
        return self

    def copy_slide(
        self,
        slide_number: int,
    ) -> Self:
        raise NotImplementedError("Copying slides is not implemented yet.")
        slide_idx = slide_number - 1
        template_slide = self.slides[slide_idx]

        try:
            blank_slide_layout = self.internal.slide_layouts[12]
        except IndexError:
            blank_slide_layout = self.internal.slide_layouts[
                len(self.internal.slide_layouts) - 1
            ]

        copied_slide = self.slides.add_slide(slide_layout=blank_slide_layout)

        for shape in template_slide.shapes:
            element = shape.element
            new_element = deepcopy(element)
            copied_slide.shapes._spTree.insert_element_before(new_element, "p:extLst")

        for _, value in six.iteritems(template_slide.part.rels):
            if "notesSlide" not in value.reltype:
                copied_slide.part.rels.add_relationship(
                    value.reltype,
                    value._target,
                    value.rId,
                )

        return self

    def move_slide(
        self,
        slide_number: int,
        to_position: int,
    ) -> Self:
        raise NotImplementedError("Moving slides is not implemented yet.")
        return self

    def reorder_slides(
        self,
        new_order: list[int],
    ) -> Self:
        raise NotImplementedError("Reordering slides is not implemented yet.")
        return self

    def insertion_spacing(
        self,
        height: Optional[Length] = None,
        width: Optional[Length] = None,
        x_shift: Optional[Length] = None,
        y_shift: Optional[Length] = None,
    ) -> dict:
        if width is None:
            if self.width is None:
                raise ValueError(
                    "Width must be specified if presentation width is not set."
                )
        if height is None:
            if self.height is None:
                raise ValueError(
                    "Height must be specified if presentation height is not set."
                )
        if x_shift is None:
            if self.width is None:
                raise ValueError(
                    "Width must be specified if presentation width is not set."
                )
            if width is not None:
                x_shift = Length.from_float(value=(self.width - width) * 0.5)
            else:
                x_shift = Length.from_float(value=0.05 * self.width)
        if y_shift is None:
            if self.height is None:
                raise ValueError(
                    "Height must be specified if presentation height is not set."
                )
            if height is not None:
                y_shift = Length.from_float(value=(self.height - height) * 0.5)
            else:
                y_shift = Length.from_float(value=0.05 * self.height)

        return dict(
            left=x_shift,
            top=y_shift,
            width=width,
            height=height,
        )

    def insert_textbox(
        self,
        slide_number: Optional[int] = None,
        height: Optional[Length] = None,
        width: Optional[Length] = None,
        x_shift: Optional[Length] = None,
        y_shift: Optional[Length] = None,
        **kwargs,
    ) -> Shape:
        slide = self.slide(
            slide_number=slide_number
            if slide_number is not None
            else len(self.slides) - 1
        )

        textbox = slide.shapes.add_textbox(
            **self.insertion_spacing(
                height=height,
                width=width,
                x_shift=x_shift,
                y_shift=y_shift,
            ),
        )

        return textbox

    def insert_text(
        self,
        text: str,
        textbox: Shape,
        bold: bool = False,
        italic: bool = False,
        underline: bool = False,
        strikethrough: bool = False,
        font_size: Optional[int] = None,
        font_family: Optional[str] = None,
        colour: Optional[Colour | str] = None,
        background_colour: Optional[Colour | str] = None,
        link: Optional[str] = None,
    ) -> Self:
        if colour is not None and not isinstance(colour, Colour):
            colour = Colour.parse(colour=colour)
        if background_colour is not None and not isinstance(background_colour, Colour):
            background_colour = Colour.parse(colour=background_colour)

        textbox.text_frame.text = text

        if font_size is not None:
            textbox.text_frame.paragraphs[0].font.size = Pt(points=font_size)
        if font_family is not None:
            textbox.text_frame.paragraphs[0].font.name = font_family
        if colour is not None:
            textbox.text_frame.paragraphs[0].font.color.rgb = RGBColor(
                r=colour.r,
                g=colour.g,
                b=colour.b,
            )
        if background_colour is not None:
            textbox.fill.solid()
            textbox.fill.fore_color.rgb = RGBColor(
                r=background_colour.r,
                g=background_colour.g,
                b=background_colour.b,
            )
        if bold:
            textbox.text_frame.paragraphs[0].font.bold = True
        if italic:
            textbox.text_frame.paragraphs[0].font.italic = True
        if underline:
            textbox.text_frame.paragraphs[0].font.underline = True
        if strikethrough:
            textbox.text_frame.paragraphs[0].font._element.attrib["strike"] = (
                "sngStrike"
            )

        if link is not None:
            raise NotImplementedError("Jyperlinks are not implemented yet.")
            from pptx.util import URI

            textbox.text_frame.paragraphs[0].hyperlink.address = URI(link)

        return self

    # TODO: Implement
    # def insert_markdown(
    #     self,
    #     text: str,
    # ) -> Self:

    # TODO: Implement
    # def insert_image(
    #     self,
    #     image_path: Path | str,
    # ) -> Self:

    # TODO: Implement
    # def insert_table(
    #     self,
    #     table: DataFrame | Styler,
    # ) -> Self:

    def save(
        self,
        file_path: Path | str,
    ) -> None:
        file_path = Path(file_path)

        return self.internal.save(
            file=str(file_path),
        )


def convert_pptx_to_pdf(
    pptx_path: str | Path,
    output_dir: str | Path = IMAGES_FOLDER.parent / "PDF",
    soffice_path: str | None = None,
) -> Path:
    pptx_path = Path(pptx_path)
    output_dir = Path(output_dir)

    if not pptx_path.exists():
        raise FileNotFoundError(f"PPTX not found: {pptx_path}")

    output_dir.mkdir(parents=True, exist_ok=True)

    if soffice_path:
        soffice = Path(soffice_path)
    else:
        soffice = shutil.which("soffice")
        if soffice:
            soffice = Path(soffice)

    if not soffice or not Path(soffice).exists():
        raise FileNotFoundError(
            "LibreOffice (soffice) not found. "
            "Install LibreOffice or pass soffice_path explicitly."
        )

    cmd = [
        str(soffice),
        "--headless",
        "--convert-to",
        "pdf",
        "--outdir",
        str(output_dir),
        str(pptx_path),
    ]

    result = subprocess.run(
        cmd,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True,
    )

    if result.returncode != 0:
        raise RuntimeError(
            "PPTX → PDF conversion failed\n"
            f"STDOUT:\n{result.stdout}\n"
            f"STDERR:\n{result.stderr}"
        )

    pdf_path = output_dir / (pptx_path.stem + ".pdf")

    if not pdf_path.exists():
        raise RuntimeError("Conversion reported success but PDF not found")

    return pdf_path
