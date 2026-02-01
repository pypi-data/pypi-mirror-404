from pptx import Presentation
from pathlib import Path
from typing import Any, Dict
import logging

logger = logging.getLogger(__name__)


async def execute(args: Dict[str, Any], context: Dict[str, Any]) -> str:
    file_path: str = args.get("file_path", "")

    path = Path(file_path)

    if not path.exists():
        return f"错误：文件不存在 {file_path}"

    if not path.is_file():
        return f"错误：{file_path} 不是文件"

    try:
        prs = Presentation(str(path))

        info: list[str] = []
        info.append(f"文件大小：{path.stat().st_size} 字节")

        slide_count = len(prs.slides)
        info.append(f"幻灯片数量：{slide_count}")

        if prs.core_properties.title:
            info.append(f"标题: {prs.core_properties.title}")
        if prs.core_properties.author:
            info.append(f"作者: {prs.core_properties.author}")
        if prs.core_properties.last_modified_by:
            info.append(f"最后修改者: {prs.core_properties.last_modified_by}")

        text_content = ""

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = f"\n--- 幻灯片 {slide_num} ---\n"

            title = ""
            if slide.shapes.title:
                title = slide.shapes.title.text.strip()
                if title:
                    slide_text += f"标题: {title}\n"

            content_count = 0
            for shape in slide.shapes:
                if shape != slide.shapes.title:
                    shape_text = getattr(shape, "text", None)
                    if shape_text:
                        text = shape_text.strip()
                        if text and len(text) > 0:
                            content_count += 1
                            slide_text += f"{text}\n"

            if not title and content_count == 0:
                slide_text += "(空白幻灯片)"

            text_content += slide_text

        if not text_content.strip():
            text_content = "(演示文稿未检测到文本内容)"

        info.append(f"\n内容预览（前 5000 字符）：\n{text_content[:5000]}")
        if len(text_content) > 5000:
            info.append(f"\n... (共 {len(text_content)} 字符)")

        return "\n".join(info)

    except Exception as e:
        logger.exception(f"解析 PowerPoint 失败: {e}")
        return f"解析 PowerPoint 失败: {e}"
