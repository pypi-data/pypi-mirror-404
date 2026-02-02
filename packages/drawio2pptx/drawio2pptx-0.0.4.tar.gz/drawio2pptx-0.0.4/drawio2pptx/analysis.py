"""
Conversion result analysis module

Compares draw.io files with generated PowerPoint files to analyze
what is correctly converted and what is not converted
"""
import logging
from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from lxml import etree as ET
import sys

from drawio2pptx.io.drawio_loader import DrawIOLoader
from drawio2pptx.logger import get_logger


def _print(*args, **kwargs):
    """Print to stdout without logging format"""
    print(*args, **kwargs, file=sys.stdout)
from drawio2pptx.mapping.shape_map import map_shape_type_to_pptx
from pptx.enum.shapes import MSO_SHAPE

NSMAP_DRAWINGML = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
NSMAP_PRESENTATIONML = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}


def format_color_for_comparison(color):
    """Convert color to string for comparison"""
    if color is None:
        return "None"
    if isinstance(color, RGBColor):
        return f"RGB({color[0]}, {color[1]}, {color[2]})"
    if isinstance(color, str):
        if color == "default":
            return "ThemeColor"
        return color
    return str(color)


def analyze_pptx_shape(shape):
    """Analyze PowerPoint shape"""
    result = {
        'shape_type': shape.shape_type,
        'fill': None,
        'stroke_color': None,
        'has_shadow': None,
        'text_frame': None,
    }
    
    # Fill
    try:
        fill = shape.fill
        fill_type = fill.type
        if fill_type is None:
            result['fill'] = "None"
        elif fill_type == 5:  # BACKGROUND (5)
            result['fill'] = "None"
        elif hasattr(fill, 'fore_color'):
            try:
                fill_color = fill.fore_color
                if hasattr(fill_color, 'rgb') and fill_color.rgb:
                    rgb = fill_color.rgb
                    result['fill'] = f"RGB({rgb.r}, {rgb.g}, {rgb.b})"
                elif hasattr(fill_color, 'theme_color'):
                    result['fill'] = f"ThemeColor({fill_color.theme_color})"
                else:
                    result['fill'] = "Unknown"
            except (TypeError, AttributeError):
                # If fore_color cannot be accessed, check from XML
                result['fill'] = None  # Check from XML later
        else:
            result['fill'] = str(fill_type)
    except Exception:
        result['fill'] = None  # Check from XML later
    
    # Check fill from XML (for ThemeColor or None cases)
    try:
        if (result['fill'] == "None" or result['fill'] is None) and hasattr(shape, '_element'):
            shape_element = shape._element
            # spPr exists as a direct child element
            sp_pr = None
            for child in shape_element:
                if child.tag.endswith('}spPr') or 'spPr' in child.tag:
                    sp_pr = child
                    break
            
            if sp_pr is not None:
                # Check noFill
                no_fill = sp_pr.find('.//a:noFill', namespaces=NSMAP_DRAWINGML)
                if no_fill is not None:
                    result['fill'] = "None"
                else:
                    solid_fill = sp_pr.find('.//a:solidFill', namespaces=NSMAP_DRAWINGML)
                    if solid_fill is not None:
                        scheme_clr = solid_fill.find('.//a:schemeClr', namespaces=NSMAP_DRAWINGML)
                        if scheme_clr is not None:
                            val = scheme_clr.get('val', '')
                            result['fill'] = f"ThemeColor({val})"
                        else:
                            srgb = solid_fill.find('.//a:srgbClr', namespaces=NSMAP_DRAWINGML)
                            if srgb is not None:
                                val = srgb.get('val', '')
                                if val and len(val) == 6:
                                    r = int(val[0:2], 16)
                                    g = int(val[2:4], 16)
                                    b = int(val[4:6], 16)
                                    result['fill'] = f"RGB({r}, {g}, {b})"
    except Exception as e:
        logger = get_logger()
        logger.debug(f"Failed to extract fill color from XML: {e}")
    
    # Set to "None" if result is None
    if result['fill'] is None:
        result['fill'] = "None"
    
    # Stroke color
    try:
        if hasattr(shape, '_element'):
            shape_element = shape._element
            ln_element = shape_element.find('.//a:ln', namespaces=NSMAP_DRAWINGML)
            if ln_element is not None:
                solid_fill = ln_element.find('.//a:solidFill', namespaces=NSMAP_DRAWINGML)
                if solid_fill is not None:
                    srgb = solid_fill.find('.//a:srgbClr', namespaces=NSMAP_DRAWINGML)
                    if srgb is not None:
                        val = srgb.get('val', '')
                        if val and len(val) == 6:
                            r = int(val[0:2], 16)
                            g = int(val[2:4], 16)
                            b = int(val[4:6], 16)
                            result['stroke_color'] = f"RGB({r}, {g}, {b})"
    except Exception as e:
        logger = get_logger()
        logger.debug(f"Failed to extract stroke color from XML: {e}")
    
    # Shadow
    try:
        if hasattr(shape, 'shadow'):
            result['has_shadow'] = shape.shadow.inherit
    except Exception as e:
        logger = get_logger()
        logger.debug(f"Failed to extract shadow information: {e}")
    
    # Text frame
    try:
        if shape.has_text_frame:
            tf = shape.text_frame
            text_runs = []
            for para in tf.paragraphs:
                for run in para.runs:
                    run_info = {
                        'text': run.text,
                        'font_name': run.font.name if run.font.name else None,
                        'font_size': run.font.size.pt if run.font.size else None,
                        'bold': run.font.bold if run.font.bold else False,
                        'italic': run.font.italic if run.font.italic else False,
                        'underline': run.font.underline if run.font.underline else False,
                        'font_color': None,
                    }
                    
                    # Get font color from XML
                    try:
                        if hasattr(run, '_r'):
                            run_element = run._r
                            r_pr = run_element.find('.//a:rPr', namespaces=NSMAP_DRAWINGML)
                            if r_pr is not None:
                                solid_fill = r_pr.find('.//a:solidFill', namespaces=NSMAP_DRAWINGML)
                                if solid_fill is not None:
                                    srgb = solid_fill.find('.//a:srgbClr', namespaces=NSMAP_DRAWINGML)
                                    if srgb is not None:
                                        val = srgb.get('val', '')
                                        if val and len(val) == 6:
                                            r = int(val[0:2], 16)
                                            g = int(val[2:4], 16)
                                            b = int(val[4:6], 16)
                                            run_info['font_color'] = f"RGB({r}, {g}, {b})"
                    except Exception as e:
                        logger = get_logger()
                        logger.debug(f"Failed to extract font color from run XML: {e}")
                    
                    text_runs.append(run_info)
            
            result['text_frame'] = {
                'vertical_anchor': str(tf.vertical_anchor),
                'runs': text_runs,
            }
    except Exception as e:
        logger = get_logger()
        logger.debug(f"Failed to extract text frame information: {e}")
    
    return result


def compare_conversion(input_path: Path, output_path: Path):
    """Compare conversion results and display analysis"""
    logger = get_logger()
    
    _print("=== Conversion Result Analysis ===")
    _print()
    
    # Load draw.io file
    loader = DrawIOLoader(logger=logger)
    diagrams = loader.load_file(input_path)
    
    if not diagrams:
        _print("Error: No diagrams found in draw.io file")
        return
    
    elements = loader.extract_elements(diagrams[0])
    shape_elements = [e for e in elements if hasattr(e, 'shape_type') and hasattr(e, 'element_type') and e.element_type == 'shape']
    connector_elements = [e for e in elements if hasattr(e, 'element_type') and e.element_type == 'connector']
    
    # Load PowerPoint file
    if not output_path.exists():
        _print(f"Error: PowerPoint file not found: {output_path}")
        return
    
    prs = Presentation(str(output_path))
    if not prs.slides:
        _print("Error: No slides found in PowerPoint file")
        return
    
    slide = prs.slides[0]
    pptx_shapes = [s for s in slide.shapes if s.shape_type == 1]  # AUTO_SHAPE
    pptx_connectors = [s for s in slide.shapes if s.shape_type == 5]  # FREEFORM (connectors)
    
    _print(f"draw.io shapes: {len(shape_elements)}")
    _print(f"PowerPoint shapes: {len(pptx_shapes)}")
    _print(f"draw.io connectors: {len(connector_elements)}")
    _print(f"PowerPoint connectors: {len(pptx_connectors)}")
    _print()
    
    # Compare each shape
    converted_features = set()
    missing_features = []
    
    for idx, (drawio_shape, pptx_shape) in enumerate(zip(shape_elements, pptx_shapes), 1):
        _print(f"--- Shape {idx} ---")
        
        # Shape type
        drawio_shape_type = drawio_shape.shape_type
        pptx_shape_type = pptx_shape.shape_type
        pptx_auto_shape_type = None
        try:
            if hasattr(pptx_shape, 'auto_shape_type'):
                pptx_auto_shape_type = pptx_shape.auto_shape_type
        except Exception as e:
            logger = get_logger()
            logger.debug(f"Failed to get auto_shape_type: {e}")
        
        auto_shape_info = f", auto_shape_type: {pptx_auto_shape_type}" if pptx_auto_shape_type else ""
        
        # Check if shape type is correctly converted
        expected_pptx_type = map_shape_type_to_pptx(drawio_shape_type)
        
        # Check if shape type is supported
        supported_types = ['rectangle', 'ellipse', 'circle', 'square', 'rect', 'rhombus', 'parallelogram']
        if drawio_shape_type.lower() not in supported_types:
            _print(f"Shape type: draw.io={drawio_shape_type}, PowerPoint={pptx_shape_type}{auto_shape_info} ✗ Unsupported (converted to {expected_pptx_type})")
            missing_features.append(f"Shape {idx}: Unsupported shape type '{drawio_shape_type}' (converted to rectangle)")
        else:
            # Compare auto_shape_type if available
            match_status = ""
            if pptx_auto_shape_type:
                if pptx_auto_shape_type == expected_pptx_type:
                    match_status = " ✓"
                    converted_features.add(f"Shape type: {drawio_shape_type}")
                else:
                    match_status = " ✗"
                    _print(f"Shape type: draw.io={drawio_shape_type}, PowerPoint={pptx_shape_type}{auto_shape_info} ✗ Mismatch (expected {expected_pptx_type}, got {pptx_auto_shape_type})")
                    missing_features.append(f"Shape {idx}: Shape type mismatch ({drawio_shape_type} -> expected {expected_pptx_type}, got {pptx_auto_shape_type})")
            else:
                # Fallback: compare shape_type values
                if pptx_shape_type == expected_pptx_type.value:
                    match_status = " ✓"
                    converted_features.add(f"Shape type: {drawio_shape_type}")
                else:
                    match_status = " ✗"
                    _print(f"Shape type: draw.io={drawio_shape_type}, PowerPoint={pptx_shape_type}{auto_shape_info} ✗ Mismatch (expected {expected_pptx_type}, got {pptx_shape_type})")
                    missing_features.append(f"Shape {idx}: Shape type mismatch ({drawio_shape_type} -> expected {expected_pptx_type}, got {pptx_shape_type})")
            
            if match_status == " ✓":
                _print(f"Shape type: draw.io={drawio_shape_type}, PowerPoint={pptx_shape_type}{auto_shape_info}{match_status}")
        
        # Text
        drawio_text = ""
        if drawio_shape.text:
            for para in drawio_shape.text:
                for run in para.runs:
                    drawio_text += run.text
        
        pptx_text = ""
        if pptx_shape.has_text_frame:
            for para in pptx_shape.text_frame.paragraphs:
                for run in para.runs:
                    pptx_text += run.text
        
        if drawio_text == pptx_text:
            _print(f"Text: draw.io='{drawio_text}', PowerPoint='{pptx_text}' ✓")
            converted_features.add("Text")
        else:
            _print(f"Text: draw.io='{drawio_text}', PowerPoint='{pptx_text}' ✗ Mismatch")
            missing_features.append(f"Shape {idx}: Text mismatch")
        
        # Fill color
        drawio_fill = format_color_for_comparison(drawio_shape.style.fill)
        pptx_result = analyze_pptx_shape(pptx_shape)
        pptx_fill = pptx_result['fill'] or "None"
        
        if drawio_fill == "ThemeColor" and "ThemeColor" in pptx_fill:
            _print(f"Fill: draw.io={drawio_fill}, PowerPoint={pptx_fill} ✓ (ThemeColor)")
            converted_features.add("Fill color (ThemeColor)")
        elif drawio_fill == pptx_fill or (drawio_fill == "None" and pptx_fill == "None"):
            _print(f"Fill: draw.io={drawio_fill}, PowerPoint={pptx_fill} ✓")
            if drawio_fill == "None":
                converted_features.add("No fill")
            else:
                converted_features.add("Fill color")
        else:
            _print(f"Fill: draw.io={drawio_fill}, PowerPoint={pptx_fill} ✗ Mismatch")
            missing_features.append(f"Shape {idx}: Fill mismatch ({drawio_fill} vs {pptx_fill})")
        
        # Stroke color
        drawio_stroke = format_color_for_comparison(drawio_shape.style.stroke)
        pptx_stroke = pptx_result['stroke_color'] or "None"
        
        if drawio_stroke == pptx_stroke or (drawio_stroke == "None" and pptx_stroke == "None"):
            _print(f"Stroke color: draw.io={drawio_stroke}, PowerPoint={pptx_stroke} ✓")
            if drawio_stroke == "None":
                converted_features.add("No stroke")
            else:
                converted_features.add("Stroke color")
        else:
            _print(f"Stroke color: draw.io={drawio_stroke}, PowerPoint={pptx_stroke} ✗ Mismatch")
            missing_features.append(f"Shape {idx}: Stroke color mismatch")
        
        # Shadow
        drawio_shadow = drawio_shape.style.has_shadow
        pptx_shadow = pptx_result['has_shadow']
        if drawio_shadow == pptx_shadow:
            _print(f"Shadow: draw.io={drawio_shadow}, PowerPoint={pptx_shadow} ✓")
            converted_features.add("Shadow" if drawio_shadow else "No shadow")
        else:
            _print(f"Shadow: draw.io={drawio_shadow}, PowerPoint={pptx_shadow} ✗ Mismatch")
            missing_features.append(f"Shape {idx}: Shadow mismatch")
        
        # Font information
        if drawio_shape.text and pptx_shape.has_text_frame:
            drawio_runs = []
            for para in drawio_shape.text:
                drawio_runs.extend(para.runs)
            
            pptx_runs = []
            for para in pptx_shape.text_frame.paragraphs:
                pptx_runs.extend(para.runs)
            
            for run_idx, (drawio_run, pptx_run) in enumerate(zip(drawio_runs, pptx_runs), 1):
                # Font color
                drawio_color = format_color_for_comparison(drawio_run.font_color)
                pptx_color = "None"
                if pptx_result['text_frame'] and pptx_result['text_frame']['runs']:
                    if run_idx <= len(pptx_result['text_frame']['runs']):
                        pptx_color = pptx_result['text_frame']['runs'][run_idx-1].get('font_color', 'None')
                
                if drawio_color == pptx_color or (drawio_color == "None" and pptx_color == "RGB(0, 0, 0)"):
                    _print(f"  Run {run_idx} - Font color: draw.io={drawio_color}, PowerPoint={pptx_color} ✓")
                    converted_features.add("Font color")
                else:
                    _print(f"  Run {run_idx} - Font color: draw.io={drawio_color}, PowerPoint={pptx_color} ✗ Mismatch")
                    missing_features.append(f"Shape {idx} Run {run_idx}: Font color mismatch")
                
                # Bold
                pptx_bold = pptx_run.font.bold if pptx_run.font.bold else False
                if drawio_run.bold == pptx_bold:
                    _print(f"  Run {run_idx} - Bold: draw.io={drawio_run.bold}, PowerPoint={pptx_bold} ✓")
                    if drawio_run.bold:
                        converted_features.add("Bold")
                else:
                    _print(f"  Run {run_idx} - Bold: draw.io={drawio_run.bold}, PowerPoint={pptx_bold} ✗ Mismatch")
                    missing_features.append(f"Shape {idx} Run {run_idx}: Bold mismatch")
                
                # Italic
                pptx_italic = pptx_run.font.italic if pptx_run.font.italic else False
                if drawio_run.italic == pptx_italic:
                    _print(f"  Run {run_idx} - Italic: draw.io={drawio_run.italic}, PowerPoint={pptx_italic} ✓")
                    if drawio_run.italic:
                        converted_features.add("Italic")
                else:
                    _print(f"  Run {run_idx} - Italic: draw.io={drawio_run.italic}, PowerPoint={pptx_italic} ✗ Mismatch")
                    missing_features.append(f"Shape {idx} Run {run_idx}: Italic mismatch")
                
                # Underline
                pptx_underline = pptx_run.font.underline if pptx_run.font.underline else False
                if drawio_run.underline == pptx_underline:
                    _print(f"  Run {run_idx} - Underline: draw.io={drawio_run.underline}, PowerPoint={pptx_underline} ✓")
                    if drawio_run.underline:
                        converted_features.add("Underline")
                else:
                    _print(f"  Run {run_idx} - Underline: draw.io={drawio_run.underline}, PowerPoint={pptx_underline} ✗ Mismatch")
                    missing_features.append(f"Shape {idx} Run {run_idx}: Underline mismatch")
                
                # Font size
                drawio_size = drawio_run.font_size
                pptx_size = pptx_run.font.size.pt if pptx_run.font.size else None
                if drawio_size == pptx_size or (drawio_size is None and pptx_size == 12.0):
                    _print(f"  Run {run_idx} - Font size: draw.io={drawio_size}, PowerPoint={pptx_size} ✓")
                    converted_features.add("Font size")
                else:
                    _print(f"  Run {run_idx} - Font size: draw.io={drawio_size}, PowerPoint={pptx_size} ✗ Mismatch")
                    missing_features.append(f"Shape {idx} Run {run_idx}: Font size mismatch")
                
                # Font name
                drawio_font = drawio_run.font_family
                pptx_font = pptx_run.font.name
                # If drawio_font is None or empty string, draw.io's default font (Helvetica) is used
                # pptx_writer.py also sets draw.io's default font, so they should match
                if drawio_font == pptx_font:
                    _print(f"  Run {run_idx} - Font name: draw.io={drawio_font}, PowerPoint={pptx_font} ✓")
                    if drawio_font:
                        converted_features.add("Font name")
                elif drawio_font is None or drawio_font == "":
                    # This should not normally occur (default font is set in drawio_loader.py)
                    # Report as mismatch just in case
                    _print(f"  Run {run_idx} - Font name: draw.io={drawio_font} (default), PowerPoint={pptx_font} ✗ Mismatch (draw.io font not specified)")
                    missing_features.append(f"Shape {idx} Run {run_idx}: Font name mismatch (draw.io font not specified)")
                else:
                    _print(f"  Run {run_idx} - Font name: draw.io={drawio_font}, PowerPoint={pptx_font} ✗ Mismatch")
                    missing_features.append(f"Shape {idx} Run {run_idx}: Font name mismatch")
        
        # Vertical alignment
        if drawio_shape.text:
            drawio_vertical = drawio_shape.text[0].vertical_align if drawio_shape.text else None
        else:
            drawio_vertical = None
        
        pptx_vertical = None
        if pptx_shape.has_text_frame:
            pptx_vertical = str(pptx_shape.text_frame.vertical_anchor)
        
        if drawio_vertical is None and "MIDDLE" in pptx_vertical:
            _print(f"Vertical alignment: draw.io={drawio_vertical}, PowerPoint={pptx_vertical} ✓ (default: middle)")
            converted_features.add("Vertical alignment (default middle)")
        elif drawio_vertical and pptx_vertical:
            _print(f"Vertical alignment: draw.io={drawio_vertical}, PowerPoint={pptx_vertical} ✓")
            converted_features.add("Vertical alignment")
        else:
            _print(f"Vertical alignment: draw.io={drawio_vertical}, PowerPoint={pptx_vertical}")
        
        _print()
    
    # Compare connectors
    if connector_elements:
        _print("=== Connectors ===")
        _print(f"Connector count: draw.io={len(connector_elements)}, PowerPoint={len(pptx_connectors)}")
        
        if len(connector_elements) == len(pptx_connectors):
            _print(f"Connector count: draw.io={len(connector_elements)}, PowerPoint={len(pptx_connectors)} ✓")
            converted_features.add("Connectors")
        else:
            _print(f"Connector count: draw.io={len(connector_elements)}, PowerPoint={len(pptx_connectors)} ✗ Mismatch")
            missing_features.append(f"Connector count mismatch ({len(connector_elements)} vs {len(pptx_connectors)})")
        
        for idx, (drawio_connector, pptx_connector) in enumerate(zip(connector_elements, pptx_connectors), 1):
            _print(f"--- Connector {idx} ---")
            
            # Edge style (curved/orthogonal are not supported, converted to straight polyline)
            drawio_edge_style = getattr(drawio_connector, 'edge_style', 'straight')
            if drawio_edge_style and drawio_edge_style != 'straight':
                _print(f"Edge style: draw.io={drawio_edge_style}, PowerPoint=straight (polyline) ✗ Unsupported (converted to straight polyline)")
                missing_features.append(f"Connector {idx}: Unsupported edge style '{drawio_edge_style}' (converted to straight polyline)")
            else:
                _print(f"Edge style: draw.io={drawio_edge_style}, PowerPoint=straight ✓")
                converted_features.add("Connector edge style")
            
            # Stroke color
            drawio_stroke = format_color_for_comparison(drawio_connector.style.stroke)
            pptx_result = analyze_pptx_shape(pptx_connector)
            pptx_stroke = pptx_result['stroke_color'] or "None"
            
            if drawio_stroke == pptx_stroke or (drawio_stroke == "None" and pptx_stroke == "None"):
                _print(f"Stroke color: draw.io={drawio_stroke}, PowerPoint={pptx_stroke} ✓")
                converted_features.add("Connector stroke color")
            else:
                _print(f"Stroke color: draw.io={drawio_stroke}, PowerPoint={pptx_stroke} ✗ Mismatch")
                missing_features.append(f"Connector {idx}: Stroke color mismatch")
            
            # Shadow
            drawio_shadow = drawio_connector.style.has_shadow
            pptx_shadow = pptx_result['has_shadow']
            if drawio_shadow == pptx_shadow:
                _print(f"Shadow: draw.io={drawio_shadow}, PowerPoint={pptx_shadow} ✓")
                converted_features.add("Connector shadow")
            else:
                _print(f"Shadow: draw.io={drawio_shadow}, PowerPoint={pptx_shadow} ✗ Mismatch")
                missing_features.append(f"Connector {idx}: Shadow mismatch")
            
            _print()
    
    # Summary
    _print("=== Converted Features ===")
    if converted_features:
        for feature in sorted(converted_features):
            _print(f"  ✓ {feature}")
    else:
        _print("  (No features converted)")
    
    _print()
    _print("=== Missing Features ===")
    if missing_features:
        for feature in missing_features:
            _print(f"  ✗ {feature}")
    else:
        _print("  ✓ All features converted successfully")
    
    _print()
