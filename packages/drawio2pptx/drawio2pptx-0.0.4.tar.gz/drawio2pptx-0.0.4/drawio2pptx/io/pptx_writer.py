"""
PowerPoint output module

Generates PowerPoint presentations from intermediate models using python-pptx + lxml
"""
from typing import List, Optional, Tuple
from lxml import etree as ET
from pptx import Presentation  # type: ignore[import]
from pptx.util import Emu, Pt  # type: ignore[import]
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR  # type: ignore[import]
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR  # type: ignore[import]
from pptx.dml.color import RGBColor  # type: ignore[import]

from ..model.intermediate import BaseElement, ShapeElement, ConnectorElement, TextElement, TextParagraph, TextRun
from ..geom.units import px_to_emu, px_to_pt, scale_font_size_for_pptx
from ..geom.transform import split_polyline_to_segments
from ..mapping.shape_map import map_shape_type_to_pptx
from ..mapping.style_map import map_arrow_type, map_arrow_type_with_size, map_dash_pattern
from ..logger import ConversionLogger
from ..fonts import replace_font, DRAWIO_DEFAULT_FONT_FAMILY
from ..config import PARALLELOGRAM_SKEW, ConversionConfig, default_config

# XML namespaces
NS_DRAWINGML = 'http://schemas.openxmlformats.org/drawingml/2006/main'
NS_PRESENTATIONML = 'http://schemas.openxmlformats.org/presentationml/2006/main'
NSMAP_DRAWINGML = {'a': NS_DRAWINGML}
NSMAP_PRESENTATIONML = {'p': NS_PRESENTATIONML}
NSMAP_BOTH = {'p': NS_PRESENTATIONML, 'a': NS_DRAWINGML}


def _a(tag_name: str) -> str:
    """Create DrawingML namespace-qualified tag name"""
    return f'{{{NS_DRAWINGML}}}{tag_name}'


def _p(tag_name: str) -> str:
    """Create PresentationML namespace-qualified tag name"""
    return f'{{{NS_PRESENTATIONML}}}{tag_name}'


class PPTXWriter:
    """PowerPoint presentation writer"""
    
    def __init__(self, logger: Optional[ConversionLogger] = None, config: Optional[ConversionConfig] = None):
        """
        Args:
            logger: ConversionLogger instance
            config: ConversionConfig instance (uses default_config if None)
        """
        self.config = config or default_config
        self.logger = logger
    
    def create_presentation(self, page_size: Optional[Tuple[float, float]] = None) -> Presentation:
        """
        Create presentation
        
        Args:
            page_size: (width, height) tuple (px), or None
        
        Returns:
            Presentation object
        """
        prs = Presentation()
        
        # Get blank layout
        blank_layout_index = 6
        try:
            blank_layout = prs.slide_layouts[blank_layout_index]
        except Exception:
            blank_layout = prs.slide_layouts[0]
        
        # Set slide size
        if page_size and page_size[0] and page_size[1]:
            try:
                prs.slide_width = px_to_emu(page_size[0])
                prs.slide_height = px_to_emu(page_size[1])
            except Exception as e:
                if self.logger:
                    self.logger.debug(f"Failed to set slide size: {e}")
        
        return prs, blank_layout

    def add_slide(self, prs: Presentation, blank_layout, elements: List[BaseElement]):
        """
        Add elements to slide
        
        Args:
            prs: Presentation object
            blank_layout: Blank layout
            elements: List of elements (sorted by Z-order; later elements are on top)
        """
        slide = prs.slides.add_slide(blank_layout)
        
        # Add elements in the provided stacking order (later = topmost in PowerPoint).
        for element in elements:
            if isinstance(element, ShapeElement):
                self._add_shape(slide, element)
            elif isinstance(element, ConnectorElement):
                self._add_connector(slide, element)
            elif isinstance(element, TextElement):
                self._add_text(slide, element)
    
    def _add_shape(self, slide, shape: ShapeElement):
        """Add shape"""
        if shape.w <= 0 or shape.h <= 0:
            return None

        if (shape.shape_type or "").lower() == "line":
            return self._add_line_shape(slide, shape)
        
        # Map shape type
        pptx_shape_type = map_shape_type_to_pptx(shape.shape_type)
        
        # For rectangles with corner radius, use ROUNDED_RECTANGLE
        if (pptx_shape_type == MSO_SHAPE.RECTANGLE and 
            shape.style.corner_radius is not None and 
            shape.style.corner_radius > 0):
            pptx_shape_type = MSO_SHAPE.ROUNDED_RECTANGLE
        
        # Create shape
        left = px_to_emu(shape.x)
        top = px_to_emu(shape.y)
        width = px_to_emu(shape.w)
        height = px_to_emu(shape.h)
        
        shp = slide.shapes.add_shape(
            pptx_shape_type,
            left, top, width, height
        )
        
        # Debug-friendly name (not visible in normal slideshow mode).
        try:
            if shape.id:
                shp.name = f"drawio2pptx:shape:{shape.id}"
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to set shape name: {e}")

        # For parallelograms, explicitly set the skew (adjust) to match connection point calculation with appearance
        try:
            if shape.shape_type and "parallelogram" in shape.shape_type.lower():
                if hasattr(shp, "adjustments") and len(shp.adjustments) > 0:
                    shp.adjustments[0] = float(PARALLELOGRAM_SKEW)
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to set parallelogram adjustments: {e}")
        
        # Set text
        if shape.text:
            margin_overrides = None
            text_direction = None
            try:
                if getattr(shape.style, "is_swimlane", False):
                    start_size = float(getattr(shape.style, "swimlane_start_size", 0.0) or 0.0)
                    if start_size > 0:
                        first_para = shape.text[0]
                        base_top = first_para.spacing_top or 0.0
                        base_left = first_para.spacing_left or 0.0
                        base_bottom = first_para.spacing_bottom or 0.0
                        base_right = first_para.spacing_right or 0.0

                        if getattr(shape.style, "swimlane_horizontal", False):
                            # Horizontal swimlane: header on the top
                            margin_overrides = (
                                base_top,
                                base_left,
                                max(shape.h - start_size + base_bottom, 0.0),
                                base_right,
                            )
                        else:
                            # Vertical swimlane: header on the left (vertical text)
                            # Use vert270 to match draw.io's left header vertical text
                            text_direction = "vert270"
                            margin_overrides = (
                                base_top,
                                base_left,
                                base_bottom,
                                max(shape.w - start_size + base_right, 0.0),
                            )
            except Exception as e:
                if self.logger:
                    self.logger.debug(f"Failed to compute swimlane header margins: {e}")

            self._set_text_frame(
                shp.text_frame,
                shape.text,
                default_highlight_color=shape.style.label_background_color,
                word_wrap=shape.style.word_wrap,
                margin_overrides_px=margin_overrides,
                text_direction=text_direction,
            )
        
        # Set fill
        fill_color = shape.style.fill
        is_swimlane = getattr(shape.style, "is_swimlane", False)
        swimlane_start = float(getattr(shape.style, "swimlane_start_size", 0) or 0)

        # Swimlane gradient only when fillColor is explicit (e.g. flowchart3). Do not apply when
        # fillColor is omitted (e.g. flowchart2 Pool/Lane) to avoid unwanted gradient.
        if is_swimlane and swimlane_start > 0 and shape.h > 0 and isinstance(fill_color, RGBColor):
            # Swimlane with explicit fillColor: header = fillColor, body = swimlaneFillColor.
            self._set_swimlane_gradient_fill_xml(shp, shape)
        else:
            if fill_color == "default":
                self._set_default_fill_xml(shp)
            elif fill_color:
                try:
                    shp.fill.solid()
                    shp.fill.fore_color.rgb = fill_color
                except Exception as e:
                    if self.logger:
                        self.logger.debug(f"Failed to set fill color: {e}")
            else:
                try:
                    shp.fill.background()
                except Exception as e:
                    if self.logger:
                        self.logger.debug(f"Failed to set background fill: {e}")

        # Gradient fill: only used for swimlane header vs body (_set_swimlane_gradient_fill_xml above).
        # Do not apply draw.io gradientColor/gradientDirection to normal shapes to avoid unwanted gradient.
        #
        # Set stroke
        if getattr(shape.style, "no_stroke", False):
            try:
                self._set_no_line_xml(shp)
            except Exception as e:
                if self.logger:
                    self.logger.debug(f"Failed to disable stroke: {e}")
        else:
            # Use black as default if stroke is None
            stroke_color = shape.style.stroke if shape.style.stroke else RGBColor(0, 0, 0)
            try:
                shp.line.fill.solid()
                self._set_stroke_color_xml(shp, stroke_color)
            except Exception as e:
                if self.logger:
                    self.logger.debug(f"Failed to set stroke color: {e}")
            
            if shape.style.stroke_width > 0:
                try:
                    shp.line.width = px_to_pt(shape.style.stroke_width)
                except Exception as e:
                    if self.logger:
                        self.logger.debug(f"Failed to set stroke width: {e}")

        # Swimlane header divider (draw.io visual split between header and content)
        try:
            if getattr(shape.style, "is_swimlane", False):
                start_size = float(getattr(shape.style, "swimlane_start_size", 0.0) or 0.0)
                draw_divider = getattr(shape.style, "swimlane_line", True)
                if start_size > 0 and draw_divider:
                    self._add_swimlane_header_divider(
                        slide=slide,
                        shape=shape,
                        stroke_color=stroke_color,
                        stroke_width_px=shape.style.stroke_width,
                    )
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to add swimlane header divider: {e}")
        
        # Shadow settings (similar to legacy: disable shadow when has_shadow is False)
        try:
            if shape.style.has_shadow:
                # Enable shadow (inherit from theme)
                shp.shadow.inherit = True
            else:
                # Disable shadow
                shp.shadow.inherit = False
                # Disable shadow via XML (similar to legacy _disable_shadow_xml)
                self._disable_shadow_xml(shp)
        except Exception:
            if not shape.style.has_shadow:
                self._disable_shadow_xml(shp)

        # Add BPMN symbol overlay (e.g., plus sign for parallel gateway)
        try:
            bpmn_symbol = getattr(shape.style, "bpmn_symbol", None)
            if bpmn_symbol and bpmn_symbol.lower() == "parallelgw":
                self._add_bpmn_parallel_gateway_symbol(slide, shape)
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to add BPMN symbol overlay: {e}")
        
        return shp

    def _add_line_shape(self, slide, shape: ShapeElement):
        """Add line shape (draw.io line vertex) as a connector."""
        # Determine orientation (default to horizontal).
        if shape.w >= shape.h:
            y = shape.y + (shape.h / 2.0)
            x1, y1 = shape.x, y
            x2, y2 = shape.x + shape.w, y
        else:
            x = shape.x + (shape.w / 2.0)
            x1, y1 = x, shape.y
            x2, y2 = x, shape.y + shape.h

        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            px_to_emu(x1),
            px_to_emu(y1),
            px_to_emu(x2),
            px_to_emu(y2),
        )

        try:
            if shape.id:
                line.name = f"drawio2pptx:line:{shape.id}"
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to set line name: {e}")

        stroke_color = shape.style.stroke if shape.style.stroke else RGBColor(0, 0, 0)
        try:
            line.line.fill.solid()
            self._set_edge_stroke_color_xml(line, stroke_color)
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to set line stroke color: {e}")

        if shape.style.stroke_width > 0:
            try:
                line.line.width = px_to_pt(shape.style.stroke_width)
            except Exception as e:
                if self.logger:
                    self.logger.debug(f"Failed to set line width: {e}")

        try:
            if shape.style.has_shadow:
                line.shadow.inherit = True
            else:
                line.shadow.inherit = False
                self._disable_shadow_xml(line)
        except Exception:
            if not shape.style.has_shadow:
                self._disable_shadow_xml(line)

        return line

    def _add_swimlane_header_divider(
        self,
        slide,
        shape: ShapeElement,
        stroke_color: RGBColor,
        stroke_width_px: float,
    ):
        """Draw the swimlane header divider line."""
        start_size = float(getattr(shape.style, "swimlane_start_size", 0.0) or 0.0)
        if start_size <= 0:
            return None

        is_horizontal = bool(getattr(shape.style, "swimlane_horizontal", False))
        if is_horizontal:
            # Header on top: horizontal divider at y + startSize
            x1 = shape.x
            y1 = shape.y + start_size
            x2 = shape.x + shape.w
            y2 = y1
        else:
            # Header on left: vertical divider at x + startSize
            x1 = shape.x + start_size
            y1 = shape.y
            x2 = x1
            y2 = shape.y + shape.h

        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            px_to_emu(x1),
            px_to_emu(y1),
            px_to_emu(x2),
            px_to_emu(y2),
        )
        try:
            if shape.id:
                line.name = f"drawio2pptx:swimlane-divider:{shape.id}"
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to set swimlane divider name: {e}")

        try:
            line.line.fill.solid()
            self._set_stroke_color_xml(line, stroke_color)
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to set swimlane divider color: {e}")

        if stroke_width_px and stroke_width_px > 0:
            try:
                line.line.width = px_to_pt(stroke_width_px)
            except Exception as e:
                if self.logger:
                    self.logger.debug(f"Failed to set swimlane divider width: {e}")

        # Shadow follows swimlane shape
        try:
            if getattr(shape.style, "has_shadow", False):
                line.shadow.inherit = True
            else:
                line.shadow.inherit = False
                self._disable_shadow_xml(line)
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to set swimlane divider shadow: {e}")

        return line
    
    def _add_connector(self, slide, connector: ConnectorElement):
        """Add connector as a single polyline shape."""
        if not connector.points or len(connector.points) < 2:
            return None

        # Simplify nearly straight polylines to avoid jitter in Freeform.
        points_px = self._simplify_polyline_points(connector.points, tol_px=0.5)
        if len(points_px) >= 2 and self._is_almost_straight(points_px, tol_px=0.5):
            return self._add_straight_connector(slide, connector, points_px)

        # Arrow settings may affect geometry (open-oval marker needs endpoint trimming).
        start_arrow = connector.style.arrow_start
        end_arrow = connector.style.arrow_end

        # PowerPoint line-end types cannot represent "open" (unfilled) oval markers.
        # When draw.io specifies startFill/endFill=0 with oval, emulate it by overlaying a small
        # ellipse outline at the endpoint, and suppress the line-end arrow.
        add_open_oval_start = self._should_emulate_open_oval_marker(start_arrow, connector.style.arrow_start_fill)
        add_open_oval_end = self._should_emulate_open_oval_marker(end_arrow, connector.style.arrow_end_fill)
        effective_start_arrow = None if add_open_oval_start else start_arrow
        effective_end_arrow = None if add_open_oval_end else end_arrow

        # Keep marker centers at the original endpoints, but trim the line so it stops at the marker boundary.
        start_marker_center_px = connector.points[0]
        end_marker_center_px = connector.points[-1]
        line_points_px = list(points_px)
        if add_open_oval_start or add_open_oval_end:
            try:
                start_trim = (
                    self._open_oval_trim_radius_px(
                        stroke_width_px=connector.style.stroke_width,
                        arrow_size_px=connector.style.arrow_start_size_px,
                    )
                    if add_open_oval_start
                    else 0.0
                )
                end_trim = (
                    self._open_oval_trim_radius_px(
                        stroke_width_px=connector.style.stroke_width,
                        arrow_size_px=connector.style.arrow_end_size_px,
                    )
                    if add_open_oval_end
                    else 0.0
                )
                line_points_px = self._trim_polyline_endpoints_px(line_points_px, start_trim, end_trim)
            except Exception:
                line_points_px = list(connector.points)
        
        # For normal lines (straight), use FreeformBuilder
        # Convert points to EMU
        points_emu = [(px_to_emu(x), px_to_emu(y)) for x, y in line_points_px]
        
        try:
            if len(points_emu) < 2:
                return None
            
            # Create polyline with FreeformBuilder
            x0, y0 = points_emu[0]
            builder = slide.shapes.build_freeform(x0, y0)
            
            if len(points_emu) > 1:
                remaining_points = points_emu[1:]
                builder.add_line_segments(remaining_points, close=False)
            
            line_shape = builder.convert_to_shape()
        except Exception:
            return None

        # Debug-friendly name (not visible in normal slideshow mode).
        try:
            if connector.id:
                line_shape.name = f"drawio2pptx:connector:{connector.id}"
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to set connector name: {e}")
        
        # Set stroke
        # Use black as default if stroke is None
        stroke_color = connector.style.stroke if connector.style.stroke else RGBColor(0, 0, 0)
        try:
            line_shape.line.fill.solid()
            self._set_edge_stroke_color_xml(line_shape, stroke_color)
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to set connector stroke color: {e}")
        
        if connector.style.stroke_width > 0:
            try:
                line_shape.line.width = px_to_pt(connector.style.stroke_width)
            except Exception as e:
                if self.logger:
                    self.logger.debug(f"Failed to set connector stroke width: {e}")
        
        # Disable fill
        try:
            line_shape.fill.background()
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to disable connector fill: {e}")
        
        # Set dash pattern
        if connector.style.dash:
            try:
                self._set_dash_pattern_xml(line_shape, connector.style.dash)
            except Exception as e:
                if self.logger:
                    self.logger.debug(f"Failed to set connector dash pattern: {e}")
        
        if effective_start_arrow or effective_end_arrow:
            self._set_arrow_heads_xml(
                line_shape,
                effective_start_arrow,
                effective_end_arrow,
                connector.style.arrow_start_fill,
                connector.style.arrow_end_fill,
                connector.style.stroke,
                connector.style.arrow_start_size_px,
                connector.style.arrow_end_size_px,
            )

        # Overlay open-oval markers after the connector so it sits on top of the line geometry.
        marker_configs = {
            "start": (add_open_oval_start, start_marker_center_px, start_arrow, connector.style.arrow_start_size_px),
            "end": (add_open_oval_end, end_marker_center_px, end_arrow, connector.style.arrow_end_size_px),
        }
        for position, (should_add, center_px, arrow_name, arrow_size_px) in marker_configs.items():
            if should_add:
                x, y = center_px
                self._add_open_oval_marker(
                    slide=slide,
                    x_px=x,
                    y_px=y,
                    stroke_color=connector.style.stroke,
                    stroke_width_px=connector.style.stroke_width,
                    arrow_name=arrow_name,
                    arrow_size_px=arrow_size_px,
                    marker_name=f"drawio2pptx:marker:open-oval:{connector.id}:{position}",
                )
        
        # Shadow settings (similar to legacy: disable shadow when has_shadow is False)
        try:
            if connector.style.has_shadow:
                # Enable shadow (inherit from theme)
                line_shape.shadow.inherit = True
            else:
                # Disable shadow
                line_shape.shadow.inherit = False
                # Disable shadow via XML (similar to legacy _disable_shadow_xml)
                self._disable_shadow_xml(line_shape)
        except Exception:
            if not connector.style.has_shadow:
                self._disable_shadow_xml(line_shape)
        
        return line_shape

    def _add_straight_connector(self, slide, connector: ConnectorElement, points_px: List[Tuple[float, float]]):
        """Add a single straight connector (line) between endpoints."""
        if len(points_px) < 2:
            return None

        # Arrow settings may affect geometry (open-oval marker needs endpoint trimming).
        start_arrow = connector.style.arrow_start
        end_arrow = connector.style.arrow_end

        add_open_oval_start = self._should_emulate_open_oval_marker(start_arrow, connector.style.arrow_start_fill)
        add_open_oval_end = self._should_emulate_open_oval_marker(end_arrow, connector.style.arrow_end_fill)
        effective_start_arrow = None if add_open_oval_start else start_arrow
        effective_end_arrow = None if add_open_oval_end else end_arrow

        start_marker_center_px = points_px[0]
        end_marker_center_px = points_px[-1]
        line_points_px = list(points_px)
        if add_open_oval_start or add_open_oval_end:
            try:
                start_trim = (
                    self._open_oval_trim_radius_px(
                        stroke_width_px=connector.style.stroke_width,
                        arrow_size_px=connector.style.arrow_start_size_px,
                    )
                    if add_open_oval_start
                    else 0.0
                )
                end_trim = (
                    self._open_oval_trim_radius_px(
                        stroke_width_px=connector.style.stroke_width,
                        arrow_size_px=connector.style.arrow_end_size_px,
                    )
                    if add_open_oval_end
                    else 0.0
                )
                line_points_px = self._trim_polyline_endpoints_px(line_points_px, start_trim, end_trim)
            except Exception:
                line_points_px = list(points_px)

        if len(line_points_px) < 2:
            return None
        (x1, y1), (x2, y2) = line_points_px[0], line_points_px[-1]
        line_shape = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            px_to_emu(x1),
            px_to_emu(y1),
            px_to_emu(x2),
            px_to_emu(y2),
        )

        try:
            if connector.id:
                line_shape.name = f"drawio2pptx:connector:{connector.id}"
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to set connector name: {e}")

        # Set stroke
        stroke_color = connector.style.stroke if connector.style.stroke else RGBColor(0, 0, 0)
        try:
            line_shape.line.fill.solid()
            self._set_edge_stroke_color_xml(line_shape, stroke_color)
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to set connector stroke color: {e}")

        if connector.style.stroke_width > 0:
            try:
                line_shape.line.width = px_to_pt(connector.style.stroke_width)
            except Exception as e:
                if self.logger:
                    self.logger.debug(f"Failed to set connector stroke width: {e}")

        # Disable fill
        try:
            line_shape.fill.background()
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to disable connector fill: {e}")

        # Set dash pattern
        if connector.style.dash:
            try:
                self._set_dash_pattern_xml(line_shape, connector.style.dash)
            except Exception as e:
                if self.logger:
                    self.logger.debug(f"Failed to set connector dash pattern: {e}")

        if effective_start_arrow or effective_end_arrow:
            self._set_arrow_heads_xml(
                line_shape,
                effective_start_arrow,
                effective_end_arrow,
                connector.style.arrow_start_fill,
                connector.style.arrow_end_fill,
                connector.style.stroke,
                connector.style.arrow_start_size_px,
                connector.style.arrow_end_size_px,
            )

        # Overlay open-oval markers after the connector so it sits on top of the line geometry.
        marker_configs = {
            "start": (add_open_oval_start, start_marker_center_px, start_arrow, connector.style.arrow_start_size_px),
            "end": (add_open_oval_end, end_marker_center_px, end_arrow, connector.style.arrow_end_size_px),
        }
        for position, (should_add, center_px, arrow_name, arrow_size_px) in marker_configs.items():
            if should_add:
                x, y = center_px
                self._add_open_oval_marker(
                    slide=slide,
                    x_px=x,
                    y_px=y,
                    stroke_color=connector.style.stroke,
                    stroke_width_px=connector.style.stroke_width,
                    arrow_name=arrow_name,
                    arrow_size_px=arrow_size_px,
                    marker_name=f"drawio2pptx:marker:open-oval:{connector.id}:{position}",
                )

        # Shadow settings
        try:
            if connector.style.has_shadow:
                line_shape.shadow.inherit = True
            else:
                line_shape.shadow.inherit = False
                self._disable_shadow_xml(line_shape)
        except Exception:
            if not connector.style.has_shadow:
                self._disable_shadow_xml(line_shape)

        return line_shape

    @staticmethod
    def _simplify_polyline_points(points: List[Tuple[float, float]], tol_px: float = 0.5) -> List[Tuple[float, float]]:
        """Drop nearly collinear points to stabilize Freeform rendering."""
        if len(points) <= 2:
            return list(points)

        def _dist_point_to_line(px, py, ax, ay, bx, by):
            dx = bx - ax
            dy = by - ay
            denom = dx * dx + dy * dy
            if denom <= 1e-9:
                return ((px - ax) ** 2 + (py - ay) ** 2) ** 0.5
            t = ((px - ax) * dx + (py - ay) * dy) / denom
            t = 0.0 if t < 0.0 else 1.0 if t > 1.0 else t
            cx = ax + t * dx
            cy = ay + t * dy
            return ((px - cx) ** 2 + (py - cy) ** 2) ** 0.5

        simplified = [points[0]]
        for i in range(1, len(points) - 1):
            ax, ay = simplified[-1]
            bx, by = points[i + 1]
            px, py = points[i]
            if _dist_point_to_line(px, py, ax, ay, bx, by) > tol_px:
                simplified.append(points[i])
        simplified.append(points[-1])
        return simplified

    @staticmethod
    def _is_almost_straight(points: List[Tuple[float, float]], tol_px: float = 0.5) -> bool:
        """Return True if all points are close to the line from start to end."""
        if len(points) <= 2:
            return True
        (ax, ay) = points[0]
        (bx, by) = points[-1]
        dx = bx - ax
        dy = by - ay
        denom = dx * dx + dy * dy
        if denom <= 1e-9:
            return True
        for (px, py) in points[1:-1]:
            t = ((px - ax) * dx + (py - ay) * dy) / denom
            cx = ax + t * dx
            cy = ay + t * dy
            dist = ((px - cx) ** 2 + (py - cy) ** 2) ** 0.5
            if dist > tol_px:
                return False
        return True

    def _add_text(self, slide, text_element: TextElement):
        """Add standalone text element."""
        if text_element.w <= 0 or text_element.h <= 0:
            return None

        left = px_to_emu(text_element.x)
        top = px_to_emu(text_element.y)
        width = px_to_emu(text_element.w)
        height = px_to_emu(text_element.h)

        tb = slide.shapes.add_textbox(left, top, width, height)
        try:
            if text_element.id:
                tb.name = f"drawio2pptx:text:{text_element.id}"
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to set text box name: {e}")

        # Make the text box background transparent.
        try:
            tb.fill.background()
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to set text box fill background: {e}")
        try:
            tb.line.fill.background()
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to set text box line background: {e}")

        if text_element.text:
            self._set_text_frame(
                tb.text_frame,
                text_element.text,
                default_highlight_color=text_element.style.label_background_color,
                word_wrap=text_element.style.word_wrap,
            )

        return tb

    def _add_bpmn_parallel_gateway_symbol(self, slide, shape: ShapeElement):
        """Add plus sign overlay for BPMN parallel gateway using lines (not text)."""
        symbol_size = min(shape.w, shape.h) * 1.0
        center_x = shape.x + shape.w / 2.0
        center_y = shape.y + shape.h / 2.0

        stroke_color = shape.style.stroke if shape.style.stroke else RGBColor(0, 0, 0)
        line_width_pt = 2.0
        half_length = symbol_size * 0.25

        h_line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            px_to_emu(center_x - half_length),
            px_to_emu(center_y),
            px_to_emu(center_x + half_length),
            px_to_emu(center_y),
        )
        v_line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            px_to_emu(center_x),
            px_to_emu(center_y - half_length),
            px_to_emu(center_x),
            px_to_emu(center_y + half_length),
        )

        try:
            h_line.line.fill.solid()
            self._set_edge_stroke_color_xml(h_line, stroke_color)
            h_line.line.width = Pt(line_width_pt)
            self._remove_arrowheads_xml(h_line)
            h_line.shadow.inherit = False
            self._disable_shadow_xml(h_line)
            self._remove_effect_ref_xml(h_line)
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to set horizontal line properties: {e}")

        try:
            if shape.id:
                h_line.name = f"drawio2pptx:bpmn-symbol-h:{shape.id}"
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to set horizontal line name: {e}")

        try:
            v_line.line.fill.solid()
            self._set_edge_stroke_color_xml(v_line, stroke_color)
            v_line.line.width = Pt(line_width_pt)
            self._remove_arrowheads_xml(v_line)
            v_line.shadow.inherit = False
            self._disable_shadow_xml(v_line)
            self._remove_effect_ref_xml(v_line)
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to set vertical line properties: {e}")

        try:
            if shape.id:
                v_line.name = f"drawio2pptx:bpmn-symbol-v:{shape.id}"
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to set vertical line name: {e}")
    
    def _add_orthogonal_connector(self, slide, connector: ConnectorElement):
        """Add polyline as straight connectors for each segment"""
        if not connector.points or len(connector.points) < 2:
            return None

        # Open-oval marker emulation needs endpoint trimming on the underlying polyline.
        add_open_oval_start = self._should_emulate_open_oval_marker(connector.style.arrow_start, connector.style.arrow_start_fill)
        add_open_oval_end = self._should_emulate_open_oval_marker(connector.style.arrow_end, connector.style.arrow_end_fill)

        start_marker_center_px = connector.points[0]
        end_marker_center_px = connector.points[-1]
        points_for_segments = list(connector.points)
        if add_open_oval_start or add_open_oval_end:
            try:
                start_trim = (
                    self._open_oval_trim_radius_px(
                        stroke_width_px=connector.style.stroke_width,
                        arrow_size_px=connector.style.arrow_start_size_px,
                    )
                    if add_open_oval_start
                    else 0.0
                )
                end_trim = (
                    self._open_oval_trim_radius_px(
                        stroke_width_px=connector.style.stroke_width,
                        arrow_size_px=connector.style.arrow_end_size_px,
                    )
                    if add_open_oval_end
                    else 0.0
                )
                points_for_segments = self._trim_polyline_endpoints_px(points_for_segments, start_trim, end_trim)
            except Exception:
                points_for_segments = list(connector.points)
        
        # Split polyline into segments
        segments = split_polyline_to_segments(points_for_segments)
        if not segments:
            return None
        
        created_shapes = []
        
        for idx, ((x1, y1), (x2, y2)) in enumerate(segments):
            try:
                x1_emu, y1_emu = px_to_emu(x1), px_to_emu(y1)
                x2_emu, y2_emu = px_to_emu(x2), px_to_emu(y2)
                
                # Create each segment as a straight connector
                conn_shape = slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT,
                    x1_emu, y1_emu, x2_emu, y2_emu
                )

                # Debug-friendly name (not visible in normal slideshow mode).
                try:
                    if connector.id:
                        conn_shape.name = f"drawio2pptx:connector:{connector.id}:seg:{idx}"
                except Exception as e:
                    if self.logger:
                        self.logger.debug(f"Failed to set connector segment name: {e}")
                
                # Set stroke
                # Use black as default if stroke is None
                stroke_color = connector.style.stroke if connector.style.stroke else RGBColor(0, 0, 0)
                try:
                    conn_shape.line.fill.solid()
                    self._set_edge_stroke_color_xml(conn_shape, stroke_color)
                except Exception as e:
                    if self.logger:
                        self.logger.debug(f"Failed to set connector segment stroke color: {e}")
                
                if connector.style.stroke_width > 0:
                    try:
                        conn_shape.line.width = px_to_pt(connector.style.stroke_width)
                    except Exception as e:
                        if self.logger:
                            self.logger.debug(f"Failed to set connector segment stroke width: {e}")
                
                # Disable fill
                try:
                    conn_shape.fill.background()
                except Exception as e:
                    if self.logger:
                        self.logger.debug(f"Failed to disable connector segment fill: {e}")
                
                # Set dash pattern
                if connector.style.dash:
                    try:
                        self._set_dash_pattern_xml(conn_shape, connector.style.dash)
                    except Exception as e:
                        if self.logger:
                            self.logger.debug(f"Failed to set connector segment dash pattern: {e}")
                
                # Arrows: start arrow on first segment, end arrow on last segment
                is_first_segment = idx == 0
                is_last_segment = idx == len(segments) - 1
                start_arrow = connector.style.arrow_start if is_first_segment else None
                end_arrow = connector.style.arrow_end if is_last_segment else None
                
                if start_arrow or end_arrow:
                    # Suppress oval line-end when it is unfilled, and emulate with an overlay marker later.
                    add_open_oval_start = self._should_emulate_open_oval_marker(
                        start_arrow, connector.style.arrow_start_fill if is_first_segment else False
                    )
                    add_open_oval_end = self._should_emulate_open_oval_marker(
                        end_arrow, connector.style.arrow_end_fill if is_last_segment else False
                    )
                    effective_start_arrow = None if add_open_oval_start else start_arrow
                    effective_end_arrow = None if add_open_oval_end else end_arrow

                    self._set_arrow_heads_xml(
                        conn_shape,
                        effective_start_arrow,
                        effective_end_arrow,
                        connector.style.arrow_start_fill if is_first_segment else False,
                        connector.style.arrow_end_fill if is_last_segment else False,
                        connector.style.stroke,
                        connector.style.arrow_start_size_px if is_first_segment else None,
                        connector.style.arrow_end_size_px if is_last_segment else None,
                    )
                
                # Shadow
                try:
                    if connector.style.has_shadow:
                        conn_shape.shadow.inherit = True
                    else:
                        conn_shape.shadow.inherit = False
                        self._disable_shadow_xml(conn_shape)
                except Exception:
                    if not connector.style.has_shadow:
                        self._disable_shadow_xml(conn_shape)
                
                created_shapes.append(conn_shape)
            except Exception as e:
                if self.logger:
                    self.logger.warning(f"Failed to create connector segment {idx}: {e}")
                continue
        
        if not created_shapes:
            return None

        # Add open-oval markers once, after all segments are created (ensures the marker is on top
        # of the line geometry for this connector).
        marker_configs = {
            "start": (add_open_oval_start, start_marker_center_px, connector.style.arrow_start, connector.style.arrow_start_size_px),
            "end": (add_open_oval_end, end_marker_center_px, connector.style.arrow_end, connector.style.arrow_end_size_px),
        }
        for position, (should_add, center_px, arrow_name, arrow_size_px) in marker_configs.items():
            if should_add:
                x, y = center_px
                self._add_open_oval_marker(
                    slide=slide,
                    x_px=x,
                    y_px=y,
                    stroke_color=connector.style.stroke,
                    stroke_width_px=connector.style.stroke_width,
                    arrow_name=arrow_name,
                    arrow_size_px=arrow_size_px,
                    marker_name=f"drawio2pptx:marker:open-oval:{connector.id}:{position}",
                )
        
        return created_shapes[0]

    @staticmethod
    def _should_emulate_open_oval_marker(arrow_name: Optional[str], fill: bool) -> bool:
        """Return True when we should emulate an open oval marker with an overlay shape.

        PowerPoint's line-end (a:headEnd/a:tailEnd) doesn't support an "unfilled oval" variant.
        draw.io represents it with startArrow/endArrow=oval and startFill/endFill=0.
        """
        if not arrow_name:
            return False
        try:
            return arrow_name.strip().lower() == "oval" and (fill is False)
        except Exception:
            return False

    def _add_open_oval_marker(
        self,
        slide,
        x_px: float,
        y_px: float,
        stroke_color: Optional[RGBColor],
        stroke_width_px: float,
        arrow_name: Optional[str],
        arrow_size_px: Optional[float],
        marker_name: str,
    ):
        """Add a small ellipse outline at the given point to emulate an unfilled oval line-end.

        The connector geometry is trimmed so the line stops at the marker boundary, so we keep the marker
        unfilled (noFill) rather than using a white "mask" fill (which breaks on non-white backgrounds).
        """
        try:
            # Prefer draw.io marker size (startSize/endSize) when present.
            # In mxGraph style, marker size values are expressed in screen px.
            # For an oval marker, treating the size as an approximate diameter matches draw.io more closely.
            #
            # Important: draw.io can omit startSize/endSize from the style string. In that case,
            # mxGraph defaults the marker size to 6. Use that default here to avoid oversizing.
            try:
                effective_size_px = float(arrow_size_px) if arrow_size_px is not None else 6.0
            except Exception:
                effective_size_px = 6.0

            base_d = max(effective_size_px, 1.0)
            d_px = max(base_d, 6.0 + float(stroke_width_px) * 1.25)

            left = px_to_emu(x_px - d_px / 2.0)
            top = px_to_emu(y_px - d_px / 2.0)
            width = px_to_emu(d_px)
            height = px_to_emu(d_px)

            marker = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
            try:
                marker.name = marker_name
            except Exception as e:
                if self.logger:
                    self.logger.debug(f"Failed to set marker name: {e}")

            # No fill (hollow marker).
            try:
                marker.fill.background()
            except Exception as e:
                if self.logger:
                    self.logger.debug(f"Failed to set marker fill background: {e}")
            try:
                self._set_no_fill_xml(marker)
            except Exception as e:
                if self.logger:
                    self.logger.debug(f"Failed to set marker no fill XML: {e}")

            # Stroke matches connector.
            # Use black as default if stroke_color is None
            effective_stroke_color = stroke_color if stroke_color is not None else RGBColor(0, 0, 0)
            try:
                marker.line.fill.solid()
                self._set_stroke_color_xml(marker, effective_stroke_color)
            except Exception as e:
                if self.logger:
                    self.logger.debug(f"Failed to set marker stroke color: {e}")

            try:
                if stroke_width_px and stroke_width_px > 0:
                    marker.line.width = px_to_pt(stroke_width_px)
            except Exception as e:
                if self.logger:
                    self.logger.debug(f"Failed to set marker stroke width: {e}")

            # Markers should not add extra visual effects.
            try:
                marker.shadow.inherit = False
                self._disable_shadow_xml(marker)
            except Exception as e:
                if self.logger:
                    self.logger.debug(f"Failed to disable marker shadow: {e}")
        except Exception:
            return None

    @staticmethod
    def _open_oval_marker_diameter_px(stroke_width_px: float, arrow_size_px: Optional[float]) -> float:
        """Return the diameter (px) for the emulated open-oval marker."""
        try:
            effective_size_px = float(arrow_size_px) if arrow_size_px is not None else 6.0
        except Exception:
            effective_size_px = 6.0
        base_d = max(effective_size_px, 1.0)
        try:
            return max(base_d, 6.0 + float(stroke_width_px) * 1.25)
        except Exception:
            return max(base_d, 6.0)

    @classmethod
    def _open_oval_trim_radius_px(cls, stroke_width_px: float, arrow_size_px: Optional[float]) -> float:
        """Return trim distance (px) so the connector line stops at the open-oval boundary."""
        d_px = cls._open_oval_marker_diameter_px(stroke_width_px=stroke_width_px, arrow_size_px=arrow_size_px)
        try:
            return max(d_px / 2.0 - float(stroke_width_px) / 2.0, 0.0)
        except Exception:
            return max(d_px / 2.0, 0.0)

    @staticmethod
    def _trim_polyline_endpoints_px(
        points: List[Tuple[float, float]],
        start_trim_px: float,
        end_trim_px: float,
    ) -> List[Tuple[float, float]]:
        """Trim polyline endpoints by the requested distances (px)."""
        if not points or len(points) < 2:
            return points

        def dist(a: Tuple[float, float], b: Tuple[float, float]) -> float:
            dx = b[0] - a[0]
            dy = b[1] - a[1]
            return (dx * dx + dy * dy) ** 0.5

        def lerp(a: Tuple[float, float], b: Tuple[float, float], t: float) -> Tuple[float, float]:
            return (a[0] + (b[0] - a[0]) * t, a[1] + (b[1] - a[1]) * t)

        pts = list(points)

        # Trim start
        remaining = max(float(start_trim_px or 0.0), 0.0)
        while remaining > 1e-6 and len(pts) >= 2:
            seg_len = dist(pts[0], pts[1])
            if seg_len < 1e-6:
                pts.pop(0)
                continue
            if seg_len <= remaining and len(pts) > 2:
                remaining -= seg_len
                pts.pop(0)
                continue
            t = min(remaining / seg_len, 0.9)
            pts[0] = lerp(pts[0], pts[1], t)
            break

        # Trim end
        remaining = max(float(end_trim_px or 0.0), 0.0)
        while remaining > 1e-6 and len(pts) >= 2:
            seg_len = dist(pts[-1], pts[-2])
            if seg_len < 1e-6:
                pts.pop(-1)
                continue
            if seg_len <= remaining and len(pts) > 2:
                remaining -= seg_len
                pts.pop(-1)
                continue
            t = min(remaining / seg_len, 0.9)
            pts[-1] = lerp(pts[-1], pts[-2], t)
            break

        if len(pts) < 2:
            return list(points)
        return pts

    def _set_no_fill_xml(self, shape):
        """Force <a:noFill/> on the shape fill (spPr) via XML."""
        try:
            if not hasattr(shape, "_element"):
                return

            shape_element = shape._element

            # spPr is often a direct child in python-pptx shapes.
            sp_pr = None
            for child in shape_element:
                if child.tag.endswith("}spPr") or "spPr" in child.tag:
                    sp_pr = child
                    break

            if sp_pr is None:
                sp_pr = shape_element.find(".//a:spPr", namespaces=NSMAP_DRAWINGML)
            if sp_pr is None:
                return

            for tag in ("noFill", "solidFill", "gradFill", "pattFill", "blipFill"):
                for elem in sp_pr.findall(f".//a:{tag}", namespaces=NSMAP_DRAWINGML):
                    try:
                        sp_pr.remove(elem)
                    except Exception as e:
                        if self.logger:
                            self.logger.debug(f"Failed to remove fill element {tag}: {e}")

            ET.SubElement(sp_pr, _a("noFill"))
        except Exception:
            return None

    def _set_no_line_xml(self, shape):
        """Force <a:noFill/> on the line (ln) via XML."""
        try:
            if not hasattr(shape, "_element"):
                return

            shape_element = shape._element
            sp_pr = None
            for child in shape_element:
                if child.tag.endswith("}spPr") or "spPr" in child.tag:
                    sp_pr = child
                    break
            if sp_pr is None:
                sp_pr = shape_element.find(".//a:spPr", namespaces=NSMAP_DRAWINGML)
            if sp_pr is None:
                return

            ln_element = sp_pr.find(".//a:ln", namespaces=NSMAP_DRAWINGML)
            if ln_element is None:
                ln_element = ET.SubElement(sp_pr, _a("ln"))

            for tag in ("noFill", "solidFill", "gradFill", "pattFill", "blipFill"):
                for elem in ln_element.findall(f".//a:{tag}", namespaces=NSMAP_DRAWINGML):
                    try:
                        ln_element.remove(elem)
                    except Exception as e:
                        if self.logger:
                            self.logger.debug(f"Failed to remove line element {tag}: {e}")

            ET.SubElement(ln_element, _a("noFill"))
        except Exception:
            return None
    
    def _set_text_frame(
        self,
        text_frame,
        paragraphs: List[TextParagraph],
        default_highlight_color: Optional[RGBColor] = None,
        word_wrap: bool = True,
        margin_overrides_px: Optional[Tuple[float, float, float, float]] = None,
        text_direction: Optional[str] = None,
    ):
        """Set text frame"""
        if not paragraphs:
            return
        
        # Configure text frame
        text_frame.word_wrap = word_wrap
        text_frame.auto_size = None
        
        # Optional text direction (used for vertical swimlane headers)
        if text_direction is not None:
            try:
                self._set_text_direction_xml(text_frame, text_direction)
            except Exception as e:
                if self.logger:
                    self.logger.debug(f"Failed to set text direction: {e}")

        # Get padding from first paragraph
        first_para = paragraphs[0]
        if margin_overrides_px is not None:
            top_px, left_px, bottom_px, right_px = margin_overrides_px
            text_frame.margin_top = px_to_emu(top_px or 0)
            text_frame.margin_left = px_to_emu(left_px or 0)
            text_frame.margin_bottom = px_to_emu(bottom_px or 0)
            text_frame.margin_right = px_to_emu(right_px or 0)
        else:
            top_px = first_para.spacing_top or 0
            left_px = first_para.spacing_left or 0
            bottom_px = first_para.spacing_bottom or 0
            right_px = first_para.spacing_right or 0
            text_frame.margin_top = px_to_emu(top_px)
            text_frame.margin_left = px_to_emu(left_px)
            text_frame.margin_bottom = px_to_emu(bottom_px)
            text_frame.margin_right = px_to_emu(right_px)
        
        margin_px = (top_px or 0, left_px or 0, bottom_px or 0, right_px or 0)
        
        # Set vertical anchor (similar to legacy: default is middle)
        VERTICAL_ALIGN_MAP = {
            "top": (MSO_ANCHOR.TOP, 't'),
            "middle": (MSO_ANCHOR.MIDDLE, 'ctr'),
            "bottom": (MSO_ANCHOR.BOTTOM, 'b'),
        }
        vertical_align = first_para.vertical_align or "middle"
        saved_vertical_anchor, anchor_value = VERTICAL_ALIGN_MAP.get(
            vertical_align, (MSO_ANCHOR.MIDDLE, 'ctr')
        )
        
        # Clear existing paragraphs
        text_frame.clear()
        
        # Remove empty paragraphs (similar to legacy: to ensure vertical_anchor works correctly)
        while len(text_frame.paragraphs) > 0 and not text_frame.paragraphs[0].text:
            try:
                first_p = text_frame.paragraphs[0]
                first_p._element.getparent().remove(first_p._element)
            except Exception:
                break
        
        # Add paragraphs
        for para_data in paragraphs:
            p = text_frame.add_paragraph()
            p.space_before = Pt(0)
            p.space_after = Pt(0)
            p.line_spacing = 1.0
            
            # Horizontal alignment
            HORIZONTAL_ALIGN_MAP = {
                "left": PP_ALIGN.LEFT,
                "center": PP_ALIGN.CENTER,
                "right": PP_ALIGN.RIGHT,
            }
            p.alignment = HORIZONTAL_ALIGN_MAP.get(
                para_data.align or "center", PP_ALIGN.CENTER
            )
            
            # Add runs
            for run_data in para_data.runs:
                run = p.add_run()
                run.text = run_data.text
                
                # Font family
                # If run_data.font_family is None or empty string, use draw.io's default font (Helvetica)
                # This ensures the same font is used in draw.io and PowerPoint
                effective_font_family = run_data.font_family or DRAWIO_DEFAULT_FONT_FAMILY
                replaced_font = replace_font(effective_font_family, config=self.config)
                if replaced_font:
                    try:
                        run.font.name = replaced_font
                    except Exception as e:
                        # Log if font setting fails (for debugging)
                        if self.logger:
                            self.logger.warning(f"Failed to set font '{replaced_font}': {e}")
                        # Continue processing even if font setting fails
                        pass
                else:
                    # If replace_font returns None, try the original font name
                    try:
                        run.font.name = effective_font_family
                    except Exception as e:
                        if self.logger:
                            self.logger.warning(f"Failed to set font '{effective_font_family}': {e}")
                        pass
                
                # Font size
                # Scale draw.io font size for PowerPoint (considering coordinate transformation)
                if run_data.font_size:
                    scaled_font_size = scale_font_size_for_pptx(run_data.font_size)
                    run.font.size = Pt(scaled_font_size)
                else:
                    # Also scale default font size (12pt)
                    scaled_default_size = scale_font_size_for_pptx(12.0)
                    run.font.size = Pt(scaled_default_size)
                
                # Font style (explicitly set)
                run.font.bold = run_data.bold
                run.font.italic = run_data.italic
                run.font.underline = run_data.underline
                
                # Font color (similar to legacy: always set via XML when font color is set)
                if run_data.font_color:
                    self._set_font_color_xml(run, run_data.font_color)
                else:
                    # Set default black color even when font color is not set (legacy behavior)
                    default_black = RGBColor(0, 0, 0)
                    self._set_font_color_xml(run, default_black)

                # draw.io labelBackgroundColor -> PPTX highlight (best effort).
                if default_highlight_color is not None:
                    self._set_highlight_color_xml(run, default_highlight_color)
        
        # Reset vertical_anchor (similar to legacy: set after clear())
        # vertical_anchor is reset after tf.clear(), so need to set again
        text_frame.vertical_anchor = saved_vertical_anchor
        
        # Set vertical anchor via XML (similar to legacy: workaround for python-pptx bug)
        self._set_vertical_anchor_xml(text_frame, anchor_value, word_wrap, margin_px)
        
        # Set spacing of all paragraphs to 0 (similar to legacy: affects vertical alignment)
        for para in text_frame.paragraphs:
            para.space_before = Pt(0)
            para.space_after = Pt(0)
    
    def _set_vertical_anchor_xml(
        self,
        text_frame,
        anchor_value: str,
        word_wrap: bool = True,
        margin_px: Optional[Tuple[float, float, float, float]] = None,
    ):
        """Set vertical anchor via XML"""
        try:
            body_pr = text_frame._element.find(f'.//{_a("bodyPr")}')
            if body_pr is not None:
                body_pr.set('anchor', anchor_value)
                if body_pr.get('anchorCtr') is not None:
                    body_pr.attrib.pop('anchorCtr', None)
                if margin_px is not None:
                    top_px, left_px, bottom_px, right_px = margin_px
                    body_pr.set('tIns', str(int(px_to_emu(top_px or 0))))
                    body_pr.set('lIns', str(int(px_to_emu(left_px or 0))))
                    body_pr.set('bIns', str(int(px_to_emu(bottom_px or 0))))
                    body_pr.set('rIns', str(int(px_to_emu(right_px or 0))))
                # Set wrap attribute based on word_wrap setting
                # 'square' enables wrapping, 'none' disables it
                wrap_value = 'square' if word_wrap else 'none'
                body_pr.set('wrap', wrap_value)
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to set vertical anchor XML: {e}")

    def _set_text_direction_xml(self, text_frame, direction: str):
        """Set text direction via XML (bodyPr@vert)."""
        if not direction:
            return
        try:
            body_pr = text_frame._element.find(f'.//{_a("bodyPr")}')
            if body_pr is None:
                return
            body_pr.set('vert', str(direction))
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to set text direction XML: {e}")
    
    def _set_font_color_xml(self, run, font_color: RGBColor):
        """Set font color via XML (similar to legacy _set_font_color)"""
        if not font_color:
            return
        
        try:
            # First try normal method to set
            try:
                run.font.color.rgb = font_color
            except Exception:
                # If setting fails, we'll use XML method directly
                pass
            
            # Always use XML method to ensure color is set correctly
            # python-pptx's RGB setting doesn't always work reliably
            try:
                run_element = getattr(run, "_element", None) or getattr(run, "_r", None)
                if run_element is not None:
                    
                    # Find rPr (run properties) element
                    r_pr = run_element.find('.//a:rPr', namespaces=NSMAP_DRAWINGML)
                    if r_pr is None:
                        # Create rPr if it doesn't exist
                        r_pr = ET.SubElement(run_element, _a('rPr'))
                    
                    # Find or create solidFill element
                    solid_fill = r_pr.find('.//a:solidFill', namespaces=NSMAP_DRAWINGML)
                    if solid_fill is None:
                        solid_fill = ET.SubElement(r_pr, _a('solidFill'))
                    
                    # Remove existing color elements
                    for color_elem in solid_fill.findall('.//a:srgbClr', namespaces=NSMAP_DRAWINGML):
                        solid_fill.remove(color_elem)
                    for color_elem in solid_fill.findall('.//a:schemeClr', namespaces=NSMAP_DRAWINGML):
                        solid_fill.remove(color_elem)
                    
                    # Add srgbClr element
                    srgb = ET.SubElement(solid_fill, _a('srgbClr'))
                    val = f"{font_color[0]:02X}{font_color[1]:02X}{font_color[2]:02X}"
                    srgb.set('val', val)
            except Exception as e:
                if self.logger:
                    self.logger.debug(f"Failed to set font color XML: {e}")
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to access run element for font color: {e}")

    def _set_highlight_color_xml(self, run, highlight_color: RGBColor) -> None:
        """Set highlight color via XML"""
        if not highlight_color:
            return
        try:
            run_element = getattr(run, "_element", None) or getattr(run, "_r", None)
            if run_element is None:
                return

            r_pr = run_element.find(".//a:rPr", namespaces=NSMAP_DRAWINGML)
            if r_pr is None:
                r_pr = ET.SubElement(run_element, _a("rPr"))

            # Remove existing highlight elements (avoid duplicates).
            for hi in r_pr.findall(".//a:highlight", namespaces=NSMAP_DRAWINGML):
                try:
                    r_pr.remove(hi)
                except Exception as e:
                    if self.logger:
                        self.logger.debug(f"Failed to remove highlight element: {e}")

            # Important: DrawingML elements are order-sensitive in many Office apps.
            # Insert <a:highlight> before font elements (<a:latin>/<a:ea>/<a:cs>) when present,
            # otherwise after fill elements when present, otherwise append.
            hi = ET.Element(_a("highlight"))
            srgb = ET.SubElement(hi, _a("srgbClr"))
            val = f"{highlight_color[0]:02X}{highlight_color[1]:02X}{highlight_color[2]:02X}"
            srgb.set("val", val)

            children = list(r_pr)
            insert_idx = len(children)

            font_tags = {_a("latin"), _a("ea"), _a("cs"), _a("sym")}
            for i, child in enumerate(children):
                if child.tag in font_tags:
                    insert_idx = i
                    break

            if insert_idx == len(children):
                fill_tags = {_a("noFill"), _a("solidFill"), _a("gradFill"), _a("pattFill"), _a("blipFill")}
                last_fill = -1
                for i, child in enumerate(children):
                    if child.tag in fill_tags:
                        last_fill = i
                if last_fill != -1:
                    insert_idx = last_fill + 1

            r_pr.insert(insert_idx, hi)
        except Exception:
            return
    
    def _set_default_fill_xml(self, shape):
        """Set default fill (white)"""
        try:
            if not hasattr(shape, '_element'):
                return
            
            shape_element = shape._element
            # spPr exists as a direct child element
            sp_pr = None
            for child in shape_element:
                if child.tag.endswith('}spPr') or 'spPr' in child.tag:
                    sp_pr = child
                    break
            
            if sp_pr is None:
                return
            
            # Remove existing fill elements (noFill, solidFill, etc.)
            for fill_elem in sp_pr.findall('.//a:noFill', namespaces=NSMAP_DRAWINGML):
                sp_pr.remove(fill_elem)
            for fill_elem in sp_pr.findall('.//a:solidFill', namespaces=NSMAP_DRAWINGML):
                sp_pr.remove(fill_elem)
            for fill_elem in sp_pr.findall('.//a:gradFill', namespaces=NSMAP_DRAWINGML):
                sp_pr.remove(fill_elem)
            
            # Create solidFill element
            solid_fill = ET.SubElement(sp_pr, _a('solidFill'))
            
            # Add srgbClr element with white color (RGB: 255, 255, 255)
            srgb_clr = ET.SubElement(solid_fill, _a('srgbClr'))
            srgb_clr.set('val', 'FFFFFF')
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to set default fill XML: {e}")

    def _set_swimlane_gradient_fill_xml(self, shp, shape: ShapeElement) -> None:
        """
        Set swimlane fill: header = fillColor, body = swimlaneFillColor.
        Uses a multi-stop gradient so only the header area is colored; body stays white/transparent.
        """
        def _rgb_to_hex(rgb: RGBColor) -> str:
            return f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"

        try:
            if not hasattr(shp, "_element"):
                return
            shape_element = shp._element
            sp_pr = None
            for child in shape_element:
                if child.tag.endswith("}spPr") or "spPr" in child.tag:
                    sp_pr = child
                    break
            if sp_pr is None:
                sp_pr = shape_element.find(".//a:spPr", namespaces=NSMAP_DRAWINGML)
            if sp_pr is None:
                return

            start_size = float(getattr(shape.style, "swimlane_start_size", 0) or 0.0)
            is_horizontal = bool(getattr(shape.style, "swimlane_horizontal", True))
            header_color = shape.style.fill
            body_color = getattr(shape.style, "swimlane_fill_color", None)
            if isinstance(header_color, RGBColor):
                header_rgb = header_color
            else:
                header_rgb = RGBColor(0xE0, 0xE0, 0xE0)
            if isinstance(body_color, RGBColor):
                body_rgb = body_color
            elif body_color in ("default", "auto"):
                body_rgb = RGBColor(0xFF, 0xFF, 0xFF)
            else:
                body_rgb = RGBColor(0xFF, 0xFF, 0xFF)

            if is_horizontal and shape.h > 0:
                ratio = min(1.0, max(0.0, start_size / shape.h))
                pos_header = int(ratio * 100000)
                ang = "5400000"  # 90 degrees, top to bottom
            else:
                if shape.w > 0:
                    ratio = min(1.0, max(0.0, start_size / shape.w))
                else:
                    ratio = 0.0
                pos_header = int(ratio * 100000)
                ang = "0"  # 0 degrees, left to right

            # Exactly 2 stops on the divider line (e.g. 50%); distance between them almost 0.
            # Before first stop = header, between = tiny blend, after second = body.
            pos_header = min(max(0, pos_header), 100000)
            pos1 = max(0, pos_header - 1)
            pos2 = min(100000, pos_header + 1)

            for tag in ("noFill", "solidFill", "gradFill", "pattFill", "blipFill"):
                for elem in sp_pr.findall(f".//a:{tag}", namespaces=NSMAP_DRAWINGML):
                    try:
                        sp_pr.remove(elem)
                    except Exception as e:
                        if self.logger:
                            self.logger.debug(f"Failed to remove fill element {tag}: {e}")

            grad_fill = ET.SubElement(sp_pr, _a("gradFill"))
            grad_fill.set("rotWithShape", "1")
            gs_lst = ET.SubElement(grad_fill, _a("gsLst"))

            # Stop 1: just above divider → header (extends to 0%)
            gs1 = ET.SubElement(gs_lst, _a("gs"))
            gs1.set("pos", str(pos1))
            ET.SubElement(gs1, _a("srgbClr")).set("val", _rgb_to_hex(header_rgb))
            # Stop 2: just below divider → body (extends to 100%); distance to pos1 ≈ 0
            gs2 = ET.SubElement(gs_lst, _a("gs"))
            gs2.set("pos", str(pos2))
            ET.SubElement(gs2, _a("srgbClr")).set("val", _rgb_to_hex(body_rgb))

            lin = ET.SubElement(grad_fill, _a("lin"))
            lin.set("ang", ang)
            lin.set("scaled", "1")
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to set swimlane gradient fill: {e}")

    def _set_linear_gradient_fill_xml(
        self,
        shape,
        base_fill: Optional[object],
        gradient_color: Optional[object],
        gradient_direction: Optional[str] = None,
    ) -> None:
        """
        Set a simple 2-stop linear gradient via DrawingML XML (<a:gradFill>).

        Notes:
            - draw.io uses style keys: fillColor (start) and gradientColor (end).
            - When either side is "default", we use white (FFFFFF) for start, light gray (E0E0E0) for end.
            - This is best-effort; PowerPoint themes may render slightly differently than draw.io.
        """

        def _rgb_to_hex(rgb: RGBColor) -> str:
            return f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"

        def _darken(rgb: RGBColor, amount: float = 0.25) -> RGBColor:
            # amount in [0, 1]; 0.25 means 25% closer to black
            r = max(0, int(round(rgb[0] * (1.0 - amount))))
            g = max(0, int(round(rgb[1] * (1.0 - amount))))
            b = max(0, int(round(rgb[2] * (1.0 - amount))))
            return RGBColor(r, g, b)

        def _direction_to_ang(direction: Optional[str]) -> str:
            # OpenXML uses 60000ths of a degree.
            DIRECTION_DEGREE_MAP = {
                "east": 0,
                "right": 0,
                "west": 180,
                "left": 180,
                "north": 270,
                "up": 270,
                "south": 90,
                "down": 90,
            }
            d = (direction or "").strip().lower()
            deg = DIRECTION_DEGREE_MAP.get(d, 90)  # default: vertical (top->bottom)
            return str(int(deg * 60000))

        try:
            if not hasattr(shape, "_element"):
                return

            shape_element = shape._element

            # spPr is often a direct child in python-pptx shapes.
            sp_pr = None
            for child in shape_element:
                if child.tag.endswith("}spPr") or "spPr" in child.tag:
                    sp_pr = child
                    break
            if sp_pr is None:
                sp_pr = shape_element.find(".//a:spPr", namespaces=NSMAP_DRAWINGML)
            if sp_pr is None:
                return

            # Remove existing fill elements to avoid duplicates.
            for tag in ("noFill", "solidFill", "gradFill", "pattFill", "blipFill"):
                for elem in sp_pr.findall(f".//a:{tag}", namespaces=NSMAP_DRAWINGML):
                    try:
                        sp_pr.remove(elem)
                    except Exception as e:
                        if self.logger:
                            self.logger.debug(f"Failed to remove gradient fill element {tag}: {e}")

            grad_fill = ET.SubElement(sp_pr, _a("gradFill"))
            grad_fill.set("rotWithShape", "1")

            gs_lst = ET.SubElement(grad_fill, _a("gsLst"))

            # Determine start (fillColor) and end (gradientColor)
            start_is_default = (base_fill in (None, "default", "auto"))
            end_is_default = (gradient_color in ("default", "auto"))

            start_rgb = base_fill if isinstance(base_fill, RGBColor) else None
            end_rgb = gradient_color if isinstance(gradient_color, RGBColor) else None

            # If gradientColor is "default" but we have an explicit fillColor, derive a darker end.
            if end_is_default and start_rgb is not None:
                end_rgb = _darken(start_rgb, 0.25)
                end_is_default = False

            # Stop 1 (pos=0)
            gs1 = ET.SubElement(gs_lst, _a("gs"))
            gs1.set("pos", "0")
            if start_rgb is not None:
                clr1 = ET.SubElement(gs1, _a("srgbClr"))
                clr1.set("val", _rgb_to_hex(start_rgb))
            else:
                # Default: white
                clr1 = ET.SubElement(gs1, _a("srgbClr"))
                clr1.set("val", "FFFFFF")

            # Stop 2 (pos=100000)
            gs2 = ET.SubElement(gs_lst, _a("gs"))
            gs2.set("pos", "100000")
            if end_rgb is not None:
                clr2 = ET.SubElement(gs2, _a("srgbClr"))
                clr2.set("val", _rgb_to_hex(end_rgb))
            elif end_is_default or start_is_default:
                # Default: slightly darker gray (for gradient effect)
                clr2 = ET.SubElement(gs2, _a("srgbClr"))
                clr2.set("val", "E0E0E0")
            elif start_rgb is not None:
                # Fallback: derive darker from start
                clr2 = ET.SubElement(gs2, _a("srgbClr"))
                clr2.set("val", _rgb_to_hex(_darken(start_rgb, 0.25)))
            else:
                # Default: slightly darker gray (for gradient effect)
                clr2 = ET.SubElement(gs2, _a("srgbClr"))
                clr2.set("val", "E0E0E0")

            lin = ET.SubElement(grad_fill, _a("lin"))
            lin.set("ang", _direction_to_ang(gradient_direction))
            lin.set("scaled", "1")
        except Exception:
            return
    
    def _set_stroke_color_xml(self, shape, stroke_color: RGBColor):
        """Set stroke color via XML"""
        try:
            if not hasattr(shape, '_element'):
                return
            
            shape_element = shape._element
            ln_element = shape_element.find('.//a:ln', namespaces=NSMAP_DRAWINGML)
            if ln_element is None:
                sp_pr = shape_element.find('.//a:spPr', namespaces=NSMAP_DRAWINGML)
                if sp_pr is not None:
                    ln_element = ET.SubElement(sp_pr, _a('ln'))
                else:
                    return
            
            solid_fill = ln_element.find('.//a:solidFill', namespaces=NSMAP_DRAWINGML)
            if solid_fill is None:
                no_fill = ln_element.find('.//a:noFill', namespaces=NSMAP_DRAWINGML)
                if no_fill is not None:
                    ln_element.remove(no_fill)
                solid_fill = ET.SubElement(ln_element, _a('solidFill'))
            
            for color_elem in solid_fill.findall('.//a:srgbClr', namespaces=NSMAP_DRAWINGML):
                solid_fill.remove(color_elem)
            
            srgb = ET.SubElement(solid_fill, _a('srgbClr'))
            val = f"{stroke_color[0]:02X}{stroke_color[1]:02X}{stroke_color[2]:02X}"
            srgb.set('val', val)
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to set stroke color XML: {e}")
    
    def _set_edge_stroke_color_xml(self, shape, stroke_color: RGBColor):
        """Set edge stroke color via XML"""
        self._set_stroke_color_xml(shape, stroke_color)
    
    def _set_dash_pattern_xml(self, shape, dash_pattern: Optional[str]):
        """Set dash pattern via XML
        
        Args:
            shape: PowerPoint shape object
            dash_pattern: draw.io dash pattern name (e.g., "dashed", "dotted", "dashDot")
        """
        try:
            if not hasattr(shape, '_element'):
                return
            
            if not dash_pattern:
                return
            
            # Map draw.io dash pattern to PowerPoint prstDash value
            prst_dash = map_dash_pattern(dash_pattern)
            if not prst_dash or prst_dash == "solid":
                return
            
            shape_element = shape._element
            ln_element = shape_element.find(f'.//a:ln', namespaces=NSMAP_DRAWINGML)
            if ln_element is None:
                sp_pr = shape_element.find('.//a:spPr', namespaces=NSMAP_DRAWINGML)
                if sp_pr is not None:
                    ln_element = ET.SubElement(sp_pr, _a('ln'))
                else:
                    return
            
            # Remove existing prstDash
            for prst_dash_elem in ln_element.findall('.//a:prstDash', namespaces=NSMAP_DRAWINGML):
                ln_element.remove(prst_dash_elem)
            
            # Add new prstDash
            prst_dash_elem = ET.SubElement(ln_element, _a('prstDash'))
            prst_dash_elem.set('val', prst_dash)
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to set dash pattern XML: {e}")
    
    def _set_arrow_heads_xml(
        self,
        shape,
        start_arrow: Optional[str],
        end_arrow: Optional[str],
        start_fill: bool,
        end_fill: bool,
        stroke_color: Optional[RGBColor],
        start_size_px: Optional[float],
        end_size_px: Optional[float],
    ):
        """Set arrows via XML

        Notes:
            DrawingML line-end elements (a:headEnd / a:tailEnd) are empty elements with attributes.
            They must NOT contain fill children (a:solidFill / a:noFill). PowerPoint uses the line
            formatting for arrow color. For "open" arrows, map to the appropriate line-end type.
        """
        try:
            if not hasattr(shape, '_element'):
                return
            
            shape_element = shape._element
            ln_element = shape_element.find(f'.//a:ln', namespaces=NSMAP_DRAWINGML)
            if ln_element is None:
                sp_pr = shape_element.find('.//a:spPr', namespaces=NSMAP_DRAWINGML)
                if sp_pr is not None:
                    ln_element = ET.SubElement(sp_pr, _a('ln'))
                else:
                    return
            
            # Remove existing arrows
            for head_end in ln_element.findall('.//a:headEnd', namespaces=NSMAP_DRAWINGML):
                ln_element.remove(head_end)
            for tail_end in ln_element.findall('.//a:tailEnd', namespaces=NSMAP_DRAWINGML):
                ln_element.remove(tail_end)
            
            # Start arrow (draw.io: source side / line start)
            if start_arrow:
                arrow_info = map_arrow_type_with_size(start_arrow, start_size_px)
                if arrow_info:
                    arrow_type, arrow_w, arrow_len = arrow_info
                    head_end = ET.SubElement(ln_element, _a('headEnd'))
                    head_end.set('type', arrow_type)
                    head_end.set('w', arrow_w)
                    head_end.set('len', arrow_len)
            
            # End arrow (draw.io: target side / line end)
            if end_arrow:
                arrow_info = map_arrow_type_with_size(end_arrow, end_size_px)
                if arrow_info:
                    arrow_type, arrow_w, arrow_len = arrow_info
                    tail_end = ET.SubElement(ln_element, _a('tailEnd'))
                    tail_end.set('type', arrow_type)
                    tail_end.set('w', arrow_w)
                    tail_end.set('len', arrow_len)
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to set arrow XML: {e}")

    def _remove_arrowheads_xml(self, shape):
        """Remove arrowheads from a connector line via XML."""
        try:
            if not hasattr(shape, '_element'):
                return

            shape_element = shape._element
            ln_element = shape_element.find('.//a:ln', namespaces=NSMAP_DRAWINGML)
            if ln_element is None:
                return

            for head_end in ln_element.findall('.//a:headEnd', namespaces=NSMAP_DRAWINGML):
                ln_element.remove(head_end)
            for tail_end in ln_element.findall('.//a:tailEnd', namespaces=NSMAP_DRAWINGML):
                ln_element.remove(tail_end)
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to remove arrowheads XML: {e}")

    def _remove_effect_ref_xml(self, shape):
        """Remove effectRef from connector style to avoid theme shadow."""
        try:
            if not hasattr(shape, '_element'):
                return
            shape_element = shape._element
            for effect_ref in list(shape_element.iter()):
                if effect_ref.tag.endswith('effectRef'):
                    parent = effect_ref.getparent()
                    if parent is not None:
                        parent.remove(effect_ref)
            for style in shape_element.findall('.//p:style', namespaces=NSMAP_BOTH):
                parent = style.getparent()
                if parent is not None:
                    parent.remove(style)
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to remove effectRef XML: {e}")
    
    def _disable_shadow_xml(self, shape):
        """Disable shadow via XML (similar to legacy _disable_shadow_xml)"""
        try:
            if not hasattr(shape, '_element'):
                return
            
            shape_element = shape._element
            
            # Find spPr (shape properties) element
            sp_pr = shape_element.find('.//a:spPr', namespaces=NSMAP_DRAWINGML)
            if sp_pr is None:
                return
            
            # Remove existing effectLst (effects list) which contains shadow
            effect_lst = sp_pr.find('.//a:effectLst', namespaces=NSMAP_DRAWINGML)
            if effect_lst is not None:
                sp_pr.remove(effect_lst)
            
            # Also remove any shadow elements
            for shadow_elem in sp_pr.findall('.//a:outerShdw', namespaces=NSMAP_DRAWINGML):
                sp_pr.remove(shadow_elem)
            for shadow_elem in sp_pr.findall('.//a:innerShdw', namespaces=NSMAP_DRAWINGML):
                sp_pr.remove(shadow_elem)
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to disable shadow XML: {e}")
