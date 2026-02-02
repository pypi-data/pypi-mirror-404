"""
Test module for PowerPoint file analysis functionality
"""
import pytest
from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor


def format_color(color_obj):
    """Convert color object to string"""
    if color_obj is None:
        return "None"
    
    try:
        # RGBColor object itself
        if isinstance(color_obj, RGBColor):
            # RGBColor has tuple-like structure, accessible by index
            if len(color_obj) >= 3:
                r, g, b = color_obj[0], color_obj[1], color_obj[2]
                return f"RGB({r}, {g}, {b}) / #{r:02X}{g:02X}{b:02X}"
            else:
                # Get from string representation
                hex_str = str(color_obj)
                if len(hex_str) == 6:
                    r = int(hex_str[0:2], 16)
                    g = int(hex_str[2:4], 16)
                    b = int(hex_str[4:6], 16)
                    return f"RGB({r}, {g}, {b}) / #{hex_str}"
                return f"RGBColor({hex_str})"
        
        # RGB color case
        if hasattr(color_obj, 'rgb'):
            try:
                rgb = color_obj.rgb
                if rgb is not None and isinstance(rgb, RGBColor):
                    return f"RGB({rgb.r}, {rgb.g}, {rgb.b}) / #{rgb.r:02X}{rgb.g:02X}{rgb.b:02X}"
            except Exception:
                pass
        
        # Theme color case
        if hasattr(color_obj, 'theme_color'):
            try:
                theme_color = color_obj.theme_color
                if theme_color is not None:
                    return f"ThemeColor({theme_color})"
            except Exception:
                pass
        
        # ColorFormat object case, get directly from XML element
        if hasattr(color_obj, '_element'):
            try:
                element = color_obj._element
                from lxml import etree
                nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                srgb = element.find('.//a:srgbClr', namespaces=nsmap)
                if srgb is None:
                    srgb = element.find('.//srgbClr')
                if srgb is not None:
                    val = srgb.get('val', '')
                    if val and len(val) == 6:
                        r = int(val[0:2], 16)
                        g = int(val[2:4], 16)
                        b = int(val[4:6], 16)
                        return f"RGB({r}, {g}, {b}) / #{val}"
                
                scheme = element.find('.//a:schemeClr', namespaces=nsmap)
                if scheme is None:
                    scheme = element.find('.//schemeClr')
                if scheme is not None:
                    val = scheme.get('val', '')
                    return f"ThemeColor({val})"
            except Exception:
                pass
        
        return f"Unknown color format"
    except Exception as e:
        return f"Error formatting color: {e}"


def analyze_shape(shape):
    """Analyze shape color information and return dictionary"""
    result = {
        'shape_type': shape.shape_type,
        'shape_id': shape.shape_id,
        'fill': None,
        'line_color': None,
        'text_frame': None,
        'lock_aspect_ratio': None
    }
    
    # Fill color
    try:
        fill = shape.fill
        if fill.type is None:
            result['fill'] = "None"
        elif hasattr(fill, 'fore_color'):
            fill_color = fill.fore_color
            result['fill'] = format_color(fill_color)
        else:
            result['fill'] = str(fill.type)
    except Exception as e:
        result['fill'] = f"Error: {e}"
    
    # Line color
    try:
        if hasattr(shape, '_element'):
            from lxml import etree
            shape_element = shape._element
            nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
            ln_element = shape_element.find('.//a:ln', namespaces=nsmap)
            if ln_element is not None:
                solid_fill = ln_element.find('.//a:solidFill', namespaces=nsmap)
                if solid_fill is not None:
                    srgb = solid_fill.find('.//a:srgbClr', namespaces=nsmap)
                    if srgb is not None:
                        val = srgb.get('val', '')
                        if val and len(val) == 6:
                            r = int(val[0:2], 16)
                            g = int(val[2:4], 16)
                            b = int(val[4:6], 16)
                            result['line_color'] = f"RGB({r}, {g}, {b}) / #{val}"
    except Exception:
        pass
    
    # Text frame
    try:
        if shape.has_text_frame:
            tf = shape.text_frame
            result['text_frame'] = {
                'paragraph_count': len(tf.paragraphs),
                'vertical_anchor': str(tf.vertical_anchor),
                'margin_top': tf.margin_top,
                'margin_left': tf.margin_left,
                'margin_bottom': tf.margin_bottom,
                'margin_right': tf.margin_right,
            }
    except Exception:
        pass
    
    # Lock aspect ratio
    try:
        if hasattr(shape, '_element'):
            from lxml import etree
            shape_element = shape._element
            nsmap = {
                'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
                'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'
            }
            nv_sp_pr = shape_element.find('.//p:nvSpPr', namespaces=nsmap)
            if nv_sp_pr is not None:
                c_nv_sp_pr = nv_sp_pr.find('.//p:cNvSpPr', namespaces=nsmap)
                if c_nv_sp_pr is not None:
                    sp_locks = c_nv_sp_pr.find('.//a:spLocks', namespaces=nsmap)
                    if sp_locks is not None:
                        no_change_aspect = sp_locks.get('noChangeAspect')
                        if no_change_aspect:
                            result['lock_aspect_ratio'] = no_change_aspect
    except Exception:
        pass
    
    return result


class TestPPTXAnalysis:
    """Test class for PowerPoint file analysis"""
    
    @pytest.fixture
    def sample_pptx_path(self):
        """Path to sample pptx file"""
        # First, look for sample/sample.pptx
        sample_pptx = Path(__file__).parent.parent / "sample" / "sample.pptx"
        if sample_pptx.exists():
            return sample_pptx
        
        # Otherwise, look for test_output.pptx
        test_output = Path(__file__).parent.parent / "test_output.pptx"
        if test_output.exists():
            return test_output
        
        # Otherwise, look for output.pptx
        output = Path(__file__).parent.parent / "output.pptx"
        if output.exists():
            return output
        
        pytest.skip("No sample PPTX file found")
    
    def test_format_color_rgb(self):
        """Test RGB color formatting"""
        rgb = RGBColor(255, 0, 0)
        result = format_color(rgb)
        assert "RGB(255, 0, 0)" in result
        assert "#FF0000" in result or "#ff0000" in result
    
    def test_format_color_none(self):
        """Test None color formatting"""
        assert format_color(None) == "None"
    
    def test_parse_pptx_file(self, sample_pptx_path):
        """Test PowerPoint file parsing"""
        prs = Presentation(str(sample_pptx_path))
        assert prs is not None
        assert len(prs.slides) > 0, "At least one slide must exist"
    
    def test_analyze_shapes(self, sample_pptx_path):
        """Test shape analysis"""
        prs = Presentation(str(sample_pptx_path))
        
        total_shapes = 0
        for slide in prs.slides:
            total_shapes += len(slide.shapes)
            
            for shape in slide.shapes:
                result = analyze_shape(shape)
                
                # Verify basic structure exists
                assert 'shape_type' in result
                assert 'shape_id' in result
                assert 'fill' in result
    
    def test_pptx_has_shapes(self, sample_pptx_path):
        """Verify PowerPoint file contains shapes"""
        prs = Presentation(str(sample_pptx_path))
        
        total_shapes = 0
        for slide in prs.slides:
            total_shapes += len(slide.shapes)
        
        assert total_shapes > 0, "At least one shape must exist"
    
    def test_pptx_writer_integration(self, sample_pptx_path):
        """Integration test with drawio2pptx PPTXWriter"""
        from drawio2pptx.io.pptx_writer import PPTXWriter
        
        prs = Presentation(str(sample_pptx_path))
        assert prs is not None
        assert len(prs.slides) > 0
    
    def test_shape_properties(self, sample_pptx_path):
        """Verify shape properties can be retrieved correctly"""
        prs = Presentation(str(sample_pptx_path))
        
        for slide in prs.slides:
            for shape in slide.shapes:
                result = analyze_shape(shape)
                
                # shape_type and shape_id always exist
                assert result['shape_type'] is not None
                assert result['shape_id'] is not None
                
                # fill is always set (even if None)
                assert result['fill'] is not None




