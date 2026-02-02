"""
Test module for comparing text positions between drawio and PowerPoint files
"""
import pytest
from pathlib import Path
from pptx import Presentation
from lxml import etree as ET
from lxml import html as lxml_html


def parse_drawio_file(path):
    """Parse drawio file and return list of diagrams"""
    tree = ET.parse(path)
    root = tree.getroot()
    
    diagrams = []
    for d in root.findall(".//diagram"):
        inner = (d.text or "").strip()
        if not inner:
            mgm = d.find(".//mxGraphModel")
            if mgm is not None:
                diagrams.append(mgm)
                continue
            continue
        
        if "&lt;" in inner or "&amp;" in inner:
            try:
                wrapped = f"<div>{inner}</div>"
                parsed = lxml_html.fromstring(wrapped)
                inner = parsed.text_content()
            except Exception:
                pass
        
        if "<mxGraphModel" in inner or "<root" in inner or "<mxCell" in inner:
            try:
                parsed = ET.fromstring(inner)
                mgm = None
                if parsed.tag.endswith("mxGraphModel") or parsed.tag == "mxGraphModel":
                    mgm = parsed
                else:
                    mgm = parsed.find(".//mxGraphModel")
                    if mgm is None and parsed.tag.endswith("root"):
                        mgm = parsed
                if mgm is not None:
                    diagrams.append(mgm)
                    continue
            except ET.ParseError:
                pass
        
        mgm_global = root.find(".//mxGraphModel")
        if mgm_global is not None:
            diagrams.append(mgm_global)
        else:
            diagrams.append(root)
    
    if not diagrams:
        mgm_global = root.find(".//mxGraphModel")
        if mgm_global is not None:
            diagrams.append(mgm_global)
        else:
            diagrams.append(root)
    
    return diagrams


def extract_style_value(style_str, key):
    """Extract value for specified key from style string"""
    if not style_str:
        return None
    for part in style_str.split(";"):
        if "=" in part:
            k, v = part.split("=", 1)
            if k.strip() == key:
                return v.strip()
    return None


def extract_drawio_text_properties(drawio_path):
    """Extract text properties from drawio file"""
    diagrams = parse_drawio_file(drawio_path)
    properties = []
    
    for diagram in diagrams:
        cells = diagram.findall(".//mxCell")
        vertex_cells = [c for c in cells if c.attrib.get("vertex") == "1"]
        
        for cell in vertex_cells:
            style = cell.attrib.get("style", "")
            if style:
                props = {
                    'fontSize': extract_style_value(style, "fontSize"),
                    'fontStyle': extract_style_value(style, "fontStyle"),
                    'align': extract_style_value(style, "align"),
                    'verticalAlign': extract_style_value(style, "verticalAlign"),
                    'spacingTop': extract_style_value(style, "spacingTop"),
                    'spacingLeft': extract_style_value(style, "spacingLeft"),
                    'spacingBottom': extract_style_value(style, "spacingBottom"),
                    'spacingRight': extract_style_value(style, "spacingRight"),
                }
                if any(props.values()):  # At least one property is set
                    properties.append(props)
    
    return properties


def extract_pptx_text_properties(pptx_path):
    """Extract text properties from PowerPoint file"""
    prs = Presentation(str(pptx_path))
    properties = []
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape.has_text_frame:
                tf = shape.text_frame
                props = {
                    'vertical_anchor': str(tf.vertical_anchor) if tf.vertical_anchor else None,
                    'margin_top': tf.margin_top,
                    'margin_bottom': tf.margin_bottom,
                    'margin_left': tf.margin_left,
                    'margin_right': tf.margin_right,
                    'auto_size': str(tf.auto_size) if tf.auto_size else None,
                    'word_wrap': tf.word_wrap,
                    'paragraphs': []
                }
                
                for paragraph in tf.paragraphs:
                    para_props = {
                        'space_before': paragraph.space_before,
                        'space_after': paragraph.space_after,
                        'line_spacing': str(paragraph.line_spacing) if paragraph.line_spacing else None,
                        'alignment': str(paragraph.alignment) if paragraph.alignment else None,
                        'runs': []
                    }
                    
                    for run in paragraph.runs:
                        run_props = {
                            'font_size': run.font.size.pt if run.font.size else None,
                            'bold': run.font.bold if hasattr(run.font, 'bold') else None,
                            'italic': run.font.italic if hasattr(run.font, 'italic') else None,
                        }
                        para_props['runs'].append(run_props)
                    
                    props['paragraphs'].append(para_props)
                
                properties.append(props)
    
    return properties


class TestTextPositionComparison:
    """Test class for comparing text positions between drawio and PowerPoint"""
    
    @pytest.fixture
    def sample_drawio_path(self):
        """Path to sample drawio file"""
        path = Path(__file__).parent.parent / "sample" / "sample.drawio"
        if not path.exists():
            pytest.skip(f"Sample file not found: {path}")
        return path
    
    @pytest.fixture
    def sample_pptx_path(self):
        """Path to sample pptx file"""
        # First, look for sample/sample.pptx
        sample_pptx = Path(__file__).parent.parent / "sample" / "sample.pptx"
        if sample_pptx.exists():
            return sample_pptx
        
        # Otherwise, look for test_output.pptx
        test_output = Path(__file__).parent.parent / "test_output.pptx"
        if test_output.exists():
            return test_output
        
        # Otherwise, look for output.pptx
        output = Path(__file__).parent.parent / "output.pptx"
        if output.exists():
            return output
        
        pytest.skip("No sample PPTX file found")
    
    def test_drawio_text_properties_extraction(self, sample_drawio_path):
        """Verify text properties can be correctly extracted from drawio file"""
        properties = extract_drawio_text_properties(sample_drawio_path)
        
        # If at least one shape exists, properties may be extracted
        # (Some shapes may not have text properties set, so empty is OK)
        assert isinstance(properties, list)
    
    def test_pptx_text_properties_extraction(self, sample_pptx_path):
        """Verify text properties can be correctly extracted from PowerPoint file"""
        properties = extract_pptx_text_properties(sample_pptx_path)
        
        # Verify at least one text frame exists
        assert len(properties) > 0
        
        # Verify each property has correct structure
        for props in properties:
            assert 'vertical_anchor' in props
            assert 'margin_top' in props
            assert 'margin_bottom' in props
            assert 'paragraphs' in props
            assert isinstance(props['paragraphs'], list)
    
    def test_text_properties_comparison(self, sample_drawio_path, sample_pptx_path):
        """Compare text properties between drawio and PowerPoint"""
        drawio_props = extract_drawio_text_properties(sample_drawio_path)
        pptx_props = extract_pptx_text_properties(sample_pptx_path)
        
        # Verify properties can be extracted from both files
        assert isinstance(drawio_props, list)
        assert isinstance(pptx_props, list)
        assert len(pptx_props) > 0
    
    def test_drawio_has_text_properties(self, sample_drawio_path):
        """Verify drawio file has text properties defined"""
        properties = extract_drawio_text_properties(sample_drawio_path)
        
        # Check if any shape has text properties set
        has_text_props = any(
            any(props.values()) for props in properties
        )
        
        # Whether properties exist depends on the file, so just check
        assert isinstance(has_text_props, bool)
    
    def test_pptx_has_text_frame_properties(self, sample_pptx_path):
        """Verify PowerPoint file has text frame properties set"""
        properties = extract_pptx_text_properties(sample_pptx_path)
        
        # Verify at least one text frame exists
        assert len(properties) > 0
        
        # Check if text frame properties are set
        has_props = any(
            props['vertical_anchor'] is not None or
            props['margin_top'] != 0 or
            props['margin_bottom'] != 0 or
            props['margin_left'] != 0 or
            props['margin_right'] != 0
            for props in properties
        )
        
        # Whether properties are set depends on the file, so just check
        assert isinstance(has_props, bool)
    
    def test_paragraph_comparison(self, sample_pptx_path):
        """Verify paragraph properties in PowerPoint file"""
        properties = extract_pptx_text_properties(sample_pptx_path)
        
        for props in properties:
            for paragraph in props['paragraphs']:
                # Verify paragraph properties are correctly extracted
                assert 'space_before' in paragraph
                assert 'space_after' in paragraph
                assert 'runs' in paragraph
                assert isinstance(paragraph['runs'], list)
    
    def test_font_properties_comparison(self, sample_pptx_path):
        """Compare font properties"""
        properties = extract_pptx_text_properties(sample_pptx_path)
        
        for props in properties:
            for paragraph in props['paragraphs']:
                for run in paragraph['runs']:
                    # Verify font properties are correctly extracted
                    assert 'font_size' in run
                    # font_size is None or numeric
                    assert run['font_size'] is None or isinstance(run['font_size'], (int, float))




