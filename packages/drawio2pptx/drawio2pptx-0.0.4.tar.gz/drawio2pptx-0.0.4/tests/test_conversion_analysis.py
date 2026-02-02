"""
Detailed analysis test for conversion results

Compare draw.io files with generated PowerPoint files to analyze
what is correctly converted and what is not converted
"""
import pytest
from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from lxml import etree as ET
from lxml import html as lxml_html

from drawio2pptx.io.drawio_loader import DrawIOLoader
from drawio2pptx.logger import get_logger


def analyze_drawio_shape(cell, style_extractor):
    """Analyze draw.io shape"""
    result = {
        'id': cell.attrib.get('id'),
        'text': cell.attrib.get('value', ''),
        'shape_type': style_extractor.extract_shape_type(cell),
        'fill_color': style_extractor.extract_fill_color(cell),
        'stroke_color': style_extractor.extract_stroke_color(cell),
        'font_color': style_extractor.extract_font_color(cell),
        'has_shadow': style_extractor.extract_shadow(cell, None),
        'text_properties': {}
    }
    
    style_str = cell.attrib.get("style", "")
    result['text_properties'] = {
        'fontSize': style_extractor.extract_style_float(style_str, "fontSize"),
        'fontFamily': style_extractor.extract_style_value(style_str, "fontFamily"),
        'fontStyle': style_extractor.extract_style_value(style_str, "fontStyle"),
        'align': style_extractor.extract_style_value(style_str, "align"),
        'verticalAlign': style_extractor.extract_style_value(style_str, "verticalAlign"),
    }
    
    return result


def analyze_pptx_shape(shape):
    """Analyze PowerPoint shape"""
    result = {
        'shape_type': shape.shape_type,
        'fill': None,
        'stroke_color': None,
        'has_shadow': None,
        'text_frame': None,
    }
    
    # Fill
    try:
        fill = shape.fill
        if fill.type is None:
            result['fill'] = "None"
        elif hasattr(fill, 'fore_color'):
            fill_color = fill.fore_color
            if hasattr(fill_color, 'rgb') and fill_color.rgb:
                rgb = fill_color.rgb
                result['fill'] = f"RGB({rgb.r}, {rgb.g}, {rgb.b})"
            elif hasattr(fill_color, 'theme_color'):
                result['fill'] = f"ThemeColor({fill_color.theme_color})"
            else:
                result['fill'] = "Unknown"
        else:
            result['fill'] = str(fill.type)
    except Exception:
        result['fill'] = "Error"
    
    # Stroke color
    try:
        if hasattr(shape, '_element'):
            shape_element = shape._element
            nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
            ln_element = shape_element.find('.//a:ln', namespaces=nsmap)
            if ln_element is not None:
                solid_fill = ln_element.find('.//a:solidFill', namespaces=nsmap)
                if solid_fill is not None:
                    srgb = solid_fill.find('.//a:srgbClr', namespaces=nsmap)
                    if srgb is not None:
                        val = srgb.get('val', '')
                        if val and len(val) == 6:
                            r = int(val[0:2], 16)
                            g = int(val[2:4], 16)
                            b = int(val[4:6], 16)
                            result['stroke_color'] = f"RGB({r}, {g}, {b})"
    except Exception:
        pass
    
    # Shadow
    try:
        if hasattr(shape, 'shadow'):
            result['has_shadow'] = shape.shadow.inherit
    except Exception:
        pass
    
    # Text frame
    try:
        if shape.has_text_frame:
            tf = shape.text_frame
            text_runs = []
            for para in tf.paragraphs:
                for run in para.runs:
                    run_info = {
                        'text': run.text,
                        'font_name': run.font.name if run.font.name else None,
                        'font_size': run.font.size.pt if run.font.size else None,
                        'bold': run.font.bold,
                        'italic': run.font.italic,
                        'underline': run.font.underline,
                        'font_color': None,
                    }
                    
                    # Get font color from XML
                    try:
                        if hasattr(run, '_r'):
                            run_element = run._r
                            nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                            r_pr = run_element.find('.//a:rPr', namespaces=nsmap)
                            if r_pr is not None:
                                solid_fill = r_pr.find('.//a:solidFill', namespaces=nsmap)
                                if solid_fill is not None:
                                    srgb = solid_fill.find('.//a:srgbClr', namespaces=nsmap)
                                    if srgb is not None:
                                        val = srgb.get('val', '')
                                        if val and len(val) == 6:
                                            r = int(val[0:2], 16)
                                            g = int(val[2:4], 16)
                                            b = int(val[4:6], 16)
                                            run_info['font_color'] = f"RGB({r}, {g}, {b})"
                    except Exception:
                        pass
                    
                    text_runs.append(run_info)
            
            result['text_frame'] = {
                'vertical_anchor': str(tf.vertical_anchor),
                'runs': text_runs,
            }
    except Exception:
        pass
    
    return result


def format_color_for_comparison(color):
    """Convert color to string for comparison"""
    if color is None:
        return "None"
    if isinstance(color, RGBColor):
        return f"RGB({color[0]}, {color[1]}, {color[2]})"
    if isinstance(color, str):
        if color == "default":
            return "ThemeColor"
        return color
    return str(color)


class TestConversionAnalysis:
    """Test class for detailed analysis of conversion results"""
    
    @pytest.fixture
    def sample_drawio_path(self):
        """Path to sample drawio file"""
        path = Path(__file__).parent.parent / "sample" / "sample.drawio"
        if not path.exists():
            pytest.skip(f"Sample file not found: {path}")
        return path
    
    @pytest.fixture
    def sample_pptx_path(self):
        """Path to sample pptx file"""
        # First, look for sample/sample.pptx
        sample_pptx = Path(__file__).parent.parent / "sample" / "sample.pptx"
        if sample_pptx.exists():
            return sample_pptx
        
        # Otherwise, look for output.pptx
        output = Path(__file__).parent.parent / "output.pptx"
        if output.exists():
            return output
        
        pytest.skip("No sample PPTX file found")
    
    def test_compare_shapes(self, sample_drawio_path, sample_pptx_path):
        """Compare shapes between draw.io and PowerPoint"""
        # Load draw.io
        logger = get_logger()
        loader = DrawIOLoader(logger=logger)
        diagrams = loader.load_file(sample_drawio_path)
        
        if not diagrams:
            pytest.skip("No diagrams found in draw.io file")
        
        elements = loader.extract_elements(diagrams[0])
        shape_elements = [e for e in elements if hasattr(e, 'shape_type')]
        
        # Load PowerPoint
        prs = Presentation(str(sample_pptx_path))
        if not prs.slides:
            pytest.skip("No slides found in PowerPoint file")
        
        slide = prs.slides[0]
        pptx_shapes = [s for s in slide.shapes if s.shape_type == 1]  # AUTO_SHAPE
        
        print(f"\n=== Conversion Result Comparison ===\n")
        print(f"draw.io shapes: {len(shape_elements)}")
        print(f"PowerPoint shapes: {len(pptx_shapes)}\n")
        
        # Compare each shape
        for idx, (drawio_shape, pptx_shape) in enumerate(zip(shape_elements, pptx_shapes), 1):
            print(f"--- Shape {idx} ---")
            
            # Text
            drawio_text = ""
            if drawio_shape.text:
                for para in drawio_shape.text:
                    for run in para.runs:
                        drawio_text += run.text
            
            pptx_text = ""
            if pptx_shape.has_text_frame:
                for para in pptx_shape.text_frame.paragraphs:
                    for run in para.runs:
                        pptx_text += run.text
            
            print(f"Text:")
            print(f"  draw.io: '{drawio_text}'")
            print(f"  PowerPoint: '{pptx_text}'")
            print(f"  ✓" if drawio_text == pptx_text else f"  ✗ (mismatch)")
            
            # Fill color
            drawio_fill = format_color_for_comparison(drawio_shape.style.fill)
            pptx_fill = "None"
            try:
                fill_type = pptx_shape.fill.type
                if fill_type:
                    if hasattr(pptx_shape.fill, 'fore_color'):
                        try:
                            fore_color = pptx_shape.fill.fore_color
                            if hasattr(fore_color, 'rgb') and fore_color.rgb:
                                rgb = fore_color.rgb
                                pptx_fill = f"RGB({rgb.r}, {rgb.g}, {rgb.b})"
                            elif hasattr(fore_color, 'theme_color'):
                                pptx_fill = f"ThemeColor({fore_color.theme_color})"
                        except (TypeError, AttributeError):
                            pass
                
                # Check from XML (for ThemeColor case)
                if pptx_fill == "None" and hasattr(pptx_shape, '_element'):
                    shape_element = pptx_shape._element
                    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                    # spPr exists as direct child element
                    sp_pr = None
                    for child in shape_element:
                        if child.tag.endswith('}spPr') or 'spPr' in child.tag:
                            sp_pr = child
                            break
                    
                    if sp_pr is not None:
                        solid_fill = sp_pr.find('.//a:solidFill', namespaces=nsmap)
                        if solid_fill is not None:
                            scheme_clr = solid_fill.find('.//a:schemeClr', namespaces=nsmap)
                            if scheme_clr is not None:
                                val = scheme_clr.get('val', '')
                                pptx_fill = f"ThemeColor({val})"
                            else:
                                srgb = solid_fill.find('.//a:srgbClr', namespaces=nsmap)
                                if srgb is not None:
                                    val = srgb.get('val', '')
                                    if val and len(val) == 6:
                                        r = int(val[0:2], 16)
                                        g = int(val[2:4], 16)
                                        b = int(val[4:6], 16)
                                        pptx_fill = f"RGB({r}, {g}, {b})"
            except Exception:
                pass
            
            print(f"Fill:")
            print(f"  draw.io: {drawio_fill}")
            print(f"  PowerPoint: {pptx_fill}")
            # For ThemeColor case, consider as match
            if drawio_fill == "ThemeColor" and "ThemeColor" in pptx_fill:
                print(f"  ✓")
            elif drawio_fill == pptx_fill or (drawio_fill == "None" and pptx_fill == "None"):
                print(f"  ✓")
            else:
                print(f"  ✗ (mismatch)")
            
            # Stroke color
            drawio_stroke = format_color_for_comparison(drawio_shape.style.stroke)
            pptx_stroke = "None"
            pptx_result = analyze_pptx_shape(pptx_shape)
            if pptx_result['stroke_color']:
                pptx_stroke = pptx_result['stroke_color']
            
            print(f"Stroke color:")
            print(f"  draw.io: {drawio_stroke}")
            print(f"  PowerPoint: {pptx_stroke}")
            
            # Shadow
            drawio_shadow = drawio_shape.style.has_shadow
            pptx_shadow = pptx_result['has_shadow']
            print(f"Shadow:")
            print(f"  draw.io: {drawio_shadow}")
            print(f"  PowerPoint: {pptx_shadow}")
            print(f"  ✓" if drawio_shadow == pptx_shadow else f"  ✗ (mismatch)")
            
            # Font information
            if drawio_shape.text and pptx_shape.has_text_frame:
                drawio_runs = []
                for para in drawio_shape.text:
                    drawio_runs.extend(para.runs)
                
                pptx_runs = []
                for para in pptx_shape.text_frame.paragraphs:
                    pptx_runs.extend(para.runs)
                
                print(f"Font information:")
                for run_idx, (drawio_run, pptx_run) in enumerate(zip(drawio_runs, pptx_runs), 1):
                    print(f"  Run {run_idx}:")
                    
                    # Font color
                    drawio_color = format_color_for_comparison(drawio_run.font_color)
                    pptx_color = "None"
                    pptx_run_info = analyze_pptx_shape(pptx_shape)
                    if pptx_run_info['text_frame'] and pptx_run_info['text_frame']['runs']:
                        if run_idx <= len(pptx_run_info['text_frame']['runs']):
                            pptx_color = pptx_run_info['text_frame']['runs'][run_idx-1].get('font_color', 'None')
                    
                    print(f"    Font color:")
                    print(f"      draw.io: {drawio_color}")
                    print(f"      PowerPoint: {pptx_color}")
                    
                    # Bold
                    pptx_bold = pptx_run.font.bold if pptx_run.font.bold else False
                    print(f"    Bold:")
                    print(f"      draw.io: {drawio_run.bold}")
                    print(f"      PowerPoint: {pptx_bold}")
                    print(f"      ✓" if drawio_run.bold == pptx_bold else f"      ✗ (mismatch)")
                    
                    # Font size
                    drawio_size = drawio_run.font_size
                    pptx_size = pptx_run.font.size.pt if pptx_run.font.size else None
                    print(f"    Font size:")
                    print(f"      draw.io: {drawio_size}")
                    print(f"      PowerPoint: {pptx_size}")
                    
                    # Font name
                    drawio_font = drawio_run.font_family
                    pptx_font = pptx_run.font.name
                    print(f"    Font name:")
                    print(f"      draw.io: {drawio_font}")
                    print(f"      PowerPoint: {pptx_font}")
            
            # Vertical alignment
            if drawio_shape.text:
                drawio_vertical = drawio_shape.text[0].vertical_align if drawio_shape.text else None
            else:
                drawio_vertical = None
            
            pptx_vertical = None
            if pptx_shape.has_text_frame:
                pptx_vertical = str(pptx_shape.text_frame.vertical_anchor)
            
            print(f"Vertical alignment:")
            print(f"  draw.io: {drawio_vertical}")
            print(f"  PowerPoint: {pptx_vertical}")
            
            print()
    
    def test_list_converted_features(self, sample_drawio_path, sample_pptx_path):
        """List features that are converted and not converted"""
        logger = get_logger()
        loader = DrawIOLoader(logger=logger)
        diagrams = loader.load_file(sample_drawio_path)
        
        if not diagrams:
            pytest.skip("No diagrams found")
        
        elements = loader.extract_elements(diagrams[0])
        
        prs = Presentation(str(sample_pptx_path))
        if not prs.slides:
            pytest.skip("No slides found")
        
        slide = prs.slides[0]
        
        converted = []
        not_converted = []
        
        # Check each feature
        for element in elements:
            if hasattr(element, 'shape_type'):
                # Shape type
                converted.append(f"Shape type: {element.shape_type}")
                
                # Fill
                if element.style.fill:
                    converted.append("Fill color")
                else:
                    converted.append("No fill")
                
                # Stroke
                if element.style.stroke:
                    converted.append("Stroke color")
                else:
                    converted.append("No stroke")
                
                # Shadow
                if element.style.has_shadow:
                    converted.append("Shadow")
                else:
                    converted.append("No shadow")
                
                # Text
                if element.text:
                    converted.append("Text")
                    for para in element.text:
                        for run in para.runs:
                            if run.bold:
                                converted.append("Bold")
                            if run.italic:
                                converted.append("Italic")
                            if run.underline:
                                converted.append("Underline")
                            if run.font_color:
                                converted.append("Font color")
                            if run.font_size:
                                converted.append("Font size")
                            if run.font_family:
                                converted.append("Font name")
        
        print("\n=== Converted Features ===")
        for feature in set(converted):
            print(f"  ✓ {feature}")
        
        print("\n=== Not Converted Features ===")
        if not_converted:
            for feature in set(not_converted):
                print(f"  ✗ {feature}")
        else:
            print("  (All features are converted)")




