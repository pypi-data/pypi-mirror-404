"""
Test module for vertical text frame settings in PowerPoint files
"""
import pytest
from pathlib import Path
from pptx import Presentation
from lxml import etree


class TestVerticalTextAlignment:
    """Test class for vertical text frame settings"""
    
    @pytest.fixture
    def sample_pptx_path(self):
        """Path to sample pptx file"""
        # First, look for sample/sample.pptx
        sample_pptx = Path(__file__).parent.parent / "sample" / "sample.pptx"
        if sample_pptx.exists():
            return sample_pptx
        
        # Otherwise, look for test_output.pptx
        test_output = Path(__file__).parent.parent / "test_output.pptx"
        if test_output.exists():
            return test_output
        
        # Otherwise, look for output.pptx
        output = Path(__file__).parent.parent / "output.pptx"
        if output.exists():
            return output
        
        pytest.skip("No sample PPTX file found")
    
    def test_text_frame_vertical_anchor(self, sample_pptx_path):
        """Verify text frame vertical anchor settings"""
        prs = Presentation(str(sample_pptx_path))
        assert len(prs.slides) > 0
        
        slide = prs.slides[0]
        assert len(slide.shapes) > 0
        
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape.has_text_frame:
                tf = shape.text_frame
                
                # Verify vertical_anchor property exists
                assert hasattr(tf, 'vertical_anchor')
                
                # Verify vertical_anchor value is valid
                # Value is one of: None, MSO_ANCHOR_TOP, MSO_ANCHOR_MIDDLE, MSO_ANCHOR_BOTTOM
                va = tf.vertical_anchor
                # vertical_anchor is None or MSO_VERTICAL_ANCHOR enum
                assert va is None or hasattr(va, 'value') or 'ANCHOR' in str(type(va))
    
    def test_text_frame_margins(self, sample_pptx_path):
        """Verify text frame margin settings"""
        prs = Presentation(str(sample_pptx_path))
        slide = prs.slides[0]
        
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape.has_text_frame:
                tf = shape.text_frame
                
                # Verify margin properties exist
                assert hasattr(tf, 'margin_top')
                assert hasattr(tf, 'margin_bottom')
                assert hasattr(tf, 'margin_left')
                assert hasattr(tf, 'margin_right')
                
                # Verify margin values are numeric (in EMU units)
                assert isinstance(tf.margin_top, (int, type(None)))
                assert isinstance(tf.margin_bottom, (int, type(None)))
                assert isinstance(tf.margin_left, (int, type(None)))
                assert isinstance(tf.margin_right, (int, type(None)))
    
    def test_text_frame_auto_size(self, sample_pptx_path):
        """Verify text frame auto size settings"""
        prs = Presentation(str(sample_pptx_path))
        slide = prs.slides[0]
        
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape.has_text_frame:
                tf = shape.text_frame
                
                # Verify auto_size property exists
                assert hasattr(tf, 'auto_size')
                
                # Verify auto_size value is valid
                auto_size = tf.auto_size
                assert auto_size is None or str(auto_size) in ['None', 'NONE', 'AUTO_SHAPE', 'TEXT_FRAME_AUTO_SIZE']
    
    def test_text_frame_word_wrap(self, sample_pptx_path):
        """Verify text frame word wrap settings"""
        prs = Presentation(str(sample_pptx_path))
        slide = prs.slides[0]
        
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape.has_text_frame:
                tf = shape.text_frame
                
                # Verify word_wrap property exists
                assert hasattr(tf, 'word_wrap')
                
                # Verify word_wrap value is boolean or None
                assert isinstance(tf.word_wrap, (bool, type(None)))
    
    def test_body_pr_xml_element(self, sample_pptx_path):
        """Verify bodyPr XML element existence and structure"""
        prs = Presentation(str(sample_pptx_path))
        slide = prs.slides[0]
        
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape.has_text_frame:
                tf = shape.text_frame
                element = tf._element
                
                # Search for bodyPr element
                nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                body_pr = element.find('.//a:bodyPr', namespaces=nsmap)
                
                # If bodyPr element exists, verify its attributes
                if body_pr is not None:
                    # Verify bodyPr element can be parsed correctly
                    xml_str = etree.tostring(body_pr, encoding='unicode')
                    assert 'bodyPr' in xml_str
    
    def test_paragraph_properties(self, sample_pptx_path):
        """Verify paragraph properties"""
        prs = Presentation(str(sample_pptx_path))
        slide = prs.slides[0]
        
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape.has_text_frame:
                tf = shape.text_frame
                
                # Verify paragraphs exist
                assert len(tf.paragraphs) > 0
                
                for paragraph in tf.paragraphs:
                    # Verify paragraph properties exist
                    assert hasattr(paragraph, 'space_before')
                    assert hasattr(paragraph, 'space_after')
                    assert hasattr(paragraph, 'line_spacing')
                    assert hasattr(paragraph, 'text')
                    assert hasattr(paragraph, 'runs')
                    
                    # Verify space_before and space_after are numeric or None
                    assert isinstance(paragraph.space_before, (int, type(None)))
                    assert isinstance(paragraph.space_after, (int, type(None)))
    
    def test_run_font_properties(self, sample_pptx_path):
        """Verify text run font properties"""
        prs = Presentation(str(sample_pptx_path))
        slide = prs.slides[0]
        
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape.has_text_frame:
                tf = shape.text_frame
                
                for paragraph in tf.paragraphs:
                    for run in paragraph.runs:
                        # Verify font properties exist
                        assert hasattr(run, 'font')
                        assert hasattr(run.font, 'size')
                        
                        # If font size is set, verify pt property exists
                        if run.font.size:
                            assert hasattr(run.font.size, 'pt')
                            assert isinstance(run.font.size.pt, (int, float))
    
    def test_shape_geometry(self, sample_pptx_path):
        """Verify shape geometry information"""
        prs = Presentation(str(sample_pptx_path))
        slide = prs.slides[0]
        
        has_valid_shape = False
        for shape in slide.shapes:
            # Verify shape geometry properties exist
            assert hasattr(shape, 'height')
            assert hasattr(shape, 'width')
            
            # Verify height and width are numeric (in EMU units)
            assert isinstance(shape.height, int)
            assert isinstance(shape.width, int)
            
            # Verify height and width are non-negative (shapes with 0 size may exist)
            assert shape.height >= 0
            assert shape.width >= 0
            
            # Verify at least one shape with positive values exists
            if shape.height > 0 and shape.width > 0:
                has_valid_shape = True
                # Verify EMU to pixel conversion works correctly
                height_px = shape.height / 9525
                width_px = shape.width / 9525
                assert height_px > 0
                assert width_px > 0
        
        # Verify at least one valid shape exists
        assert has_valid_shape, "At least one shape with positive size must exist"




