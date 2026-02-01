import asyncio
import csv
import json
import math
import re
import textwrap
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path
from typing import Any, Dict, List, Union

try:  # Python 3.11+
    from typing import LiteralString
except ImportError:  # 3.9–3.10
    from typing_extensions import LiteralString

import numpy as np
import pdfplumber
from docx import Document
from PIL import Image
from pptx import Presentation
from projectdavid_common import UtilsInterface

log = UtilsInterface.LoggingUtility()


def latlon_to_unit_vec(lat: float, lon: float) -> List[float]:
    """Convert geographic lat/lon (deg) to a 3-D unit vector for Qdrant."""
    lat_r = math.radians(lat)
    lon_r = math.radians(lon)
    return [
        math.cos(lat_r) * math.cos(lon_r),
        math.cos(lat_r) * math.sin(lon_r),
        math.sin(lat_r),
    ]


class FileProcessor:
    # ------------------------------------------------------------------ #
    #  Construction
    # ------------------------------------------------------------------ #
    def __init__(
        self,
        *,
        max_workers: int = 4,
        chunk_size: int = 512,
        use_gpu: bool = True,
        use_ocr: bool = True,
        use_detection: bool = False,
        image_model_name: str = "ViT-H-14",
        caption_model_name: str = "Salesforce/blip2-flan-t5-xl",
    ):
        # Configuration
        self._use_gpu = use_gpu
        self._max_workers = max_workers
        self._requested_chunk_size = chunk_size
        self._image_model_name = image_model_name
        self._executor = ThreadPoolExecutor(max_workers=max_workers)

        # Lazy ML Stack Attributes
        self._device = None
        self._torch_dtype = None
        self._embedding_model = None
        self._clip_model = None
        self._clip_preprocess = None
        self._clip_tokenizer = None

        # Lazy Token Limits
        self._effective_max_length = None
        self._chunk_size = None

        log.info("Initialized Multimodal Lazy-Loaded FileProcessor")

    def _ensure_ml_stack(self):
        """Lazy-loads Torch, CLIP, and SentenceTransformers only when needed."""
        if self._embedding_model is None:
            try:
                import open_clip
                import torch
                from sentence_transformers import SentenceTransformer

                # 1. Setup Device
                if self._use_gpu and torch.cuda.is_available():
                    self._device = torch.device("cuda")
                    self._torch_dtype = torch.float16
                else:
                    self._device = torch.device("cpu")
                    self._torch_dtype = torch.float32

                # 2. Setup Text Embedder
                self.embedding_model_name = "paraphrase-MiniLM-L6-v2"
                self._embedding_model = SentenceTransformer(self.embedding_model_name)
                self._embedding_model.to(str(self._device))

                # 3. Setup CLIP
                # Note: We use the provided image_model_name (default ViT-H-14)
                self._clip_model, _, self._clip_preprocess = (
                    open_clip.create_model_and_transforms(
                        self._image_model_name, pretrained="laion2b_s32b_b79k"
                    )
                )
                self._clip_model.to(self._device)
                self._clip_tokenizer = open_clip.get_tokenizer(self._image_model_name)

                # 4. Calculate limits
                max_seq_length = self._embedding_model.get_max_seq_length()
                special_tokens_count = 2
                self._effective_max_length = max_seq_length - special_tokens_count
                self._chunk_size = min(
                    self._requested_chunk_size, self._effective_max_length * 4
                )

                log.info("ML Stack loaded (device=%s)", self._device)

            except ImportError as e:
                log.error(f"ML Stack failed to load: {e}")
                raise ImportError(
                    "This feature requires heavy ML binaries. "
                    "Please install the vector stack: pip install projectdavid[vector]"
                )
        return self._embedding_model

    @property
    def chunk_size(self):
        if self._chunk_size is None:
            self._ensure_ml_stack()
        return self._chunk_size

    @property
    def effective_max_length(self):
        if self._effective_max_length is None:
            self._ensure_ml_stack()
        return self._effective_max_length

    # ------------------------------------------------------------------ #
    #  Public Embedders
    # ------------------------------------------------------------------ #
    def encode_text(self, text: Union[str, List[str]]) -> np.ndarray:
        model = self._ensure_ml_stack()
        single = isinstance(text, str)
        out = model.encode(
            text,
            convert_to_numpy=True,
            normalize_embeddings=True,
            show_progress_bar=False,
        )
        return out if not single else out[0]

    def encode_image(self, img: Image.Image) -> np.ndarray:
        import torch

        self._ensure_ml_stack()
        with torch.no_grad():
            tensor = self._clip_preprocess(img).unsqueeze(0).to(self._device)
            feat = self._clip_model.encode_image(tensor).squeeze()
            feat = feat / feat.norm()
            return feat.float().cpu().numpy()

    def encode_clip_text(self, text: Union[str, List[str]]) -> np.ndarray:
        import torch

        self._ensure_ml_stack()
        with torch.no_grad():
            toks = (
                self._clip_tokenizer(text)
                if isinstance(text, str)
                else self._clip_tokenizer(text)  # Adjusted for clip tokenizer behavior
            )
            tensor = toks.to(self._device)
            feat = self._clip_model.encode_text(tensor).squeeze()
            if feat.dim() > 1:  # Handle batch
                feat = feat / feat.norm(dim=-1, keepdim=True)
            else:
                feat = feat / feat.norm()
            return feat.float().cpu().numpy()

    # ------------------------------------------------------------------ #
    #  Generic validators / Type Detection
    # ------------------------------------------------------------------ #
    def validate_file(self, file_path: Path):
        max_size = 100 * 1024 * 1024
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        if file_path.stat().st_size > max_size:
            raise ValueError(f"{file_path.name} exceeds 100MB limit")

    def _detect_file_type(self, file_path: Path) -> str:
        suffix = file_path.suffix.lower()
        if suffix == ".pdf":
            return "pdf"
        if suffix == ".csv":
            return "csv"
        if suffix == ".json":
            return "json"
        if suffix in {".doc", ".docx", ".pptx"}:
            return "office"
        if suffix in {".jpg", ".jpeg", ".png", ".webp", ".bmp", ".gif", ".tiff"}:
            return "image"

        text_exts = {
            ".txt",
            ".md",
            ".rst",
            ".c",
            ".cpp",
            ".cs",
            ".go",
            ".java",
            ".js",
            ".ts",
            ".py",
            ".html",
            ".css",
        }
        if suffix in text_exts:
            return "text"
        raise ValueError(f"Unsupported file type: {file_path.name}")

    async def process_file(self, file_path: Union[str, Path]) -> Dict[str, Any]:
        path = Path(file_path)
        self.validate_file(path)
        ftype = self._detect_file_type(path)
        return await getattr(self, f"_process_{ftype}")(path)

    # ------------------------------------------------------------------ #
    #  Processors (PDF, Text, CSV, Office, JSON, Image)
    # ------------------------------------------------------------------ #
    async def _process_pdf(self, file_path: Path) -> Dict[str, Any]:
        page_chunks, doc_meta = await self._extract_text(file_path)
        all_chunks, line_data = [], []
        for page_text, page_num, line_nums in page_chunks:
            lines = page_text.split("\n")
            buf, buf_lines, length = [], [], 0
            for line, ln in zip(lines, line_nums):
                l = len(line) + 1
                if length + l <= self.chunk_size:
                    buf.append(line)
                    buf_lines.append(ln)
                    length += l
                else:
                    if buf:
                        all_chunks.append("\n".join(buf))
                        line_data.append({"page": page_num, "lines": buf_lines})
                        buf, buf_lines, length = [], [], 0
                    for piece in self._split_oversized_chunk(line):
                        all_chunks.append(piece)
                        line_data.append({"page": page_num, "lines": [ln]})
            if buf:
                all_chunks.append("\n".join(buf))
                line_data.append({"page": page_num, "lines": buf_lines})

        vectors = await asyncio.gather(
            *[self._encode_chunk_async(c) for c in all_chunks]
        )
        return {
            "content": "\n\n".join(all_chunks),
            "metadata": {
                **doc_meta,
                "source": str(file_path),
                "chunks": len(all_chunks),
                "type": "pdf",
            },
            "chunks": all_chunks,
            "vectors": [v.tolist() for v in vectors],
            "line_data": line_data,
        }

    async def _process_text(self, file_path: Path) -> Dict[str, Any]:
        text, extra_meta, _ = await self._extract_text(file_path)
        chunks = self._chunk_text(text)
        vectors = await asyncio.gather(*[self._encode_chunk_async(c) for c in chunks])
        return {
            "content": text,
            "metadata": {
                **extra_meta,
                "source": str(file_path),
                "chunks": len(chunks),
                "type": "text",
            },
            "chunks": chunks,
            "vectors": [v.tolist() for v in vectors],
        }

    async def _process_csv(
        self, file_path: Path, text_field: str = "description"
    ) -> Dict[str, Any]:
        texts, metas = [], []
        with file_path.open(newline="", encoding="utf-8") as f:
            reader = csv.DictReader(f)
            for row in reader:
                txt = row.get(text_field, "").strip()
                if not txt:
                    continue
                texts.append(txt)
                metas.append({k: v for k, v in row.items() if k != text_field and v})
        vectors = await asyncio.gather(*[self._encode_chunk_async(t) for t in texts])
        return {
            "content": None,
            "metadata": {"source": str(file_path), "rows": len(texts), "type": "csv"},
            "chunks": texts,
            "vectors": [v.tolist() for v in vectors],
            "csv_row_metadata": metas,
        }

    async def _process_office(self, file_path: Path) -> Dict[str, Any]:
        loop = asyncio.get_event_loop()
        method = (
            self._read_docx
            if file_path.suffix.lower() in {".doc", ".docx"}
            else self._read_pptx
        )
        text = await loop.run_in_executor(self._executor, method, file_path)
        chunks = self._chunk_text(text)
        vectors = await asyncio.gather(*[self._encode_chunk_async(c) for c in chunks])
        return {
            "content": text,
            "metadata": {
                "source": str(file_path),
                "chunks": len(chunks),
                "type": "office",
            },
            "chunks": chunks,
            "vectors": [v.tolist() for v in vectors],
        }

    async def _process_json(self, file_path: Path) -> Dict[str, Any]:
        text = await asyncio.get_event_loop().run_in_executor(
            self._executor, self._read_json, file_path
        )
        chunks = self._chunk_text(text)
        vectors = await asyncio.gather(*[self._encode_chunk_async(c) for c in chunks])
        return {
            "content": text,
            "metadata": {
                "source": str(file_path),
                "chunks": len(chunks),
                "type": "json",
            },
            "chunks": chunks,
            "vectors": [v.tolist() for v in vectors],
        }

    async def _process_image(self, file_path: Path) -> Dict[str, Any]:
        """Handles image embedding via CLIP."""
        img = Image.open(file_path).convert("RGB")
        vector = await asyncio.get_event_loop().run_in_executor(
            self._executor, self.encode_image, img
        )
        return {
            "content": None,
            "metadata": {"source": str(file_path), "type": "image"},
            "chunks": [],
            "vectors": [vector.tolist()],
        }

    # ------------------------------------------------------------------ #
    #  Extraction/Read Helpers
    # ------------------------------------------------------------------ #
    async def _extract_text(self, file_path: Path):
        loop = asyncio.get_event_loop()
        if file_path.suffix.lower() == ".pdf":
            return await loop.run_in_executor(
                self._executor, self._extract_pdf_text, file_path
            )
        text = await loop.run_in_executor(
            self._executor, self._read_text_file, file_path
        )
        return text, {}, []

    def _extract_pdf_text(self, file_path: Path):
        page_chunks, meta = [], {}
        with pdfplumber.open(file_path) as pdf:
            meta.update(
                {
                    "author": pdf.metadata.get("Author", ""),
                    "title": pdf.metadata.get("Title", file_path.stem),
                    "page_count": len(pdf.pages),
                }
            )
            for i, page in enumerate(pdf.pages, start=1):
                lines = page.extract_text_lines()
                sorted_lines = sorted(lines, key=lambda x: x["top"])
                txts, nums = [], []
                for ln_idx, line in enumerate(sorted_lines, start=1):
                    t = line.get("text", "").strip()
                    if t:
                        txts.append(t)
                        nums.append(ln_idx)
                if txts:
                    page_chunks.append(("\n".join(txts), i, nums))
        return page_chunks, meta

    def _read_text_file(self, file_path: Path) -> str:
        try:
            return file_path.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            return file_path.read_text(encoding="latin-1")

    def _read_docx(self, path: Path) -> str:
        return "\n".join(p.text for p in Document(path).paragraphs if p.text.strip())

    def _read_pptx(self, path: Path) -> str:
        slides = []
        for slide in Presentation(path).slides:
            slides.append(
                "\n".join(
                    filter(
                        None, [sh.text for sh in slide.shapes if hasattr(sh, "text")]
                    )
                )
            )
        return "\n\n".join(slides)

    def _read_json(self, path: Path) -> str:
        obj = json.loads(path.read_text(encoding="utf-8"))
        return "\n".join(
            textwrap.wrap(json.dumps(obj, indent=2, ensure_ascii=False), width=120)
        )

    async def _encode_chunk_async(self, chunk: str) -> np.ndarray:
        model = self._ensure_ml_stack()
        return await asyncio.get_event_loop().run_in_executor(
            self._executor,
            lambda: model.encode(
                [chunk],
                convert_to_numpy=True,
                normalize_embeddings=True,
                show_progress_bar=False,
            )[0],
        )

    # ------------------------------------------------------------------ #
    #  Chunking Logic
    # ------------------------------------------------------------------ #
    def _chunk_text(self, text: str) -> List[str]:
        sentences = re.split(r"(?<=[\.!?])\s+", text)
        chunks, buf, length = [], [], 0
        for sent in sentences:
            slen = len(sent) + 1
            if length + slen <= self.chunk_size:
                buf.append(sent)
                length += slen
            else:
                if buf:
                    chunks.append(" ".join(buf))
                while len(sent) > self.chunk_size:
                    part, sent = sent[: self.chunk_size], sent[self.chunk_size :]
                    chunks.append(part)
                buf, length = [sent], len(sent)
        if buf:
            chunks.append(" ".join(buf))
        return chunks

    def _split_oversized_chunk(self, chunk: str, tokens: List[str] = None) -> List[str]:
        model = self._ensure_ml_stack()
        if tokens is None:
            tokens = model.tokenizer.tokenize(chunk)
        out = []
        for i in range(0, len(tokens), self.effective_max_length):
            seg = tokens[i : i + self.effective_max_length]
            out.append(model.tokenizer.convert_tokens_to_string(seg))
        return out
