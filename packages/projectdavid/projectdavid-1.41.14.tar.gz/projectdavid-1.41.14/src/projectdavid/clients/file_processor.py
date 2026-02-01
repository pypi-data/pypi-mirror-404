import asyncio
import csv
import json
import re
import textwrap
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path
from typing import Any, Dict, List, Tuple, Union

try:  # Python 3.11+
    from typing import LiteralString
except ImportError:  # 3.9–3.10
    from typing_extensions import LiteralString

import numpy as np
import pdfplumber
from docx import Document
from pptx import Presentation
from projectdavid_common import UtilsInterface

log = UtilsInterface.LoggingUtility()


class FileProcessor:
    # ------------------------------------------------------------------ #
    #  Construction
    # ------------------------------------------------------------------ #
    def __init__(self, max_workers: int = 4, chunk_size: int = 512):
        self.embedding_model_name = "paraphrase-MiniLM-L6-v2"
        self._embedding_model = None
        self._executor = ThreadPoolExecutor(max_workers=max_workers)

        # Lazy-initialized attributes
        self._requested_chunk_size = chunk_size
        self._max_seq_length = None
        self._effective_max_length = None
        self._chunk_size = None

        log.info("Initialized Lazy-Loaded FileProcessor")

    def _ensure_model(self):
        """
        Internal helper to load the model and calculate limits only once.
        This prevents heavy imports (scipy, torch) until actually needed.
        """
        if self._embedding_model is None:
            try:
                from sentence_transformers import SentenceTransformer

                log.info(f"Lazy-loading model: {self.embedding_model_name}")

                self._embedding_model = SentenceTransformer(self.embedding_model_name)

                # Ported Limit Calculations
                self._max_seq_length = self._embedding_model.get_max_seq_length()
                special_tokens_count = 2
                self._effective_max_length = self._max_seq_length - special_tokens_count
                self._chunk_size = min(
                    self._requested_chunk_size, self._effective_max_length * 4
                )

            except ImportError:
                log.error(
                    "sentence-transformers not found. Ensure 'pip install projectdavid[vision]' is installed."
                )
                raise ImportError(
                    "Model-based features require 'sentence-transformers'. Install with [vision] extra."
                )
        return self._embedding_model

    # Properties to maintain access to derived attributes
    @property
    def chunk_size(self):
        if self._chunk_size is None:
            self._ensure_model()
        return self._chunk_size

    @property
    def effective_max_length(self):
        if self._effective_max_length is None:
            self._ensure_model()
        return self._effective_max_length

    # ------------------------------------------------------------------ #
    #  Embeddings
    # ------------------------------------------------------------------ #
    def encode_text(self, text: str):
        model = self._ensure_model()
        return model.encode(
            [text],
            convert_to_numpy=True,
            truncate="model_max_length",
            normalize_embeddings=True,
        )[0]

    async def _encode_chunk_async(self, chunk: str) -> np.ndarray:
        model = self._ensure_model()
        return await asyncio.get_event_loop().run_in_executor(
            self._executor,
            lambda: model.encode(
                [chunk],
                convert_to_numpy=True,
                truncate="model_max_length",
                normalize_embeddings=True,
                show_progress_bar=False,
            )[0],
        )

    # ------------------------------------------------------------------ #
    #  Generic validators
    # ------------------------------------------------------------------ #
    def validate_file(self, file_path: Path):
        """Ensure file exists and is under 100 MB."""
        max_size = 100 * 1024 * 1024
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        if file_path.stat().st_size > max_size:
            mb = max_size // (1024 * 1024)
            raise ValueError(f"{file_path.name} > {mb} MB limit")

    # ------------------------------------------------------------------ #
    #  File-type detection
    # ------------------------------------------------------------------ #
    def _detect_file_type(self, file_path: Path) -> str:
        suffix = file_path.suffix.lower()
        if suffix == ".pdf":
            return "pdf"
        if suffix == ".csv":
            return "csv"
        if suffix == ".json":
            return "json"
        if suffix in {".doc", ".docx", ".pptx"}:
            return "office"

        text_exts = {
            ".txt",
            ".md",
            ".rst",
            ".c",
            ".cpp",
            ".cs",
            ".go",
            ".java",
            ".js",
            ".ts",
            ".php",
            ".py",
            ".rb",
            ".sh",
            ".tex",
            ".html",
            ".css",
        }
        if suffix in text_exts:
            return "text"
        raise ValueError(f"Unsupported file type: {file_path.name} (ext={suffix})")

    # ------------------------------------------------------------------ #
    #  Public entry-point
    # ------------------------------------------------------------------ #
    async def process_file(self, file_path: Union[str, Path]) -> Dict[str, Any]:
        """Validate → detect → dispatch to the appropriate processor."""
        file_path = Path(file_path)
        self.validate_file(file_path)
        ftype = self._detect_file_type(file_path)

        dispatch_map = {
            "pdf": self._process_pdf,
            "text": self._process_text,
            "csv": self._process_csv,
            "office": self._process_office,
            "json": self._process_json,
        }
        if ftype not in dispatch_map:
            raise ValueError(f"Unsupported file type: {file_path.suffix}")

        return await dispatch_map[ftype](file_path)

    # ------------------------------------------------------------------ #
    #  PDF
    # ------------------------------------------------------------------ #
    async def _process_pdf(self, file_path: Path) -> Dict[str, Any]:
        page_chunks, doc_meta = await self._extract_text(file_path)
        all_chunks, line_data = [], []

        for page_text, page_num, line_nums in page_chunks:
            lines = page_text.split("\n")
            buf, buf_lines, length = [], [], 0
            for line, ln in zip(lines, line_nums):
                l = len(line) + 1
                if length + l <= self.chunk_size:
                    buf.append(line)
                    buf_lines.append(ln)
                    length += l
                else:
                    if buf:
                        all_chunks.append("\n".join(buf))
                        line_data.append({"page": page_num, "lines": buf_lines})
                        buf, buf_lines, length = [], [], 0
                    for piece in self._split_oversized_chunk(line):
                        all_chunks.append(piece)
                        line_data.append({"page": page_num, "lines": [ln]})
            if buf:
                all_chunks.append("\n".join(buf))
                line_data.append({"page": page_num, "lines": buf_lines})

        vectors = await asyncio.gather(
            *[self._encode_chunk_async(c) for c in all_chunks]
        )
        return {
            "content": "\n\n".join(all_chunks),
            "metadata": {
                **doc_meta,
                "source": str(file_path),
                "chunks": len(all_chunks),
                "type": "pdf",
            },
            "chunks": all_chunks,
            "vectors": [v.tolist() for v in vectors],
            "line_data": line_data,
        }

    # ------------------------------------------------------------------ #
    #  Plain-text / code / markup
    # ------------------------------------------------------------------ #
    async def _process_text(self, file_path: Path) -> Dict[str, Any]:
        text, extra_meta, _ = await self._extract_text(file_path)
        chunks = self._chunk_text(text)
        vectors = await asyncio.gather(*[self._encode_chunk_async(c) for c in chunks])
        return {
            "content": text,
            "metadata": {
                **extra_meta,
                "source": str(file_path),
                "chunks": len(chunks),
                "type": "text",
            },
            "chunks": chunks,
            "vectors": [v.tolist() for v in vectors],
        }

    # ------------------------------------------------------------------ #
    #  CSV
    # ------------------------------------------------------------------ #
    async def _process_csv(
        self, file_path: Path, text_field: str = "description"
    ) -> Dict[str, Any]:
        texts, metas = [], []
        with file_path.open(newline="", encoding="utf-8") as f:
            reader = csv.DictReader(f)
            for row in reader:
                txt = row.get(text_field, "").strip()
                if not txt:
                    continue
                texts.append(txt)
                metas.append({k: v for k, v in row.items() if k != text_field and v})

        vectors = await asyncio.gather(*[self._encode_chunk_async(t) for t in texts])
        return {
            "content": None,
            "metadata": {"source": str(file_path), "rows": len(texts), "type": "csv"},
            "chunks": texts,
            "vectors": [v.tolist() for v in vectors],
            "csv_row_metadata": metas,
        }

    # ------------------------------------------------------------------ #
    #  Office docs (.doc/.docx/.pptx)
    # ------------------------------------------------------------------ #
    async def _process_office(self, file_path: Path) -> Dict[str, Any]:
        loop = asyncio.get_event_loop()
        if file_path.suffix.lower() in {".doc", ".docx"}:
            text = await loop.run_in_executor(
                self._executor, self._read_docx, file_path
            )
        else:  # .pptx
            text = await loop.run_in_executor(
                self._executor, self._read_pptx, file_path
            )

        chunks = self._chunk_text(text)
        vectors = await asyncio.gather(*[self._encode_chunk_async(c) for c in chunks])
        return {
            "content": text,
            "metadata": {
                "source": str(file_path),
                "chunks": len(chunks),
                "type": "office",
            },
            "chunks": chunks,
            "vectors": [v.tolist() for v in vectors],
        }

    # ------------------------------------------------------------------ #
    #  JSON
    # ------------------------------------------------------------------ #
    async def _process_json(self, file_path: Path) -> Dict[str, Any]:
        text = await asyncio.get_event_loop().run_in_executor(
            self._executor, self._read_json, file_path
        )
        chunks = self._chunk_text(text)
        vectors = await asyncio.gather(*[self._encode_chunk_async(c) for c in chunks])
        return {
            "content": text,
            "metadata": {
                "source": str(file_path),
                "chunks": len(chunks),
                "type": "json",
            },
            "chunks": chunks,
            "vectors": [v.tolist() for v in vectors],
        }

    # ------------------------------------------------------------------ #
    #  Shared helpers
    # ------------------------------------------------------------------ #
    async def _extract_text(self, file_path: Path) -> Union[
        Tuple[List[Tuple[str, int, List[int]]], Dict[str, Any]],
        Tuple[str, Dict[str, Any], List[int]],
    ]:
        loop = asyncio.get_event_loop()
        if file_path.suffix.lower() == ".pdf":
            return await loop.run_in_executor(
                self._executor, self._extract_pdf_text, file_path
            )
        else:
            text = await loop.run_in_executor(
                self._executor, self._read_text_file, file_path
            )
            return text, {}, []

    def _extract_pdf_text(self, file_path: Path):
        page_chunks, meta = [], {}
        with pdfplumber.open(file_path) as pdf:
            meta.update(
                {
                    "author": pdf.metadata.get("Author", ""),
                    "title": pdf.metadata.get("Title", file_path.stem),
                    "page_count": len(pdf.pages),
                }
            )
            for i, page in enumerate(pdf.pages, start=1):
                lines = page.extract_text_lines()
                sorted_lines = sorted(lines, key=lambda x: x["top"])
                txts, nums = [], []
                for ln_idx, L in enumerate(sorted_lines, start=1):
                    t = L.get("text", "").strip()
                    if t:
                        txts.append(t)
                        nums.append(ln_idx)
                if txts:
                    page_chunks.append(("\n".join(txts), i, nums))
        return page_chunks, meta

    def _read_text_file(self, file_path: Path) -> str:
        try:
            return file_path.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            return file_path.read_text(encoding="latin-1")

    def _read_docx(self, path: Path) -> str:
        doc = Document(path)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

    def _read_pptx(self, path: Path) -> str:
        prs = Presentation(path)
        slides = []
        for slide in prs.slides:
            chunks = [sh.text for sh in slide.shapes if hasattr(sh, "text")]
            slides.append("\n".join(filter(None, chunks)))
        return "\n\n".join(slides)

    def _read_json(self, path: Path) -> str:
        obj = json.loads(path.read_text(encoding="utf-8"))
        pretty = json.dumps(obj, indent=2, ensure_ascii=False)
        return "\n".join(textwrap.wrap(pretty, width=120))

    # ------------------------------------------------------------------ #
    #  Text chunking helpers
    # ------------------------------------------------------------------ #
    def _chunk_text(self, text: str) -> List[str]:
        sentences = re.split(r"(?<=[\.!?])\s+", text)
        chunks, buf, length = [], [], 0
        for sent in sentences:
            slen = len(sent) + 1
            if length + slen <= self.chunk_size:
                buf.append(sent)
                length += slen
            else:
                if buf:
                    chunks.append(" ".join(buf))
                    buf, length = [], 0
                while len(sent) > self.chunk_size:
                    part, sent = sent[: self.chunk_size], sent[self.chunk_size :]
                    chunks.append(part)
                buf, length = [sent], len(sent)
        if buf:
            chunks.append(" ".join(buf))
        return chunks

    def _split_oversized_chunk(self, chunk: str, tokens: List[str] = None) -> List[str]:
        model = self._ensure_model()  # Ensure model is loaded to access tokenizer
        if tokens is None:
            tokens = model.tokenizer.tokenize(chunk)
        out = []
        for i in range(0, len(tokens), self.effective_max_length):
            seg = tokens[i : i + self.effective_max_length]
            out.append(model.tokenizer.convert_tokens_to_string(seg))
        return out
