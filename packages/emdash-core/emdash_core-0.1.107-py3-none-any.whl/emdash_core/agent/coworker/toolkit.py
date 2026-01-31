"""Toolkit for the CoworkerAgent.

The CoworkerToolkit provides tools for general-purpose assistance
WITHOUT coding capabilities. It focuses on:
- Web research (using the real WebTool with DuckDuckGo)
- Note-taking and memory (coworker-specific)
- Task management (same as coding agent)
- Collaboration and clarification
"""

from typing import Optional
from pathlib import Path

from ..toolkits.base import BaseToolkit
from ..tools.base import BaseTool, ToolResult


# =============================================================================
# Coworker-specific tools (notes, ideation)
# =============================================================================

class SaveNoteTool(BaseTool):
    """Save a note to the session memory."""

    name = "save_note"
    description = "Save a note to memory for later recall. Use for storing important information, decisions, or findings."

    def __init__(self, notes_storage: list):
        # Don't call super().__init__() - we don't need database connection
        self._notes = notes_storage

    def get_schema(self) -> dict:
        return self._make_schema(
            properties={
                "title": {
                    "type": "string",
                    "description": "Short title for the note",
                },
                "content": {
                    "type": "string",
                    "description": "The note content",
                },
                "tags": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "Optional tags for categorization",
                },
            },
            required=["title", "content"],
        )

    def execute(self, title: str, content: str, tags: list = None) -> ToolResult:
        note = {
            "id": len(self._notes) + 1,
            "title": title,
            "content": content,
            "tags": tags or [],
        }
        self._notes.append(note)
        return ToolResult(
            success=True,
            data={
                "note_id": note["id"],
                "message": f"Note '{title}' saved successfully",
            },
        )


class RecallNotesTool(BaseTool):
    """Recall notes from session memory."""

    name = "recall_notes"
    description = "Recall saved notes. Can filter by tag or search in content."

    def __init__(self, notes_storage: list):
        # Don't call super().__init__() - we don't need database connection
        self._notes = notes_storage

    def get_schema(self) -> dict:
        return self._make_schema(
            properties={
                "tag": {
                    "type": "string",
                    "description": "Filter by tag (optional)",
                },
                "search": {
                    "type": "string",
                    "description": "Search in note titles and content (optional)",
                },
            },
        )

    def execute(self, tag: str = None, search: str = None) -> ToolResult:
        results = self._notes

        if tag:
            results = [n for n in results if tag in n.get("tags", [])]

        if search:
            search_lower = search.lower()
            results = [
                n for n in results
                if search_lower in n["title"].lower() or search_lower in n["content"].lower()
            ]

        return ToolResult(
            success=True,
            data={
                "notes": results,
                "count": len(results),
            },
        )


class SummarizeTool(BaseTool):
    """Summarize text content."""

    name = "summarize"
    description = "Summarize provided text into key points."

    def __init__(self):
        # Don't call super().__init__() - we don't need database connection
        pass

    def get_schema(self) -> dict:
        return self._make_schema(
            properties={
                "text": {
                    "type": "string",
                    "description": "The text to summarize",
                },
                "style": {
                    "type": "string",
                    "enum": ["bullets", "paragraph", "brief"],
                    "description": "Summary style (default: bullets)",
                    "default": "bullets",
                },
            },
            required=["text"],
        )

    def execute(self, text: str, style: str = "bullets") -> ToolResult:
        return ToolResult(
            success=True,
            data={
                "original_length": len(text),
                "style": style,
                "instruction": f"Please summarize the following in {style} format:\n\n{text[:2000]}",
            },
        )


class BrainstormTool(BaseTool):
    """Brainstorm ideas on a topic."""

    name = "brainstorm"
    description = "Generate ideas on a topic with optional constraints."

    def __init__(self):
        # Don't call super().__init__() - we don't need database connection
        pass

    def get_schema(self) -> dict:
        return self._make_schema(
            properties={
                "topic": {
                    "type": "string",
                    "description": "The topic to brainstorm about",
                },
                "num_ideas": {
                    "type": "integer",
                    "description": "Number of ideas to generate (default 5)",
                    "default": 5,
                },
                "constraints": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "Constraints to consider",
                },
            },
            required=["topic"],
        )

    def execute(self, topic: str, num_ideas: int = 5, constraints: list = None) -> ToolResult:
        return ToolResult(
            success=True,
            data={
                "topic": topic,
                "num_ideas": num_ideas,
                "constraints": constraints or [],
                "instruction": f"Brainstorm {num_ideas} ideas about: {topic}",
            },
        )


class PresentOptionsTool(BaseTool):
    """Present multiple options for the user to choose from."""

    name = "present_options"
    description = "Present a set of options with pros/cons for user decision."

    def __init__(self):
        # Don't call super().__init__() - we don't need database connection
        pass

    def get_schema(self) -> dict:
        return self._make_schema(
            properties={
                "context": {
                    "type": "string",
                    "description": "Context for the decision",
                },
                "options": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "properties": {
                            "name": {"type": "string"},
                            "description": {"type": "string"},
                            "pros": {"type": "array", "items": {"type": "string"}},
                            "cons": {"type": "array", "items": {"type": "string"}},
                        },
                        "required": ["name", "description"],
                    },
                    "description": "The options to present",
                },
            },
            required=["context", "options"],
        )

    def execute(self, context: str, options: list) -> ToolResult:
        return ToolResult(
            success=True,
            data={
                "status": "awaiting_choices",
                "context": context,
                "questions": [
                    {
                        "question": f"Choose an option for: {context}",
                        "options": [opt.get("name", f"Option {i+1}") for i, opt in enumerate(options)],
                        "details": options,
                    }
                ],
            },
        )


class CreateDocumentTool(BaseTool):
    """Create Word documents (.docx) from markdown or structured content."""

    name = "create_document"
    description = (
        "Create a Word document (.docx) from markdown content. "
        "Supports headings, paragraphs, bullet lists, numbered lists, bold, and italic text. "
        "Perfect for reports, memos, proposals, and other business documents."
    )

    def __init__(self):
        # Don't call super().__init__() - we don't need database connection
        pass

    def get_schema(self) -> dict:
        return self._make_schema(
            properties={
                "filename": {
                    "type": "string",
                    "description": "Output filename (will add .docx if not present)",
                },
                "content": {
                    "type": "string",
                    "description": "Document content in Markdown format. Supports # headings, **bold**, *italic*, - bullets, 1. numbered lists",
                },
                "title": {
                    "type": "string",
                    "description": "Document title (optional, added as metadata)",
                },
                "author": {
                    "type": "string",
                    "description": "Author name (optional, added as metadata)",
                },
            },
            required=["filename", "content"],
        )

    def execute(
        self,
        filename: str,
        content: str,
        title: str = "",
        author: str = "",
    ) -> ToolResult:
        try:
            from docx import Document
            from docx.shared import Pt, Inches
            from docx.enum.text import WD_ALIGN_PARAGRAPH
        except ImportError:
            return ToolResult.error_result(
                "python-docx is not installed. Install it with: pip install python-docx",
                suggestions=["pip install python-docx"],
            )

        if not filename.endswith(".docx"):
            filename = f"{filename}.docx"

        try:
            doc = Document()

            # Set metadata
            if title:
                doc.core_properties.title = title
            if author:
                doc.core_properties.author = author

            # Parse markdown-like content
            lines = content.split("\n")
            in_list = False
            list_type = None

            for line in lines:
                stripped = line.strip()
                if not stripped:
                    in_list = False
                    continue

                # Headings
                if stripped.startswith("# "):
                    doc.add_heading(stripped[2:], level=1)
                    in_list = False
                elif stripped.startswith("## "):
                    doc.add_heading(stripped[3:], level=2)
                    in_list = False
                elif stripped.startswith("### "):
                    doc.add_heading(stripped[4:], level=3)
                    in_list = False
                # Bullet lists
                elif stripped.startswith("- ") or stripped.startswith("* "):
                    text = stripped[2:]
                    p = doc.add_paragraph(style="List Bullet")
                    self._add_formatted_text(p, text)
                    in_list = True
                    list_type = "bullet"
                # Numbered lists
                elif len(stripped) > 2 and stripped[0].isdigit() and stripped[1] == ".":
                    text = stripped[2:].strip()
                    p = doc.add_paragraph(style="List Number")
                    self._add_formatted_text(p, text)
                    in_list = True
                    list_type = "number"
                # Regular paragraph
                else:
                    p = doc.add_paragraph()
                    self._add_formatted_text(p, stripped)
                    in_list = False

            doc.save(filename)

            return ToolResult(
                success=True,
                data={
                    "filename": filename,
                    "message": f"Created document '{filename}'",
                },
            )

        except Exception as e:
            return ToolResult.error_result(f"Failed to create document: {str(e)}")

    def _add_formatted_text(self, paragraph, text: str):
        """Add text with basic markdown formatting (bold, italic)."""
        import re

        # Simple pattern matching for **bold** and *italic*
        parts = re.split(r"(\*\*[^*]+\*\*|\*[^*]+\*)", text)

        for part in parts:
            if part.startswith("**") and part.endswith("**"):
                run = paragraph.add_run(part[2:-2])
                run.bold = True
            elif part.startswith("*") and part.endswith("*"):
                run = paragraph.add_run(part[1:-1])
                run.italic = True
            else:
                paragraph.add_run(part)


class CreateSpreadsheetTool(BaseTool):
    """Create Excel spreadsheets (.xlsx) with data and formulas."""

    name = "create_spreadsheet"
    description = (
        "Create an Excel spreadsheet (.xlsx) with data, formulas, and basic formatting. "
        "Supports multiple sheets, formulas, and cell formatting. "
        "Perfect for data analysis, budgets, and reports."
    )

    def __init__(self):
        # Don't call super().__init__() - we don't need database connection
        pass

    def get_schema(self) -> dict:
        return self._make_schema(
            properties={
                "filename": {
                    "type": "string",
                    "description": "Output filename (will add .xlsx if not present)",
                },
                "sheets": {
                    "type": "array",
                    "description": "List of sheets to create",
                    "items": {
                        "type": "object",
                        "properties": {
                            "name": {
                                "type": "string",
                                "description": "Sheet name",
                            },
                            "headers": {
                                "type": "array",
                                "items": {"type": "string"},
                                "description": "Column headers (first row)",
                            },
                            "rows": {
                                "type": "array",
                                "items": {
                                    "type": "array",
                                    "items": {},
                                },
                                "description": "Data rows (can include numbers, strings, or formulas starting with =)",
                            },
                        },
                        "required": ["name", "headers", "rows"],
                    },
                },
            },
            required=["filename", "sheets"],
        )

    def execute(self, filename: str, sheets: list) -> ToolResult:
        try:
            from openpyxl import Workbook
            from openpyxl.styles import Font, PatternFill, Alignment
        except ImportError:
            return ToolResult.error_result(
                "openpyxl is not installed. Install it with: pip install openpyxl",
                suggestions=["pip install openpyxl"],
            )

        if not filename.endswith(".xlsx"):
            filename = f"{filename}.xlsx"

        try:
            wb = Workbook()
            # Remove default sheet
            default_sheet = wb.active

            for i, sheet_data in enumerate(sheets):
                sheet_name = sheet_data.get("name", f"Sheet{i+1}")
                headers = sheet_data.get("headers", [])
                rows = sheet_data.get("rows", [])

                if i == 0:
                    ws = default_sheet
                    ws.title = sheet_name
                else:
                    ws = wb.create_sheet(title=sheet_name)

                # Add headers with formatting
                header_font = Font(bold=True)
                header_fill = PatternFill(start_color="DAEEF3", end_color="DAEEF3", fill_type="solid")

                for col, header in enumerate(headers, 1):
                    cell = ws.cell(row=1, column=col, value=header)
                    cell.font = header_font
                    cell.fill = header_fill
                    cell.alignment = Alignment(horizontal="center")

                # Add data rows
                for row_idx, row_data in enumerate(rows, 2):
                    for col_idx, value in enumerate(row_data, 1):
                        cell = ws.cell(row=row_idx, column=col_idx)
                        # Handle formulas
                        if isinstance(value, str) and value.startswith("="):
                            cell.value = value
                        else:
                            cell.value = value

                # Auto-adjust column widths
                for col in ws.columns:
                    max_length = 0
                    column = col[0].column_letter
                    for cell in col:
                        try:
                            if cell.value:
                                max_length = max(max_length, len(str(cell.value)))
                        except:
                            pass
                    ws.column_dimensions[column].width = min(max_length + 2, 50)

            wb.save(filename)

            return ToolResult(
                success=True,
                data={
                    "filename": filename,
                    "sheets_count": len(sheets),
                    "message": f"Created spreadsheet '{filename}' with {len(sheets)} sheet(s)",
                },
            )

        except Exception as e:
            return ToolResult.error_result(f"Failed to create spreadsheet: {str(e)}")


class CalculateTool(BaseTool):
    """Evaluate mathematical expressions and perform calculations."""

    name = "calculate"
    description = (
        "Evaluate mathematical expressions. Supports basic arithmetic (+, -, *, /, **), "
        "percentages, unit conversions, and common math functions (sqrt, sin, cos, log, etc.). "
        "Use for calculations, data analysis, and number crunching."
    )

    def __init__(self):
        # Don't call super().__init__() - we don't need database connection
        pass

    def get_schema(self) -> dict:
        return self._make_schema(
            properties={
                "expression": {
                    "type": "string",
                    "description": "Mathematical expression to evaluate (e.g., '2 + 2', '100 * 1.15', 'sqrt(16)')",
                },
                "variables": {
                    "type": "object",
                    "description": "Optional variables to use in the expression (e.g., {'x': 10, 'y': 20})",
                    "additionalProperties": {"type": "number"},
                },
            },
            required=["expression"],
        )

    def execute(self, expression: str, variables: dict = None) -> ToolResult:
        import math
        import re

        # Safe math functions
        safe_funcs = {
            "abs": abs,
            "round": round,
            "min": min,
            "max": max,
            "sum": sum,
            "sqrt": math.sqrt,
            "sin": math.sin,
            "cos": math.cos,
            "tan": math.tan,
            "log": math.log,
            "log10": math.log10,
            "exp": math.exp,
            "pow": pow,
            "ceil": math.ceil,
            "floor": math.floor,
            "pi": math.pi,
            "e": math.e,
        }

        # Add user variables
        if variables:
            safe_funcs.update(variables)

        # Handle percentage syntax (e.g., "100 + 15%" -> "100 * 1.15")
        expression = re.sub(r"(\d+(?:\.\d+)?)\s*%", r"(\1/100)", expression)

        # Validate expression (only allow safe characters)
        allowed_chars = set("0123456789+-*/.()%, abcdefghijklmnopqrstuvwxyzABCDEFGHIJKLMNOPQRSTUVWXYZ_")
        if not all(c in allowed_chars for c in expression):
            return ToolResult.error_result("Expression contains invalid characters")

        try:
            # Evaluate safely
            result = eval(expression, {"__builtins__": {}}, safe_funcs)

            return ToolResult(
                success=True,
                data={
                    "expression": expression,
                    "result": result,
                    "formatted": f"{expression} = {result}",
                },
            )

        except ZeroDivisionError:
            return ToolResult.error_result("Division by zero")
        except Exception as e:
            return ToolResult.error_result(f"Calculation error: {str(e)}")


class SendEmailTool(BaseTool):
    """Send emails via SMTP."""

    name = "send_email"
    description = (
        "Send an email via SMTP. Requires SMTP configuration via environment variables "
        "(SMTP_HOST, SMTP_PORT, SMTP_USER, SMTP_PASSWORD, SMTP_FROM). "
        "For Gmail, use smtp.gmail.com:587 with an app password."
    )

    def __init__(self):
        # Don't call super().__init__() - we don't need database connection
        pass

    def get_schema(self) -> dict:
        return self._make_schema(
            properties={
                "to": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "Recipient email addresses",
                },
                "subject": {
                    "type": "string",
                    "description": "Email subject line",
                },
                "body": {
                    "type": "string",
                    "description": "Email body (plain text or HTML)",
                },
                "html": {
                    "type": "boolean",
                    "description": "If true, body is treated as HTML (default: false)",
                    "default": False,
                },
                "cc": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "CC recipients (optional)",
                },
            },
            required=["to", "subject", "body"],
        )

    def execute(
        self,
        to: list,
        subject: str,
        body: str,
        html: bool = False,
        cc: list = None,
    ) -> ToolResult:
        import os
        import smtplib
        from email.mime.text import MIMEText
        from email.mime.multipart import MIMEMultipart

        # Get SMTP config from environment
        smtp_host = os.getenv("SMTP_HOST")
        smtp_port = int(os.getenv("SMTP_PORT", "587"))
        smtp_user = os.getenv("SMTP_USER")
        smtp_password = os.getenv("SMTP_PASSWORD")
        smtp_from = os.getenv("SMTP_FROM", smtp_user)

        if not all([smtp_host, smtp_user, smtp_password]):
            return ToolResult.error_result(
                "SMTP not configured. Set SMTP_HOST, SMTP_USER, SMTP_PASSWORD environment variables.",
                suggestions=[
                    "export SMTP_HOST=smtp.gmail.com",
                    "export SMTP_PORT=587",
                    "export SMTP_USER=your@email.com",
                    "export SMTP_PASSWORD=your-app-password",
                ],
            )

        try:
            # Create message
            msg = MIMEMultipart("alternative")
            msg["Subject"] = subject
            msg["From"] = smtp_from
            msg["To"] = ", ".join(to)
            if cc:
                msg["Cc"] = ", ".join(cc)

            # Attach body
            content_type = "html" if html else "plain"
            msg.attach(MIMEText(body, content_type))

            # Send
            all_recipients = to + (cc or [])

            with smtplib.SMTP(smtp_host, smtp_port) as server:
                server.starttls()
                server.login(smtp_user, smtp_password)
                server.sendmail(smtp_from, all_recipients, msg.as_string())

            return ToolResult(
                success=True,
                data={
                    "message": f"Email sent to {', '.join(to)}",
                    "subject": subject,
                    "recipients": len(all_recipients),
                },
            )

        except smtplib.SMTPAuthenticationError:
            return ToolResult.error_result("SMTP authentication failed. Check credentials.")
        except Exception as e:
            return ToolResult.error_result(f"Failed to send email: {str(e)}")


class GmailSendTool(BaseTool):
    """Send emails via Gmail API using OAuth authentication."""

    name = "gmail_send"
    description = (
        "Send an email using Gmail API with OAuth. Requires Google authentication "
        "(run /auth google login first). This is the preferred way to send emails "
        "if you have authenticated with Google. Supports file attachments."
    )

    def __init__(self):
        pass

    def get_schema(self) -> dict:
        return self._make_schema(
            properties={
                "to": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "Recipient email addresses",
                },
                "subject": {
                    "type": "string",
                    "description": "Email subject line",
                },
                "body": {
                    "type": "string",
                    "description": "Email body (plain text or HTML)",
                },
                "cc": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "CC recipients (optional)",
                },
                "attachments": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "File paths to attach (optional). Supports any file type.",
                },
                "html": {
                    "type": "boolean",
                    "description": "If true, body is treated as HTML (optional, default false)",
                },
            },
            required=["to", "subject", "body"],
        )

    def execute(
        self,
        to: list,
        subject: str,
        body: str,
        cc: list = None,
        attachments: list = None,
        html: bool = False,
    ) -> ToolResult:
        import base64
        import mimetypes
        import os
        from email.mime.text import MIMEText
        from email.mime.multipart import MIMEMultipart
        from email.mime.base import MIMEBase
        from email import encoders
        import requests

        from ...auth.google import get_google_access_token, is_google_authenticated

        if not is_google_authenticated():
            return ToolResult.error_result(
                "Not authenticated with Google. Run /auth google login first.",
                suggestions=["/auth google login"],
            )

        access_token = get_google_access_token()
        if not access_token:
            return ToolResult.error_result(
                "Failed to get Google access token. Try /auth google login again.",
                suggestions=["/auth google login"],
            )

        try:
            # Create the email message
            if attachments:
                # Use multipart for attachments
                message = MIMEMultipart()
                message["to"] = ", ".join(to)
                message["subject"] = subject
                if cc:
                    message["cc"] = ", ".join(cc)

                # Add body
                body_type = "html" if html else "plain"
                message.attach(MIMEText(body, body_type))

                # Add attachments
                attached_files = []
                for file_path in attachments:
                    if not os.path.exists(file_path):
                        return ToolResult.error_result(
                            f"Attachment not found: {file_path}",
                            suggestions=["Check the file path exists"],
                        )

                    # Guess the MIME type
                    mime_type, _ = mimetypes.guess_type(file_path)
                    if mime_type is None:
                        mime_type = "application/octet-stream"

                    main_type, sub_type = mime_type.split("/", 1)

                    with open(file_path, "rb") as f:
                        file_data = f.read()

                    attachment = MIMEBase(main_type, sub_type)
                    attachment.set_payload(file_data)
                    encoders.encode_base64(attachment)

                    filename = os.path.basename(file_path)
                    attachment.add_header(
                        "Content-Disposition",
                        "attachment",
                        filename=filename,
                    )
                    message.attach(attachment)
                    attached_files.append(filename)
            else:
                # Simple text message (no attachments)
                body_type = "html" if html else "plain"
                message = MIMEText(body, body_type)
                message["to"] = ", ".join(to)
                message["subject"] = subject
                if cc:
                    message["cc"] = ", ".join(cc)
                attached_files = []

            # Encode the message in base64url format
            raw_message = base64.urlsafe_b64encode(message.as_bytes()).decode("utf-8")

            # Send via Gmail API
            response = requests.post(
                "https://gmail.googleapis.com/gmail/v1/users/me/messages/send",
                headers={
                    "Authorization": f"Bearer {access_token}",
                    "Content-Type": "application/json",
                },
                json={"raw": raw_message},
                timeout=60,  # Longer timeout for attachments
            )

            if response.status_code == 200:
                result = response.json()
                msg = f"Email sent to {', '.join(to)}"
                if attached_files:
                    msg += f" with {len(attached_files)} attachment(s): {', '.join(attached_files)}"
                return ToolResult(
                    success=True,
                    data={
                        "message": msg,
                        "subject": subject,
                        "message_id": result.get("id"),
                        "thread_id": result.get("threadId"),
                        "attachments": attached_files,
                    },
                )
            elif response.status_code == 401:
                return ToolResult.error_result(
                    "Google token expired. Run /auth google login to re-authenticate.",
                    suggestions=["/auth google login"],
                )
            else:
                error_data = response.json()
                error_msg = error_data.get("error", {}).get("message", response.text)
                return ToolResult.error_result(f"Gmail API error: {error_msg}")

        except Exception as e:
            return ToolResult.error_result(f"Failed to send email: {str(e)}")


class CreateImageTool(BaseTool):
    """Generate images using AI (Fireworks flux-kontext-max model)."""

    name = "create_image"
    description = (
        "Generate an image using AI from a text prompt. "
        "Uses the Fireworks flux-kontext-max model for high-quality image generation. "
        "Supports various aspect ratios (1:1, 16:9, 9:16, 4:3, 3:4). "
        "Returns the path to the saved image file. "
        "Requires FIREWORKS_API_KEY environment variable."
    )

    def __init__(self):
        # Don't call super().__init__() - we don't need database connection
        pass

    def get_schema(self) -> dict:
        return self._make_schema(
            properties={
                "prompt": {
                    "type": "string",
                    "description": "Text description of the image to generate (be detailed and specific)",
                },
                "filename": {
                    "type": "string",
                    "description": "Output filename without extension (will add .jpg). If not provided, uses a timestamp.",
                },
                "aspect_ratio": {
                    "type": "string",
                    "enum": ["1:1", "16:9", "9:16", "4:3", "3:4"],
                    "description": "Aspect ratio for the generated image (default: 1:1)",
                    "default": "1:1",
                },
            },
            required=["prompt"],
        )

    def execute(
        self,
        prompt: str,
        filename: str = None,
        aspect_ratio: str = "1:1",
    ) -> ToolResult:
        import os
        import time
        import requests
        import base64
        from datetime import datetime

        # Get API key from environment
        api_key = os.getenv("FIREWORKS_API_KEY")
        if not api_key:
            return ToolResult.error_result(
                "FIREWORKS_API_KEY environment variable is not set. "
                "Get your API key from https://fireworks.ai and set it in .env",
                suggestions=["export FIREWORKS_API_KEY=your-api-key"],
            )

        # Generate filename if not provided
        if not filename:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"generated_image_{timestamp}"

        # Ensure filename has .jpg extension
        if not filename.endswith(".jpg") and not filename.endswith(".jpeg"):
            filename = f"{filename}.jpg"

        # Step 1: Submit the generation request
        url = "https://api.fireworks.ai/inference/v1/workflows/accounts/fireworks/models/flux-kontext-max"
        headers = {
            "Content-Type": "application/json",
            "Accept": "application/json",
            "Authorization": f"Bearer {api_key}",
        }
        data = {
            "prompt": prompt,
            "aspect_ratio": aspect_ratio,
            "seed": -1,  # Random seed
        }

        try:
            response = requests.post(url, headers=headers, json=data, timeout=30)
            response.raise_for_status()
            result = response.json()
            request_id = result.get("request_id")

            if not request_id:
                return ToolResult.error_result(
                    f"No request ID returned from Fireworks API: {result}"
                )

        except requests.exceptions.RequestException as e:
            return ToolResult.error_result(f"Failed to submit image generation request: {str(e)}")

        # Step 2: Poll for the result
        result_endpoint = f"{url}/get_result"
        max_attempts = 120  # 2 minutes max
        poll_interval = 1  # 1 second between polls

        for attempt in range(max_attempts):
            time.sleep(poll_interval)

            try:
                result_response = requests.post(
                    result_endpoint,
                    headers={
                        "Content-Type": "application/json",
                        "Accept": "image/jpeg",
                        "Authorization": f"Bearer {api_key}",
                    },
                    json={"id": request_id},
                    timeout=30,
                )

                poll_result = result_response.json()
                status = poll_result.get("status")

                if status in ["Ready", "Complete", "Finished"]:
                    image_data = poll_result.get("result", {}).get("sample")

                    if isinstance(image_data, str) and image_data.startswith("http"):
                        # Download from URL
                        image_response = requests.get(image_data, timeout=60)
                        image_response.raise_for_status()
                        with open(filename, "wb") as f:
                            f.write(image_response.content)
                    elif image_data:
                        # Base64 data
                        with open(filename, "wb") as f:
                            f.write(base64.b64decode(image_data))
                    else:
                        return ToolResult.error_result("No image data in response")

                    # Get absolute path for clarity
                    abs_path = os.path.abspath(filename)

                    return ToolResult(
                        success=True,
                        data={
                            "filename": filename,
                            "absolute_path": abs_path,
                            "prompt": prompt,
                            "aspect_ratio": aspect_ratio,
                            "message": f"Image generated and saved to '{filename}'",
                        },
                    )

                if status in ["Failed", "Error"]:
                    error_details = poll_result.get("details", poll_result.get("error", "Unknown error"))
                    return ToolResult.error_result(f"Image generation failed: {error_details}")

                # Still processing, continue polling
                if attempt % 10 == 0 and attempt > 0:
                    # Log progress every 10 seconds
                    pass

            except requests.exceptions.RequestException as e:
                # Network error during polling, continue trying
                if attempt >= max_attempts - 1:
                    return ToolResult.error_result(f"Polling failed after {max_attempts} attempts: {str(e)}")
                continue

        return ToolResult.error_result(
            f"Image generation timed out after {max_attempts} seconds. "
            f"Request ID: {request_id}"
        )


class CreatePresentationTool(BaseTool):
    """Create PowerPoint presentations with slides."""

    name = "create_presentation"
    description = (
        "Create a PowerPoint (.pptx) presentation with multiple slides. "
        "Supports title slides, content slides with bullet points, and section headers. "
        "Perfect for creating reports, proposals, or any structured presentation."
    )

    def __init__(self):
        # Don't call super().__init__() - we don't need database connection
        pass

    def get_schema(self) -> dict:
        return self._make_schema(
            properties={
                "filename": {
                    "type": "string",
                    "description": "Output filename (will add .pptx if not present). Saved to current directory.",
                },
                "title": {
                    "type": "string",
                    "description": "Presentation title (shown on title slide)",
                },
                "subtitle": {
                    "type": "string",
                    "description": "Presentation subtitle (optional, shown on title slide)",
                },
                "author": {
                    "type": "string",
                    "description": "Author name (optional, added to metadata)",
                },
                "slides": {
                    "type": "array",
                    "description": "List of slides to add after the title slide",
                    "items": {
                        "type": "object",
                        "properties": {
                            "layout": {
                                "type": "string",
                                "enum": ["title_content", "section", "title_only", "blank", "two_column"],
                                "description": "Slide layout type",
                            },
                            "title": {
                                "type": "string",
                                "description": "Slide title",
                            },
                            "content": {
                                "type": "array",
                                "items": {"type": "string"},
                                "description": "Bullet points or paragraphs for the slide",
                            },
                            "left_content": {
                                "type": "array",
                                "items": {"type": "string"},
                                "description": "Left column content (for two_column layout)",
                            },
                            "right_content": {
                                "type": "array",
                                "items": {"type": "string"},
                                "description": "Right column content (for two_column layout)",
                            },
                            "notes": {
                                "type": "string",
                                "description": "Speaker notes for this slide",
                            },
                        },
                        "required": ["layout", "title"],
                    },
                },
            },
            required=["filename", "title", "slides"],
        )

    def execute(
        self,
        filename: str,
        title: str,
        slides: list,
        subtitle: str = "",
        author: str = "",
    ) -> ToolResult:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
        except ImportError:
            return ToolResult.error_result(
                "python-pptx is not installed. Install it with: pip install python-pptx",
                suggestions=["pip install python-pptx"],
            )

        # Ensure filename has .pptx extension
        if not filename.endswith(".pptx"):
            filename = f"{filename}.pptx"

        try:
            # Create presentation
            prs = Presentation()
            prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
            prs.slide_height = Inches(7.5)

            # Set metadata
            if author:
                prs.core_properties.author = author
            prs.core_properties.title = title

            # Add title slide
            title_slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(title_slide_layout)

            # Add title text box
            left = Inches(0.5)
            top = Inches(2.5)
            width = Inches(12.333)
            height = Inches(1.5)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(44)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

            # Add subtitle if provided
            if subtitle:
                left = Inches(0.5)
                top = Inches(4.2)
                width = Inches(12.333)
                height = Inches(1)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = subtitle
                p.font.size = Pt(24)
                p.alignment = PP_ALIGN.CENTER

            # Add content slides
            slides_added = 0
            for slide_data in slides:
                layout = slide_data.get("layout", "title_content")
                slide_title = slide_data.get("title", "")
                content = slide_data.get("content", [])
                notes = slide_data.get("notes", "")
                left_content = slide_data.get("left_content", [])
                right_content = slide_data.get("right_content", [])

                # Use blank layout and add custom shapes
                blank_layout = prs.slide_layouts[6]
                new_slide = prs.slides.add_slide(blank_layout)

                if layout == "section":
                    # Section header - centered title
                    left = Inches(0.5)
                    top = Inches(3)
                    width = Inches(12.333)
                    height = Inches(1.5)
                    txBox = new_slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = slide_title
                    p.font.size = Pt(40)
                    p.font.bold = True
                    p.alignment = PP_ALIGN.CENTER

                elif layout == "two_column":
                    # Title
                    self._add_title(new_slide, slide_title)

                    # Left column
                    left = Inches(0.5)
                    top = Inches(1.8)
                    width = Inches(5.9)
                    height = Inches(5)
                    self._add_content_box(new_slide, left, top, width, height, left_content)

                    # Right column
                    left = Inches(6.9)
                    self._add_content_box(new_slide, left, top, width, height, right_content)

                elif layout == "title_only":
                    # Just title, no content
                    self._add_title(new_slide, slide_title)

                elif layout == "blank":
                    # Blank slide - add content if provided
                    if content:
                        left = Inches(0.5)
                        top = Inches(0.5)
                        width = Inches(12.333)
                        height = Inches(6.5)
                        self._add_content_box(new_slide, left, top, width, height, content)

                else:  # title_content (default)
                    # Title at top
                    self._add_title(new_slide, slide_title)

                    # Content below
                    if content:
                        left = Inches(0.5)
                        top = Inches(1.8)
                        width = Inches(12.333)
                        height = Inches(5)
                        self._add_content_box(new_slide, left, top, width, height, content)

                # Add speaker notes
                if notes:
                    notes_slide = new_slide.notes_slide
                    notes_slide.notes_text_frame.text = notes

                slides_added += 1

            # Save presentation
            prs.save(filename)

            return ToolResult(
                success=True,
                data={
                    "filename": filename,
                    "slides_count": slides_added + 1,  # +1 for title slide
                    "message": f"Created presentation '{filename}' with {slides_added + 1} slides",
                },
            )

        except Exception as e:
            return ToolResult.error_result(f"Failed to create presentation: {str(e)}")

    def _add_title(self, slide, title_text: str):
        """Add a title to a slide."""
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        left = Inches(0.5)
        top = Inches(0.4)
        width = Inches(12.333)
        height = Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(32)
        p.font.bold = True

    def _add_content_box(self, slide, left, top, width, height, content: list):
        """Add a content box with bullet points."""
        from pptx.util import Pt
        from pptx.enum.text import PP_ALIGN

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, item in enumerate(content):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}" if not item.startswith("•") else item
            p.font.size = Pt(18)
            p.space_after = Pt(12)


# =============================================================================
# CoworkerToolkit
# =============================================================================

class CoworkerToolkit(BaseToolkit):
    """Toolkit for general-purpose coworker agents.

    Provides tools for:
    - Web research (search, fetch) - SHARED with coding agent
    - Note-taking (save, recall) - COWORKER ONLY
    - Task management (todos) - SHARED with coding agent
    - Collaboration (clarification, options)
    - Ideation (brainstorm, summarize) - COWORKER ONLY
    - Document creation (PPTX, DOCX, XLSX) - COWORKER ONLY (built-in)
    - Calculations (math, percentages) - COWORKER ONLY
    - Email sending (SMTP) - COWORKER ONLY
    - Skills (skill, list_skills) - SHARED via BaseToolkit
    - MCP servers - SHARED via BaseToolkit

    Does NOT provide:
    - File system access (except document creation)
    - Code execution
    - System commands

    Built-in Tools:
        - create_presentation: PowerPoint (.pptx)
        - create_document: Word documents (.docx)
        - create_spreadsheet: Excel spreadsheets (.xlsx)
        - create_image: AI image generation (Fireworks flux-kontext-max)
        - calculate: Math expressions and calculations
        - send_email: Send emails via SMTP

    MCP Server Options (more powerful, requires setup):
        Documents:
        - ms-office: Advanced DOCX/PPTX/XLSX (Docker)
        - markitdown: Convert any document to Markdown
        - pdf-forms: Fill PDF forms

        Productivity:
        - google-workspace: Gmail, Calendar, Drive, Docs, Sheets
        - memory (mem0): Long-term memory across sessions
        - replicate: AI image generation
        - puppeteer: Screenshots and web automation
        - slack: Slack messaging

        Configure MCP servers in .emdash/mcp.json for advanced features.

    Inherits extensibility features from BaseToolkit:
    - MCP servers from .emdash/mcp.json
    - Skills from .emdash/skills/
    - Rules from .emdash/rules/ (via load_rules())
    """

    def __init__(
        self,
        repo_root: Optional[Path] = None,
        enable_skills: bool = True,
        enable_mcp_config: bool = True,
        mcp_config_path: Optional[Path] = None,
    ):
        """Initialize the coworker toolkit.

        Args:
            repo_root: Root directory (optional for coworker, used by BaseToolkit)
            enable_skills: If True, register skill tools (default: True)
            enable_mcp_config: If True, load MCP servers from config (default: True)
            mcp_config_path: Path to MCP config file (defaults to .emdash/mcp.json)
        """
        # Shared storage for notes - must be initialized before super().__init__
        # because _register_tools() is called by BaseToolkit.__init__
        self._notes: list[dict] = []

        # Initialize base class with shared extensibility features
        # This will:
        # 1. Call _register_tools() to register our tools
        # 2. Register skill tools if enable_skills=True
        # 3. Load MCP servers from config if enable_mcp_config=True
        super().__init__(
            repo_root=repo_root or Path.cwd(),
            enable_skills=enable_skills,
            enable_mcp_config=enable_mcp_config,
            mcp_config_path=mcp_config_path,
        )

    def get_tools(self) -> list[BaseTool]:
        """Return list of coworker tools.

        Uses shared tools from the coding agent where appropriate:
        - WebTool (real DuckDuckGo search + BeautifulSoup fetch)
        - WriteTodoTool, UpdateTodoListTool, AskChoiceQuestionsTool
        - TaskTool (for spawning sub-agents: Researcher, GeneralPlanner)

        Returns:
            List of BaseTool instances for coworker functionality
        """
        # Import shared tools from the coding agent
        from ..tools.web import WebTool
        from ..tools.task import TaskTool
        from ..tools.tasks import (
            WriteTodoTool,
            UpdateTodoListTool,
            AskChoiceQuestionsTool,
        )

        return [
            # Web tools - SHARED (real implementation with DuckDuckGo)
            # Pass connection=False to skip database connection (not needed)
            WebTool(connection=False),

            # Note tools - COWORKER ONLY
            SaveNoteTool(self._notes),
            RecallNotesTool(self._notes),

            # Ideation tools - COWORKER ONLY
            SummarizeTool(),
            BrainstormTool(),

            # Collaboration tools - COWORKER ONLY
            PresentOptionsTool(),

            # Document creation tools - COWORKER ONLY
            CreatePresentationTool(),
            CreateDocumentTool(),
            CreateSpreadsheetTool(),

            # Image generation - COWORKER ONLY
            CreateImageTool(),

            # Utility tools - COWORKER ONLY
            CalculateTool(),
            # SendEmailTool(),  # SMTP-based (disabled - use GmailSendTool instead)
            GmailSendTool(),  # OAuth-based (requires /auth google login)

            # Task tools - SHARED with coding agent
            WriteTodoTool(),
            UpdateTodoListTool(),
            AskChoiceQuestionsTool(),

            # Sub-agent spawning - can spawn Researcher and GeneralPlanner
            # Pass connection=False to skip database connection (not needed)
            TaskTool(repo_root=self.repo_root, connection=False),
        ]

    def get_notes(self) -> list[dict]:
        """Get all saved notes."""
        return self._notes.copy()

    def get_todos(self) -> list[dict]:
        """Get all todos from TaskState (shared with coding agent)."""
        from ..tools.tasks import TaskState
        return TaskState.get_instance().get_all_tasks()

    def reset_session(self) -> None:
        """Reset session state."""
        self._notes.clear()
        # Also reset shared task state
        from ..tools.tasks import TaskState
        TaskState.reset()
