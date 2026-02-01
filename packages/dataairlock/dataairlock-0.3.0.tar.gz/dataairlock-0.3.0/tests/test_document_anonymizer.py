"""ドキュメント匿名化テスト"""

import tempfile
from pathlib import Path

import pytest
from docx import Document
from pptx import Presentation
from typer.testing import CliRunner

from dataairlock.cli import app
from dataairlock.document_anonymizer import (
    DocumentAnonymizer,
    DocumentPIIResult,
    anonymize_document,
    deanonymize_document,
    scan_document,
)
from dataairlock.anonymizer import load_mapping, save_mapping


runner = CliRunner()


@pytest.fixture
def sample_docx(tmp_path):
    """サンプルWordファイルを作成"""
    docx_path = tmp_path / "test_doc.docx"
    doc = Document()
    doc.add_heading("患者情報", 0)
    doc.add_paragraph("患者名: 山田太郎")
    doc.add_paragraph("電話番号: 090-1234-5678")
    doc.add_paragraph("メール: yamada@example.com")
    doc.add_paragraph("住所: 東京都新宿区西新宿1-1-1")
    doc.add_paragraph("生年月日: 1990/01/15")
    doc.save(docx_path)
    return docx_path


@pytest.fixture
def sample_docx_with_table(tmp_path):
    """テーブル付きWordファイルを作成"""
    docx_path = tmp_path / "test_doc_table.docx"
    doc = Document()
    doc.add_heading("患者一覧", 0)

    table = doc.add_table(rows=3, cols=3)
    table.cell(0, 0).text = "氏名"
    table.cell(0, 1).text = "電話番号"
    table.cell(0, 2).text = "メール"
    table.cell(1, 0).text = "山田太郎"
    table.cell(1, 1).text = "090-1111-2222"
    table.cell(1, 2).text = "yamada@test.com"
    table.cell(2, 0).text = "鈴木花子"
    table.cell(2, 1).text = "080-3333-4444"
    table.cell(2, 2).text = "suzuki@test.com"

    doc.save(docx_path)
    return docx_path


@pytest.fixture
def sample_pptx(tmp_path):
    """サンプルPowerPointファイルを作成"""
    pptx_path = tmp_path / "test_pres.pptx"
    prs = Presentation()

    # スライド1: タイトル
    slide_layout = prs.slide_layouts[5]  # 空白スライド
    slide = prs.slides.add_slide(slide_layout)
    shape = slide.shapes.add_textbox(
        left=914400, top=914400, width=5486400, height=914400  # 1inch = 914400 EMU
    )
    shape.text_frame.paragraphs[0].text = "患者情報レポート"

    # スライド2: 内容
    slide2 = prs.slides.add_slide(slide_layout)
    shape2 = slide2.shapes.add_textbox(
        left=914400, top=914400, width=5486400, height=3657600
    )
    tf = shape2.text_frame
    tf.paragraphs[0].text = "担当者: 佐藤一郎"
    p1 = tf.add_paragraph()
    p1.text = "連絡先: 03-1234-5678"
    p2 = tf.add_paragraph()
    p2.text = "メール: sato@hospital.jp"

    prs.save(pptx_path)
    return pptx_path


@pytest.fixture
def sample_docx_no_pii(tmp_path):
    """PII無しWordファイルを作成"""
    docx_path = tmp_path / "no_pii.docx"
    doc = Document()
    doc.add_heading("製品仕様書", 0)
    doc.add_paragraph("製品コード: PRD-001")
    doc.add_paragraph("カテゴリ: 電子機器")
    doc.add_paragraph("価格: 10,000円")
    doc.save(docx_path)
    return docx_path


class TestDocumentAnonymizer:
    """DocumentAnonymizerのテスト"""

    def test_scan_docx(self, sample_docx):
        """Word文書のスキャン"""
        result = scan_document(sample_docx)

        assert result.file_type == "docx"
        assert result.total_matches > 0
        assert len(result.pii_by_type) > 0

    def test_scan_docx_with_table(self, sample_docx_with_table):
        """テーブル付きWord文書のスキャン"""
        result = scan_document(sample_docx_with_table)

        assert result.file_type == "docx"
        assert result.total_matches > 0
        # 電話番号とメールを検出
        assert "電話番号" in result.pii_by_type or "メールアドレス" in result.pii_by_type

    def test_scan_pptx(self, sample_pptx):
        """PowerPointのスキャン"""
        result = scan_document(sample_pptx)

        assert result.file_type == "pptx"
        assert result.total_matches > 0

    def test_scan_no_pii(self, sample_docx_no_pii):
        """PII無しファイルのスキャン"""
        result = scan_document(sample_docx_no_pii)

        assert result.file_type == "docx"
        assert result.total_matches == 0

    def test_anonymize_docx(self, sample_docx, tmp_path):
        """Word文書の匿名化"""
        output_path = tmp_path / "anonymized.docx"
        result, mapping = anonymize_document(sample_docx, output_path, "replace")

        assert output_path.exists()
        assert result.total_matches > 0
        assert "values" in mapping
        assert len(mapping["values"]) > 0

        # 匿名化されたファイルを読み込んで確認
        doc = Document(output_path)
        full_text = "\n".join([p.text for p in doc.paragraphs])
        # セマンティックID（PHONE_, EMAIL_, ADDR_など）が含まれていることを確認
        assert any(prefix in full_text for prefix in ["PHONE_", "EMAIL_", "ADDR_", "BIRTHDATE_"])
        # 元の電話番号がないことを確認
        assert "090-1234-5678" not in full_text

    def test_anonymize_docx_generalize(self, sample_docx, tmp_path):
        """一般化戦略での匿名化"""
        output_path = tmp_path / "generalized.docx"
        result, mapping = anonymize_document(sample_docx, output_path, "generalize")

        assert output_path.exists()
        assert result.total_matches > 0

    def test_anonymize_pptx(self, sample_pptx, tmp_path):
        """PowerPointの匿名化"""
        output_path = tmp_path / "anonymized.pptx"
        result, mapping = anonymize_document(sample_pptx, output_path, "replace")

        assert output_path.exists()
        assert result.total_matches > 0
        assert "values" in mapping

    def test_deanonymize_docx(self, sample_docx, tmp_path):
        """Word文書の復元"""
        anonymized_path = tmp_path / "anonymized.docx"
        restored_path = tmp_path / "restored.docx"

        # 匿名化
        result, mapping = anonymize_document(sample_docx, anonymized_path, "replace")

        # 復元
        deanonymize_document(anonymized_path, restored_path, mapping)

        assert restored_path.exists()

        # 復元されたファイルを確認
        doc = Document(restored_path)
        full_text = "\n".join([p.text for p in doc.paragraphs])
        # 元の電話番号が復元されていることを確認
        assert "090-1234-5678" in full_text

    def test_deanonymize_pptx(self, sample_pptx, tmp_path):
        """PowerPointの復元"""
        anonymized_path = tmp_path / "anonymized.pptx"
        restored_path = tmp_path / "restored.pptx"

        # 匿名化
        result, mapping = anonymize_document(sample_pptx, anonymized_path, "replace")

        # 復元
        deanonymize_document(anonymized_path, restored_path, mapping)

        assert restored_path.exists()


class TestScanDocCommand:
    """scan-doc コマンドのテスト"""

    def test_scan_doc_basic(self, sample_docx):
        """基本的なスキャン"""
        result = runner.invoke(app, ["scan-doc", str(sample_docx)])

        assert result.exit_code == 0
        assert "個人情報を検出" in result.output or "PII" in result.output

    def test_scan_doc_pptx(self, sample_pptx):
        """PowerPointのスキャン"""
        result = runner.invoke(app, ["scan-doc", str(sample_pptx)])

        assert result.exit_code == 0

    def test_scan_doc_no_pii(self, sample_docx_no_pii):
        """PII無しファイル"""
        result = runner.invoke(app, ["scan-doc", str(sample_docx_no_pii)])

        assert result.exit_code == 0
        assert "検出されませんでした" in result.output

    def test_scan_doc_not_found(self, tmp_path):
        """存在しないファイル"""
        result = runner.invoke(app, ["scan-doc", str(tmp_path / "nonexistent.docx")])

        assert result.exit_code == 1
        assert "ファイルが見つかりません" in result.output

    def test_scan_doc_unsupported_format(self, tmp_path):
        """サポートされていないフォーマット"""
        txt_file = tmp_path / "test.txt"
        txt_file.write_text("test content")

        result = runner.invoke(app, ["scan-doc", str(txt_file)])

        assert result.exit_code == 1
        assert "サポートされていない" in result.output


class TestAnonymizeDocCommand:
    """anonymize-doc コマンドのテスト"""

    def test_anonymize_doc_basic(self, sample_docx, tmp_path):
        """基本的な匿名化"""
        output_path = tmp_path / "output" / "anonymized.docx"

        result = runner.invoke(app, [
            "anonymize-doc",
            str(sample_docx),
            "-o", str(output_path),
            "-p", "testpassword123",
        ])

        assert result.exit_code == 0
        assert output_path.exists()
        assert (tmp_path / "output" / "anonymized.mapping.enc").exists()

    def test_anonymize_doc_pptx(self, sample_pptx, tmp_path):
        """PowerPointの匿名化"""
        output_path = tmp_path / "output" / "anonymized.pptx"

        result = runner.invoke(app, [
            "anonymize-doc",
            str(sample_pptx),
            "-o", str(output_path),
            "-p", "testpassword123",
        ])

        assert result.exit_code == 0
        assert output_path.exists()

    def test_anonymize_doc_generalize(self, sample_docx, tmp_path):
        """一般化戦略"""
        output_path = tmp_path / "output" / "generalized.docx"

        result = runner.invoke(app, [
            "anonymize-doc",
            str(sample_docx),
            "-o", str(output_path),
            "-p", "testpassword123",
            "-s", "generalize",
        ])

        assert result.exit_code == 0


class TestRestoreDocCommand:
    """restore-doc コマンドのテスト"""

    def test_restore_doc_basic(self, sample_docx, tmp_path):
        """基本的な復元"""
        # まず匿名化
        anonymized_path = tmp_path / "anonymized.docx"
        mapping_path = tmp_path / "anonymized.mapping.enc"

        result, mapping = anonymize_document(sample_docx, anonymized_path, "replace")
        save_mapping(mapping, mapping_path, "testpassword123")

        # 復元
        restored_path = tmp_path / "restored.docx"
        result = runner.invoke(app, [
            "restore-doc",
            str(anonymized_path),
            "-m", str(mapping_path),
            "-o", str(restored_path),
            "-p", "testpassword123",
        ])

        assert result.exit_code == 0
        assert restored_path.exists()


class TestMappingPersistence:
    """マッピング保存・復元のテスト"""

    def test_mapping_save_load(self, sample_docx, tmp_path):
        """マッピングの保存と読み込み"""
        anonymized_path = tmp_path / "anonymized.docx"
        mapping_path = tmp_path / "test.mapping.enc"
        password = "testpassword123"

        # 匿名化してマッピング取得
        result, mapping = anonymize_document(sample_docx, anonymized_path, "replace")

        # マッピング保存
        save_mapping(mapping, mapping_path, password)
        assert mapping_path.exists()

        # マッピング読み込み
        loaded_mapping = load_mapping(mapping_path, password)
        assert loaded_mapping["values"] == mapping["values"]

    def test_mapping_wrong_password(self, sample_docx, tmp_path):
        """間違ったパスワード"""
        anonymized_path = tmp_path / "anonymized.docx"
        mapping_path = tmp_path / "test.mapping.enc"

        result, mapping = anonymize_document(sample_docx, anonymized_path, "replace")
        save_mapping(mapping, mapping_path, "correctpassword")

        # 間違ったパスワードで読み込み
        with pytest.raises(ValueError):
            load_mapping(mapping_path, "wrongpassword")
