"""Word/PowerPointドキュメント匿名化モジュール"""

from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Any, Literal

from docx import Document
from docx.text.paragraph import Paragraph
from docx.table import Table
from pptx import Presentation
from pptx.shapes.base import BaseShape
from pptx.util import Inches

from .anonymizer import (
    PIIDetector,
    PIIType,
    PIIValueResult,
    SEMANTIC_PREFIXES,
    save_mapping,
    load_mapping,
    _derive_key_from_password,
    _generalize_birthdate,
    _generalize_address,
    _generalize_age,
    generate_session_id,
)


@dataclass
class TextMatch:
    """テキスト内のPIIマッチ情報"""
    original: str
    replacement: str
    pii_type: PIIType
    start: int
    end: int


@dataclass
class DocumentPIIResult:
    """ドキュメント全体のPII検出結果"""
    file_path: str
    file_type: str
    total_matches: int
    pii_by_type: dict[str, int] = field(default_factory=dict)
    sample_matches: list[TextMatch] = field(default_factory=list)


class DocumentAnonymizer:
    """Word/PowerPointドキュメント匿名化クラス"""

    def __init__(self):
        self.detector = PIIDetector()
        self._value_mapping: dict[str, str] = {}  # original -> replacement
        self._reverse_mapping: dict[str, str] = {}  # replacement -> original
        self._type_mapping: dict[str, str] = {}  # original -> pii_type
        self._type_counters: dict[PIIType, int] = {}  # PIIType別のカウンター
        self._session_id: str | None = None  # セッションID

    def _generate_replacement(
        self,
        value: str,
        pii_type: PIIType,
        strategy: Literal["replace", "generalize"] = "replace",
    ) -> str:
        """値に対する置換文字列を生成"""
        if value in self._value_mapping:
            return self._value_mapping[value]

        if strategy == "generalize":
            if pii_type == PIIType.BIRTHDATE:
                replacement = _generalize_birthdate(value)
            elif pii_type == PIIType.ADDRESS:
                replacement = _generalize_address(value)
            elif pii_type == PIIType.AGE:
                replacement = _generalize_age(value)
            else:
                # generalizeでも対応できない場合はセマンティックIDにフォールバック
                prefix = SEMANTIC_PREFIXES.get(pii_type, "ID")
                counter = self._type_counters.get(pii_type, 0) + 1
                self._type_counters[pii_type] = counter
                if self._session_id:
                    replacement = f"{prefix}_{counter:03d}_{self._session_id}"
                else:
                    replacement = f"{prefix}_{counter:03d}"
        else:
            # セマンティックIDを使用
            prefix = SEMANTIC_PREFIXES.get(pii_type, "ID")
            counter = self._type_counters.get(pii_type, 0) + 1
            self._type_counters[pii_type] = counter
            if self._session_id:
                replacement = f"{prefix}_{counter:03d}_{self._session_id}"
            else:
                replacement = f"{prefix}_{counter:03d}"

        self._value_mapping[value] = replacement
        self._reverse_mapping[replacement] = value
        self._type_mapping[value] = pii_type.value
        return replacement

    def _remove_overlapping_matches(self, pii_results: list) -> list:
        """重複するマッチを除去（長いマッチを優先）"""
        if not pii_results:
            return []

        # 開始位置でソートし、同じ位置なら長いものを優先
        sorted_results = sorted(pii_results, key=lambda x: (x.start, -(x.end - x.start)))

        filtered = []
        last_end = -1

        for pii in sorted_results:
            if pii.start >= last_end:
                filtered.append(pii)
                last_end = pii.end

        return filtered

    def _anonymize_text(
        self,
        text: str,
        strategy: Literal["replace", "generalize"] = "replace",
    ) -> tuple[str, list[TextMatch]]:
        """テキスト内のPIIを匿名化"""
        if not text:
            return text, []

        matches: list[TextMatch] = []
        pii_results = self.detector._detect_pii_in_text(text)

        if not pii_results:
            return text, []

        # 重複を除去（長いマッチを優先）
        pii_results = self._remove_overlapping_matches(pii_results)

        # 後ろから置換（位置がずれないように）
        pii_results_sorted = sorted(pii_results, key=lambda x: x.start, reverse=True)
        result_text = text

        for pii in pii_results_sorted:
            original = pii.value
            replacement = self._generate_replacement(original, pii.pii_type, strategy)
            result_text = result_text[:pii.start] + replacement + result_text[pii.end:]
            matches.append(TextMatch(
                original=original,
                replacement=replacement,
                pii_type=pii.pii_type,
                start=pii.start,
                end=pii.end,
            ))

        return result_text, matches

    def _deanonymize_text(self, text: str, mapping: dict) -> str:
        """テキスト内の匿名化を解除"""
        if not text:
            return text

        result_text = text
        values_mapping = mapping.get("values", {})

        # 逆マッピングを使用
        for original, replacement in values_mapping.items():
            if replacement in result_text:
                result_text = result_text.replace(replacement, original)

        return result_text

    # ==================== Word (.docx) ====================

    def scan_docx(self, file_path: str | Path) -> DocumentPIIResult:
        """Wordファイル内のPIIをスキャン"""
        file_path = Path(file_path)
        doc = Document(file_path)
        total_matches = 0
        pii_by_type: dict[str, int] = {}
        sample_matches: list[TextMatch] = []

        # 本文段落
        for para in doc.paragraphs:
            pii_results = self.detector._detect_pii_in_text(para.text)
            for pii in pii_results:
                total_matches += 1
                type_name = pii.pii_type.value
                pii_by_type[type_name] = pii_by_type.get(type_name, 0) + 1
                if len(sample_matches) < 10:
                    sample_matches.append(TextMatch(
                        original=pii.value,
                        replacement="",
                        pii_type=pii.pii_type,
                        start=pii.start,
                        end=pii.end,
                    ))

        # テーブル
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for para in cell.paragraphs:
                        pii_results = self.detector._detect_pii_in_text(para.text)
                        for pii in pii_results:
                            total_matches += 1
                            type_name = pii.pii_type.value
                            pii_by_type[type_name] = pii_by_type.get(type_name, 0) + 1
                            if len(sample_matches) < 10:
                                sample_matches.append(TextMatch(
                                    original=pii.value,
                                    replacement="",
                                    pii_type=pii.pii_type,
                                    start=pii.start,
                                    end=pii.end,
                                ))

        # ヘッダー・フッター
        for section in doc.sections:
            for header in [section.header, section.first_page_header, section.even_page_header]:
                if header:
                    for para in header.paragraphs:
                        pii_results = self.detector._detect_pii_in_text(para.text)
                        for pii in pii_results:
                            total_matches += 1
                            type_name = pii.pii_type.value
                            pii_by_type[type_name] = pii_by_type.get(type_name, 0) + 1

            for footer in [section.footer, section.first_page_footer, section.even_page_footer]:
                if footer:
                    for para in footer.paragraphs:
                        pii_results = self.detector._detect_pii_in_text(para.text)
                        for pii in pii_results:
                            total_matches += 1
                            type_name = pii.pii_type.value
                            pii_by_type[type_name] = pii_by_type.get(type_name, 0) + 1

        return DocumentPIIResult(
            file_path=str(file_path),
            file_type="docx",
            total_matches=total_matches,
            pii_by_type=pii_by_type,
            sample_matches=sample_matches,
        )

    def _anonymize_paragraph(
        self,
        para: Paragraph,
        strategy: Literal["replace", "generalize"] = "replace",
    ) -> list[TextMatch]:
        """段落内のテキストを匿名化（書式を保持）"""
        matches: list[TextMatch] = []

        # 各runの処理
        for run in para.runs:
            if run.text:
                new_text, run_matches = self._anonymize_text(run.text, strategy)
                run.text = new_text
                matches.extend(run_matches)

        return matches

    def anonymize_docx(
        self,
        input_path: str | Path,
        output_path: str | Path,
        strategy: Literal["replace", "generalize"] = "replace",
    ) -> tuple[DocumentPIIResult, dict]:
        """Wordファイルを匿名化"""
        input_path = Path(input_path)
        output_path = Path(output_path)

        # マッピングをリセット
        self._value_mapping = {}
        self._reverse_mapping = {}
        self._type_mapping = {}
        self._type_counters = {}
        self._session_id = generate_session_id()

        doc = Document(input_path)
        all_matches: list[TextMatch] = []
        pii_by_type: dict[str, int] = {}

        # 本文段落
        for para in doc.paragraphs:
            matches = self._anonymize_paragraph(para, strategy)
            all_matches.extend(matches)

        # テーブル
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for para in cell.paragraphs:
                        matches = self._anonymize_paragraph(para, strategy)
                        all_matches.extend(matches)

        # ヘッダー・フッター
        for section in doc.sections:
            for header in [section.header, section.first_page_header, section.even_page_header]:
                if header:
                    for para in header.paragraphs:
                        matches = self._anonymize_paragraph(para, strategy)
                        all_matches.extend(matches)

            for footer in [section.footer, section.first_page_footer, section.even_page_footer]:
                if footer:
                    for para in footer.paragraphs:
                        matches = self._anonymize_paragraph(para, strategy)
                        all_matches.extend(matches)

        # PIIタイプ別集計
        for match in all_matches:
            type_name = match.pii_type.value
            pii_by_type[type_name] = pii_by_type.get(type_name, 0) + 1

        # 保存
        doc.save(output_path)

        # マッピング生成
        mapping = {
            "metadata": {
                "created_at": datetime.now().isoformat(),
                "strategy": strategy,
                "original_file": str(input_path),
                "file_type": "docx",
                "total_replacements": len(all_matches),
                "session_id": self._session_id,
            },
            "values": self._value_mapping.copy(),
            "types": self._type_mapping.copy(),
        }

        result = DocumentPIIResult(
            file_path=str(input_path),
            file_type="docx",
            total_matches=len(all_matches),
            pii_by_type=pii_by_type,
            sample_matches=all_matches[:10],
        )

        return result, mapping

    def deanonymize_docx(
        self,
        input_path: str | Path,
        output_path: str | Path,
        mapping: dict,
    ) -> None:
        """Wordファイルの匿名化を解除"""
        input_path = Path(input_path)
        output_path = Path(output_path)

        doc = Document(input_path)

        # 本文段落
        for para in doc.paragraphs:
            for run in para.runs:
                if run.text:
                    run.text = self._deanonymize_text(run.text, mapping)

        # テーブル
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for para in cell.paragraphs:
                        for run in para.runs:
                            if run.text:
                                run.text = self._deanonymize_text(run.text, mapping)

        # ヘッダー・フッター
        for section in doc.sections:
            for header in [section.header, section.first_page_header, section.even_page_header]:
                if header:
                    for para in header.paragraphs:
                        for run in para.runs:
                            if run.text:
                                run.text = self._deanonymize_text(run.text, mapping)

            for footer in [section.footer, section.first_page_footer, section.even_page_footer]:
                if footer:
                    for para in footer.paragraphs:
                        for run in para.runs:
                            if run.text:
                                run.text = self._deanonymize_text(run.text, mapping)

        doc.save(output_path)

    # ==================== PowerPoint (.pptx) ====================

    def _get_text_from_shape(self, shape: BaseShape) -> str:
        """シェイプからテキストを取得"""
        text_parts = []
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text:
                        text_parts.append(run.text)
        return "".join(text_parts)

    def scan_pptx(self, file_path: str | Path) -> DocumentPIIResult:
        """PowerPointファイル内のPIIをスキャン"""
        file_path = Path(file_path)
        prs = Presentation(file_path)
        total_matches = 0
        pii_by_type: dict[str, int] = {}
        sample_matches: list[TextMatch] = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        full_text = "".join(run.text for run in para.runs if run.text)
                        pii_results = self.detector._detect_pii_in_text(full_text)
                        for pii in pii_results:
                            total_matches += 1
                            type_name = pii.pii_type.value
                            pii_by_type[type_name] = pii_by_type.get(type_name, 0) + 1
                            if len(sample_matches) < 10:
                                sample_matches.append(TextMatch(
                                    original=pii.value,
                                    replacement="",
                                    pii_type=pii.pii_type,
                                    start=pii.start,
                                    end=pii.end,
                                ))

                # テーブルの処理
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            if cell.text_frame:
                                for para in cell.text_frame.paragraphs:
                                    full_text = "".join(run.text for run in para.runs if run.text)
                                    pii_results = self.detector._detect_pii_in_text(full_text)
                                    for pii in pii_results:
                                        total_matches += 1
                                        type_name = pii.pii_type.value
                                        pii_by_type[type_name] = pii_by_type.get(type_name, 0) + 1
                                        if len(sample_matches) < 10:
                                            sample_matches.append(TextMatch(
                                                original=pii.value,
                                                replacement="",
                                                pii_type=pii.pii_type,
                                                start=pii.start,
                                                end=pii.end,
                                            ))

            # ノート（発表者ノート）の処理
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                for shape in notes_slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            full_text = "".join(run.text for run in para.runs if run.text)
                            pii_results = self.detector._detect_pii_in_text(full_text)
                            for pii in pii_results:
                                total_matches += 1
                                type_name = pii.pii_type.value
                                pii_by_type[type_name] = pii_by_type.get(type_name, 0) + 1

        return DocumentPIIResult(
            file_path=str(file_path),
            file_type="pptx",
            total_matches=total_matches,
            pii_by_type=pii_by_type,
            sample_matches=sample_matches,
        )

    def _anonymize_text_frame(
        self,
        text_frame,
        strategy: Literal["replace", "generalize"] = "replace",
    ) -> list[TextMatch]:
        """テキストフレーム内を匿名化"""
        all_matches: list[TextMatch] = []

        for para in text_frame.paragraphs:
            for run in para.runs:
                if run.text:
                    new_text, matches = self._anonymize_text(run.text, strategy)
                    run.text = new_text
                    all_matches.extend(matches)

        return all_matches

    def anonymize_pptx(
        self,
        input_path: str | Path,
        output_path: str | Path,
        strategy: Literal["replace", "generalize"] = "replace",
    ) -> tuple[DocumentPIIResult, dict]:
        """PowerPointファイルを匿名化"""
        input_path = Path(input_path)
        output_path = Path(output_path)

        # マッピングをリセット
        self._value_mapping = {}
        self._reverse_mapping = {}
        self._type_mapping = {}
        self._type_counters = {}
        self._session_id = generate_session_id()

        prs = Presentation(input_path)
        all_matches: list[TextMatch] = []
        pii_by_type: dict[str, int] = {}

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    matches = self._anonymize_text_frame(shape.text_frame, strategy)
                    all_matches.extend(matches)

                # テーブルの処理
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            if cell.text_frame:
                                matches = self._anonymize_text_frame(cell.text_frame, strategy)
                                all_matches.extend(matches)

            # ノートの処理
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                for shape in notes_slide.shapes:
                    if shape.has_text_frame:
                        matches = self._anonymize_text_frame(shape.text_frame, strategy)
                        all_matches.extend(matches)

        # PIIタイプ別集計
        for match in all_matches:
            type_name = match.pii_type.value
            pii_by_type[type_name] = pii_by_type.get(type_name, 0) + 1

        # 保存
        prs.save(output_path)

        # マッピング生成
        mapping = {
            "metadata": {
                "created_at": datetime.now().isoformat(),
                "strategy": strategy,
                "original_file": str(input_path),
                "file_type": "pptx",
                "total_replacements": len(all_matches),
                "session_id": self._session_id,
            },
            "values": self._value_mapping.copy(),
            "types": self._type_mapping.copy(),
        }

        result = DocumentPIIResult(
            file_path=str(input_path),
            file_type="pptx",
            total_matches=len(all_matches),
            pii_by_type=pii_by_type,
            sample_matches=all_matches[:10],
        )

        return result, mapping

    def deanonymize_pptx(
        self,
        input_path: str | Path,
        output_path: str | Path,
        mapping: dict,
    ) -> None:
        """PowerPointファイルの匿名化を解除"""
        input_path = Path(input_path)
        output_path = Path(output_path)

        prs = Presentation(input_path)

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.text:
                                run.text = self._deanonymize_text(run.text, mapping)

                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            if cell.text_frame:
                                for para in cell.text_frame.paragraphs:
                                    for run in para.runs:
                                        if run.text:
                                            run.text = self._deanonymize_text(run.text, mapping)

            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                for shape in notes_slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                if run.text:
                                    run.text = self._deanonymize_text(run.text, mapping)

        prs.save(output_path)


# ==================== 便利関数 ====================

def scan_document(file_path: str | Path) -> DocumentPIIResult:
    """ドキュメントファイル内のPIIをスキャン"""
    file_path = Path(file_path)
    anonymizer = DocumentAnonymizer()

    suffix = file_path.suffix.lower()
    if suffix == ".docx":
        return anonymizer.scan_docx(file_path)
    elif suffix == ".pptx":
        return anonymizer.scan_pptx(file_path)
    else:
        raise ValueError(f"サポートされていないファイル形式: {suffix}")


def anonymize_document(
    input_path: str | Path,
    output_path: str | Path,
    strategy: Literal["replace", "generalize"] = "replace",
) -> tuple[DocumentPIIResult, dict]:
    """ドキュメントファイルを匿名化"""
    input_path = Path(input_path)
    output_path = Path(output_path)
    anonymizer = DocumentAnonymizer()

    suffix = input_path.suffix.lower()
    if suffix == ".docx":
        return anonymizer.anonymize_docx(input_path, output_path, strategy)
    elif suffix == ".pptx":
        return anonymizer.anonymize_pptx(input_path, output_path, strategy)
    else:
        raise ValueError(f"サポートされていないファイル形式: {suffix}")


def deanonymize_document(
    input_path: str | Path,
    output_path: str | Path,
    mapping: dict,
) -> None:
    """ドキュメントファイルの匿名化を解除"""
    input_path = Path(input_path)
    output_path = Path(output_path)
    anonymizer = DocumentAnonymizer()

    suffix = input_path.suffix.lower()
    if suffix == ".docx":
        anonymizer.deanonymize_docx(input_path, output_path, mapping)
    elif suffix == ".pptx":
        anonymizer.deanonymize_pptx(input_path, output_path, mapping)
    else:
        raise ValueError(f"サポートされていないファイル形式: {suffix}")
