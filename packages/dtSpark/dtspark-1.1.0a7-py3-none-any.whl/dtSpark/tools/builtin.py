"""
Built-in Tools module for providing default tool capabilities.

This module provides built-in tools that are always available to the LLM,
such as date/time information with timezone awareness and filesystem access.


"""

import logging
import os
import base64
import fnmatch
from pathlib import Path
from datetime import datetime
from typing import Dict, List, Any, Optional
from zoneinfo import ZoneInfo, available_timezones

# Common error/description string constants (SonarCloud S1192)
_ERR_FS_NOT_CONFIGURED = "Filesystem tools not configured"
_ERR_FILE_PATH_REQUIRED = "File path is required"
_ERR_DOC_NOT_CONFIGURED = "Document tools not configured"
_ERR_DOC_NOT_ENABLED = "Document tools are not enabled"
_ERR_WRITE_MODE_REQUIRED = "Write operations require access_mode: read_write"
_ERR_OUTPUT_PATH_REQUIRED = "Output file path is required"
_ERR_ARCHIVE_NOT_CONFIGURED = "Archive tools not configured"
_ERR_ARCHIVE_NOT_ENABLED = "Archive tools are not enabled"
_ERR_ARCHIVE_PATH_REQUIRED = "Archive path is required"
_DESC_ARCHIVE_PATH = "Path to the archive file"
_TAR_GZ = 'tar.gz'
_TAR_BZ2 = 'tar.bz2'
_TAR_BZ2_MODE = 'r:bz2'
_TAR_OPEN_MODES = {_TAR_GZ: 'r:gz', _TAR_BZ2: _TAR_BZ2_MODE}



def get_builtin_tools(config: Optional[Dict[str, Any]] = None) -> List[Dict[str, Any]]:
    """
    Get the list of built-in tool definitions.

    Args:
        config: Optional configuration dictionary containing embedded_tools settings

    Returns:
        List of tool definitions in Claude API format
    """
    tools = [
        {
            "name": "get_current_datetime",
            "description": "Get the current date and time with timezone awareness. "
                          "Returns the current datetime in ISO 8601 format. "
                          "Optionally specify a timezone to get the time in that zone.",
            "input_schema": {
                "type": "object",
                "properties": {
                    "timezone": {
                        "type": "string",
                        "description": "Optional timezone identifier (e.g., 'Australia/Sydney', 'America/New_York', 'UTC'). "
                                      "If not provided, uses the system's local timezone.",
                        "default": None
                    },
                    "format": {
                        "type": "string",
                        "description": "Optional format for the datetime output. Options: 'iso' (ISO 8601), 'human' (human-readable). "
                                      "Default is 'iso'.",
                        "enum": ["iso", "human"],
                        "default": "iso"
                    }
                },
                "required": []
            }
        }
    ]

    # Add filesystem tools if enabled
    if config.get('embedded_tools', None):
        fs_config = config.get('embedded_tools', {}).get('filesystem', {})
        if fs_config.get('enabled', False):
            fs_tools = _get_filesystem_tools(fs_config)
            tools.extend(fs_tools)
            logging.info(f"Embedded filesystem tools enabled: {len(fs_tools)} tools added")

        # Add document tools if enabled
        doc_config = config.get('embedded_tools', {}).get('documents', {})
        if doc_config.get('enabled', False):
            doc_tools = _get_document_tools(doc_config)
            tools.extend(doc_tools)
            logging.info(f"Embedded document tools enabled: {len(doc_tools)} tools added")

        # Add archive tools if enabled
        archive_config = config.get('embedded_tools', {}).get('archives', {})
        if archive_config.get('enabled', False):
            archive_tools = _get_archive_tools(archive_config)
            tools.extend(archive_tools)
            logging.info(f"Embedded archive tools enabled: {len(archive_tools)} tools added")

    return tools


def execute_builtin_tool(tool_name: str, tool_input: Dict[str, Any],
                        config: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
    """
    Execute a built-in tool.

    Args:
        tool_name: Name of the tool to execute
        tool_input: Input parameters for the tool
        config: Optional configuration dictionary for filesystem tools

    Returns:
        Dictionary containing:
        - success: Boolean indicating if execution was successful
        - result: The tool execution result (if successful)
        - error: Error message (if failed)
    """
    try:
        if tool_name == "get_current_datetime":
            return _execute_get_current_datetime(tool_input)

        # Filesystem tools
        elif tool_name == "list_files_recursive":
            return _execute_list_files_recursive(tool_input, config)
        elif tool_name == "search_files":
            return _execute_search_files(tool_input, config)
        elif tool_name == "read_file_text":
            return _execute_read_file_text(tool_input, config)
        elif tool_name == "read_file_binary":
            return _execute_read_file_binary(tool_input, config)
        elif tool_name == "write_file":
            return _execute_write_file(tool_input, config)
        elif tool_name == "create_directories":
            return _execute_create_directories(tool_input, config)

        # Document tools
        elif tool_name == "get_file_info":
            return _execute_get_file_info(tool_input, config)
        elif tool_name == "read_word_document":
            return _execute_read_word_document(tool_input, config)
        elif tool_name == "read_excel_document":
            return _execute_read_excel_document(tool_input, config)
        elif tool_name == "read_powerpoint_document":
            return _execute_read_powerpoint_document(tool_input, config)
        elif tool_name == "read_pdf_document":
            return _execute_read_pdf_document(tool_input, config)
        elif tool_name == "create_word_document":
            return _execute_create_word_document(tool_input, config)
        elif tool_name == "create_excel_document":
            return _execute_create_excel_document(tool_input, config)
        elif tool_name == "create_powerpoint_document":
            return _execute_create_powerpoint_document(tool_input, config)

        # Archive tools
        elif tool_name == "list_archive_contents":
            return _execute_list_archive_contents(tool_input, config)
        elif tool_name == "read_archive_file":
            return _execute_read_archive_file(tool_input, config)
        elif tool_name == "extract_archive":
            return _execute_extract_archive(tool_input, config)

        else:
            return {
                "success": False,
                "error": f"Unknown built-in tool: {tool_name}"
            }
    except Exception as e:
        logging.error(f"Error executing built-in tool {tool_name}: {e}")
        return {
            "success": False,
            "error": str(e)
        }


def _execute_get_current_datetime(tool_input: Dict[str, Any]) -> Dict[str, Any]:
    """
    Execute the get_current_datetime tool.

    Args:
        tool_input: Dictionary containing optional 'timezone' and 'format' keys

    Returns:
        Dictionary with success status and datetime result
    """
    timezone_str = tool_input.get("timezone")
    output_format = tool_input.get("format", "iso")

    try:
        # Get current datetime
        if timezone_str:
            # Validate timezone
            if timezone_str not in available_timezones():
                return {
                    "success": False,
                    "error": f"Invalid timezone: {timezone_str}. Use a valid IANA timezone identifier."
                }

            # Get datetime in specified timezone
            tz = ZoneInfo(timezone_str)
            now = datetime.now(tz)
        else:
            # Get local datetime with system timezone
            now = datetime.now().astimezone()

        # Format output
        if output_format == "human":
            # Human-readable format
            result = {
                "datetime": now.strftime("%A, %d %B %Y at %I:%M:%S %p"),
                "timezone": now.strftime("%Z (UTC%z)"),
                "iso_format": now.isoformat()
            }
        else:
            # ISO 8601 format (default)
            result = {
                "datetime": now.isoformat(),
                "timezone": str(now.tzinfo),
                "timezone_offset": now.strftime("%z"),
                "unix_timestamp": int(now.timestamp())
            }

        logging.info(f"Built-in tool get_current_datetime executed: timezone={timezone_str or 'local'}, format={output_format}")

        return {
            "success": True,
            "result": result
        }

    except Exception as e:
        logging.error(f"Error in get_current_datetime: {e}")
        return {
            "success": False,
            "error": str(e)
        }


def get_available_timezones() -> List[str]:
    """
    Get a list of all available timezone identifiers.

    Returns:
        Sorted list of timezone identifiers
    """
    return sorted(available_timezones())


def validate_timezone(timezone_str: str) -> bool:
    """
    Validate if a timezone string is valid.

    Args:
        timezone_str: Timezone identifier to validate

    Returns:
        True if valid, False otherwise
    """
    return timezone_str in available_timezones()


# ============================================================================
# Filesystem Tools
# ============================================================================

def _get_filesystem_tools(fs_config: Dict[str, Any]) -> List[Dict[str, Any]]:
    """
    Get filesystem tool definitions based on configuration.

    Args:
        fs_config: Filesystem configuration dictionary

    Returns:
        List of filesystem tool definitions
    """
    access_mode = fs_config.get('access_mode', 'read')
    allowed_path = fs_config.get('allowed_path', '.')

    # Read-only tools (always included when filesystem is enabled)
    tools = [
        {
            "name": "list_files_recursive",
            "description": f"List all files and directories recursively within the allowed path ({allowed_path}). "
                          "Returns a structured list of all files with their paths, sizes, and modification times. "
                          "Useful for understanding directory structure and finding files.",
            "input_schema": {
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": f"Optional subdirectory within {allowed_path} to list. "
                                      "If not provided, lists from the root of the allowed path.",
                        "default": ""
                    },
                    "include_hidden": {
                        "type": "boolean",
                        "description": "Include hidden files and directories (those starting with '.')",
                        "default": False
                    }
                },
                "required": []
            }
        },
        {
            "name": "search_files",
            "description": f"Search for files by filename within the allowed path ({allowed_path}). "
                          "Supports wildcards (* for any characters, ? for single character). "
                          "Returns matching file paths.",
            "input_schema": {
                "type": "object",
                "properties": {
                    "pattern": {
                        "type": "string",
                        "description": "Search pattern for filename. Supports wildcards: * (any characters), ? (single character). "
                                      "Examples: '*.py' (all Python files), 'test_*.py' (test files), 'config.???' (config with 3-char extension)",
                    },
                    "case_sensitive": {
                        "type": "boolean",
                        "description": "Whether the search should be case-sensitive",
                        "default": False
                    }
                },
                "required": ["pattern"]
            }
        },
        {
            "name": "read_file_text",
            "description": f"Read the contents of a text file within the allowed path ({allowed_path}). "
                          "Attempts to decode the file as UTF-8 text. Use read_file_binary for non-text files.",
            "input_schema": {
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": "Path to the file to read (relative to allowed path or absolute within allowed path)",
                    }
                },
                "required": ["path"]
            }
        },
        {
            "name": "read_file_binary",
            "description": f"Read the contents of a file as binary data within the allowed path ({allowed_path}). "
                          "Returns base64-encoded binary content. Use for images, PDFs, or other non-text files.",
            "input_schema": {
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": "Path to the file to read (relative to allowed path or absolute within allowed path)",
                    },
                    "max_size_mb": {
                        "type": "number",
                        "description": "Maximum file size in MB to read (default: 10MB). Prevents reading very large files.",
                        "default": 10
                    }
                },
                "required": ["path"]
            }
        }
    ]

    # Write tools (only added if access_mode is read_write)
    if access_mode == 'read_write':
        tools.extend([
            {
                "name": "write_file",
                "description": f"Write content to a file within the allowed path ({allowed_path}). "
                              "Creates the file if it doesn't exist, or overwrites if it exists. "
                              "Parent directories must already exist (use create_directories first if needed).",
                "input_schema": {
                    "type": "object",
                    "properties": {
                        "path": {
                            "type": "string",
                            "description": "Path to the file to write (relative to allowed path or absolute within allowed path)",
                        },
                        "content": {
                            "type": "string",
                            "description": "Content to write to the file",
                        },
                        "encoding": {
                            "type": "string",
                            "description": "Text encoding to use (default: utf-8)",
                            "default": "utf-8"
                        }
                    },
                    "required": ["path", "content"]
                }
            },
            {
                "name": "create_directories",
                "description": f"Create one or more nested directories within the allowed path ({allowed_path}). "
                              "Creates all intermediate directories as needed (like 'mkdir -p'). "
                              "Safe to call even if directories already exist.",
                "input_schema": {
                    "type": "object",
                    "properties": {
                        "path": {
                            "type": "string",
                            "description": "Directory path to create (relative to allowed path or absolute within allowed path). "
                                          "Can include multiple nested levels (e.g., 'data/processed/reports')",
                        }
                    },
                    "required": ["path"]
                }
            }
        ])

    return tools


def _validate_path(file_path: str, allowed_path: str) -> Dict[str, Any]:
    """
    Validate that a file path is within the allowed directory.

    Args:
        file_path: File path to validate
        allowed_path: Root path that file must be within

    Returns:
        Dictionary with:
        - valid: Boolean indicating if path is valid
        - resolved_path: Absolute resolved path (if valid)
        - error: Error message (if invalid)
    """
    try:
        # Resolve allowed path to absolute
        allowed_abs = Path(allowed_path).resolve()

        # Handle empty file_path (means root of allowed path)
        if not file_path or file_path == '.':
            return {
                "valid": True,
                "resolved_path": str(allowed_abs),
                "error": None
            }

        # Resolve file path
        # If file_path is absolute, use it directly; otherwise treat as relative to allowed_path
        if Path(file_path).is_absolute():
            file_abs = Path(file_path).resolve()
        else:
            file_abs = (allowed_abs / file_path).resolve()

        # Check if file path is within allowed path
        try:
            file_abs.relative_to(allowed_abs)
        except ValueError:
            return {
                "valid": False,
                "resolved_path": None,
                "error": f"Access denied: Path '{file_path}' is outside allowed directory '{allowed_path}'"
            }

        return {
            "valid": True,
            "resolved_path": str(file_abs),
            "error": None
        }

    except Exception as e:
        return {
            "valid": False,
            "resolved_path": None,
            "error": f"Invalid path: {str(e)}"
        }


def _execute_list_files_recursive(tool_input: Dict[str, Any],
                                  config: Optional[Dict[str, Any]]) -> Dict[str, Any]:
    """
    Execute the list_files_recursive tool.

    Args:
        tool_input: Tool input parameters
        config: Configuration dictionary

    Returns:
        Dictionary with success status and file listing
    """
    if not config.get('embedded_tools'):
        return {"success": False, "error": _ERR_FS_NOT_CONFIGURED}

    fs_config = config.get('embedded_tools', {}).get('filesystem', {})
    allowed_path = fs_config.get('allowed_path', '.')

    # Get parameters
    sub_path = tool_input.get('path', '')
    include_hidden = tool_input.get('include_hidden', False)

    # Validate path
    validation = _validate_path(sub_path, allowed_path)
    if not validation['valid']:
        return {"success": False, "error": validation['error']}

    root_path = Path(validation['resolved_path'])

    # Check if path exists
    if not root_path.exists():
        return {"success": False, "error": f"Path does not exist: {sub_path}"}

    if not root_path.is_dir():
        return {"success": False, "error": f"Path is not a directory: {sub_path}"}

    # Collect all files and directories
    files = []
    directories = []

    try:
        for item in root_path.rglob('*'):
            # Skip hidden files if not requested
            if not include_hidden and any(part.startswith('.') for part in item.parts):
                continue

            # Get relative path from root
            rel_path = item.relative_to(root_path)

            if item.is_file():
                files.append({
                    "path": str(rel_path),
                    "full_path": str(item),
                    "size_bytes": item.stat().st_size,
                    "modified": datetime.fromtimestamp(item.stat().st_mtime).isoformat(),
                    "type": "file"
                })
            elif item.is_dir():
                directories.append({
                    "path": str(rel_path),
                    "full_path": str(item),
                    "type": "directory"
                })

        result = {
            "root_path": str(root_path),
            "total_files": len(files),
            "total_directories": len(directories),
            "files": sorted(files, key=lambda x: x['path']),
            "directories": sorted(directories, key=lambda x: x['path'])
        }

        logging.info(f"Listed {len(files)} files and {len(directories)} directories from {root_path}")
        return {"success": True, "result": result}

    except Exception as e:
        logging.error(f"Error listing files: {e}")
        return {"success": False, "error": str(e)}


def _execute_search_files(tool_input: Dict[str, Any],
                          config: Optional[Dict[str, Any]]) -> Dict[str, Any]:
    """
    Execute the search_files tool.

    Args:
        tool_input: Tool input parameters
        config: Configuration dictionary

    Returns:
        Dictionary with success status and search results
    """
    if not config.get('embedded_tools'):
        return {"success": False, "error": _ERR_FS_NOT_CONFIGURED}

    fs_config = config.get('embedded_tools', {}).get('filesystem', {})
    allowed_path = fs_config.get('allowed_path', '.')

    # Get parameters
    pattern = tool_input.get('pattern')
    case_sensitive = tool_input.get('case_sensitive', False)

    if not pattern:
        return {"success": False, "error": "Search pattern is required"}

    # Validate path
    validation = _validate_path('', allowed_path)
    if not validation['valid']:
        return {"success": False, "error": validation['error']}

    root_path = Path(validation['resolved_path'])

    # Search for matching files
    matches = []

    try:
        for item in root_path.rglob('*'):
            if item.is_file():
                filename = item.name

                # Apply pattern matching
                if case_sensitive:
                    match = fnmatch.fnmatch(filename, pattern)
                else:
                    match = fnmatch.fnmatch(filename.lower(), pattern.lower())

                if match:
                    rel_path = item.relative_to(root_path)
                    matches.append({
                        "filename": filename,
                        "path": str(rel_path),
                        "full_path": str(item),
                        "size_bytes": item.stat().st_size,
                        "modified": datetime.fromtimestamp(item.stat().st_mtime).isoformat()
                    })

        result = {
            "pattern": pattern,
            "total_matches": len(matches),
            "matches": sorted(matches, key=lambda x: x['path'])
        }

        logging.info(f"Search for '{pattern}' found {len(matches)} matches")
        return {"success": True, "result": result}

    except Exception as e:
        logging.error(f"Error searching files: {e}")
        return {"success": False, "error": str(e)}


def _execute_read_file_text(tool_input: Dict[str, Any],
                            config: Optional[Dict[str, Any]]) -> Dict[str, Any]:
    """
    Execute the read_file_text tool.

    Args:
        tool_input: Tool input parameters
        config: Configuration dictionary

    Returns:
        Dictionary with success status and file content
    """
    if not config.get('embedded_tools'):
        return {"success": False, "error": _ERR_FS_NOT_CONFIGURED}

    fs_config = config.get('embedded_tools', {}).get('filesystem', {})
    allowed_path = fs_config.get('allowed_path', '.')

    # Get parameters
    file_path = tool_input.get('path')

    if not file_path:
        return {"success": False, "error": _ERR_FILE_PATH_REQUIRED}

    # Validate path
    validation = _validate_path(file_path, allowed_path)
    if not validation['valid']:
        return {"success": False, "error": validation['error']}

    full_path = Path(validation['resolved_path'])

    # Check if file exists
    if not full_path.exists():
        return {"success": False, "error": f"File does not exist: {file_path}"}

    if not full_path.is_file():
        return {"success": False, "error": f"Path is not a file: {file_path}"}

    # Read file as text
    try:
        with open(full_path, 'r', encoding='utf-8') as f:
            content = f.read()

        result = {
            "path": file_path,
            "full_path": str(full_path),
            "content": content,
            "size_bytes": full_path.stat().st_size,
            "encoding": "utf-8"
        }

        logging.info(f"Read text file: {file_path} ({result['size_bytes']} bytes)")
        return {"success": True, "result": result}

    except UnicodeDecodeError:
        return {
            "success": False,
            "error": f"File is not valid UTF-8 text. Use read_file_binary instead: {file_path}"
        }
    except Exception as e:
        logging.error(f"Error reading file {file_path}: {e}")
        return {"success": False, "error": str(e)}


def _execute_read_file_binary(tool_input: Dict[str, Any],
                              config: Optional[Dict[str, Any]]) -> Dict[str, Any]:
    """
    Execute the read_file_binary tool.

    Args:
        tool_input: Tool input parameters
        config: Configuration dictionary

    Returns:
        Dictionary with success status and base64-encoded content
    """
    if not config.get('embedded_tools'):
        return {"success": False, "error": _ERR_FS_NOT_CONFIGURED}

    fs_config = config.get('embedded_tools', {}).get('filesystem', {})
    allowed_path = fs_config.get('allowed_path', '.')

    # Get parameters
    file_path = tool_input.get('path')
    max_size_mb = tool_input.get('max_size_mb', 10)

    if not file_path:
        return {"success": False, "error": _ERR_FILE_PATH_REQUIRED}

    # Validate path
    validation = _validate_path(file_path, allowed_path)
    if not validation['valid']:
        return {"success": False, "error": validation['error']}

    full_path = Path(validation['resolved_path'])

    # Check if file exists
    if not full_path.exists():
        return {"success": False, "error": f"File does not exist: {file_path}"}

    if not full_path.is_file():
        return {"success": False, "error": f"Path is not a file: {file_path}"}

    # Check file size
    file_size = full_path.stat().st_size
    max_size_bytes = max_size_mb * 1024 * 1024

    if file_size > max_size_bytes:
        return {
            "success": False,
            "error": f"File size ({file_size / 1024 / 1024:.2f} MB) exceeds maximum ({max_size_mb} MB)"
        }

    # Read file as binary
    try:
        with open(full_path, 'rb') as f:
            binary_content = f.read()

        # Encode as base64
        base64_content = base64.b64encode(binary_content).decode('utf-8')

        result = {
            "path": file_path,
            "full_path": str(full_path),
            "content_base64": base64_content,
            "size_bytes": file_size
        }

        logging.info(f"Read binary file: {file_path} ({file_size} bytes)")
        return {"success": True, "result": result}

    except Exception as e:
        logging.error(f"Error reading binary file {file_path}: {e}")
        return {"success": False, "error": str(e)}


def _execute_write_file(tool_input: Dict[str, Any],
                       config: Optional[Dict[str, Any]]) -> Dict[str, Any]:
    """
    Execute the write_file tool.

    Args:
        tool_input: Tool input parameters
        config: Configuration dictionary

    Returns:
        Dictionary with success status
    """
    logging.debug(f"write_file called with config keys: {list(config.keys()) if config else 'None'}")

    if not config.get('embedded_tools'):
        logging.warning("write_file failed: embedded_tools not in config")
        return {"success": False, "error": _ERR_FS_NOT_CONFIGURED}

    fs_config = config.get('embedded_tools', {}).get('filesystem', {})
    allowed_path = fs_config.get('allowed_path', '.')
    access_mode = fs_config.get('access_mode', 'read')

    logging.debug(f"write_file fs_config: allowed_path={allowed_path}, access_mode={access_mode}")

    # Check if write access is enabled
    if access_mode != 'read_write':
        logging.warning(f"write_file failed: access_mode is '{access_mode}', not 'read_write'")
        return {
            "success": False,
            "error": "Write operations are disabled. Set access_mode to 'read_write' in configuration."
        }

    # Get parameters
    file_path = tool_input.get('path')
    content = tool_input.get('content')
    encoding = tool_input.get('encoding', 'utf-8')

    logging.debug(f"write_file params: path={file_path}, content_len={len(content) if content else 0}")

    if not file_path:
        logging.warning("write_file failed: no file path provided")
        return {"success": False, "error": _ERR_FILE_PATH_REQUIRED}

    if content is None:
        logging.warning("write_file failed: no content provided")
        return {"success": False, "error": "Content is required"}

    # Validate path
    validation = _validate_path(file_path, allowed_path)
    if not validation['valid']:
        logging.warning(f"write_file failed: path validation error: {validation['error']}")
        return {"success": False, "error": validation['error']}

    full_path = Path(validation['resolved_path'])
    logging.debug(f"write_file resolved path: {full_path}")

    # Check if parent directory exists
    if not full_path.parent.exists():
        logging.warning(f"write_file failed: parent directory does not exist: {full_path.parent}")
        return {
            "success": False,
            "error": f"Parent directory does not exist: {full_path.parent}. Use create_directories first."
        }

    # Write file
    try:
        logging.debug(f"write_file: attempting to write {len(content)} chars to {full_path}")
        with open(full_path, 'w', encoding=encoding) as f:
            f.write(content)

        result = {
            "path": file_path,
            "full_path": str(full_path),
            "size_bytes": full_path.stat().st_size,
            "encoding": encoding
        }

        logging.info(f"Wrote file: {file_path} ({result['size_bytes']} bytes)")
        return {"success": True, "result": result}

    except Exception as e:
        logging.error(f"Error writing file {file_path}: {e}")
        return {"success": False, "error": str(e)}


def _execute_create_directories(tool_input: Dict[str, Any],
                                config: Optional[Dict[str, Any]]) -> Dict[str, Any]:
    """
    Execute the create_directories tool.

    Args:
        tool_input: Tool input parameters
        config: Configuration dictionary

    Returns:
        Dictionary with success status
    """
    if not config.get('embedded_tools'):
        return {"success": False, "error": _ERR_FS_NOT_CONFIGURED}

    fs_config = config.get('embedded_tools', {}).get('filesystem', {})
    allowed_path = fs_config.get('allowed_path', '.')
    access_mode = fs_config.get('access_mode', 'read')

    # Check if write access is enabled
    if access_mode != 'read_write':
        return {
            "success": False,
            "error": "Write operations are disabled. Set access_mode to 'read_write' in configuration."
        }

    # Get parameters
    dir_path = tool_input.get('path')

    if not dir_path:
        return {"success": False, "error": "Directory path is required"}

    # Validate path
    validation = _validate_path(dir_path, allowed_path)
    if not validation['valid']:
        return {"success": False, "error": validation['error']}

    full_path = Path(validation['resolved_path'])

    # Create directories
    try:
        full_path.mkdir(parents=True, exist_ok=True)

        result = {
            "path": dir_path,
            "full_path": str(full_path),
            "created": not full_path.exists() or len(list(full_path.iterdir())) == 0
        }

        logging.info(f"Created directories: {dir_path}")
        return {"success": True, "result": result}

    except Exception as e:
        logging.error(f"Error creating directories {dir_path}: {e}")
        return {"success": False, "error": str(e)}


# ============================================================================
# Document Tools
# ============================================================================

def _get_document_tools(doc_config: Dict[str, Any]) -> List[Dict[str, Any]]:
    """
    Get document tool definitions based on configuration.

    Args:
        doc_config: Document tools configuration dictionary

    Returns:
        List of document tool definitions
    """
    access_mode = doc_config.get('access_mode', 'read')
    allowed_path = doc_config.get('allowed_path', '.')

    # File info tool (always included)
    tools = [
        {
            "name": "get_file_info",
            "description": f"Get detailed file information including type, size, MIME type, and extension. "
                          f"Works for any file within the allowed path ({allowed_path}). "
                          "Useful for determining how to process a file before reading it.",
            "input_schema": {
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": "Path to the file (relative to allowed path or absolute within allowed path)"
                    }
                },
                "required": ["path"]
            }
        }
    ]

    # Read-only document tools
    tools.extend([
        {
            "name": "read_word_document",
            "description": f"Extract text content from Microsoft Word documents (.docx) within the allowed path ({allowed_path}). "
                          "Returns the document text organised by paragraphs. Also extracts headings and tables if present.",
            "input_schema": {
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": "Path to the .docx file"
                    },
                    "include_tables": {
                        "type": "boolean",
                        "description": "Include table content in the output",
                        "default": True
                    },
                    "include_headers_footers": {
                        "type": "boolean",
                        "description": "Include header and footer content",
                        "default": False
                    }
                },
                "required": ["path"]
            }
        },
        {
            "name": "read_excel_document",
            "description": f"Extract data from Microsoft Excel documents (.xlsx) within the allowed path ({allowed_path}). "
                          "Returns spreadsheet data as structured JSON. Can read specific sheets or all sheets.",
            "input_schema": {
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": "Path to the .xlsx file"
                    },
                    "sheet_name": {
                        "type": "string",
                        "description": "Specific sheet name to read. If not provided, reads the active sheet."
                    },
                    "include_all_sheets": {
                        "type": "boolean",
                        "description": "Read all sheets in the workbook",
                        "default": False
                    },
                    "max_rows": {
                        "type": "integer",
                        "description": "Maximum number of rows to read (0 = use config default)",
                        "default": 0
                    }
                },
                "required": ["path"]
            }
        },
        {
            "name": "read_powerpoint_document",
            "description": f"Extract text content from Microsoft PowerPoint documents (.pptx) within the allowed path ({allowed_path}). "
                          "Returns text organised by slide, including titles, body text, and notes.",
            "input_schema": {
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": "Path to the .pptx file"
                    },
                    "include_notes": {
                        "type": "boolean",
                        "description": "Include speaker notes in the output",
                        "default": True
                    }
                },
                "required": ["path"]
            }
        },
        {
            "name": "read_pdf_document",
            "description": f"Extract text content from PDF documents within the allowed path ({allowed_path}). "
                          "Returns text organised by page. Can extract metadata and specific pages.",
            "input_schema": {
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": "Path to the .pdf file"
                    },
                    "page_numbers": {
                        "type": "array",
                        "items": {"type": "integer"},
                        "description": "Specific page numbers to extract (1-indexed). If not provided, extracts all pages."
                    },
                    "include_metadata": {
                        "type": "boolean",
                        "description": "Include document metadata (author, title, etc.)",
                        "default": True
                    }
                },
                "required": ["path"]
            }
        }
    ])

    # Write/create tools (only if access_mode is read_write)
    if access_mode == 'read_write':
        tools.extend([
            {
                "name": "create_word_document",
                "description": f"Create a Microsoft Word document (.docx) within the allowed path ({allowed_path}). "
                              "Supports creating from scratch with structured content, or using a template with placeholder replacement. "
                              "When using a template, placeholders in the format {{{{placeholder_name}}}} will be replaced with provided values.",
                "input_schema": {
                    "type": "object",
                    "properties": {
                        "path": {
                            "type": "string",
                            "description": "Path for the output .docx file"
                        },
                        "content": {
                            "type": "object",
                            "description": "Document content structure",
                            "properties": {
                                "title": {"type": "string", "description": "Document title"},
                                "paragraphs": {
                                    "type": "array",
                                    "items": {
                                        "type": "object",
                                        "properties": {
                                            "text": {"type": "string"},
                                            "style": {"type": "string", "description": "Style: Normal, Heading 1, Heading 2, Heading 3, Title"}
                                        }
                                    },
                                    "description": "List of paragraphs with optional styles"
                                }
                            }
                        },
                        "template_path": {
                            "type": "string",
                            "description": "Path to a .docx template file. If provided, placeholders will be replaced."
                        },
                        "placeholders": {
                            "type": "object",
                            "description": "Dictionary of placeholder names to values for template replacement",
                            "additionalProperties": {"type": "string"}
                        }
                    },
                    "required": ["path"]
                }
            },
            {
                "name": "create_excel_document",
                "description": f"Create a Microsoft Excel document (.xlsx) within the allowed path ({allowed_path}). "
                              "Creates spreadsheets from structured data. Supports multiple sheets.",
                "input_schema": {
                    "type": "object",
                    "properties": {
                        "path": {
                            "type": "string",
                            "description": "Path for the output .xlsx file"
                        },
                        "sheets": {
                            "type": "array",
                            "items": {
                                "type": "object",
                                "properties": {
                                    "name": {"type": "string", "description": "Sheet name"},
                                    "headers": {"type": "array", "items": {"type": "string"}, "description": "Column headers"},
                                    "data": {
                                        "type": "array",
                                        "items": {"type": "array"},
                                        "description": "2D array of cell values (rows of columns)"
                                    }
                                },
                                "required": ["name", "data"]
                            },
                            "description": "List of sheets to create"
                        }
                    },
                    "required": ["path", "sheets"]
                }
            },
            {
                "name": "create_powerpoint_document",
                "description": f"Create a Microsoft PowerPoint document (.pptx) within the allowed path ({allowed_path}). "
                              "Creates presentations with title and content slides. Supports templates with placeholder replacement.",
                "input_schema": {
                    "type": "object",
                    "properties": {
                        "path": {
                            "type": "string",
                            "description": "Path for the output .pptx file"
                        },
                        "slides": {
                            "type": "array",
                            "items": {
                                "type": "object",
                                "properties": {
                                    "layout": {
                                        "type": "string",
                                        "description": "Slide layout: title, title_content, content, blank"
                                    },
                                    "title": {"type": "string", "description": "Slide title"},
                                    "content": {
                                        "type": "array",
                                        "items": {"type": "string"},
                                        "description": "Bullet points or paragraphs"
                                    },
                                    "notes": {"type": "string", "description": "Speaker notes"}
                                }
                            },
                            "description": "List of slides to create"
                        },
                        "template_path": {
                            "type": "string",
                            "description": "Path to a .pptx template file"
                        },
                        "placeholders": {
                            "type": "object",
                            "description": "Dictionary of placeholder names to values for template replacement",
                            "additionalProperties": {"type": "string"}
                        }
                    },
                    "required": ["path", "slides"]
                }
            }
        ])

    return tools


def _execute_get_file_info(tool_input: Dict[str, Any],
                           config: Optional[Dict[str, Any]]) -> Dict[str, Any]:
    """Execute the get_file_info tool."""
    if not config.get('embedded_tools'):
        return {"success": False, "error": _ERR_DOC_NOT_CONFIGURED}

    doc_config = config.get('embedded_tools', {}).get('documents', {})
    if not doc_config.get('enabled', False):
        return {"success": False, "error": _ERR_DOC_NOT_ENABLED}

    allowed_path = doc_config.get('allowed_path', '.')
    file_path = tool_input.get('path')

    if not file_path:
        return {"success": False, "error": _ERR_FILE_PATH_REQUIRED}

    validation = _validate_path(file_path, allowed_path)
    if not validation['valid']:
        return {"success": False, "error": validation['error']}

    full_path = Path(validation['resolved_path'])

    if not full_path.exists():
        return {"success": False, "error": f"File does not exist: {file_path}"}

    if not full_path.is_file():
        return {"success": False, "error": f"Path is not a file: {file_path}"}

    try:
        import mimetypes
        stat_info = full_path.stat()

        # Try to get MIME type
        mime_type, _ = mimetypes.guess_type(str(full_path))

        # Try python-magic for more accurate detection
        try:
            import magic
            mime_type_magic = magic.from_file(str(full_path), mime=True)
            if mime_type_magic:
                mime_type = mime_type_magic
        except ImportError:
            pass
        except Exception:
            pass

        result = {
            "path": file_path,
            "full_path": str(full_path),
            "filename": full_path.name,
            "extension": full_path.suffix.lower(),
            "mime_type": mime_type or "application/octet-stream",
            "size_bytes": stat_info.st_size,
            "size_human": _format_size(stat_info.st_size),
            "modified": datetime.fromtimestamp(stat_info.st_mtime).isoformat(),
            "created": datetime.fromtimestamp(stat_info.st_ctime).isoformat()
        }

        logging.info(f"Got file info: {file_path}")
        return {"success": True, "result": result}

    except Exception as e:
        logging.error(f"Error getting file info {file_path}: {e}")
        return {"success": False, "error": str(e)}


def _format_size(size_bytes: int) -> str:
    """Format file size in human-readable format."""
    for unit in ['B', 'KB', 'MB', 'GB', 'TB']:
        if size_bytes < 1024.0:
            return f"{size_bytes:.2f} {unit}"
        size_bytes /= 1024.0
    return f"{size_bytes:.2f} PB"


def _execute_read_word_document(tool_input: Dict[str, Any],
                                config: Optional[Dict[str, Any]]) -> Dict[str, Any]:
    """Execute the read_word_document tool."""
    if not config.get('embedded_tools'):
        return {"success": False, "error": _ERR_DOC_NOT_CONFIGURED}

    doc_config = config.get('embedded_tools', {}).get('documents', {})
    if not doc_config.get('enabled', False):
        return {"success": False, "error": _ERR_DOC_NOT_ENABLED}

    allowed_path = doc_config.get('allowed_path', '.')
    max_size_mb = doc_config.get('max_file_size_mb', 50)

    file_path = tool_input.get('path')
    include_tables = tool_input.get('include_tables', True)
    include_headers_footers = tool_input.get('include_headers_footers', False)

    if not file_path:
        return {"success": False, "error": _ERR_FILE_PATH_REQUIRED}

    validation = _validate_path(file_path, allowed_path)
    if not validation['valid']:
        return {"success": False, "error": validation['error']}

    full_path = Path(validation['resolved_path'])

    if not full_path.exists():
        return {"success": False, "error": f"File does not exist: {file_path}"}

    if full_path.suffix.lower() != '.docx':
        return {"success": False, "error": f"File is not a Word document (.docx): {file_path}"}

    if full_path.stat().st_size > max_size_mb * 1024 * 1024:
        return {"success": False, "error": f"File exceeds maximum size of {max_size_mb} MB"}

    try:
        from docx import Document
        doc = Document(str(full_path))

        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append({
                    "text": para.text,
                    "style": para.style.name if para.style else "Normal"
                })

        tables = []
        if include_tables:
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                if table_data:
                    tables.append(table_data)

        headers_footers = []
        if include_headers_footers:
            for section in doc.sections:
                if section.header and section.header.paragraphs:
                    for para in section.header.paragraphs:
                        if para.text.strip():
                            headers_footers.append({"type": "header", "text": para.text})
                if section.footer and section.footer.paragraphs:
                    for para in section.footer.paragraphs:
                        if para.text.strip():
                            headers_footers.append({"type": "footer", "text": para.text})

        result = {
            "path": file_path,
            "full_path": str(full_path),
            "paragraph_count": len(paragraphs),
            "paragraphs": paragraphs,
            "table_count": len(tables),
            "tables": tables if tables else None,
            "headers_footers": headers_footers if headers_footers else None
        }

        logging.info(f"Read Word document: {file_path} ({len(paragraphs)} paragraphs, {len(tables)} tables)")
        return {"success": True, "result": result}

    except Exception as e:
        logging.error(f"Error reading Word document {file_path}: {e}")
        return {"success": False, "error": str(e)}


def _execute_read_excel_document(tool_input: Dict[str, Any],
                                 config: Optional[Dict[str, Any]]) -> Dict[str, Any]:
    """Execute the read_excel_document tool."""
    if not config.get('embedded_tools'):
        return {"success": False, "error": _ERR_DOC_NOT_CONFIGURED}

    doc_config = config.get('embedded_tools', {}).get('documents', {})
    if not doc_config.get('enabled', False):
        return {"success": False, "error": _ERR_DOC_NOT_ENABLED}

    allowed_path = doc_config.get('allowed_path', '.')
    max_size_mb = doc_config.get('max_file_size_mb', 50)
    default_max_rows = doc_config.get('reading', {}).get('max_excel_rows', 10000)

    file_path = tool_input.get('path')
    sheet_name = tool_input.get('sheet_name')
    include_all_sheets = tool_input.get('include_all_sheets', False)
    max_rows = tool_input.get('max_rows', 0) or default_max_rows

    if not file_path:
        return {"success": False, "error": _ERR_FILE_PATH_REQUIRED}

    validation = _validate_path(file_path, allowed_path)
    if not validation['valid']:
        return {"success": False, "error": validation['error']}

    full_path = Path(validation['resolved_path'])

    if not full_path.exists():
        return {"success": False, "error": f"File does not exist: {file_path}"}

    if full_path.suffix.lower() != '.xlsx':
        return {"success": False, "error": f"File is not an Excel document (.xlsx): {file_path}"}

    if full_path.stat().st_size > max_size_mb * 1024 * 1024:
        return {"success": False, "error": f"File exceeds maximum size of {max_size_mb} MB"}

    try:
        from openpyxl import load_workbook
        wb = load_workbook(str(full_path), read_only=True, data_only=True)

        sheets_data = {}
        sheet_names = wb.sheetnames

        if include_all_sheets:
            sheets_to_read = sheet_names
        elif sheet_name:
            if sheet_name not in sheet_names:
                return {"success": False, "error": f"Sheet '{sheet_name}' not found. Available: {sheet_names}"}
            sheets_to_read = [sheet_name]
        else:
            sheets_to_read = [wb.active.title] if wb.active else sheet_names[:1]

        for sname in sheets_to_read:
            ws = wb[sname]
            rows = []
            row_count = 0
            for row in ws.iter_rows(values_only=True):
                if row_count >= max_rows:
                    break
                rows.append(list(row))
                row_count += 1

            sheets_data[sname] = {
                "rows": rows,
                "row_count": len(rows),
                "truncated": row_count >= max_rows
            }

        wb.close()

        result = {
            "path": file_path,
            "full_path": str(full_path),
            "sheet_names": sheet_names,
            "sheets_read": list(sheets_data.keys()),
            "data": sheets_data,
            "max_rows_limit": max_rows
        }

        logging.info(f"Read Excel document: {file_path} ({len(sheets_data)} sheets)")
        return {"success": True, "result": result}

    except Exception as e:
        logging.error(f"Error reading Excel document {file_path}: {e}")
        return {"success": False, "error": str(e)}


def _execute_read_powerpoint_document(tool_input: Dict[str, Any],
                                      config: Optional[Dict[str, Any]]) -> Dict[str, Any]:
    """Execute the read_powerpoint_document tool."""
    if not config.get('embedded_tools'):
        return {"success": False, "error": _ERR_DOC_NOT_CONFIGURED}

    doc_config = config.get('embedded_tools', {}).get('documents', {})
    if not doc_config.get('enabled', False):
        return {"success": False, "error": _ERR_DOC_NOT_ENABLED}

    allowed_path = doc_config.get('allowed_path', '.')
    max_size_mb = doc_config.get('max_file_size_mb', 50)

    file_path = tool_input.get('path')
    include_notes = tool_input.get('include_notes', True)

    if not file_path:
        return {"success": False, "error": _ERR_FILE_PATH_REQUIRED}

    validation = _validate_path(file_path, allowed_path)
    if not validation['valid']:
        return {"success": False, "error": validation['error']}

    full_path = Path(validation['resolved_path'])

    if not full_path.exists():
        return {"success": False, "error": f"File does not exist: {file_path}"}

    if full_path.suffix.lower() != '.pptx':
        return {"success": False, "error": f"File is not a PowerPoint document (.pptx): {file_path}"}

    if full_path.stat().st_size > max_size_mb * 1024 * 1024:
        return {"success": False, "error": f"File exceeds maximum size of {max_size_mb} MB"}

    try:
        from pptx import Presentation
        prs = Presentation(str(full_path))

        slides = []
        for idx, slide in enumerate(prs.slides, 1):
            slide_data = {
                "slide_number": idx,
                "title": None,
                "content": []
            }

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            if shape.is_placeholder and hasattr(shape, 'placeholder_format'):
                                if shape.placeholder_format.type == 1:  # Title
                                    slide_data["title"] = text
                                else:
                                    slide_data["content"].append(text)
                            else:
                                slide_data["content"].append(text)

            if include_notes and slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame:
                    notes_text = notes_frame.text.strip()
                    if notes_text:
                        slide_data["notes"] = notes_text

            slides.append(slide_data)

        result = {
            "path": file_path,
            "full_path": str(full_path),
            "slide_count": len(slides),
            "slides": slides
        }

        logging.info(f"Read PowerPoint document: {file_path} ({len(slides)} slides)")
        return {"success": True, "result": result}

    except Exception as e:
        logging.error(f"Error reading PowerPoint document {file_path}: {e}")
        return {"success": False, "error": str(e)}


def _execute_read_pdf_document(tool_input: Dict[str, Any],
                               config: Optional[Dict[str, Any]]) -> Dict[str, Any]:
    """Execute the read_pdf_document tool."""
    if not config.get('embedded_tools'):
        return {"success": False, "error": _ERR_DOC_NOT_CONFIGURED}

    doc_config = config.get('embedded_tools', {}).get('documents', {})
    if not doc_config.get('enabled', False):
        return {"success": False, "error": _ERR_DOC_NOT_ENABLED}

    allowed_path = doc_config.get('allowed_path', '.')
    max_size_mb = doc_config.get('max_file_size_mb', 50)
    max_pages = doc_config.get('reading', {}).get('max_pdf_pages', 100)

    file_path = tool_input.get('path')
    page_numbers = tool_input.get('page_numbers')
    include_metadata = tool_input.get('include_metadata', True)

    if not file_path:
        return {"success": False, "error": _ERR_FILE_PATH_REQUIRED}

    validation = _validate_path(file_path, allowed_path)
    if not validation['valid']:
        return {"success": False, "error": validation['error']}

    full_path = Path(validation['resolved_path'])

    if not full_path.exists():
        return {"success": False, "error": f"File does not exist: {file_path}"}

    if full_path.suffix.lower() != '.pdf':
        return {"success": False, "error": f"File is not a PDF document (.pdf): {file_path}"}

    if full_path.stat().st_size > max_size_mb * 1024 * 1024:
        return {"success": False, "error": f"File exceeds maximum size of {max_size_mb} MB"}

    try:
        import pdfplumber

        pages_data = []
        metadata = None

        with pdfplumber.open(str(full_path)) as pdf:
            total_pages = len(pdf.pages)

            if include_metadata:
                metadata = pdf.metadata

            # Determine which pages to extract
            if page_numbers:
                pages_to_read = [p - 1 for p in page_numbers if 0 < p <= total_pages]
            else:
                pages_to_read = list(range(min(total_pages, max_pages)))

            for page_idx in pages_to_read:
                page = pdf.pages[page_idx]
                text = page.extract_text() or ""
                pages_data.append({
                    "page_number": page_idx + 1,
                    "text": text,
                    "width": page.width,
                    "height": page.height
                })

        result = {
            "path": file_path,
            "full_path": str(full_path),
            "total_pages": total_pages,
            "pages_extracted": len(pages_data),
            "pages": pages_data,
            "truncated": len(pages_data) < total_pages and not page_numbers,
            "metadata": metadata
        }

        logging.info(f"Read PDF document: {file_path} ({len(pages_data)} pages)")
        return {"success": True, "result": result}

    except Exception as e:
        logging.error(f"Error reading PDF document {file_path}: {e}")
        return {"success": False, "error": str(e)}


def _execute_create_word_document(tool_input: Dict[str, Any],
                                  config: Optional[Dict[str, Any]]) -> Dict[str, Any]:
    """Execute the create_word_document tool."""
    if not config.get('embedded_tools'):
        return {"success": False, "error": _ERR_DOC_NOT_CONFIGURED}

    doc_config = config.get('embedded_tools', {}).get('documents', {})
    if not doc_config.get('enabled', False):
        return {"success": False, "error": _ERR_DOC_NOT_ENABLED}

    if doc_config.get('access_mode', 'read') != 'read_write':
        return {"success": False, "error": _ERR_WRITE_MODE_REQUIRED}

    allowed_path = doc_config.get('allowed_path', '.')
    templates_path = doc_config.get('creation', {}).get('templates_path')

    file_path = tool_input.get('path')
    content = tool_input.get('content', {})
    template_path = tool_input.get('template_path')
    placeholders = tool_input.get('placeholders', {})

    if not file_path:
        return {"success": False, "error": _ERR_OUTPUT_PATH_REQUIRED}

    validation = _validate_path(file_path, allowed_path)
    if not validation['valid']:
        return {"success": False, "error": validation['error']}

    full_path = Path(validation['resolved_path'])

    if not full_path.parent.exists():
        return {"success": False, "error": f"Parent directory does not exist: {full_path.parent}"}

    try:
        from docx import Document

        # Use template if provided
        if template_path:
            # Validate template path
            if templates_path:
                template_full = Path(templates_path) / template_path
            else:
                template_validation = _validate_path(template_path, allowed_path)
                if not template_validation['valid']:
                    return {"success": False, "error": f"Template path error: {template_validation['error']}"}
                template_full = Path(template_validation['resolved_path'])

            if not template_full.exists():
                return {"success": False, "error": f"Template does not exist: {template_path}"}

            doc = Document(str(template_full))

            # Replace placeholders in paragraphs
            for para in doc.paragraphs:
                for key, value in placeholders.items():
                    if f"{{{{{key}}}}}" in para.text:
                        for run in para.runs:
                            run.text = run.text.replace(f"{{{{{key}}}}}", str(value))

            # Replace placeholders in tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        for key, value in placeholders.items():
                            if f"{{{{{key}}}}}" in cell.text:
                                cell.text = cell.text.replace(f"{{{{{key}}}}}", str(value))

        else:
            doc = Document()

            # Add title if provided
            if content.get('title'):
                doc.add_heading(content['title'], 0)

            # Define valid built-in styles that python-docx supports
            valid_styles = {
                'Normal', 'Title', 'Subtitle', 'Quote', 'Intense Quote',
                'List Paragraph', 'List Bullet', 'List Number',
                'Heading 1', 'Heading 2', 'Heading 3', 'Heading 4',
                'Heading 5', 'Heading 6', 'Heading 7', 'Heading 8', 'Heading 9',
                'Body Text', 'Body Text 2', 'Body Text 3',
                'Caption', 'Macro Text', 'No Spacing'
            }
            invalid_styles_used = set()

            # Add paragraphs
            for para_data in content.get('paragraphs', []):
                text = para_data.get('text', '')
                style = para_data.get('style', 'Normal')
                if style.startswith('Heading'):
                    level = int(style.split()[-1]) if style.split()[-1].isdigit() else 1
                    doc.add_heading(text, level)
                else:
                    # Validate style - fall back to Normal if invalid
                    if style not in valid_styles:
                        invalid_styles_used.add(style)
                        style = 'Normal'
                    doc.add_paragraph(text, style=style)

            if invalid_styles_used:
                logging.warning(f"Invalid styles replaced with 'Normal': {', '.join(sorted(invalid_styles_used))}")

        doc.save(str(full_path))

        result = {
            "path": file_path,
            "full_path": str(full_path),
            "size_bytes": full_path.stat().st_size,
            "used_template": template_path is not None,
            "placeholders_replaced": list(placeholders.keys()) if placeholders else []
        }

        logging.info(f"Created Word document: {file_path}")
        return {"success": True, "result": result}

    except Exception as e:
        logging.error(f"Error creating Word document {file_path}: {e}")
        return {"success": False, "error": str(e)}


def _execute_create_excel_document(tool_input: Dict[str, Any],
                                   config: Optional[Dict[str, Any]]) -> Dict[str, Any]:
    """Execute the create_excel_document tool."""
    if not config.get('embedded_tools'):
        return {"success": False, "error": _ERR_DOC_NOT_CONFIGURED}

    doc_config = config.get('embedded_tools', {}).get('documents', {})
    if not doc_config.get('enabled', False):
        return {"success": False, "error": _ERR_DOC_NOT_ENABLED}

    if doc_config.get('access_mode', 'read') != 'read_write':
        return {"success": False, "error": _ERR_WRITE_MODE_REQUIRED}

    allowed_path = doc_config.get('allowed_path', '.')

    file_path = tool_input.get('path')
    sheets = tool_input.get('sheets', [])

    if not file_path:
        return {"success": False, "error": _ERR_OUTPUT_PATH_REQUIRED}

    if not sheets:
        return {"success": False, "error": "At least one sheet is required"}

    validation = _validate_path(file_path, allowed_path)
    if not validation['valid']:
        return {"success": False, "error": validation['error']}

    full_path = Path(validation['resolved_path'])

    if not full_path.parent.exists():
        return {"success": False, "error": f"Parent directory does not exist: {full_path.parent}"}

    try:
        from openpyxl import Workbook

        wb = Workbook()
        # Remove default sheet
        if 'Sheet' in wb.sheetnames:
            del wb['Sheet']

        for sheet_data in sheets:
            sheet_name = sheet_data.get('name', 'Sheet')
            headers = sheet_data.get('headers', [])
            data = sheet_data.get('data', [])

            ws = wb.create_sheet(title=sheet_name)

            # Add headers if provided
            if headers:
                for col, header in enumerate(headers, 1):
                    ws.cell(row=1, column=col, value=header)
                start_row = 2
            else:
                start_row = 1

            # Add data
            for row_idx, row_data in enumerate(data, start_row):
                for col_idx, value in enumerate(row_data, 1):
                    ws.cell(row=row_idx, column=col_idx, value=value)

        wb.save(str(full_path))

        result = {
            "path": file_path,
            "full_path": str(full_path),
            "size_bytes": full_path.stat().st_size,
            "sheets_created": [s.get('name', 'Sheet') for s in sheets]
        }

        logging.info(f"Created Excel document: {file_path}")
        return {"success": True, "result": result}

    except Exception as e:
        logging.error(f"Error creating Excel document {file_path}: {e}")
        return {"success": False, "error": str(e)}


def _execute_create_powerpoint_document(tool_input: Dict[str, Any],
                                        config: Optional[Dict[str, Any]]) -> Dict[str, Any]:
    """Execute the create_powerpoint_document tool."""
    if not config.get('embedded_tools'):
        return {"success": False, "error": _ERR_DOC_NOT_CONFIGURED}

    doc_config = config.get('embedded_tools', {}).get('documents', {})
    if not doc_config.get('enabled', False):
        return {"success": False, "error": _ERR_DOC_NOT_ENABLED}

    if doc_config.get('access_mode', 'read') != 'read_write':
        return {"success": False, "error": _ERR_WRITE_MODE_REQUIRED}

    allowed_path = doc_config.get('allowed_path', '.')
    templates_path = doc_config.get('creation', {}).get('templates_path')

    file_path = tool_input.get('path')
    slides_data = tool_input.get('slides', [])
    template_path = tool_input.get('template_path')
    placeholders = tool_input.get('placeholders', {})

    if not file_path:
        return {"success": False, "error": _ERR_OUTPUT_PATH_REQUIRED}

    if not slides_data and not template_path:
        return {"success": False, "error": "Either slides or template_path is required"}

    validation = _validate_path(file_path, allowed_path)
    if not validation['valid']:
        return {"success": False, "error": validation['error']}

    full_path = Path(validation['resolved_path'])

    if not full_path.parent.exists():
        return {"success": False, "error": f"Parent directory does not exist: {full_path.parent}"}

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        # Use template if provided
        if template_path:
            if templates_path:
                template_full = Path(templates_path) / template_path
            else:
                template_validation = _validate_path(template_path, allowed_path)
                if not template_validation['valid']:
                    return {"success": False, "error": f"Template path error: {template_validation['error']}"}
                template_full = Path(template_validation['resolved_path'])

            if not template_full.exists():
                return {"success": False, "error": f"Template does not exist: {template_path}"}

            prs = Presentation(str(template_full))

            # Replace placeholders in existing slides
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                for key, value in placeholders.items():
                                    if f"{{{{{key}}}}}" in run.text:
                                        run.text = run.text.replace(f"{{{{{key}}}}}", str(value))
        else:
            prs = Presentation()

        # Add new slides
        for slide_data in slides_data:
            layout_name = slide_data.get('layout', 'title_content')
            title = slide_data.get('title', '')
            content = slide_data.get('content', [])
            notes = slide_data.get('notes', '')

            # Map layout names to indices
            layout_map = {
                'title': 0,
                'title_content': 1,
                'content': 5,
                'blank': 6
            }
            layout_idx = layout_map.get(layout_name, 1)

            if layout_idx < len(prs.slide_layouts):
                slide_layout = prs.slide_layouts[layout_idx]
            else:
                slide_layout = prs.slide_layouts[0]

            slide = prs.slides.add_slide(slide_layout)

            # Set title
            if title and slide.shapes.title:
                slide.shapes.title.text = title

            # Add content
            if content:
                for shape in slide.shapes:
                    if shape.has_text_frame and shape != slide.shapes.title:
                        tf = shape.text_frame
                        tf.clear()
                        for i, text in enumerate(content):
                            if i == 0:
                                tf.paragraphs[0].text = text
                            else:
                                p = tf.add_paragraph()
                                p.text = text
                        break

            # Add notes
            if notes:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = notes

        prs.save(str(full_path))

        result = {
            "path": file_path,
            "full_path": str(full_path),
            "size_bytes": full_path.stat().st_size,
            "slides_added": len(slides_data),
            "used_template": template_path is not None,
            "placeholders_replaced": list(placeholders.keys()) if placeholders else []
        }

        logging.info(f"Created PowerPoint document: {file_path}")
        return {"success": True, "result": result}

    except Exception as e:
        logging.error(f"Error creating PowerPoint document {file_path}: {e}")
        return {"success": False, "error": str(e)}


# ============================================================================
# Archive Tools
# ============================================================================

def _get_archive_tools(archive_config: Dict[str, Any]) -> List[Dict[str, Any]]:
    """
    Get archive tool definitions based on configuration.

    Args:
        archive_config: Archive tools configuration dictionary

    Returns:
        List of archive tool definitions
    """
    access_mode = archive_config.get('access_mode', 'read')
    allowed_path = archive_config.get('allowed_path', '.')

    tools = [
        {
            "name": "list_archive_contents",
            "description": f"List the contents of an archive file within the allowed path ({allowed_path}). "
                          "Supports .zip, .tar, .tar.gz, .tgz, and .tar.bz2 files. "
                          "Returns file names, sizes, and modification times.",
            "input_schema": {
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": _DESC_ARCHIVE_PATH
                    }
                },
                "required": ["path"]
            }
        },
        {
            "name": "read_archive_file",
            "description": f"Read a specific file from within an archive without extracting to disk ({allowed_path}). "
                          "Returns the file content. Text files are returned as strings, binary files as base64.",
            "input_schema": {
                "type": "object",
                "properties": {
                    "archive_path": {
                        "type": "string",
                        "description": _DESC_ARCHIVE_PATH
                    },
                    "file_path": {
                        "type": "string",
                        "description": "Path of the file within the archive to read"
                    },
                    "encoding": {
                        "type": "string",
                        "description": "Text encoding to use when reading as text (default: utf-8)",
                        "default": "utf-8"
                    },
                    "as_binary": {
                        "type": "boolean",
                        "description": "Force reading as binary (returns base64)",
                        "default": False
                    }
                },
                "required": ["archive_path", "file_path"]
            }
        }
    ]

    # Extract tool (only if access_mode is read_write)
    if access_mode == 'read_write':
        tools.append({
            "name": "extract_archive",
            "description": f"Extract an archive file to a destination directory within the allowed path ({allowed_path}). "
                          "Supports .zip, .tar, .tar.gz, .tgz, and .tar.bz2 files. "
                          "Can extract all files or specific files only.",
            "input_schema": {
                "type": "object",
                "properties": {
                    "archive_path": {
                        "type": "string",
                        "description": _DESC_ARCHIVE_PATH
                    },
                    "destination": {
                        "type": "string",
                        "description": "Destination directory for extracted files"
                    },
                    "files": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "Specific files to extract. If not provided, extracts all files."
                    },
                    "overwrite": {
                        "type": "boolean",
                        "description": "Overwrite existing files",
                        "default": False
                    }
                },
                "required": ["archive_path", "destination"]
            }
        })

    return tools


def _get_archive_type(file_path: Path) -> Optional[str]:
    """Determine archive type from file extension."""
    suffix = file_path.suffix.lower()
    name = file_path.name.lower()

    if suffix == '.zip':
        return 'zip'
    elif suffix == '.tar':
        return 'tar'
    elif name.endswith('.tar.gz') or suffix == '.tgz':
        return _TAR_GZ
    elif name.endswith('.tar.bz2'):
        return _TAR_BZ2
    return None


def _execute_list_archive_contents(tool_input: Dict[str, Any],
                                   config: Optional[Dict[str, Any]]) -> Dict[str, Any]:
    """Execute the list_archive_contents tool."""
    import zipfile
    import tarfile

    if not config.get('embedded_tools'):
        return {"success": False, "error": _ERR_ARCHIVE_NOT_CONFIGURED}

    archive_config = config.get('embedded_tools', {}).get('archives', {})
    if not archive_config.get('enabled', False):
        return {"success": False, "error": _ERR_ARCHIVE_NOT_ENABLED}

    allowed_path = archive_config.get('allowed_path', '.')
    max_size_mb = archive_config.get('max_file_size_mb', 100)
    max_files = archive_config.get('max_files_to_list', 1000)

    file_path = tool_input.get('path')

    if not file_path:
        return {"success": False, "error": _ERR_ARCHIVE_PATH_REQUIRED}

    validation = _validate_path(file_path, allowed_path)
    if not validation['valid']:
        return {"success": False, "error": validation['error']}

    full_path = Path(validation['resolved_path'])

    if not full_path.exists():
        return {"success": False, "error": f"Archive does not exist: {file_path}"}

    if full_path.stat().st_size > max_size_mb * 1024 * 1024:
        return {"success": False, "error": f"Archive exceeds maximum size of {max_size_mb} MB"}

    archive_type = _get_archive_type(full_path)
    if not archive_type:
        return {"success": False, "error": f"Unsupported archive format: {full_path.suffix}"}

    try:
        files = []

        if archive_type == 'zip':
            with zipfile.ZipFile(str(full_path), 'r') as zf:
                for info in zf.infolist()[:max_files]:
                    files.append({
                        "path": info.filename,
                        "size_bytes": info.file_size,
                        "compressed_size": info.compress_size,
                        "is_directory": info.is_dir(),
                        "modified": datetime(*info.date_time).isoformat() if info.date_time else None
                    })
        else:
            mode = _TAR_OPEN_MODES.get(archive_type, 'r')
            with tarfile.open(str(full_path), mode) as tf:
                count = 0
                for member in tf:
                    if count >= max_files:
                        break
                    files.append({
                        "path": member.name,
                        "size_bytes": member.size,
                        "is_directory": member.isdir(),
                        "modified": datetime.fromtimestamp(member.mtime).isoformat() if member.mtime else None
                    })
                    count += 1

        result = {
            "path": file_path,
            "full_path": str(full_path),
            "archive_type": archive_type,
            "total_files": len(files),
            "truncated": len(files) >= max_files,
            "files": files
        }

        logging.info(f"Listed archive contents: {file_path} ({len(files)} files)")
        return {"success": True, "result": result}

    except Exception as e:
        logging.error(f"Error listing archive {file_path}: {e}")
        return {"success": False, "error": str(e)}


def _execute_read_archive_file(tool_input: Dict[str, Any],
                               config: Optional[Dict[str, Any]]) -> Dict[str, Any]:
    """Execute the read_archive_file tool."""
    import zipfile
    import tarfile

    if not config.get('embedded_tools'):
        return {"success": False, "error": _ERR_ARCHIVE_NOT_CONFIGURED}

    archive_config = config.get('embedded_tools', {}).get('archives', {})
    if not archive_config.get('enabled', False):
        return {"success": False, "error": _ERR_ARCHIVE_NOT_ENABLED}

    allowed_path = archive_config.get('allowed_path', '.')
    max_size_mb = archive_config.get('max_file_size_mb', 100)

    archive_path = tool_input.get('archive_path')
    file_path = tool_input.get('file_path')
    encoding = tool_input.get('encoding', 'utf-8')
    as_binary = tool_input.get('as_binary', False)

    if not archive_path:
        return {"success": False, "error": _ERR_ARCHIVE_PATH_REQUIRED}

    if not file_path:
        return {"success": False, "error": "File path within archive is required"}

    validation = _validate_path(archive_path, allowed_path)
    if not validation['valid']:
        return {"success": False, "error": validation['error']}

    full_path = Path(validation['resolved_path'])

    if not full_path.exists():
        return {"success": False, "error": f"Archive does not exist: {archive_path}"}

    if full_path.stat().st_size > max_size_mb * 1024 * 1024:
        return {"success": False, "error": f"Archive exceeds maximum size of {max_size_mb} MB"}

    archive_type = _get_archive_type(full_path)
    if not archive_type:
        return {"success": False, "error": f"Unsupported archive format: {full_path.suffix}"}

    try:
        content = None

        if archive_type == 'zip':
            with zipfile.ZipFile(str(full_path), 'r') as zf:
                if file_path not in zf.namelist():
                    return {"success": False, "error": f"File not found in archive: {file_path}"}
                content = zf.read(file_path)
        else:
            mode = _TAR_OPEN_MODES.get(archive_type, 'r')
            with tarfile.open(str(full_path), mode) as tf:
                try:
                    member = tf.getmember(file_path)
                    f = tf.extractfile(member)
                    if f:
                        content = f.read()
                    else:
                        return {"success": False, "error": f"Cannot read directory: {file_path}"}
                except KeyError:
                    return {"success": False, "error": f"File not found in archive: {file_path}"}

        # Try to decode as text unless binary requested
        if as_binary:
            result = {
                "archive_path": archive_path,
                "file_path": file_path,
                "content_base64": base64.b64encode(content).decode('utf-8'),
                "size_bytes": len(content),
                "is_binary": True
            }
        else:
            try:
                text_content = content.decode(encoding)
                result = {
                    "archive_path": archive_path,
                    "file_path": file_path,
                    "content": text_content,
                    "size_bytes": len(content),
                    "encoding": encoding,
                    "is_binary": False
                }
            except UnicodeDecodeError:
                result = {
                    "archive_path": archive_path,
                    "file_path": file_path,
                    "content_base64": base64.b64encode(content).decode('utf-8'),
                    "size_bytes": len(content),
                    "is_binary": True,
                    "note": f"Could not decode as {encoding}, returned as base64"
                }

        logging.info(f"Read file from archive: {archive_path}/{file_path}")
        return {"success": True, "result": result}

    except Exception as e:
        logging.error(f"Error reading from archive {archive_path}: {e}")
        return {"success": False, "error": str(e)}


def _execute_extract_archive(tool_input: Dict[str, Any],
                             config: Optional[Dict[str, Any]]) -> Dict[str, Any]:
    """Execute the extract_archive tool."""
    import zipfile
    import tarfile

    if not config.get('embedded_tools'):
        return {"success": False, "error": _ERR_ARCHIVE_NOT_CONFIGURED}

    archive_config = config.get('embedded_tools', {}).get('archives', {})
    if not archive_config.get('enabled', False):
        return {"success": False, "error": _ERR_ARCHIVE_NOT_ENABLED}

    if archive_config.get('access_mode', 'read') != 'read_write':
        return {"success": False, "error": "Extract operations require access_mode: read_write"}

    allowed_path = archive_config.get('allowed_path', '.')
    max_size_mb = archive_config.get('max_file_size_mb', 100)

    archive_path = tool_input.get('archive_path')
    destination = tool_input.get('destination')
    files_to_extract = tool_input.get('files')
    overwrite = tool_input.get('overwrite', False)

    if not archive_path:
        return {"success": False, "error": _ERR_ARCHIVE_PATH_REQUIRED}

    if not destination:
        return {"success": False, "error": "Destination directory is required"}

    # Validate archive path
    archive_validation = _validate_path(archive_path, allowed_path)
    if not archive_validation['valid']:
        return {"success": False, "error": archive_validation['error']}

    full_archive_path = Path(archive_validation['resolved_path'])

    # Validate destination path
    dest_validation = _validate_path(destination, allowed_path)
    if not dest_validation['valid']:
        return {"success": False, "error": dest_validation['error']}

    full_dest_path = Path(dest_validation['resolved_path'])

    if not full_archive_path.exists():
        return {"success": False, "error": f"Archive does not exist: {archive_path}"}

    if full_archive_path.stat().st_size > max_size_mb * 1024 * 1024:
        return {"success": False, "error": f"Archive exceeds maximum size of {max_size_mb} MB"}

    archive_type = _get_archive_type(full_archive_path)
    if not archive_type:
        return {"success": False, "error": f"Unsupported archive format: {full_archive_path.suffix}"}

    try:
        # Create destination directory
        full_dest_path.mkdir(parents=True, exist_ok=True)

        extracted_files = []

        if archive_type == 'zip':
            with zipfile.ZipFile(str(full_archive_path), 'r') as zf:
                members = files_to_extract if files_to_extract else zf.namelist()
                for member in members:
                    if member in zf.namelist():
                        dest_file = full_dest_path / member
                        if dest_file.exists() and not overwrite:
                            continue
                        zf.extract(member, str(full_dest_path))
                        extracted_files.append(member)
        else:
            mode = _TAR_OPEN_MODES.get(archive_type, 'r')
            with tarfile.open(str(full_archive_path), mode) as tf:
                if files_to_extract:
                    members = [tf.getmember(f) for f in files_to_extract if f in tf.getnames()]
                else:
                    members = tf.getmembers()

                for member in members:
                    dest_file = full_dest_path / member.name
                    if dest_file.exists() and not overwrite:
                        continue
                    tf.extract(member, str(full_dest_path))
                    extracted_files.append(member.name)

        result = {
            "archive_path": archive_path,
            "destination": destination,
            "full_destination": str(full_dest_path),
            "files_extracted": len(extracted_files),
            "extracted": extracted_files
        }

        logging.info(f"Extracted archive: {archive_path} -> {destination} ({len(extracted_files)} files)")
        return {"success": True, "result": result}

    except Exception as e:
        logging.error(f"Error extracting archive {archive_path}: {e}")
        return {"success": False, "error": str(e)}
