from aiecs.tools import register_tool
from aiecs.tools.base_tool import BaseTool
from pydantic import BaseModel, field_validator, Field
from pydantic_settings import BaseSettings, SettingsConfigDict
from pptx.util import Inches
from pptx import Presentation
from docx.shared import Pt
from docx import Document as DocxDocument
from tika import parser  # type: ignore[import-untyped]
import os
import logging
import warnings
from typing import List, Dict, Optional, Any

import pandas as pd  # type: ignore[import-untyped]
import pdfplumber
import pytesseract  # type: ignore[import-untyped]
from PIL import Image

# Tika log path will be configured via Config class

# Suppress pkg_resources deprecation warning from tika
warnings.filterwarnings("ignore", category=UserWarning, module="tika")


# Module-level default configuration for validators
_DEFAULT_MAX_FILE_SIZE_MB = 100
_DEFAULT_ALLOWED_EXTENSIONS = [
    ".docx",
    ".pptx",
    ".xlsx",
    ".pdf",
    ".png",
    ".jpg",
    ".jpeg",
    ".tiff",
    ".bmp",
    ".gif",
]

# Exceptions


class OfficeToolError(Exception):
    """Base exception for OfficeTool errors."""


class InputValidationError(OfficeToolError):
    """Raised when input validation fails."""


class FileOperationError(OfficeToolError):
    """Raised when file operations fail."""


class SecurityError(OfficeToolError):
    """Raised for security-related issues."""


class ContentValidationError(OfficeToolError):
    """Raised when document content validation fails."""


# Base schema for common fields


class BaseFileSchema(BaseModel):
    file_path: Optional[str] = None
    output_path: Optional[str] = None
    image_path: Optional[str] = None

    @field_validator("file_path", "output_path", "image_path")
    def validate_path(cls, v: Optional[str], field) -> Optional[str]:
        """Validate file paths for existence, size, extension, and path traversal."""
        if not v:
            return v
        abs_path = os.path.abspath(os.path.normpath(v))
        # Check for path traversal
        if ".." in v or "~" in v or "%" in v:
            raise SecurityError(f"Path traversal attempt detected: {v}")
        # Ensure path is in allowed directories
        base_dir = os.path.abspath(os.getcwd())
        allowed_dirs = [os.path.abspath(os.path.normpath(d)) for d in ["/tmp", "./data", "./uploads"]]
        if not abs_path.startswith(base_dir) and not any(abs_path.startswith(d) for d in allowed_dirs):
            raise SecurityError(f"Path not in allowed directories: {abs_path}")
        # Check extension
        ext = os.path.splitext(abs_path)[1].lower()
        if ext not in _DEFAULT_ALLOWED_EXTENSIONS:
            raise SecurityError(f"Extension '{ext}' not allowed for '{field.field_name}', expected {_DEFAULT_ALLOWED_EXTENSIONS}")
        # Check file existence and size for input paths
        if field.field_name == "file_path":
            if not os.path.isfile(abs_path):
                raise FileOperationError(f"{field.field_name}: File not found: {abs_path}")
            size_mb = os.path.getsize(abs_path) / (1024 * 1024)
            if size_mb > _DEFAULT_MAX_FILE_SIZE_MB:
                raise FileOperationError(f"{field.field_name}: File too large: {size_mb:.1f}MB, max {_DEFAULT_MAX_FILE_SIZE_MB}MB")
        # Check for existing output paths
        elif field.field_name == "output_path" and os.path.exists(abs_path):
            raise FileOperationError(f"{field.field_name}: File already exists: {abs_path}")
        return abs_path


# Schemas for operations - moved to OfficeTool class as inner classes


@register_tool("office")
class OfficeTool(BaseTool):
    """
    Office document processing tool supporting:
      - read_docx: Read content from DOCX files.
      - write_docx: Write content to DOCX files.
      - read_pptx: Read content from PPTX files.
      - write_pptx: Write content to PPTX files.
      - read_xlsx: Read content from XLSX files.
      - write_xlsx: Write content to XLSX files.
      - extract_text: Extract text from various file formats.

    Inherits from BaseTool to leverage ToolExecutor for caching, concurrency, and error handling.
    """

    # Configuration schema
    class Config(BaseSettings):
        """Configuration for the office tool
        
        Automatically reads from environment variables with OFFICE_TOOL_ prefix.
        Example: OFFICE_TOOL_MAX_FILE_SIZE_MB -> max_file_size_mb
        """

        model_config = SettingsConfigDict(env_prefix="OFFICE_TOOL_")

        max_file_size_mb: int = Field(default=100, description="Maximum file size in megabytes")
        default_font: str = Field(default="Arial", description="Default font for documents")
        default_font_size: int = Field(default=12, description="Default font size in points")
        tika_log_path: str = Field(
            default=os.path.expanduser("~/.cache/tika"),
            description="Tika log directory path",
        )
        allowed_extensions: List[str] = Field(
            default=[
                ".docx",
                ".pptx",
                ".xlsx",
                ".pdf",
                ".png",
                ".jpg",
                ".jpeg",
                ".tiff",
                ".bmp",
                ".gif",
            ],
            description="Allowed document file extensions",
        )

    # Schema definitions
    class Read_docxSchema(BaseFileSchema):
        """Schema for read_docx operation"""

        file_path: str = Field(description="Path to the DOCX file to read")
        include_tables: bool = Field(default=False, description="Whether to include table data in the output. If True, tables are included as nested lists")

    class Write_docxSchema(BaseFileSchema):
        """Schema for write_docx operation"""

        text: str = Field(description="Text content to write to the DOCX file")
        output_path: str = Field(description="Path where the DOCX file will be saved")
        table_data: Optional[List[List[str]]] = Field(default=None, description="Optional table data to include in the document. Each inner list represents a row, each string represents a cell")

    class Read_pptxSchema(BaseFileSchema):
        """Schema for read_pptx operation"""

        file_path: str = Field(description="Path to the PPTX file to read")

    class Write_pptxSchema(BaseFileSchema):
        """Schema for write_pptx operation"""

        slides: List[str] = Field(description="List of slide content strings. Each string becomes a slide")
        output_path: str = Field(description="Path where the PPTX file will be saved")
        image_path: Optional[str] = Field(default=None, description="Optional path to an image file to include on the first slide")

    class Read_xlsxSchema(BaseFileSchema):
        """Schema for read_xlsx operation"""

        file_path: str = Field(description="Path to the XLSX file to read")
        sheet_name: Optional[str] = Field(default=None, description="Optional name of the sheet to read. If None, reads the first sheet")

    class Write_xlsxSchema(BaseFileSchema):
        """Schema for write_xlsx operation"""

        data: List[Dict[str, Any]] = Field(description="List of dictionaries representing Excel rows. Each dictionary key becomes a column header, values become cell data")
        output_path: str = Field(description="Path where the XLSX file will be saved")
        sheet_name: str = Field(default="Sheet1", description="Name of the Excel sheet to create")

    class Extract_textSchema(BaseFileSchema):
        """Schema for extract_text operation"""

        file_path: str = Field(description="Path to the file to extract text from. Supports DOCX, PPTX, XLSX, PDF, and image formats")

    def __init__(self, config: Optional[Dict[str, Any]] = None, **kwargs):
        """
        Initialize OfficeTool with configuration.

        Configuration is automatically loaded by BaseTool from:
        1. Explicit config dict (highest priority)
        2. YAML config files (config/tools/office_tool.yaml)
        3. Environment variables (via dotenv from .env files)
        4. Tool defaults (lowest priority)

        Args:
            config (Dict, optional): Configuration overrides for OfficeTool.
            **kwargs: Additional arguments passed to BaseTool (e.g., tool_name)

        Raises:
            ValueError: If config contains invalid settings.
        """
        super().__init__(config, **kwargs)

        # Configuration is automatically loaded by BaseTool into self._config_obj
        # Access config via self._config_obj (BaseSettings instance)
        self.config = self._config_obj if self._config_obj else self.Config()
        
        # Configure Tika log path from config
        os.environ["TIKA_LOG_PATH"] = self.config.tika_log_path
        os.makedirs(self.config.tika_log_path, exist_ok=True)

        self.logger = logging.getLogger(__name__)
        if not self.logger.handlers:
            handler = logging.StreamHandler()
            handler.setFormatter(logging.Formatter("%(asctime)s %(levelname)s %(message)s"))
            self.logger.addHandler(handler)
        self.logger.setLevel(logging.INFO)

    def _validate_document(self, file_path: str, file_type: str) -> None:
        """
        Validate document structure before processing.

        Args:
            file_path (str): Path to the document file.
            file_type (str): Type of document ('docx', 'pptx', 'xlsx', 'pdf', 'image').

        Raises:
            ContentValidationError: If document structure is invalid.
        """
        try:
            if file_type == "docx":
                doc = DocxDocument(file_path)
                if not hasattr(doc, "paragraphs"):
                    raise ContentValidationError("Invalid DOCX structure")
            elif file_type == "pptx":
                prs = Presentation(file_path)
                if not hasattr(prs, "slides"):
                    raise ContentValidationError("Invalid PPTX structure")
            elif file_type == "xlsx":
                # Just validate that file can be read - don't care about return
                # type
                pd.read_excel(file_path, nrows=5)
            elif file_type == "pdf":
                with pdfplumber.open(file_path) as pdf:
                    if len(pdf.pages) == 0:
                        raise ContentValidationError("PDF has no pages")
            elif file_type == "image":
                img = Image.open(file_path)
                img.verify()  # Verify it's a valid image
            else:
                # Use tika as fallback for other formats
                parsed = parser.from_file(file_path)
                if not parsed or not parsed.get("content"):
                    raise ContentValidationError("Unable to parse file content")
        except Exception as e:
            raise ContentValidationError(f"Invalid {file_type.upper()} file: {str(e)}")

    def _sanitize_text(self, text: str) -> str:
        """
        Sanitize text to remove potentially harmful control characters.

        Args:
            text (str): Input text.

        Returns:
            str: Sanitized text.
        """
        if not text:
            return ""
        return "".join(char for char in text if ord(char) >= 32 or char in "\n\r\t")

    def _sanitize_table_data(self, table_data: Optional[List[List[str]]]) -> Optional[List[List[str]]]:
        """
        Sanitize table data to remove harmful content.

        Args:
            table_data (Optional[List[List[str]]]): Table data to sanitize.

        Returns:
            Optional[List[List[str]]]: Sanitized table data.
        """
        if not table_data:
            return None
        return [[self._sanitize_text(str(cell)) for cell in row] for row in table_data]

    def _sanitize_data(self, data_list: List[Dict]) -> List[Dict]:
        """
        Sanitize Excel data to remove harmful content and enforce limits.

        Args:
            data_list (List[Dict]): List of dictionaries to sanitize.

        Returns:
            List[Dict]: Sanitized data.
        """
        if not data_list:
            return []
        sanitized = []
        for item in data_list:
            clean_item = {}
            for k, v in item.items():
                # Excel key limit with sanitization
                clean_key = self._sanitize_text(str(k))[:255]
                if isinstance(v, str):
                    clean_value = self._sanitize_text(v)[:32767]  # Excel cell limit
                else:
                    clean_value = v
                clean_item[clean_key] = clean_value
            sanitized.append(clean_item)
        return sanitized

    def _extract_pdf_text(self, file_path: str) -> str:
        """
        Extract text from PDF using pdfplumber.

        Args:
            file_path (str): Path to the PDF file.

        Returns:
            str: Extracted text content.

        Raises:
            FileOperationError: If PDF text extraction fails.
        """
        try:
            text_content = []
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text_content.append(page_text)
            return "\n".join(text_content)
        except Exception as e:
            raise FileOperationError(f"Failed to extract PDF text: {str(e)}")

    def _extract_image_text(self, file_path: str) -> str:
        """
        Extract text from image using pytesseract OCR.

        Args:
            file_path (str): Path to the image file.

        Returns:
            str: Extracted text content.

        Raises:
            FileOperationError: If image text extraction fails.
        """
        try:
            image: Image.Image = Image.open(file_path)
            # Convert to RGB if necessary
            if image.mode != "RGB":
                image = image.convert("RGB")
            text = pytesseract.image_to_string(image, lang="eng+chi_sim")
            return text.strip()
        except Exception as e:
            raise FileOperationError(f"Failed to extract image text: {str(e)}")

    def _extract_tika_text(self, file_path: str) -> str:
        """
        Extract text using Apache Tika as fallback.

        Args:
            file_path (str): Path to the file.

        Returns:
            str: Extracted text content.

        Raises:
            FileOperationError: If Tika text extraction fails.
        """
        try:
            parsed = parser.from_file(file_path)
            content = parsed.get("content", "")
            return content.strip() if content else ""
        except Exception as e:
            raise FileOperationError(f"Failed to extract text with Tika: {str(e)}")

    def read_docx(self, file_path: str, include_tables: bool = False) -> Dict[str, Any]:
        """
        Read content from a DOCX file.

        Args:
            file_path (str): Path to the DOCX file.
            include_tables (bool): Whether to include table data.

        Returns:
            Dict[str, Any]: Document content {'paragraphs': List[str], 'tables': Optional[List[List[List[str]]]]}.

        Raises:
            FileOperationError: If file cannot be read.
            ContentValidationError: If document structure is invalid.
        """
        try:
            self._validate_document(file_path, "docx")
            doc = DocxDocument(file_path)
            paras = [p.text for p in doc.paragraphs if p.text.strip()]
            tables = None
            if include_tables:
                tables = [[[cell.text for cell in row.cells] for row in table.rows] for table in doc.tables]
            return {"paragraphs": paras, "tables": tables}
        except ContentValidationError:
            raise
        except Exception as e:
            raise FileOperationError(f"Failed to read DOCX: {str(e)}")

    def write_docx(
        self,
        text: str,
        output_path: str,
        table_data: Optional[List[List[str]]] = None,
    ) -> Dict[str, Any]:
        """
        Write content to a DOCX file.

        Args:
            text (str): Text content to write.
            output_path (str): Path to save the DOCX file.
            table_data (Optional[List[List[str]]]): Table data to include.

        Returns:
            Dict[str, Any]: Status {'success': bool, 'file_path': str}.

        Raises:
            FileOperationError: If file cannot be written.
        """
        try:
            sanitized_text = self._sanitize_text(text)
            sanitized_table_data = self._sanitize_table_data(table_data)
            doc = DocxDocument()
            style = doc.styles["Normal"]
            style.font.name = self.config.default_font
            style.font.size = Pt(self.config.default_font_size)
            for line in sanitized_text.splitlines():
                doc.add_paragraph(line)
            if sanitized_table_data and sanitized_table_data[0]:
                # Find maximum number of columns to handle irregular table data
                max_cols = max(len(row) for row in sanitized_table_data)
                table = doc.add_table(rows=len(sanitized_table_data), cols=max_cols)
                for i, row in enumerate(sanitized_table_data):
                    for j in range(max_cols):
                        if j < len(row):
                            table.rows[i].cells[j].text = str(row[j])
                        else:
                            # Empty cell for missing data
                            table.rows[i].cells[j].text = ""
            doc.save(output_path)
            return {"success": True, "file_path": output_path}
        except Exception as e:
            raise FileOperationError(f"Failed to write DOCX: {str(e)}")

    def read_pptx(self, file_path: str) -> List[str]:
        """
        Read content from a PPTX file.

        Args:
            file_path (str): Path to the PPTX file.

        Returns:
            List[str]: List of text content from slides.

        Raises:
            FileOperationError: If file cannot be read.
            ContentValidationError: If document structure is invalid.
        """
        try:
            self._validate_document(file_path, "pptx")
            prs = Presentation(file_path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        txt = shape.text.strip()
                        if txt:
                            texts.append(txt)
            return texts
        except ContentValidationError:
            raise
        except Exception as e:
            raise FileOperationError(f"Failed to read PPTX: {str(e)}")

    def write_pptx(
        self,
        slides: List[str],
        output_path: str,
        image_path: Optional[str] = None,
    ) -> Dict[str, Any]:
        """
        Write content to a PPTX file.

        Args:
            slides (List[str]): List of slide contents.
            output_path (str): Path to save the PPTX file.
            image_path (Optional[str]): Path to an image to include on the first slide.

        Returns:
            Dict[str, Any]: Status {'success': bool, 'file_path': str}.

        Raises:
            FileOperationError: If file cannot be written.
        """
        try:
            sanitized_slides = [self._sanitize_text(slide) for slide in slides]
            prs = Presentation()
            blank = prs.slide_layouts[6]
            for idx, content in enumerate(sanitized_slides):
                slide = prs.slides.add_slide(blank)
                box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
                tf = box.text_frame
                lines = content.splitlines()
                if lines:
                    # Set text for the first paragraph (which already exists)
                    tf.text = lines[0]
                    # Add additional paragraphs for remaining lines
                    for line in lines[1:]:
                        p = tf.add_paragraph()
                        p.text = line
                if idx == 0 and image_path:
                    try:
                        slide.shapes.add_picture(image_path, Inches(1), Inches(6), Inches(4))
                    except Exception as img_err:
                        self.logger.warning(f"Could not add image to slide: {img_err}")
            prs.save(output_path)
            return {"success": True, "file_path": output_path}
        except Exception as e:
            raise FileOperationError(f"Failed to write PPTX: {str(e)}")

    def read_xlsx(self, file_path: str, sheet_name: Optional[str] = None) -> List[Dict]:
        """
        Read content from an XLSX file.

        Args:
            file_path (str): Path to the XLSX file.
            sheet_name (Optional[str]): Name of the sheet to read.

        Returns:
            List[Dict]: List of dictionaries representing Excel data.

        Raises:
            FileOperationError: If file cannot be read.
            ContentValidationError: If document structure is invalid.
        """
        try:
            self._validate_document(file_path, "xlsx")
            data = pd.read_excel(file_path, sheet_name=sheet_name)

            # Handle different return types from pd.read_excel()
            if isinstance(data, pd.DataFrame):
                # Single sheet or specific sheet requested
                return data.to_dict(orient="records")
            elif isinstance(data, dict):
                # Multiple sheets returned as dict - use the first sheet
                first_sheet_name = list(data.keys())[0]
                first_df = data[first_sheet_name]
                return first_df.to_dict(orient="records")
            else:
                raise FileOperationError("Unexpected data type returned from Excel file")

        except ContentValidationError:
            raise
        except Exception as e:
            raise FileOperationError(f"Failed to read XLSX: {str(e)}")

    def write_xlsx(self, data: List[Dict], output_path: str, sheet_name: str = "Sheet1") -> Dict[str, Any]:
        """
        Write content to an XLSX file.

        Args:
            data (List[Dict]): Data to write.
            output_path (str): Path to save the XLSX file.
            sheet_name (str): Name of the sheet.

        Returns:
            Dict[str, Any]: Status {'success': bool, 'file_path': str}.

        Raises:
            FileOperationError: If file cannot be written.
        """
        try:
            sanitized_data = self._sanitize_data(data)
            if not sanitized_data:
                pd.DataFrame().to_excel(output_path, index=False, sheet_name=sheet_name)
            else:
                pd.DataFrame(sanitized_data).to_excel(output_path, index=False, sheet_name=sheet_name)
            return {"success": True, "file_path": output_path}
        except Exception as e:
            raise FileOperationError(f"Failed to write XLSX: {str(e)}")

    def extract_text(self, file_path: str) -> str:
        """
        Extract text from various file formats using combination library approach.

        Args:
            file_path (str): Path to the file.

        Returns:
            str: Extracted text content.

        Raises:
            FileOperationError: If text extraction fails.
            ContentValidationError: If document structure is invalid.
        """
        try:
            file_ext = os.path.splitext(file_path)[1].lower()

            # Determine file type and validate
            if file_ext == ".pdf":
                file_type = "pdf"
            elif file_ext == ".docx":
                file_type = "docx"
            elif file_ext == ".pptx":
                file_type = "pptx"
            elif file_ext == ".xlsx":
                file_type = "xlsx"
            elif file_ext in [
                ".png",
                ".jpg",
                ".jpeg",
                ".tiff",
                ".bmp",
                ".gif",
            ]:
                file_type = "image"
            else:
                file_type = "other"

            # Validate document structure
            self._validate_document(file_path, file_type)

            # Extract text based on file type
            if file_type == "pdf":
                return self._sanitize_text(self._extract_pdf_text(file_path))
            elif file_type == "docx":
                doc = DocxDocument(file_path)
                paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                return self._sanitize_text("\n".join(paragraphs))
            elif file_type == "pptx":
                prs = Presentation(file_path)
                texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            texts.append(shape.text)
                return self._sanitize_text("\n".join(texts))
            elif file_type == "xlsx":
                data = pd.read_excel(file_path)
                # Handle different return types from pd.read_excel()
                if isinstance(data, pd.DataFrame):
                    return self._sanitize_text(data.to_string(index=False))
                elif isinstance(data, dict):
                    # Multiple sheets returned as dict - use the first sheet
                    first_sheet_name = list(data.keys())[0]
                    first_df = data[first_sheet_name]
                    return self._sanitize_text(first_df.to_string(index=False))
                else:
                    # Fallback for unexpected data types
                    return self._sanitize_text("")
            elif file_type == "image":
                return self._sanitize_text(self._extract_image_text(file_path))
            else:
                # Use Tika as fallback for other formats
                return self._sanitize_text(self._extract_tika_text(file_path))

        except ContentValidationError:
            raise
        except Exception as e:
            raise FileOperationError(f"Failed to extract text: {str(e)}")
