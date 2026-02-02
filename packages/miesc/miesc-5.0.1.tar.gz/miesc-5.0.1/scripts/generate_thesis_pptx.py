#!/usr/bin/env python3
"""
MIESC Thesis Defense Presentation Generator
Generates a professional PPTX for the thesis defense on December 18, 2025.

Author: Fernando Boiero
Date: December 2025
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
from pptx.dml.color import RGBColor
import os

# Color scheme - Professional dark theme
COLORS = {
    'bg_dark': RGBColor(26, 26, 46),       # #1a1a2e
    'bg_darker': RGBColor(15, 15, 35),     # #0f0f23
    'cyan': RGBColor(0, 212, 255),         # #00d4ff
    'green': RGBColor(0, 255, 157),        # #00ff9d
    'red': RGBColor(255, 107, 107),        # #ff6b6b
    'white': RGBColor(234, 234, 234),      # #eaeaea
    'yellow': RGBColor(255, 215, 0),       # Gold
}

def set_slide_background(slide, color):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs, title, subtitle=""):
    """Add a title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_background(slide, COLORS['bg_darker'])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['cyan']
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        # Subtitle
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = COLORS['green']
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, content_lines, bullet_color=None):
    """Add a content slide with bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_background(slide, COLORS['bg_dark'])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['cyan']

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Check if it's a bullet point
        if line.startswith('- ') or line.startswith('• '):
            p.text = "• " + line[2:]
            p.level = 0
        elif line.startswith('  - ') or line.startswith('  • '):
            p.text = "  • " + line[4:]
            p.level = 1
        else:
            p.text = line

        p.font.size = Pt(20)
        p.font.color.rgb = bullet_color or COLORS['white']
        p.space_after = Pt(8)

    return slide

def add_table_slide(prs, title, headers, rows):
    """Add a slide with a table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_background(slide, COLORS['bg_dark'])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['cyan']

    # Table
    cols = len(headers)
    table_rows = len(rows) + 1
    table = slide.shapes.add_table(table_rows, cols, Inches(0.5), Inches(1.3), Inches(9), Inches(0.5 * table_rows)).table

    # Headers
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].font.color.rgb = COLORS['cyan']

    # Data rows
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].font.color.rgb = COLORS['white']

    return slide

def add_metrics_slide(prs, title, metrics):
    """Add a slide with highlighted metrics."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_background(slide, COLORS['bg_dark'])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['cyan']

    # Metrics grid (2 columns)
    col_width = 4.25
    row_height = 1.2
    start_y = 1.5

    for i, (metric_name, metric_value, color) in enumerate(metrics):
        col = i % 2
        row = i // 2
        x = 0.5 + (col * col_width)
        y = start_y + (row * row_height)

        # Metric box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(4), Inches(1))
        box.fill.solid()
        box.fill.fore_color.rgb = COLORS['bg_darker']
        box.line.color.rgb = color
        box.line.width = Pt(2)

        # Metric name
        name_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.1), Inches(3.8), Inches(0.4))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = metric_name
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['white']

        # Metric value
        val_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.5), Inches(3.8), Inches(0.5))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = metric_value
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = color

    return slide

def add_image_slide(prs, title, image_path, subtitle=""):
    """Add a slide with an image - white background for visibility."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    # Use WHITE background so figures with dark text are visible
    set_slide_background(slide, RGBColor(255, 255, 255))

    # Title - dark blue on white
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue

    # Image
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(0.3), Inches(0.9), width=Inches(9.4))
    else:
        # Placeholder text if image not found
        msg_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
        tf = msg_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"[Imagen: {os.path.basename(image_path)}]"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(200, 0, 0)  # Red
        p.alignment = PP_ALIGN.CENTER

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = RGBColor(80, 80, 80)  # Dark gray
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_code_slide(prs, title, code, language="python"):
    """Add a slide with code."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_background(slide, COLORS['bg_dark'])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['cyan']

    # Code box
    code_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.0), Inches(9), Inches(5.5))
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = RGBColor(30, 30, 30)
    code_box.line.color.rgb = COLORS['green']

    # Code text
    text_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.6), Inches(5.1))
    tf = text_box.text_frame
    tf.word_wrap = False

    for i, line in enumerate(code.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.name = "Consolas"
        p.font.color.rgb = COLORS['green']

    return slide

def generate_presentation():
    """Generate the complete thesis defense presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ==================== SLIDE 1: Title ====================
    slide = add_title_slide(prs, "MIESC", "Multi-layer Intelligent Evaluation for Smart Contracts")
    # Add subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Un Enfoque de Ciberdefensa para la Seguridad de Contratos Inteligentes"
    p.font.size = Pt(20)
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    # ==================== SLIDE 2: Author Data ====================
    add_content_slide(prs, "Datos de la Tesis", [
        "Autor: Ing. Fernando Boiero",
        "Email: fboiero@frvm.utn.edu.ar",
        "",
        "Programa: Maestria en Ciberdefensa",
        "Institucion: Universidad de la Defensa Nacional (UNDEF)",
        "Instituto Universitario Aeronautico (IUA) - Cordoba",
        "",
        "Director: Mg. Eduardo Casanovas",
        "Fecha de Defensa: 18 de Diciembre 2025"
    ])

    # ==================== SLIDE 3: Agenda ====================
    add_content_slide(prs, "Agenda", [
        "1. Contexto y Motivacion",
        "2. Problema de Investigacion",
        "3. Marco Teorico",
        "4. Solucion Propuesta: MIESC",
        "5. Resultados Experimentales",
        "6. Demo en Vivo",
        "7. Conclusiones",
        "8. Trabajo Post-Tesis (para contexto)",
        "9. Trabajos Futuros",
        "10. Preguntas"
    ])

    # ==================== SLIDE 4: Section - Context ====================
    add_title_slide(prs, "1. Contexto y Motivacion", "El Ciberespacio como Dominio Estrategico")

    # ==================== SLIDE 5: The Threat ====================
    add_table_slide(prs, "La Amenaza es Real: $7.8+ MIL MILLONES perdidos",
        ["Ano", "Incidente", "Perdida", "Vulnerabilidad"],
        [
            ["2016", "The DAO", "$60M", "Reentrancy"],
            ["2017", "Parity Wallet", "$280M", "Access Control"],
            ["2022", "Wormhole", "$320M", "Signature Bypass"],
            ["2022", "Ronin Bridge", "$625M", "Key Compromise"],
            ["2023", "Euler Finance", "$197M", "Flash Loan"]
        ]
    )

    # ==================== SLIDE 6: Threat Taxonomy ====================
    add_image_slide(prs, "Taxonomia de Amenazas a Smart Contracts",
                    "thesis_generator/figuras/fig_08_threat_taxonomy.png",
                    "Figura: Clasificacion de vulnerabilidades por categoria")

    # ==================== SLIDE 7: Fragmentation Problem ====================
    add_content_slide(prs, "El Problema: Fragmentacion", [
        "Heterogeneidad de Enfoques:",
        "- Analisis estatico (Slither, Solhint)",
        "- Fuzzers (Echidna, Medusa)",
        "- Ejecucion simbolica (Mythril)",
        "- Verificacion formal (Certora)",
        "- IA (GPTScan)",
        "",
        "Problemas Criticos:",
        "- Salidas incompatibles",
        "- Nomenclaturas diferentes",
        "- Cobertura incompleta",
        "- NINGUNA herramienta detecta >70%"
    ])

    # ==================== SLIDE 7: Data Sovereignty ====================
    add_content_slide(prs, "El Problema: Soberania de Datos", [
        "Riesgo de Confidencialidad con IA comercial:",
        "",
        "- Propiedad intelectual: Codigo fuente enviado a terceros",
        "- Dependencia externa: Perdida de capacidad operativa",
        "- Cumplimiento normativo: GDPR, LGPD, regulaciones",
        "- Trazabilidad: Imposibilidad de auditar el procesamiento",
        "",
        "En ciberdefensa, la confidencialidad del codigo es CRITICA"
    ], COLORS['red'])

    # ==================== SLIDE 8: Section - Problem ====================
    add_title_slide(prs, "2. Planteamiento del Problema", "")

    # ==================== SLIDE 9: Specific Problems ====================
    add_content_slide(prs, "Problemas Especificos", [
        "P1: No existe un framework que integre coherentemente las",
        "    principales herramientas de analisis de seguridad.",
        "",
        "P2: Las salidas de herramientas existentes utilizan",
        "    nomenclaturas y formatos incompatibles.",
        "",
        "P3: Las soluciones con IA dependen de servicios externos,",
        "    comprometiendo la confidencialidad.",
        "",
        "P4: No existe una arquitectura que aplique Defense-in-Depth",
        "    a la seguridad de smart contracts."
    ])

    # ==================== SLIDE 10: Objectives ====================
    add_content_slide(prs, "Objetivos", [
        "Objetivo General:",
        "Desarrollar un framework de codigo abierto que integre multiples",
        "herramientas de analisis en una arquitectura de defensa en",
        "profundidad, con IA soberana.",
        "",
        "Objetivos Especificos:",
        "1. Integrar 25 herramientas en 7 capas de defensa",
        "2. Normalizar salidas a taxonomias estandar (SWC/CWE/OWASP)",
        "3. Implementar backend de IA 100% local (Ollama)",
        "4. Cumplir estandares Digital Public Good",
        "5. Integrar con asistentes IA via MCP Protocol"
    ])

    # ==================== SLIDE 11: Section - Theory ====================
    add_title_slide(prs, "3. Marco Teorico", "")

    # ==================== SLIDE 12: Theoretical Foundations ====================
    add_content_slide(prs, "Fundamentos Teoricos", [
        "Defense-in-Depth (Saltzer & Schroeder, 1975):",
        "\"Multiples capas de defensa independientes, de modo que",
        "la falla de una no comprometa la seguridad total\"",
        "",
        "Multi-Tool Analysis (Durieux et al., 2020):",
        "\"La combinacion de herramientas complementarias mejora",
        "significativamente la deteccion de vulnerabilidades\"",
        "",
        "Taxonomias: SWC Registry (37), CWE (900+), OWASP SC Top 10"
    ])

    # ==================== SLIDE 13: Section - Solution ====================
    add_title_slide(prs, "4. Solucion Propuesta", "MIESC v4.0.0")

    # ==================== SLIDE 14: MIESC Overview ====================
    add_metrics_slide(prs, "MIESC: Vision General", [
        ("Herramientas Integradas", "25", COLORS['cyan']),
        ("Capas de Defensa", "7", COLORS['cyan']),
        ("Recall Vulnerabilidades", "100%", COLORS['green']),
        ("Mejora vs Individual", "+40.8%", COLORS['green']),
        ("Costo Operativo", "$0", COLORS['yellow']),
        ("IA Soberana", "Ollama Local", COLORS['yellow'])
    ])

    # ==================== SLIDE 15: 7 Layer Architecture ====================
    add_image_slide(prs, "Arquitectura Defense-in-Depth de 7 Capas",
                    "thesis_generator/figuras/fig_01_defense_in_depth.png",
                    "Figura: Arquitectura de defensa en profundidad de MIESC")

    # ==================== SLIDE 16: Adapter Pattern ====================
    add_image_slide(prs, "Patron Adapter: Integracion Unificada",
                    "thesis_generator/figuras/fig_02_adapter_pattern.png",
                    "Figura: Diagrama de clases del patron Adapter")

    # ==================== SLIDE 17: Layer 1 ====================
    add_content_slide(prs, "Capa 1: Analisis Estatico", [
        "Herramientas: Slither, Solhint, Securify2, Semgrep",
        "",
        "Capacidades:",
        "- 90+ detectores de vulnerabilidades",
        "- Analisis de flujo de datos",
        "- Deteccion de patrones inseguros",
        "- Verificacion de mejores practicas",
        "",
        "Tiempo: ~5 segundos por contrato"
    ])

    # ==================== SLIDE 17: Layer 2 ====================
    add_content_slide(prs, "Capa 2: Testing Dinamico (Fuzzing)", [
        "Herramientas: Echidna, Medusa, Foundry Fuzz, DogeFuzz",
        "",
        "Capacidades:",
        "- Fuzzing basado en propiedades",
        "- Coverage-guided testing",
        "- Generacion automatica de inputs",
        "- Deteccion de invariantes violados",
        "",
        "Mejora v4.0: DogeFuzz - AFL-style, 3x mas rapido"
    ])

    # ==================== SLIDE 18: Layer 3 ====================
    add_content_slide(prs, "Capa 3: Ejecucion Simbolica", [
        "Herramientas: Mythril, Manticore, Halmos, Oyente",
        "",
        "Capacidades:",
        "- Exploracion exhaustiva de paths",
        "- Deteccion de condiciones de overflow",
        "- Verificacion de assertions",
        "- Analisis de dependencias",
        "",
        "Tiempo: 1-5 minutos (cuello de botella principal)"
    ])

    # ==================== SLIDE 19: Layers 4-5 ====================
    add_content_slide(prs, "Capas 4-5: Verificacion Formal", [
        "Herramientas: SMTChecker, Certora, Scribble, Halmos",
        "",
        "Capacidades:",
        "- Verificacion matematica de propiedades",
        "- Deteccion de violaciones de invariantes",
        "- Pruebas formales de correctitud",
        "",
        "Mejora v4.0: PropertyGPT",
        "- 80% recall en propiedades ground-truth",
        "- +700% adopcion de verificacion formal"
    ])

    # ==================== SLIDE 20: Layers 6-7 ====================
    add_content_slide(prs, "Capas 6-7: Analisis con IA", [
        "Herramientas: SmartLLM, GPTScan, LLMSmartAudit, ThreatModel",
        "",
        "Capacidades:",
        "- Correlacion de hallazgos",
        "- Analisis semantico profundo",
        "- Modelado de amenazas",
        "- Recomendaciones de remediacion",
        "",
        "Mejora v4.0: RAG + Verificator",
        "- +17% precision (75% -> 88%)",
        "- -52% falsos positivos"
    ])

    # ==================== SLIDE 21: RAG Architecture ====================
    add_image_slide(prs, "Arquitectura RAG en SmartLLM",
                    "thesis_generator/figuras/fig_10_rag_architecture.png",
                    "Figura: Retrieval Augmented Generation para analisis de vulnerabilidades")

    # ==================== SLIDE 22: Normalization Flow ====================
    add_image_slide(prs, "Flujo de Normalizacion SWC/CWE/OWASP",
                    "thesis_generator/figuras/fig_03_normalization_flow.png",
                    "Figura: Proceso de normalizacion de hallazgos")

    # ==================== SLIDE 23: Adapter Pattern ====================
    add_code_slide(prs, "Patron Adapter: Integracion Unificada", """
class ToolAdapter(ABC):
    @abstractmethod
    def analyze(self, contract_path: str) -> ToolResult:
        pass

    @abstractmethod
    def is_available(self) -> bool:
        pass

    @abstractmethod
    def get_metadata(self) -> ToolMetadata:
        pass

# 25 adapters implementados siguiendo este patron
""")

    # ==================== SLIDE 24: Sovereign AI ====================
    add_content_slide(prs, "IA Soberana con Ollama", [
        "CODIGO NUNCA SALE DE TU MAQUINA",
        "",
        "Problema con APIs comerciales:",
        "- OpenAI: codigo enviado a servidores USA",
        "- Costo: $0.03-$0.10 por analisis",
        "",
        "Solucion MIESC:",
        "- Ollama local (deepseek-coder, codellama)",
        "- Procesamiento 100% on-premise",
        "- Costo: $0.00",
        "- Cumple DPGA Standard"
    ])

    # ==================== SLIDE 25: Sovereign AI Diagram ====================
    add_image_slide(prs, "Arquitectura de IA Soberana",
                    "thesis_generator/figuras/fig_06_sovereign_ai.png",
                    "Figura: Codigo procesado localmente, sin envio a APIs externas")

    # ==================== SLIDE 26: MCP Integration ====================
    add_code_slide(prs, "Integracion MCP (Model Context Protocol)", """
{
  "mcpServers": {
    "miesc": {
      "url": "http://localhost:8080/mcp/jsonrpc"
    }
  }
}

Endpoints disponibles:
- run_audit: Ejecutar auditoria
- correlate_findings: Correlacionar hallazgos
- map_compliance: Mapear a estandares
- generate_report: Generar reportes
""")

    # ==================== SLIDE 27: MCP Architecture ====================
    add_image_slide(prs, "Arquitectura de Integracion MCP",
                    "thesis_generator/figuras/fig_04_mcp_architecture.png",
                    "Figura: Model Context Protocol - Integracion con Claude Desktop")

    # ==================== SLIDE 28: Section - Results ====================
    add_title_slide(prs, "5. Resultados Experimentales", "")

    # ==================== SLIDE 25: Methodology ====================
    add_content_slide(prs, "Metodologia de Evaluacion", [
        "Tipo de estudio: Evaluacion comparativa con benchmark controlado",
        "",
        "Preguntas de investigacion:",
        "- RQ1: Integracion exitosa de 25 herramientas?",
        "- RQ2: Mejora deteccion vs herramientas individuales?",
        "- RQ3: Reduccion efectiva de duplicados?",
        "- RQ4: Viabilidad para produccion?",
        "",
        "Corpus: 4 contratos, 14 vulnerabilidades conocidas, 7 categorias SWC"
    ])

    # ==================== SLIDE 26: RQ1 ====================
    add_table_slide(prs, "RQ1: Integracion de Herramientas - 100% (25/25)",
        ["Capa", "Herramientas", "Estado"],
        [
            ["1 - Estatico", "Slither, Solhint, Securify2, Semgrep", "4/4"],
            ["2 - Fuzzing", "Echidna, Medusa, Foundry, DogeFuzz", "4/4"],
            ["3 - Simbolico", "Mythril, Manticore, Halmos, Oyente", "4/4"],
            ["4-5 - Formal", "SMTChecker, Certora, Scribble, PropertyGPT", "4/4"],
            ["6-7 - IA", "SmartLLM, GPTScan, ThreatModel, etc.", "9/9"]
        ]
    )

    # ==================== SLIDE 27: RQ2 ====================
    add_table_slide(prs, "RQ2: Mejora de Deteccion - +40.8%",
        ["Herramienta", "Precision", "Recall", "F1-Score"],
        [
            ["Slither (individual)", "74%", "66%", "0.70"],
            ["Mythril (individual)", "68%", "59%", "0.63"],
            ["Echidna (individual)", "71%", "62%", "0.66"],
            ["MIESC (7 capas)", "94.5%", "92.8%", "0.936"]
        ]
    )

    # ==================== SLIDE 28: RQ3 ====================
    add_metrics_slide(prs, "RQ3: Reduccion de Duplicados - 66%", [
        ("Hallazgos Brutos (raw)", "147", COLORS['red']),
        ("Hallazgos Unicos", "50", COLORS['green']),
        ("Duplicados Eliminados", "97", COLORS['yellow']),
        ("Reduccion", "66%", COLORS['green']),
        ("Precision Normalizacion", "97.1%", COLORS['cyan']),
        ("Taxonomias", "SWC/CWE/OWASP", COLORS['cyan'])
    ])

    # ==================== SLIDE 29: RQ4 ====================
    add_metrics_slide(prs, "RQ4: Viabilidad en Produccion", [
        ("Tiempo por Contrato", "~2 min", COLORS['cyan']),
        ("Costo por Auditoria", "$0.00", COLORS['green']),
        ("Tests Unitarios", "1,277 pasando", COLORS['green']),
        ("Cobertura Codigo", "59.42%", COLORS['yellow']),
        ("Compliance Index", "91.4%", COLORS['cyan']),
        ("DPG Application", "GID0092948", COLORS['yellow'])
    ])

    # ==================== SLIDE 30: Comparison Chart ====================
    add_image_slide(prs, "Comparativa de Rendimiento",
                    "thesis_generator/figuras/fig_07_comparison.png",
                    "Figura: MIESC vs herramientas individuales (Precision, Recall, F1)")

    # ==================== SLIDE 31: Severity Distribution ====================
    add_image_slide(prs, "Distribucion de Hallazgos por Severidad",
                    "thesis_generator/figuras/fig_05_severity_distribution.png",
                    "Figura: Clasificacion de vulnerabilidades por severidad y capa")

    # ==================== SLIDE 32: Execution Timeline ====================
    add_image_slide(prs, "Timeline de Ejecucion Paralela",
                    "thesis_generator/figuras/fig_09_execution_timeline.png",
                    "Figura: Ejecucion paralela de capas para optimizar tiempos")

    # ==================== SLIDE 33: Section - Demo ====================
    add_title_slide(prs, "6. Demo en Vivo", "Ejecutando MIESC contra VulnerableBank.sol")

    # ==================== SLIDE 32: Demo Code ====================
    add_code_slide(prs, "Demo: Vulnerabilidad Detectada", """
function withdraw() public {
    uint256 balance = balances[msg.sender];
    require(balance > 0, "No balance");

    // VULNERABILIDAD: External call ANTES de state update
    (bool success, ) = msg.sender.call{value: balance}("");
    require(success, "Transfer failed");

    // State update DESPUES - puede ser re-entered!
    balances[msg.sender] = 0;
}

// SWC-107: Reentrancy | Severidad: CRITICAL
""")

    # ==================== SLIDE 33: Section - Conclusions ====================
    add_title_slide(prs, "7. Conclusiones", "")

    # ==================== SLIDE 34: Objectives Achieved ====================
    add_table_slide(prs, "Objetivos Alcanzados",
        ["Objetivo", "Indicador", "Resultado", "Estado"],
        [
            ["Integrar herramientas", "25 operativas", "25/25 (100%)", "CUMPLIDO"],
            ["Defensa en profundidad", "7 capas", "7 implementadas", "CUMPLIDO"],
            ["Normalizar resultados", "Mapeo SWC/CWE/OWASP", "97.1% precision", "CUMPLIDO"],
            ["Eliminar dep. comerciales", "Costo $0", "$0/auditoria", "CUMPLIDO"],
            ["Mejorar deteccion", ">20% recall", "+40.8%", "SUPERADO"],
            ["Reducir duplicados", ">40%", "66%", "SUPERADO"]
        ]
    )

    # ==================== SLIDE 35: Main Contributions ====================
    add_content_slide(prs, "Contribuciones Principales", [
        "1. Arquitectura de 7 Capas:",
        "   Primera implementacion de Defense-in-Depth para smart contracts",
        "",
        "2. Protocolo ToolAdapter:",
        "   Interfaz unificada para herramientas heterogeneas",
        "",
        "3. Normalizacion Triple:",
        "   Mapeo automatico SWC/CWE/OWASP",
        "",
        "4. IA Soberana:",
        "   Eliminacion de dependencias de APIs comerciales",
        "",
        "5. MCP Server:",
        "   Primera herramienta de auditoria con Model Context Protocol"
    ])

    # ==================== SLIDE 36: Limitations ====================
    add_content_slide(prs, "Limitaciones", [
        "Tecnicas:",
        "- Escalabilidad en contratos >1000 LOC",
        "- Calidad depende del modelo LLM",
        "- Vulnerabilidades de logica de negocio (flash loans, MEV)",
        "",
        "Metodologicas:",
        "- Corpus de prueba limitado (4 contratos)",
        "- Ausencia de validacion en mainnet",
        "- Metricas de IA subjetivas"
    ])

    # ==================== SLIDE 37: Section - Post-Thesis Work ====================
    add_title_slide(prs, "Trabajo Posterior a la Tesis", "Desarrollado entre Octubre y Diciembre 2025")

    # ==================== SLIDE 38: Post-Thesis Overview ====================
    add_content_slide(prs, "Trabajo Post-Tesis (v4.1.0 - v4.2.0)", [
        "IMPORTANTE: Lo siguiente NO forma parte de la tesis entregada",
        "el 23 de Octubre 2025. Se presenta para contexto.",
        "",
        "Nuevas Capas implementadas:",
        "- Capa 8: Analisis de Seguridad DeFi (flash loans, oracles, MEV)",
        "- Capa 9: Seguridad de Dependencias (CVE database, supply chain)",
        "",
        "Trabajos Futuros implementados:",
        "- TF-5.2: GitHub Action CI/CD",
        "- TF-5.3: WebSocket Streaming",
        "- TF-5.4: Dashboard Streamlit mejorado"
    ], COLORS['yellow'])

    # ==================== SLIDE 39: Post-Thesis Features ====================
    add_table_slide(prs, "Funcionalidades Post-Tesis",
        ["Feature", "Descripcion", "Version"],
        [
            ["Multi-format Exporters", "SARIF, SonarQube, Checkmarx, Markdown", "v4.2.0"],
            ["Prometheus Metrics", "Observabilidad con metricas", "v4.2.0"],
            ["WebSocket Real-Time", "16 tipos de eventos en tiempo real", "v4.1.0"],
            ["OpenAPI 3.1.0", "33 endpoints, 61 schemas", "v4.1.0+"],
            ["DPGA Application", "Digital Public Good (GID0092948)", "v4.2.0"]
        ]
    )

    # ==================== SLIDE 40: Metrics Evolution ====================
    add_table_slide(prs, "Evolucion de Metricas",
        ["Metrica", "Tesis v4.0.0", "Actual v4.2.0"],
        [
            ["Herramientas", "25", "25+ (Capas 8-9)"],
            ["Capas de defensa", "7", "9"],
            ["Tests unitarios", "~100", "1,277"],
            ["Lineas de codigo", "~30K", "~51K"],
            ["OpenAPI Endpoints", "-", "33"]
        ]
    )

    # ==================== SLIDE 41: Section - Future Work ====================
    add_title_slide(prs, "9. Trabajos Futuros", "Lineas de investigacion propuestas")

    # ==================== SLIDE 42: Future Lines ====================
    add_table_slide(prs, "Lineas de Investigacion Futuras",
        ["Linea", "Descripcion", "Impacto"],
        [
            ["Multi-chain", "Soporte para Solana, Cairo, Soroban", "Alto"],
            ["Fine-tuning LLM", "Modelos especializados en smart contracts", "Alto"],
            ["Runtime Monitoring", "Deteccion en tiempo real post-deployment", "Alto"],
            ["Automated Patching", "Generacion automatica de parches", "Medio"],
            ["IDE Integration", "Extension VSCode", "Medio"]
        ]
    )

    # ==================== SLIDE 39: Thanks ====================
    slide = add_title_slide(prs, "Gracias", "Preguntas?")

    # ==================== SLIDE 40: References ====================
    add_content_slide(prs, "Referencias Principales", [
        "- Atzei, N., et al. (2017). \"A Survey of Attacks on Ethereum",
        "  Smart Contracts\"",
        "",
        "- Durieux, T., et al. (2020). \"Empirical Review of Automated",
        "  Analysis Tools on 47,587 Contracts\"",
        "",
        "- Saltzer, J. & Schroeder, M. (1975). \"The Protection of",
        "  Information in Computer Systems\"",
        "",
        "- Anthropic (2024). \"Model Context Protocol Specification\"",
        "",
        "- DPGA (2023). \"Digital Public Goods Standard\""
    ])

    # ==================== SLIDE 41: Contact ====================
    add_content_slide(prs, "Contacto", [
        "Fernando Boiero",
        "- Email: fboiero@frvm.utn.edu.ar",
        "- GitHub: github.com/fboiero/MIESC",
        "- Documentacion: fboiero.github.io/MIESC",
        "",
        "MIESC v4.0.0",
        "- Licencia: AGPL-3.0",
        "- DPG Application: GID0092948"
    ])

    # ==================== SLIDE 42: Final ====================
    slide = add_title_slide(prs, "MIESC", "Securing the Future of Smart Contracts")
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Defense-in-Depth for the Blockchain Era"
    p.font.size = Pt(20)
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    # Save presentation
    output_path = "docs/MIESC_Defensa_Tesis_18Dic2025.pptx"
    prs.save(output_path)
    print(f"Presentacion generada: {output_path}")
    print(f"Total slides: {len(prs.slides)}")
    return output_path

if __name__ == "__main__":
    generate_presentation()
